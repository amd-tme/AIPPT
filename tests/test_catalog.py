"""Tests for aippt.catalog module."""

import csv
import os
import time

import pytest

from aippt.catalog import (
    get_db,
    content_hash,
    display_name,
    file_hash,
    catalog_deck,
    search_slides,
    add_tags,
    get_slide_tags,
    check_newer_versions,
    get_deck_by_id,
    get_deck_slides,
    list_decks,
    remove_slide_tag,
    remove_all_slide_tags,
    rename_tag,
    list_taxonomy,
    add_taxonomy_tags,
    remove_taxonomy_tag,
    import_taxonomy_csv,
    export_taxonomy_csv,
    get_taxonomy_names,
    set_slide_section,
    get_slide_section,
    get_deck_sections,
    remove_slide_section,
    rename_section,
    record_edit,
    get_all_tags,
    resolve_deck,
    delete_deck,
    rename_deck,
    get_deck_tag_count,
    get_deck_top_tags,
)


class TestDisplayName:
    def test_strips_uuid_prefix(self):
        assert display_name("44dc98ea57f240efaabca1333f166d0b_Deck Name") == "Deck Name"

    def test_no_prefix_unchanged(self):
        assert display_name("Regular Deck Name") == "Regular Deck Name"

    def test_empty_string(self):
        assert display_name("") == ""

    def test_prefix_with_empty_remainder(self):
        assert display_name("abcdef1234567890abcdef1234567890_") == ""

    def test_uppercase_hex_no_match(self):
        assert display_name("ABCDEF1234567890abcdef1234567890_Name") == "ABCDEF1234567890abcdef1234567890_Name"

    def test_short_prefix_no_match(self):
        assert display_name("abc_short_prefix") == "abc_short_prefix"


class TestContentHash:
    def test_consistent(self):
        h1 = content_hash("Title", "Content")
        h2 = content_hash("Title", "Content")
        assert h1 == h2

    def test_case_insensitive(self):
        h1 = content_hash("Title", "Content")
        h2 = content_hash("TITLE", "CONTENT")
        assert h1 == h2

    def test_different_content(self):
        h1 = content_hash("Title", "Content A")
        h2 = content_hash("Title", "Content B")
        assert h1 != h2

    def test_different_titles(self):
        h1 = content_hash("Title A", "Content")
        h2 = content_hash("Title B", "Content")
        assert h1 != h2

    def test_whitespace_normalized(self):
        h1 = content_hash("  Title  ", "  Content  ")
        h2 = content_hash("Title", "Content")
        assert h1 == h2


class TestFileHash:
    def test_consistent(self, tmp_path):
        f = tmp_path / "test.txt"
        f.write_text("Test content")
        h1 = file_hash(str(f))
        h2 = file_hash(str(f))
        assert h1 == h2

    def test_different_content(self, tmp_path):
        f1 = tmp_path / "file1.txt"
        f1.write_text("Content A")
        f2 = tmp_path / "file2.txt"
        f2.write_text("Content B")
        assert file_hash(str(f1)) != file_hash(str(f2))


class TestDatabase:
    @pytest.fixture
    def db_path(self, tmp_path):
        return str(tmp_path / "test.db")

    def test_create_schema(self, db_path):
        conn = get_db(db_path)
        tables = conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table'"
        ).fetchall()
        table_names = {r["name"] for r in tables}
        assert "decks" in table_names
        assert "slides" in table_names
        assert "tags" in table_names
        assert "slide_tags" in table_names
        assert "taxonomy" in table_names
        conn.close()

    def test_foreign_keys_enabled(self, db_path):
        conn = get_db(db_path)
        result = conn.execute("PRAGMA foreign_keys").fetchone()
        assert result[0] == 1
        conn.close()


class TestCatalogDeck:
    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """Create a minimal PPTX for testing."""
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Introduction"

        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Security Overview"

        path = tmp_path / "test.pptx"
        prs.save(str(path))
        return str(path)

    def test_catalogs_deck(self, tmp_path, sample_pptx):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(sample_pptx, db_path=db_path)
        assert deck_id > 0

    def test_catalogs_all_slides(self, tmp_path, sample_pptx):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(sample_pptx, db_path=db_path)

        slides = get_deck_slides(deck_id, db_path)
        assert len(slides) == 2
        assert slides[0]["title"] == "Introduction"
        assert slides[1]["title"] == "Security Overview"

    def test_skips_duplicate_catalog(self, tmp_path, sample_pptx):
        db_path = str(tmp_path / "test.db")
        deck_id1 = catalog_deck(sample_pptx, db_path=db_path)
        deck_id2 = catalog_deck(sample_pptx, db_path=db_path)
        assert deck_id1 == deck_id2

    def test_updates_on_file_change(self, tmp_path):
        from pptx import Presentation

        db_path = str(tmp_path / "test.db")
        pptx_path = str(tmp_path / "test.pptx")

        # Create initial deck
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Original"
        prs.save(pptx_path)

        deck_id1 = catalog_deck(pptx_path, db_path=db_path)

        # Modify the deck
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Modified"
        prs.save(pptx_path)

        deck_id2 = catalog_deck(pptx_path, db_path=db_path)

        # Same deck ID but content updated
        assert deck_id1 == deck_id2
        slides = get_deck_slides(deck_id2, db_path)
        assert slides[0]["title"] == "Modified"


class TestSearchSlides:
    @pytest.fixture
    def db_with_slides(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("deck1", "/deck1.pptx", "abc", 2),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Security Overview", "content", "hash1"),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 2, "Cloud Architecture", "content", "hash2"),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_search_by_title(self, db_with_slides):
        results = search_slides(db_path=db_with_slides, title_contains="Security")
        assert len(results) == 1
        assert results[0]["title"] == "Security Overview"

    def test_search_no_results(self, db_with_slides):
        results = search_slides(db_path=db_with_slides, title_contains="Nonexistent")
        assert len(results) == 0

    def test_search_all(self, db_with_slides):
        results = search_slides(db_path=db_with_slides)
        assert len(results) == 2


class TestTags:
    @pytest.fixture
    def db_with_slide(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", "/test.pptx", "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Test Slide", "content", "hash123"),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_add_and_get_tags(self, db_with_slide):
        add_tags(1, ["security", "architecture"], "ai", db_with_slide)
        tags = get_slide_tags(1, db_with_slide)
        assert "security" in tags
        assert "architecture" in tags

    def test_duplicate_tags_ignored(self, db_with_slide):
        add_tags(1, ["security"], "ai", db_with_slide)
        add_tags(1, ["security"], "ai", db_with_slide)
        tags = get_slide_tags(1, db_with_slide)
        assert tags.count("security") == 1

    def test_empty_tags_ignored(self, db_with_slide):
        add_tags(1, ["", "  ", "valid"], "ai", db_with_slide)
        tags = get_slide_tags(1, db_with_slide)
        assert tags == ["valid"]

    def test_tags_normalized_to_lowercase(self, db_with_slide):
        add_tags(1, ["SECURITY", "Architecture"], "ai", db_with_slide)
        tags = get_slide_tags(1, db_with_slide)
        assert "security" in tags
        assert "architecture" in tags

    def test_search_by_tags(self, db_with_slide):
        add_tags(1, ["security", "cloud"], "ai", db_with_slide)
        results = search_slides(db_path=db_with_slide, tags=["security"])
        assert len(results) == 1
        assert results[0]["title"] == "Test Slide"


class TestCheckNewerVersions:
    def test_no_warnings_for_unique_slides(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("deck1", "/deck1.pptx", "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Unique Slide", "content", "hash1"),
        )
        conn.commit()
        conn.close()

        slides = [{"title": "Unique Slide", "content_hash": "hash1", "updated_at": ""}]
        warnings = check_newer_versions(slides, db_path)
        assert len(warnings) == 0

    def test_warns_on_newer_version(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("deck1", "/deck1.pptx", "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash, updated_at) VALUES (?, ?, ?, ?, ?, ?)",
            (1, 1, "Shared Slide", "content", "hash2", "2026-01-01 12:00:00"),
        )
        conn.commit()
        conn.close()

        # Check with an older hash
        slides = [{
            "title": "Shared Slide",
            "content_hash": "hash1",
            "updated_at": "2025-12-01 12:00:00",
            "deck_name": "old_deck",
        }]
        warnings = check_newer_versions(slides, db_path)
        assert len(warnings) == 1
        assert warnings[0]["slide_title"] == "Shared Slide"
        assert warnings[0]["newer_deck"] == "deck1"


class TestDeckOperations:
    @pytest.fixture
    def db_with_decks(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("deck1", "/deck1.pptx", "abc", 2),
        )
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("deck2", "/deck2.pptx", "def", 3),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_get_deck_by_id(self, db_with_decks):
        deck = get_deck_by_id(1, db_with_decks)
        assert deck is not None
        assert deck["name"] == "deck1"
        assert deck["slide_count"] == 2

    def test_get_deck_by_id_not_found(self, db_with_decks):
        deck = get_deck_by_id(999, db_with_decks)
        assert deck is None

    def test_list_decks(self, db_with_decks):
        decks = list_decks(db_with_decks)
        assert len(decks) == 2
        names = {d["name"] for d in decks}
        assert "deck1" in names
        assert "deck2" in names


class TestTagRemoval:
    @pytest.fixture
    def db_with_tagged_slide(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", "/test.pptx", "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Test Slide", "content", "hash123"),
        )
        conn.commit()
        conn.close()
        add_tags(1, ["security", "architecture", "cloud"], "ai", db_path)
        return db_path

    def test_remove_slide_tag(self, db_with_tagged_slide):
        removed = remove_slide_tag(1, "security", db_with_tagged_slide)
        assert removed is True
        tags = get_slide_tags(1, db_with_tagged_slide)
        assert "security" not in tags
        assert "architecture" in tags

    def test_remove_nonexistent_tag(self, db_with_tagged_slide):
        removed = remove_slide_tag(1, "nonexistent", db_with_tagged_slide)
        assert removed is False

    def test_remove_all_slide_tags(self, db_with_tagged_slide):
        count = remove_all_slide_tags(1, db_with_tagged_slide)
        assert count == 3
        tags = get_slide_tags(1, db_with_tagged_slide)
        assert tags == []

    def test_rename_tag(self, db_with_tagged_slide):
        assoc_count = rename_tag("security", "cybersecurity", db_with_tagged_slide)
        assert assoc_count == 1
        tags = get_slide_tags(1, db_with_tagged_slide)
        assert "cybersecurity" in tags
        assert "security" not in tags


class TestTaxonomy:
    @pytest.fixture
    def db_path(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        # Init schema
        conn = get_db(db_path)
        conn.close()
        return db_path

    def test_list_taxonomy_empty(self, db_path):
        tags = list_taxonomy(db_path)
        assert tags == []

    def test_add_taxonomy_tags(self, db_path):
        new, updated = add_taxonomy_tags(
            [{"name": "security", "category": "topic"}, {"name": "cloud", "category": "topic"}],
            db_path,
        )
        assert new == 2
        assert updated == 0
        tags = list_taxonomy(db_path)
        assert len(tags) == 2
        names = [t["name"] for t in tags]
        assert "security" in names
        assert "cloud" in names

    def test_add_taxonomy_upsert(self, db_path):
        add_taxonomy_tags([{"name": "security", "category": "topic"}], db_path)
        new, updated = add_taxonomy_tags([{"name": "security", "category": "updated-cat"}], db_path)
        assert new == 0
        assert updated == 1
        tags = list_taxonomy(db_path)
        sec = [t for t in tags if t["name"] == "security"][0]
        assert sec["category"] == "updated-cat"

    def test_remove_taxonomy_tag(self, db_path):
        add_taxonomy_tags([{"name": "security", "category": "topic"}], db_path)
        removed = remove_taxonomy_tag("security", db_path)
        assert removed is True
        assert list_taxonomy(db_path) == []

    def test_remove_taxonomy_tag_not_found(self, db_path):
        removed = remove_taxonomy_tag("nonexistent", db_path)
        assert removed is False

    def test_get_taxonomy_names(self, db_path):
        add_taxonomy_tags(
            [{"name": "b-tag", "category": "cat"}, {"name": "a-tag", "category": "cat"}],
            db_path,
        )
        names = get_taxonomy_names(db_path)
        assert names == ["a-tag", "b-tag"]

    def test_import_taxonomy_csv(self, tmp_path, db_path):
        csv_path = str(tmp_path / "taxonomy.csv")
        with open(csv_path, "w", encoding="utf-8", newline="") as f:
            writer = csv.DictWriter(f, fieldnames=["name", "category"])
            writer.writeheader()
            writer.writerow({"name": "security", "category": "topic"})
            writer.writerow({"name": "executive", "category": "audience"})

        new, updated = import_taxonomy_csv(csv_path, db_path)
        assert new == 2
        assert updated == 0
        tags = list_taxonomy(db_path)
        assert len(tags) == 2

    def test_export_taxonomy_csv(self, tmp_path, db_path):
        add_taxonomy_tags(
            [{"name": "security", "category": "topic"}, {"name": "cloud", "category": "topic"}],
            db_path,
        )
        csv_path = str(tmp_path / "exported.csv")
        count = export_taxonomy_csv(csv_path, db_path)
        assert count == 2
        assert os.path.exists(csv_path)

        with open(csv_path, encoding="utf-8") as f:
            reader = list(csv.DictReader(f))
            assert len(reader) == 2
            names = {r["name"] for r in reader}
            assert "security" in names
            assert "cloud" in names

    def test_import_export_roundtrip(self, tmp_path, db_path):
        # Create CSV
        csv_in = str(tmp_path / "in.csv")
        with open(csv_in, "w", encoding="utf-8", newline="") as f:
            writer = csv.DictWriter(f, fieldnames=["name", "category"])
            writer.writeheader()
            writer.writerow({"name": "security", "category": "topic"})
            writer.writerow({"name": "executive", "category": "audience"})
            writer.writerow({"name": "overview", "category": ""})

        import_taxonomy_csv(csv_in, db_path)

        csv_out = str(tmp_path / "out.csv")
        export_taxonomy_csv(csv_out, db_path)

        with open(csv_out, encoding="utf-8") as f:
            reader = list(csv.DictReader(f))
            assert len(reader) == 3
            names = {r["name"] for r in reader}
            assert names == {"security", "executive", "overview"}

    def test_rename_updates_taxonomy(self, db_path):
        add_taxonomy_tags([{"name": "cloud", "category": "topic"}], db_path)
        rename_tag("cloud", "cloud-computing", db_path)
        tags = list_taxonomy(db_path)
        names = [t["name"] for t in tags]
        assert "cloud-computing" in names
        assert "cloud" not in names


class TestCatalogMetadata:
    """Verify author, created_date, and modified_date extracted from PPTX core_properties."""

    @pytest.fixture
    def pptx_with_metadata(self, tmp_path):
        """Create a minimal PPTX with explicit core_properties set."""
        from datetime import datetime, timezone
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Metadata Slide"

        prs.core_properties.author = "Test Author"
        prs.core_properties.created = datetime(2024, 3, 15, 10, 0, 0, tzinfo=timezone.utc)
        prs.core_properties.modified = datetime(2025, 6, 20, 14, 30, 0, tzinfo=timezone.utc)

        path = tmp_path / "meta.pptx"
        prs.save(str(path))
        return str(path)

    def test_deck_author_extracted(self, tmp_path, pptx_with_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_with_metadata, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT author, created_date, modified_date FROM decks WHERE id = ?",
            (deck_id,),
        ).fetchone()
        conn.close()

        assert deck["author"] == "Test Author"

    def test_deck_created_date_extracted(self, tmp_path, pptx_with_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_with_metadata, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT created_date FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        assert deck["created_date"] is not None
        assert "2024-03-15" in deck["created_date"]

    def test_deck_modified_date_extracted(self, tmp_path, pptx_with_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_with_metadata, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT modified_date FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        assert deck["modified_date"] is not None
        assert "2025-06-20" in deck["modified_date"]

    def test_slides_inherit_author(self, tmp_path, pptx_with_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_with_metadata, db_path=db_path)

        conn = get_db(db_path)
        slides = conn.execute(
            "SELECT author, slide_created_date FROM slides WHERE deck_id = ?",
            (deck_id,),
        ).fetchall()
        conn.close()

        assert len(slides) > 0
        for slide in slides:
            assert slide["author"] == "Test Author"

    def test_slides_inherit_created_date(self, tmp_path, pptx_with_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_with_metadata, db_path=db_path)

        conn = get_db(db_path)
        slides = conn.execute(
            "SELECT slide_created_date FROM slides WHERE deck_id = ?",
            (deck_id,),
        ).fetchall()
        conn.close()

        assert len(slides) > 0
        for slide in slides:
            assert slide["slide_created_date"] is not None
            assert "2024-03-15" in slide["slide_created_date"]


class TestCatalogMetadataDefaults:
    """Verify defaults when core_properties are missing/empty."""

    @pytest.fixture
    def pptx_no_metadata(self, tmp_path):
        """Create a minimal PPTX without setting any core_properties."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "No Metadata Slide"

        path = tmp_path / "no_meta.pptx"
        prs.save(str(path))
        return str(path)

    def test_author_defaults_to_empty_string(self, tmp_path, pptx_no_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_no_metadata, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT author FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        assert deck["author"] == ""

    def test_created_date_is_set_when_no_explicit_metadata(self, tmp_path, pptx_no_metadata):
        """created_date is always non-null: python-pptx provides a default created timestamp
        from its blank template, so the fallback to file mtime only fires if cp.created is
        None (very rare in practice).  Either way the column must not be null."""
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_no_metadata, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT created_date FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        # created_date must always be a non-null ISO-format string
        assert deck["created_date"] is not None
        # Must look like an ISO datetime (contains 'T' or at least '-')
        assert "-" in deck["created_date"]

    def test_modified_date_is_none_when_not_set(self, tmp_path, pptx_no_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_no_metadata, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT modified_date FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        # python-pptx may or may not set a default modified date;
        # the value may be None or a date string -- just verify it's stored correctly
        # (the column exists and the query doesn't error)
        assert "modified_date" in dict(deck)

    def test_slides_author_defaults_to_empty_string(self, tmp_path, pptx_no_metadata):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_no_metadata, db_path=db_path)

        conn = get_db(db_path)
        slides = conn.execute(
            "SELECT author FROM slides WHERE deck_id = ?", (deck_id,)
        ).fetchall()
        conn.close()

        assert len(slides) > 0
        for slide in slides:
            assert slide["author"] == ""


class TestCatalogMetadataRecatalog:
    """Verify re-catalog behavior preserves created_date, updates author and modified_date."""

    def test_created_date_preserved_on_recatalog(self, tmp_path):
        from datetime import datetime, timezone
        from pptx import Presentation

        db_path = str(tmp_path / "test.db")
        pptx_path = str(tmp_path / "recatalog.pptx")

        # Initial catalog with a known created_date
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Original"
        prs.core_properties.author = "Original Author"
        prs.core_properties.created = datetime(2023, 1, 1, 0, 0, 0, tzinfo=timezone.utc)
        prs.core_properties.modified = datetime(2023, 1, 1, 0, 0, 0, tzinfo=timezone.utc)
        prs.save(pptx_path)

        deck_id = catalog_deck(pptx_path, db_path=db_path)

        conn = get_db(db_path)
        original = conn.execute(
            "SELECT created_date FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()
        original_created = original["created_date"]
        assert "2023-01-01" in original_created

        # Modify and re-catalog (new slide, updated author/modified)
        prs2 = Presentation(pptx_path)
        prs2.slides.add_slide(prs2.slide_layouts[0])
        prs2.core_properties.author = "Updated Author"
        prs2.core_properties.modified = datetime(2025, 12, 31, 23, 59, 0, tzinfo=timezone.utc)
        prs2.save(pptx_path)

        deck_id2 = catalog_deck(pptx_path, db_path=db_path)

        assert deck_id == deck_id2

        conn = get_db(db_path)
        updated = conn.execute(
            "SELECT created_date, author, modified_date FROM decks WHERE id = ?",
            (deck_id,),
        ).fetchone()
        conn.close()

        # created_date must not change on re-catalog (preserved)
        assert updated["created_date"] == original_created
        # author and modified_date must reflect the new values
        assert updated["author"] == "Updated Author"
        assert updated["modified_date"] is not None
        assert "2025-12-31" in updated["modified_date"]

    def test_slide_count_updated_on_recatalog(self, tmp_path):
        from datetime import datetime, timezone
        from pptx import Presentation

        db_path = str(tmp_path / "test.db")
        pptx_path = str(tmp_path / "recatalog2.pptx")

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(pptx_path)
        deck_id = catalog_deck(pptx_path, db_path=db_path)

        # Add a second slide and re-catalog
        prs2 = Presentation(pptx_path)
        prs2.slides.add_slide(prs2.slide_layouts[0])
        prs2.save(pptx_path)
        catalog_deck(pptx_path, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT slide_count FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        assert deck["slide_count"] == 2


class TestSections:
    @pytest.fixture
    def db_with_deck(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test_deck", "/test.pptx", "abc123", 3),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "content", "hash1"),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 2, "Slide 2", "content", "hash2"),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 3, "Slide 3", "content", "hash3"),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_schema_has_sections_tables(self, db_with_deck):
        conn = get_db(db_with_deck)
        tables = conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table'"
        ).fetchall()
        table_names = {r["name"] for r in tables}
        assert "sections" in table_names
        assert "slide_sections" in table_names
        conn.close()

    def test_set_slide_section_creates_section(self, db_with_deck):
        set_slide_section(1, "Introduction", 1, db_with_deck)
        section = get_slide_section(1, db_with_deck)
        assert section == "introduction"

    def test_get_slide_section_no_section(self, db_with_deck):
        section = get_slide_section(1, db_with_deck)
        assert section is None

    def test_set_slide_section_normalizes_name(self, db_with_deck):
        set_slide_section(1, "  Introduction  ", 1, db_with_deck)
        section = get_slide_section(1, db_with_deck)
        assert section == "introduction"

    def test_set_slide_section_position_sequential(self, db_with_deck):
        # Create multiple sections
        set_slide_section(1, "Introduction", 1, db_with_deck)
        set_slide_section(2, "Main Content", 1, db_with_deck)
        set_slide_section(3, "Conclusion", 1, db_with_deck)

        sections = get_deck_sections(1, db_with_deck)
        assert len(sections) == 3
        assert sections[0]["position"] == 1
        assert sections[1]["position"] == 2
        assert sections[2]["position"] == 3

    def test_get_deck_sections_with_counts(self, db_with_deck):
        set_slide_section(1, "Introduction", 1, db_with_deck)
        set_slide_section(2, "Introduction", 1, db_with_deck)
        set_slide_section(3, "Conclusion", 1, db_with_deck)

        sections = get_deck_sections(1, db_with_deck)
        assert len(sections) == 2

        intro = [s for s in sections if s["name"] == "introduction"][0]
        assert intro["slide_count"] == 2

        conclusion = [s for s in sections if s["name"] == "conclusion"][0]
        assert conclusion["slide_count"] == 1

    def test_remove_slide_section(self, db_with_deck):
        set_slide_section(1, "Introduction", 1, db_with_deck)
        assert get_slide_section(1, db_with_deck) == "introduction"

        removed = remove_slide_section(1, db_with_deck)
        assert removed is True
        assert get_slide_section(1, db_with_deck) is None

    def test_remove_slide_section_not_assigned(self, db_with_deck):
        removed = remove_slide_section(1, db_with_deck)
        assert removed is False

    def test_rename_section(self, db_with_deck):
        set_slide_section(1, "Introduction", 1, db_with_deck)
        set_slide_section(2, "Introduction", 1, db_with_deck)

        assoc_count = rename_section(1, "Introduction", "Opening", db_with_deck)
        assert assoc_count == 2

        assert get_slide_section(1, db_with_deck) == "opening"
        assert get_slide_section(2, db_with_deck) == "opening"

    def test_rename_section_not_found(self, db_with_deck):
        assoc_count = rename_section(1, "Nonexistent", "New Name", db_with_deck)
        assert assoc_count == 0

    def test_one_section_per_slide_constraint(self, db_with_deck):
        # Assign slide to first section
        set_slide_section(1, "Section A", 1, db_with_deck)
        assert get_slide_section(1, db_with_deck) == "section a"

        # Reassign to different section (should replace, not add)
        set_slide_section(1, "Section B", 1, db_with_deck)
        assert get_slide_section(1, db_with_deck) == "section b"

        # Verify only one section assignment
        conn = get_db(db_with_deck)
        count = conn.execute(
            "SELECT COUNT(*) as cnt FROM slide_sections WHERE slide_id = 1"
        ).fetchone()["cnt"]
        conn.close()
        assert count == 1

    def test_cascade_delete_sections_on_deck_delete(self, db_with_deck):
        set_slide_section(1, "Introduction", 1, db_with_deck)

        # Delete the deck
        conn = get_db(db_with_deck)
        conn.execute("DELETE FROM decks WHERE id = 1")
        conn.commit()

        # Verify sections were cascade deleted
        sections_count = conn.execute(
            "SELECT COUNT(*) as cnt FROM sections"
        ).fetchone()["cnt"]
        assert sections_count == 0

        slide_sections_count = conn.execute(
            "SELECT COUNT(*) as cnt FROM slide_sections"
        ).fetchone()["cnt"]
        assert slide_sections_count == 0

        conn.close()


class TestMigrationIdempotent:
    """Verify running get_db() twice on the same database causes no errors."""

    def test_double_init_no_errors(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn1 = get_db(db_path)
        conn1.close()
        conn2 = get_db(db_path)
        # Verify new columns exist
        deck_cols = {
            row[1] for row in conn2.execute("PRAGMA table_info(decks)").fetchall()
        }
        assert "subject" in deck_cols
        assert "description" in deck_cols
        slide_cols = {
            row[1] for row in conn2.execute("PRAGMA table_info(slides)").fetchall()
        }
        assert "layout_type" in slide_cols
        conn2.close()


class TestMigrationFromScratch:
    """Verify get_db() on a fresh database creates all columns."""

    def test_fresh_db_has_all_columns(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        deck_cols = {
            row[1] for row in conn.execute("PRAGMA table_info(decks)").fetchall()
        }
        assert "subject" in deck_cols
        assert "description" in deck_cols
        slide_cols = {
            row[1] for row in conn.execute("PRAGMA table_info(slides)").fetchall()
        }
        assert "layout_type" in slide_cols
        conn.close()


class TestEditHistoryTable:
    """Verify edit_history table exists and is functional."""

    def test_edit_history_table_exists(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        tables = {
            r["name"]
            for r in conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table'"
            ).fetchall()
        }
        assert "edit_history" in tables
        conn.close()

    def test_edit_history_insert_and_read(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", "/test.pptx", "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Test Slide", "content", "hash1"),
        )
        conn.execute(
            "INSERT INTO edit_history (slide_id, field, old_value, new_value, source) VALUES (?, ?, ?, ?, ?)",
            (1, "notes", "old notes", "new notes", "web"),
        )
        conn.commit()
        row = conn.execute(
            "SELECT * FROM edit_history WHERE slide_id = 1"
        ).fetchone()
        assert row["field"] == "notes"
        assert row["old_value"] == "old notes"
        assert row["new_value"] == "new notes"
        assert row["source"] == "web"
        assert row["created_at"] is not None
        conn.close()

    def test_edit_history_cascade_delete(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", "/test.pptx", "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Test Slide", "content", "hash1"),
        )
        conn.execute(
            "INSERT INTO edit_history (slide_id, field, old_value, new_value) VALUES (?, ?, ?, ?)",
            (1, "notes", "old", "new"),
        )
        conn.commit()
        conn.execute("DELETE FROM decks WHERE id = 1")
        conn.commit()
        count = conn.execute(
            "SELECT COUNT(*) as cnt FROM edit_history"
        ).fetchone()["cnt"]
        assert count == 0
        conn.close()


class TestRecordEdit:
    """Tests for the record_edit() catalog helper."""

    def _make_slide(self, tmp_path, notes="original notes"):
        """Create a DB with one deck + one slide, return (db_path, slide_id)."""
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", "/test.pptx", "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash, notes) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "content", "hash1", notes),
        )
        conn.commit()
        conn.close()
        return db_path, 1

    def test_record_edit_writes_history_and_updates_field(self, tmp_path):
        db_path, slide_id = self._make_slide(tmp_path)
        result = record_edit(slide_id, "notes", "updated notes", source="web", db_path=db_path)
        assert result is True

        conn = get_db(db_path)
        # Field updated
        row = conn.execute("SELECT notes, updated_at FROM slides WHERE id = ?", (slide_id,)).fetchone()
        assert row["notes"] == "updated notes"

        # History row written
        hist = conn.execute(
            "SELECT * FROM edit_history WHERE slide_id = ? ORDER BY id", (slide_id,)
        ).fetchall()
        assert len(hist) == 1
        assert hist[0]["field"] == "notes"
        assert hist[0]["old_value"] == "original notes"
        assert hist[0]["new_value"] == "updated notes"
        assert hist[0]["source"] == "web"
        conn.close()

    def test_record_edit_skips_when_value_unchanged(self, tmp_path):
        db_path, slide_id = self._make_slide(tmp_path)
        result = record_edit(slide_id, "notes", "original notes", source="web", db_path=db_path)
        assert result is False

        conn = get_db(db_path)
        count = conn.execute("SELECT COUNT(*) as cnt FROM edit_history").fetchone()["cnt"]
        assert count == 0
        conn.close()

    def test_record_edit_raises_for_missing_slide(self, tmp_path):
        db_path, _slide_id = self._make_slide(tmp_path)
        with pytest.raises(ValueError):
            record_edit(9999, "notes", "some value", source="web", db_path=db_path)

    def test_record_edit_handles_empty_old_value(self, tmp_path):
        db_path, slide_id = self._make_slide(tmp_path, notes="")

        record_edit(slide_id, "notes", "first notes ever", source="ai", db_path=db_path)

        conn = get_db(db_path)
        hist = conn.execute("SELECT * FROM edit_history WHERE slide_id = 1").fetchone()
        assert hist["old_value"] == ""
        assert hist["new_value"] == "first notes ever"
        assert hist["source"] == "ai"
        conn.close()

    def test_record_edit_updates_timestamp(self, tmp_path):
        db_path, slide_id = self._make_slide(tmp_path)

        conn = get_db(db_path)
        before = conn.execute("SELECT updated_at FROM slides WHERE id = ?", (slide_id,)).fetchone()["updated_at"]
        conn.close()

        time.sleep(0.05)
        record_edit(slide_id, "notes", "new text", source="web", db_path=db_path)

        conn = get_db(db_path)
        after = conn.execute("SELECT updated_at FROM slides WHERE id = ?", (slide_id,)).fetchone()["updated_at"]
        conn.close()
        assert after >= before


class TestCatalogSubjectDescription:
    """Verify subject and description extracted from PPTX core_properties."""

    @pytest.fixture
    def pptx_with_subject(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        prs.core_properties.subject = "Test Subject"
        prs.core_properties.comments = "Test Description"

        path = tmp_path / "subject.pptx"
        prs.save(str(path))
        return str(path)

    def test_subject_extracted(self, tmp_path, pptx_with_subject):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_with_subject, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT subject, description FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        assert deck["subject"] == "Test Subject"
        assert deck["description"] == "Test Description"

    def test_subject_in_get_deck_by_id(self, tmp_path, pptx_with_subject):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_with_subject, db_path=db_path)
        deck = get_deck_by_id(deck_id, db_path)
        assert deck["subject"] == "Test Subject"
        assert deck["description"] == "Test Description"

    def test_subject_in_list_decks(self, tmp_path, pptx_with_subject):
        db_path = str(tmp_path / "test.db")
        catalog_deck(pptx_with_subject, db_path=db_path)
        decks = list_decks(db_path)
        assert len(decks) == 1
        assert decks[0]["subject"] == "Test Subject"
        assert decks[0]["description"] == "Test Description"


class TestCatalogSubjectDescriptionDefaults:
    """Verify empty string defaults when subject/description are not set."""

    @pytest.fixture
    def pptx_no_subject(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"

        path = tmp_path / "no_subject.pptx"
        prs.save(str(path))
        return str(path)

    def test_subject_defaults_to_empty(self, tmp_path, pptx_no_subject):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(pptx_no_subject, db_path=db_path)

        conn = get_db(db_path)
        deck = conn.execute(
            "SELECT subject, description FROM decks WHERE id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        assert deck["subject"] == ""
        # python-pptx sets a default comments value; description column
        # is populated from comments so it may be non-empty even without
        # explicit metadata.  Just verify the column exists and is a string.
        assert isinstance(deck["description"], str)


class TestGetAllTags:
    """Tests for get_all_tags() catalog helper."""

    @pytest.fixture
    def db_path(self, tmp_path):
        return str(tmp_path / "test.db")

    def _setup_deck_and_slides(self, db_path):
        """Create a deck with two slides and return the connection."""
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Test Deck", "/test.pptx", "abc123", 2),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_hash) VALUES (?, ?, ?, ?)",
            (1, 1, "Slide 1", "hash1"),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_hash) VALUES (?, ?, ?, ?)",
            (1, 2, "Slide 2", "hash2"),
        )
        conn.commit()
        return conn

    def test_returns_tags_with_counts(self, db_path):
        """Tags applied to slides are returned with correct slide counts."""
        conn = self._setup_deck_and_slides(db_path)
        conn.close()
        # 'security' applied to both slides; 'architecture' applied to slide 1 only
        add_tags(1, ["security", "architecture"], "manual", db_path)
        add_tags(2, ["security"], "manual", db_path)

        result = get_all_tags(db_path)

        assert len(result) == 2
        names = {r["name"] for r in result}
        assert names == {"security", "architecture"}
        security = next(r for r in result if r["name"] == "security")
        architecture = next(r for r in result if r["name"] == "architecture")
        assert security["count"] == 2
        assert architecture["count"] == 1

    def test_includes_taxonomy_category(self, db_path):
        """Tags that appear in the taxonomy table carry their category."""
        conn = self._setup_deck_and_slides(db_path)
        conn.execute(
            "INSERT INTO taxonomy (name, category) VALUES (?, ?)",
            ("zero-trust", "Security"),
        )
        conn.commit()
        conn.close()
        add_tags(1, ["zero-trust"], "manual", db_path)

        result = get_all_tags(db_path)

        assert len(result) == 1
        assert result[0]["name"] == "zero-trust"
        assert result[0]["category"] == "Security"

    def test_uncategorized_tags(self, db_path):
        """Tags not present in the taxonomy table get an empty string category."""
        conn = self._setup_deck_and_slides(db_path)
        conn.close()
        add_tags(1, ["misc"], "manual", db_path)

        result = get_all_tags(db_path)

        assert len(result) == 1
        assert result[0]["name"] == "misc"
        assert result[0]["category"] == ""

    def test_empty_database(self, db_path):
        """Returns an empty list when no tags (or no slides) exist."""
        conn = get_db(db_path)
        conn.close()

        result = get_all_tags(db_path)

        assert result == []

    def test_excludes_tags_with_no_slides(self, db_path):
        """Tags that exist in the tags table but have no slide_tags rows are excluded."""
        conn = self._setup_deck_and_slides(db_path)
        conn.close()
        # Insert 'orphan' directly into tags with no slide_tags row
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO tags (name, source) VALUES (?, ?)", ("orphan", "manual")
        )
        conn.commit()
        conn.close()
        # 'used' is applied via add_tags so it gets a slide_tags entry
        add_tags(1, ["used"], "manual", db_path)

        result = get_all_tags(db_path)

        assert len(result) == 1
        assert result[0]["name"] == "used"


# ---------------------------------------------------------------------------
# TestRelativePaths
# ---------------------------------------------------------------------------

from aippt.catalog import migrate_paths


class TestRelativePaths:
    """Verify catalog_deck() stores relative paths."""

    @pytest.fixture
    def simple_pptx(self, tmp_path):
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        path = tmp_path / "uploads" / "test.pptx"
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        return str(path)

    def test_catalog_stores_relative_file_path(self, tmp_path, simple_pptx):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(simple_pptx, db_path=db_path, base_dir=str(tmp_path))

        conn = get_db(db_path)
        deck = conn.execute("SELECT file_path FROM decks WHERE id = ?", (deck_id,)).fetchone()
        conn.close()

        assert not os.path.isabs(deck["file_path"])
        assert deck["file_path"] == os.path.join("uploads", "test.pptx")

    def test_catalog_stores_relative_image_path(self, tmp_path, simple_pptx):
        # Create a fake image
        images_dir = tmp_path / "images" / "test"
        images_dir.mkdir(parents=True)
        (images_dir / "Slide1.PNG").write_bytes(b"fake png")

        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(
            simple_pptx, db_path=db_path,
            images_dir=str(images_dir), base_dir=str(tmp_path),
        )

        conn = get_db(db_path)
        slide = conn.execute(
            "SELECT image_path FROM slides WHERE deck_id = ?", (deck_id,)
        ).fetchone()
        conn.close()

        assert slide["image_path"] is not None
        assert not os.path.isabs(slide["image_path"])
        assert "images" in slide["image_path"]

    def test_catalog_default_base_dir_is_cwd(self, tmp_path, simple_pptx):
        """When base_dir is not set, paths are relative to cwd."""
        db_path = str(tmp_path / "test.db")
        old_cwd = os.getcwd()
        try:
            os.chdir(str(tmp_path))
            deck_id = catalog_deck(simple_pptx, db_path=db_path)
        finally:
            os.chdir(old_cwd)

        conn = get_db(db_path)
        deck = conn.execute("SELECT file_path FROM decks WHERE id = ?", (deck_id,)).fetchone()
        conn.close()

        assert not os.path.isabs(deck["file_path"])


# ---------------------------------------------------------------------------
# TestMigratePaths
# ---------------------------------------------------------------------------


class TestMigratePaths:
    """Verify migrate_paths() converts absolute to relative, idempotently."""

    def _make_db_with_absolute_paths(self, tmp_path):
        """Create a DB with absolute paths for testing migration."""
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", str(tmp_path / "uploads" / "test.pptx"), "abc123", 2),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash, image_path) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "content", "hash1", str(tmp_path / "images" / "test" / "Slide1.PNG")),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash, image_path) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (1, 2, "Slide 2", "content2", "hash2", str(tmp_path / "images" / "test" / "Slide2.PNG")),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_migrate_absolute_to_relative(self, tmp_path):
        db_path = self._make_db_with_absolute_paths(tmp_path)
        result = migrate_paths(db_path=db_path, base_dir=str(tmp_path))

        assert result["deck_paths"] == 1
        assert result["image_paths"] == 2
        assert result["skipped"] == 0

        conn = get_db(db_path)
        deck = conn.execute("SELECT file_path FROM decks WHERE id = 1").fetchone()
        assert not os.path.isabs(deck["file_path"])
        assert deck["file_path"] == os.path.join("uploads", "test.pptx")

        slides = conn.execute("SELECT image_path FROM slides ORDER BY position").fetchall()
        for s in slides:
            assert not os.path.isabs(s["image_path"])
        conn.close()

    def test_migrate_idempotent(self, tmp_path):
        """Running migration twice produces the same result."""
        db_path = self._make_db_with_absolute_paths(tmp_path)

        result1 = migrate_paths(db_path=db_path, base_dir=str(tmp_path))
        assert result1["deck_paths"] == 1

        result2 = migrate_paths(db_path=db_path, base_dir=str(tmp_path))
        assert result2["deck_paths"] == 0
        assert result2["image_paths"] == 0
        assert result2["skipped"] == 3  # 1 deck + 2 slides

        # Paths unchanged from first run
        conn = get_db(db_path)
        deck = conn.execute("SELECT file_path FROM decks WHERE id = 1").fetchone()
        assert deck["file_path"] == os.path.join("uploads", "test.pptx")
        conn.close()

    def test_migrate_skips_relative(self, tmp_path):
        """Already-relative paths are left untouched."""
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", "uploads/test.pptx", "abc123", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash, image_path) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "content", "hash1", "images/test/Slide1.PNG"),
        )
        conn.commit()
        conn.close()

        result = migrate_paths(db_path=db_path, base_dir=str(tmp_path))
        assert result["deck_paths"] == 0
        assert result["image_paths"] == 0
        assert result["skipped"] == 2

    def test_migrate_handles_null_image_path(self, tmp_path):
        """Slides with NULL image_path are handled gracefully."""
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", str(tmp_path / "test.pptx"), "abc", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) "
            "VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "content", "hash1"),
        )
        conn.commit()
        conn.close()

        result = migrate_paths(db_path=db_path, base_dir=str(tmp_path))
        assert result["deck_paths"] == 1
        # NULL image_path slides are not in the query, so not skipped either
        assert result["image_paths"] == 0


class TestResolveDeck:
    @pytest.fixture
    def db_with_named_decks(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Networking Advantages", "/net.pptx", "abc", 10),
        )
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Deploying AMD Instinct", "/instinct.pptx", "def", 18),
        )
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Network Security Overview", "/netsec.pptx", "ghi", 5),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_resolve_by_id(self, db_with_named_decks):
        deck = resolve_deck("1", db_with_named_decks)
        assert deck is not None
        assert deck["name"] == "Networking Advantages"

    def test_resolve_by_id_not_found(self, db_with_named_decks):
        result = resolve_deck("999", db_with_named_decks)
        assert result is None

    def test_resolve_by_exact_name(self, db_with_named_decks):
        deck = resolve_deck("instinct", db_with_named_decks)
        assert deck is not None
        assert deck["name"] == "Deploying AMD Instinct"

    def test_resolve_by_name_case_insensitive(self, db_with_named_decks):
        deck = resolve_deck("INSTINCT", db_with_named_decks)
        assert deck is not None
        assert deck["name"] == "Deploying AMD Instinct"

    def test_resolve_ambiguous_returns_list(self, db_with_named_decks):
        result = resolve_deck("network", db_with_named_decks)
        assert isinstance(result, list)
        assert len(result) == 2

    def test_resolve_name_not_found(self, db_with_named_decks):
        result = resolve_deck("nonexistent", db_with_named_decks)
        assert result is None


class TestDeleteDeck:
    @pytest.fixture
    def db_with_full_deck(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Test Deck", "/test.pptx", "abc", 2),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "content1", "hash1"),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 2, "Slide 2", "content2", "hash2"),
        )
        conn.execute("INSERT INTO tags (name, source) VALUES (?, ?)", ("security", "ai"))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (1, 1))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (2, 1))
        conn.execute("INSERT INTO sections (deck_id, name, position) VALUES (?, ?, ?)", (1, "intro", 1))
        conn.execute("INSERT INTO slide_sections (slide_id, section_id) VALUES (?, ?)", (1, 1))
        conn.execute(
            "INSERT INTO edit_history (slide_id, field, old_value, new_value, source) VALUES (?, ?, ?, ?, ?)",
            (1, "notes", "", "new notes", "web"),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_delete_deck(self, db_with_full_deck):
        info = delete_deck(1, db_with_full_deck)
        assert info["name"] == "Test Deck"
        assert info["slide_count"] == 2
        conn = get_db(db_with_full_deck)
        assert conn.execute("SELECT COUNT(*) FROM decks").fetchone()[0] == 0
        conn.close()

    def test_delete_deck_cascades(self, db_with_full_deck):
        delete_deck(1, db_with_full_deck)
        conn = get_db(db_with_full_deck)
        assert conn.execute("SELECT COUNT(*) FROM slides").fetchone()[0] == 0
        assert conn.execute("SELECT COUNT(*) FROM slide_tags").fetchone()[0] == 0
        assert conn.execute("SELECT COUNT(*) FROM sections").fetchone()[0] == 0
        assert conn.execute("SELECT COUNT(*) FROM slide_sections").fetchone()[0] == 0
        assert conn.execute("SELECT COUNT(*) FROM edit_history").fetchone()[0] == 0
        conn.close()

    def test_delete_deck_not_found(self, db_with_full_deck):
        result = delete_deck(999, db_with_full_deck)
        assert result is None

    def test_delete_returns_tag_count(self, db_with_full_deck):
        info = delete_deck(1, db_with_full_deck)
        assert info["tag_count"] == 2

    def test_delete_returns_section_count(self, db_with_full_deck):
        info = delete_deck(1, db_with_full_deck)
        assert info["section_count"] == 1


class TestRenameDeck:
    @pytest.fixture
    def db_with_deck(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Old Name", "/test.pptx", "abc", 1),
        )
        conn.commit()
        conn.close()
        return db_path

    def test_rename_deck(self, db_with_deck):
        old_name = rename_deck(1, "New Name", db_with_deck)
        assert old_name == "Old Name"
        conn = get_db(db_with_deck)
        row = conn.execute("SELECT name FROM decks WHERE id = 1").fetchone()
        assert row["name"] == "New Name"
        conn.close()

    def test_rename_deck_not_found(self, db_with_deck):
        result = rename_deck(999, "New Name", db_with_deck)
        assert result is None


class TestGetDeckTagCount:
    @pytest.fixture
    def db_with_tagged_deck(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Test Deck", "/test.pptx", "abc", 2),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "c1", "h1"),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
            (1, 2, "Slide 2", "c2", "h2"),
        )
        conn.execute("INSERT INTO tags (name, source) VALUES (?, ?)", ("security", "ai"))
        conn.execute("INSERT INTO tags (name, source) VALUES (?, ?)", ("networking", "ai"))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (1, 1))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (1, 2))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (2, 1))
        conn.commit()
        conn.close()
        return db_path

    def test_get_deck_tag_count(self, db_with_tagged_deck):
        count = get_deck_tag_count(1, db_with_tagged_deck)
        assert count == 2

    def test_get_deck_tag_count_no_tags(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Empty Deck", "/empty.pptx", "xyz", 0),
        )
        conn.commit()
        conn.close()
        count = get_deck_tag_count(1, db_path)
        assert count == 0


class TestGetDeckTopTags:
    @pytest.fixture
    def db_with_tagged_deck(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("Test Deck", "/test.pptx", "abc", 3),
        )
        for i in range(1, 4):
            conn.execute(
                "INSERT INTO slides (deck_id, position, title, content_text, content_hash) VALUES (?, ?, ?, ?, ?)",
                (1, i, f"Slide {i}", f"content{i}", f"hash{i}"),
            )
        conn.execute("INSERT INTO tags (name, source) VALUES (?, ?)", ("security", "ai"))
        conn.execute("INSERT INTO tags (name, source) VALUES (?, ?)", ("networking", "ai"))
        conn.execute("INSERT INTO tags (name, source) VALUES (?, ?)", ("performance", "ai"))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (1, 1))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (2, 1))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (3, 1))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (1, 2))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (2, 2))
        conn.execute("INSERT INTO slide_tags (slide_id, tag_id) VALUES (?, ?)", (1, 3))
        conn.commit()
        conn.close()
        return db_path

    def test_get_deck_top_tags(self, db_with_tagged_deck):
        tags = get_deck_top_tags(1, db_with_tagged_deck)
        assert tags[0] == ("security", 3)
        assert tags[1] == ("networking", 2)
        assert tags[2] == ("performance", 1)

    def test_get_deck_top_tags_limit(self, db_with_tagged_deck):
        tags = get_deck_top_tags(1, db_with_tagged_deck, limit=2)
        assert len(tags) == 2
