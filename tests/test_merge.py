"""Tests for lib/merge.py — PPTX merge utility."""

import os
import sys
import tempfile
import pytest

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# Add lib/ to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'lib'))

from merge import merge_decks, _renumber_slides, _copy_slide, _get_blank_layout


# ---------------------------------------------------------------------------
# Fixtures: create test PPTX files
# ---------------------------------------------------------------------------

def _make_test_pptx(path, num_slides, slide_width=None, slide_height=None,
                    add_numbers=False, start_num=1, add_images=False):
    """Create a test PPTX with N slides, optional slide numbers and images.

    Each slide has a title text box and optionally a slide number text box
    in the bottom-left corner (matching the pptxgenjs-helpers footer pattern).
    """
    prs = Presentation()
    if slide_width:
        prs.slide_width = Inches(slide_width)
    if slide_height:
        prs.slide_height = Inches(slide_height)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # typically blank

    for i in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        slide_num = start_num + i

        # Add title text box
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(10), Inches(0.7)
        )
        tf = txBox.text_frame
        tf.text = f"Slide {slide_num} Title"
        tf.paragraphs[0].font.size = Pt(28)

        # Add content text box
        txBox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(10), Inches(4)
        )
        tf2 = txBox2.text_frame
        tf2.text = f"Content for slide {slide_num}"

        # Add slide number (bottom-left, small text box)
        if add_numbers:
            numBox = slide.shapes.add_textbox(
                Inches(0.2),                 # left edge
                Inches(7.05),                # near bottom (in 7.5" slide)
                Inches(0.5),                 # narrow
                Inches(0.3),                 # short
            )
            numTf = numBox.text_frame
            run = numTf.paragraphs[0].add_run()
            run.text = str(slide_num)
            run.font.size = Pt(10)

        # Add a simple shape (rectangle) to test shape copying
        slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(1), Inches(5), Inches(2), Inches(1)
        )

    prs.save(path)
    return path


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test files."""
    with tempfile.TemporaryDirectory() as d:
        yield d


@pytest.fixture
def chunk_files(temp_dir):
    """Create 3 test PPTX chunk files."""
    paths = []
    for i, (num_slides, start) in enumerate([(3, 1), (4, 4), (3, 8)], 1):
        path = os.path.join(temp_dir, f'chunk-{i}.pptx')
        _make_test_pptx(path, num_slides, slide_width=13.33, slide_height=7.5,
                        add_numbers=True, start_num=start)
        paths.append(path)
    return paths


@pytest.fixture
def chunk_files_no_numbers(temp_dir):
    """Create 2 test PPTX chunk files without slide numbers."""
    paths = []
    for i, num_slides in enumerate([5, 5], 1):
        path = os.path.join(temp_dir, f'chunk-{i}.pptx')
        _make_test_pptx(path, num_slides, slide_width=13.33, slide_height=7.5)
        paths.append(path)
    return paths


# ---------------------------------------------------------------------------
# TestMergeSlideCount
# ---------------------------------------------------------------------------

class TestMergeSlideCount:
    """Test that merge preserves total slide count."""

    def test_merge_two_chunks(self, temp_dir):
        """Merging 2 chunks preserves total slide count."""
        chunk1 = _make_test_pptx(
            os.path.join(temp_dir, 'a.pptx'), 3,
            slide_width=13.33, slide_height=7.5
        )
        chunk2 = _make_test_pptx(
            os.path.join(temp_dir, 'b.pptx'), 4,
            slide_width=13.33, slide_height=7.5
        )
        output = os.path.join(temp_dir, 'merged.pptx')
        result = merge_decks([chunk1, chunk2], output)

        assert result['slide_count'] == 7
        assert result['chunk_counts'] == [3, 4]
        assert os.path.isfile(output)

    def test_merge_three_chunks(self, chunk_files, temp_dir):
        """Merging 3 chunks preserves total slide count."""
        output = os.path.join(temp_dir, 'merged.pptx')
        result = merge_decks(chunk_files, output)

        assert result['slide_count'] == 10  # 3 + 4 + 3
        assert result['chunk_counts'] == [3, 4, 3]

    def test_merge_single_chunk(self, temp_dir):
        """Merging a single chunk copies it as-is."""
        chunk = _make_test_pptx(
            os.path.join(temp_dir, 'single.pptx'), 5,
            slide_width=13.33, slide_height=7.5
        )
        output = os.path.join(temp_dir, 'merged.pptx')
        result = merge_decks([chunk], output)

        assert result['slide_count'] == 5
        assert result['chunk_counts'] == [5]

    def test_merged_file_valid_pptx(self, chunk_files, temp_dir):
        """Merged file opens as a valid PPTX."""
        output = os.path.join(temp_dir, 'merged.pptx')
        merge_decks(chunk_files, output)

        # Should open without error
        prs = Presentation(output)
        assert len(prs.slides) == 10


# ---------------------------------------------------------------------------
# TestMergeSlideOrder
# ---------------------------------------------------------------------------

class TestMergeSlideOrder:
    """Test that slides maintain their order after merge."""

    def test_slide_titles_ordered(self, chunk_files, temp_dir):
        """Slide titles appear in chunk order."""
        output = os.path.join(temp_dir, 'merged.pptx')
        merge_decks(chunk_files, output)

        prs = Presentation(output)
        titles = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if 'Title' in text:
                        titles.append(text)
                        break

        # Titles should be in order: Slide 1 Title, Slide 2 Title, ..., Slide 10 Title
        expected = [f"Slide {i} Title" for i in range(1, 11)]
        assert titles == expected

    def test_shapes_preserved(self, chunk_files, temp_dir):
        """All shapes are preserved in merged slides."""
        output = os.path.join(temp_dir, 'merged.pptx')
        merge_decks(chunk_files, output)

        prs = Presentation(output)
        for slide in prs.slides:
            # Each test slide has: title textbox, content textbox, number textbox, rectangle
            # The blank layout may add placeholders, which are removed by _copy_slide
            shape_count = len(slide.shapes)
            assert shape_count >= 3  # at minimum: title, content, rectangle


# ---------------------------------------------------------------------------
# TestRenumbering
# ---------------------------------------------------------------------------

class TestRenumbering:
    """Test slide number renumbering."""

    def test_renumber_after_merge(self, chunk_files, temp_dir):
        """Slide numbers are renumbered sequentially after merge."""
        output = os.path.join(temp_dir, 'merged.pptx')
        merge_decks(chunk_files, output, renumber=True)

        prs = Presentation(output)
        slide_numbers = _extract_slide_numbers(prs)

        # Should be renumbered 1-10
        expected = list(range(1, 11))
        assert slide_numbers == expected

    def test_no_renumber_preserves_original(self, chunk_files, temp_dir):
        """With renumber=False, original numbers are preserved."""
        output = os.path.join(temp_dir, 'merged.pptx')
        merge_decks(chunk_files, output, renumber=False)

        prs = Presentation(output)
        slide_numbers = _extract_slide_numbers(prs)

        # Original numbers: chunk1=[1,2,3], chunk2=[4,5,6,7], chunk3=[8,9,10]
        expected = list(range(1, 11))
        assert slide_numbers == expected

    def test_renumber_no_numbers_no_crash(self, chunk_files_no_numbers, temp_dir):
        """Renumbering works even when slides have no number boxes."""
        output = os.path.join(temp_dir, 'merged.pptx')
        # Should not raise
        result = merge_decks(chunk_files_no_numbers, output, renumber=True)
        assert result['slide_count'] == 10


def _extract_slide_numbers(prs):
    """Extract slide numbers from number text boxes in each slide."""
    slide_height = prs.slide_height or Inches(7.5)
    bottom_threshold = slide_height * 0.80
    left_threshold = (prs.slide_width or Inches(13.33)) * 0.20

    numbers = []
    for slide in prs.slides:
        found = None
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.is_placeholder:
                continue
            if shape.top is None or shape.left is None:
                continue
            if shape.top < bottom_threshold or shape.left > left_threshold:
                continue
            text = shape.text_frame.text.strip()
            if text and text.isdigit():
                found = int(text)
                break
        if found is not None:
            numbers.append(found)
    return numbers


# ---------------------------------------------------------------------------
# TestErrorHandling
# ---------------------------------------------------------------------------

class TestErrorHandling:
    """Test error handling in merge operations."""

    def test_empty_chunk_list(self):
        """Merging empty list raises ValueError."""
        with pytest.raises(ValueError, match="No chunk paths"):
            merge_decks([], 'output.pptx')

    def test_missing_file(self, temp_dir):
        """Missing chunk file raises FileNotFoundError."""
        with pytest.raises(FileNotFoundError, match="Chunk file not found"):
            merge_decks([os.path.join(temp_dir, 'nonexistent.pptx')],
                        os.path.join(temp_dir, 'output.pptx'))

    def test_output_dir_created(self, temp_dir):
        """Output directory is created if it doesn't exist."""
        chunk = _make_test_pptx(
            os.path.join(temp_dir, 'a.pptx'), 2,
            slide_width=13.33, slide_height=7.5
        )
        output = os.path.join(temp_dir, 'subdir', 'nested', 'merged.pptx')
        result = merge_decks([chunk], output)
        assert os.path.isfile(output)
