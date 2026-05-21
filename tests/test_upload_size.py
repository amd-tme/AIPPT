"""Tests for the upload size limit middleware, handler backstop, and config loader."""
from __future__ import annotations

import os
from pathlib import Path
from unittest.mock import patch

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt.catalog import catalog_deck, file_hash, get_deck_by_file_hash
from aippt.config import DEFAULT_MAX_UPLOAD_MB, load_upload_config
from aippt.web.app import create_app


@pytest.fixture
def deck_path(tmp_path):
    """Create a minimal PPTX with one slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    p = str(tmp_path / "test.pptx")
    prs.save(p)
    return p


# ---------------------------------------------------------------------------
# Config loader
# ---------------------------------------------------------------------------


class TestLoadUploadConfig:
    """load_upload_config — gateway.yaml `upload.max_size_mb` resolution."""

    def test_default_when_path_none(self):
        assert load_upload_config(None) == DEFAULT_MAX_UPLOAD_MB * 1024 * 1024

    def test_default_when_file_missing(self, tmp_path):
        missing = tmp_path / "nonexistent.yaml"
        assert load_upload_config(str(missing)) == DEFAULT_MAX_UPLOAD_MB * 1024 * 1024

    def test_default_when_block_absent(self, tmp_path):
        cfg = tmp_path / "gateway.yaml"
        cfg.write_text("gateway:\n  base_url: 'http://x'\n", encoding="utf-8")
        assert load_upload_config(str(cfg)) == DEFAULT_MAX_UPLOAD_MB * 1024 * 1024

    def test_reads_max_size_mb(self, tmp_path):
        cfg = tmp_path / "gateway.yaml"
        cfg.write_text("upload:\n  max_size_mb: 25\n", encoding="utf-8")
        assert load_upload_config(str(cfg)) == 25 * 1024 * 1024

    def test_string_value_falls_back_to_default(self, tmp_path):
        cfg = tmp_path / "gateway.yaml"
        cfg.write_text("upload:\n  max_size_mb: 'fifty'\n", encoding="utf-8")
        # Coerces non-numeric to default rather than crashing the server.
        assert load_upload_config(str(cfg)) == DEFAULT_MAX_UPLOAD_MB * 1024 * 1024

    def test_zero_or_negative_falls_back(self, tmp_path):
        cfg = tmp_path / "gateway.yaml"
        cfg.write_text("upload:\n  max_size_mb: 0\n", encoding="utf-8")
        assert load_upload_config(str(cfg)) == DEFAULT_MAX_UPLOAD_MB * 1024 * 1024


# ---------------------------------------------------------------------------
# Middleware (Content-Length pre-check)
# ---------------------------------------------------------------------------


@pytest.fixture
def small_limit_client(tmp_path, deck_path):
    """TestClient with a 1 KB upload cap so we can easily exceed it."""
    db_path = str(tmp_path / "test.db")
    uploads_dir = str(tmp_path / "uploads")
    catalog_deck(deck_path, db_path=db_path)
    app = create_app(
        db_path=db_path, uploads_dir=uploads_dir, view_only=False,
        # 1 MB rounded to bytes via the helper path — but max_upload_mb=1 is
        # the smallest integer accepted; tests upload <1 MB anyway. For the
        # over-limit case we patch app.state.max_upload_bytes to 1024.
        max_upload_mb=1,
    )
    app.state.max_upload_bytes = 1024
    return TestClient(app)


class TestUploadSizeMiddleware:
    """Middleware rejects requests whose Content-Length exceeds the cap."""

    def test_under_limit_passes_through(self, small_limit_client):
        # A GET that's not an upload path should always pass through.
        resp = small_limit_client.get("/api/config")
        assert resp.status_code == 200

    def test_oversized_post_returns_413(self, small_limit_client):
        # Send 4 KB of bytes against a 1 KB cap. Content-Length header is
        # auto-set by requests. Hits middleware before the route handler runs.
        big = b"x" * 4096
        resp = small_limit_client.post(
            "/api/decks/upload-stream",
            files={"file": ("big.pptx", big, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert resp.status_code == 413
        body = resp.json()
        assert "max_bytes" in body
        assert body["max_bytes"] == 1024
        # Either observed_bytes (declared in CL) shows up or the message is
        # the generic one; both shapes are acceptable.
        assert "exceeds" in body["error"]

    def test_non_upload_path_not_subject_to_limit(self, small_limit_client):
        # /api/decks/{id}/slides is not an upload route — even a giant query
        # string would pass the middleware.
        resp = small_limit_client.get("/api/decks/1/slides")
        # Whether the deck exists or not, we must not see a 413.
        assert resp.status_code != 413


# ---------------------------------------------------------------------------
# /api/config carries max_upload_bytes
# ---------------------------------------------------------------------------


class TestConfigEndpointUpload:
    def test_config_exposes_max_upload_bytes(self, small_limit_client):
        resp = small_limit_client.get("/api/config")
        assert resp.status_code == 200
        body = resp.json()
        assert body["max_upload_bytes"] == 1024


# ---------------------------------------------------------------------------
# /api/decks/by-hash/{sha256}
# ---------------------------------------------------------------------------


@pytest.fixture
def hash_client(tmp_path, deck_path):
    """TestClient with a pre-cataloged deck so by-hash lookups can hit."""
    db_path = str(tmp_path / "test.db")
    uploads_dir = str(tmp_path / "uploads")
    catalog_deck(deck_path, db_path=db_path)
    app = create_app(db_path=db_path, uploads_dir=uploads_dir, view_only=False)
    client = TestClient(app)
    client.deck_path = deck_path  # type: ignore[attr-defined]
    return client


class TestDeckByHashRoute:
    def test_returns_200_for_known_hash(self, hash_client):
        sha = file_hash(hash_client.deck_path)  # type: ignore[attr-defined]
        resp = hash_client.get(f"/api/decks/by-hash/{sha}")
        assert resp.status_code == 200
        body = resp.json()
        assert body["file_hash"] == sha
        assert body["slide_count"] == 1
        assert "display_name" in body
        assert "id" in body

    def test_returns_404_for_unknown_hash(self, hash_client):
        unknown = "0" * 64
        resp = hash_client.get(f"/api/decks/by-hash/{unknown}")
        assert resp.status_code == 404
        assert "error" in resp.json()

    def test_returns_400_for_malformed_hash(self, hash_client):
        # Too short, wrong charset, etc.
        for bad in ("notahex", "ABC", "g" * 64, "0" * 63):
            resp = hash_client.get(f"/api/decks/by-hash/{bad}")
            assert resp.status_code == 400, f"expected 400 for {bad}, got {resp.status_code}"


# ---------------------------------------------------------------------------
# catalog.get_deck_by_file_hash unit
# ---------------------------------------------------------------------------


class TestGetDeckByFileHash:
    def test_returns_dict_when_present(self, tmp_path, deck_path):
        db_path = str(tmp_path / "test.db")
        catalog_deck(deck_path, db_path=db_path)
        sha = file_hash(deck_path)
        deck = get_deck_by_file_hash(sha, db_path=db_path)
        assert deck is not None
        assert deck["file_hash"] == sha

    def test_returns_none_when_absent(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        # Catalog initializes the schema; we don't need a deck.
        from aippt.catalog import get_db
        get_db(db_path).close()
        assert get_deck_by_file_hash("0" * 64, db_path=db_path) is None
