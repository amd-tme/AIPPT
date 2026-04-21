"""Tests for PowerPoint section support."""
import os
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from aippt.sections import read_sections, write_sections, Section


@pytest.fixture
def temp_pptx():
    """Create a temporary PPTX file for testing."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        # Create a simple presentation with 3 slides
        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            if slide.shapes.title:
                slide.shapes.title.text = f"Slide {i+1}"
        prs.save(f.name)
        yield f.name
    os.unlink(f.name)


def test_read_sections_no_sections(temp_pptx):
    """Test reading from PPTX without sections returns empty list."""
    prs = Presentation(temp_pptx)
    sections = read_sections(prs)
    assert sections == []


def test_write_and_read_sections(temp_pptx):
    """Test writing sections and reading them back."""
    prs = Presentation(temp_pptx)

    # Get slide IDs from the presentation
    slide_ids = [slide.slide_id for slide in prs.slides]

    # Create sections
    sections_to_write = [
        Section(name="Introduction", slide_ids=[slide_ids[0]]),
        Section(name="Main Content", slide_ids=[slide_ids[1], slide_ids[2]]),
    ]

    # Write sections
    write_sections(prs, sections_to_write)
    prs.save(temp_pptx)

    # Read back and verify
    prs2 = Presentation(temp_pptx)
    sections_read = read_sections(prs2)

    assert len(sections_read) == 2
    assert sections_read[0].name == "Introduction"
    assert sections_read[0].slide_ids == [slide_ids[0]]
    assert sections_read[1].name == "Main Content"
    assert sections_read[1].slide_ids == [slide_ids[1], slide_ids[2]]


def test_write_sections_validates_slide_ids(temp_pptx):
    """Test that write_sections rejects invalid slide IDs."""
    prs = Presentation(temp_pptx)

    # Create section with invalid slide ID
    invalid_sections = [
        Section(name="Bad Section", slide_ids=[99999]),
    ]

    with pytest.raises(ValueError, match="invalid slide ID"):
        write_sections(prs, invalid_sections)


def test_write_empty_sections(temp_pptx):
    """Test writing empty section list does nothing."""
    prs = Presentation(temp_pptx)
    write_sections(prs, [])
    prs.save(temp_pptx)

    # Verify no sections were created
    prs2 = Presentation(temp_pptx)
    sections = read_sections(prs2)
    assert sections == []


def test_round_trip_consistency(temp_pptx):
    """Test read → write → read produces consistent results."""
    prs = Presentation(temp_pptx)
    slide_ids = [slide.slide_id for slide in prs.slides]

    original_sections = [
        Section(name="Section A", slide_ids=[slide_ids[0]]),
        Section(name="Section B", slide_ids=[slide_ids[1]]),
        Section(name="Section C", slide_ids=[slide_ids[2]]),
    ]

    # First write
    write_sections(prs, original_sections)
    prs.save(temp_pptx)

    # Read back
    prs2 = Presentation(temp_pptx)
    read_sections_1 = read_sections(prs2)

    # Write again with the read sections
    write_sections(prs2, read_sections_1)
    prs2.save(temp_pptx)

    # Read again
    prs3 = Presentation(temp_pptx)
    read_sections_2 = read_sections(prs3)

    # Verify consistency
    assert len(read_sections_1) == len(original_sections)
    assert len(read_sections_2) == len(original_sections)

    for i in range(len(original_sections)):
        assert read_sections_1[i].name == original_sections[i].name
        assert read_sections_1[i].slide_ids == original_sections[i].slide_ids
        assert read_sections_2[i].name == original_sections[i].name
        assert read_sections_2[i].slide_ids == original_sections[i].slide_ids


def test_section_with_empty_name(temp_pptx):
    """Test that sections can have empty names (edge case)."""
    prs = Presentation(temp_pptx)
    slide_ids = [slide.slide_id for slide in prs.slides]

    sections = [
        Section(name="", slide_ids=[slide_ids[0]]),
    ]

    write_sections(prs, sections)
    prs.save(temp_pptx)

    prs2 = Presentation(temp_pptx)
    sections_read = read_sections(prs2)

    assert len(sections_read) == 1
    assert sections_read[0].name == ""
    assert sections_read[0].slide_ids == [slide_ids[0]]
