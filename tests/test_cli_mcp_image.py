"""Tests for MCP image generation in the CLI pipeline."""

import asyncio
import base64
import os
import struct
import zlib
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation

from aippt.cli import build_parser
from aippt.builder import build_slide, BuildContext


def _make_tiny_png():
    """Create a minimal valid 1x1 PNG file."""
    sig = b'\x89PNG\r\n\x1a\n'
    ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
    ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
    raw = zlib.compress(b'\x00\x00\x00\x00')
    idat_crc = zlib.crc32(b'IDAT' + raw) & 0xffffffff
    idat = struct.pack('>I', len(raw)) + b'IDAT' + raw + struct.pack('>I', idat_crc)
    iend_crc = zlib.crc32(b'IEND') & 0xffffffff
    iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
    return sig + ihdr + idat + iend


class TestBuildParserMCPArgs:
    """Verify new CLI arguments are registered."""

    def test_image_gen_mcp_choice(self):
        parser = build_parser()
        args = parser.parse_args(["create", "o.md", "t.pptx", "out.pptx",
                                  "--image-gen", "mcp"])
        assert args.image_gen == "mcp"

    def test_classification_default(self):
        parser = build_parser()
        args = parser.parse_args(["create", "o.md", "t.pptx", "out.pptx",
                                  "--image-gen", "mcp"])
        assert args.classification == "internal"

    def test_classification_custom(self):
        parser = build_parser()
        args = parser.parse_args(["create", "o.md", "t.pptx", "out.pptx",
                                  "--image-gen", "mcp", "--classification", "confidential"])
        assert args.classification == "confidential"

    def test_mcp_server_default(self):
        parser = build_parser()
        args = parser.parse_args(["create", "o.md", "t.pptx", "out.pptx",
                                  "--image-gen", "mcp"])
        assert args.mcp_server == "txt2img"

    def test_mcp_server_custom(self):
        parser = build_parser()
        args = parser.parse_args(["create", "o.md", "t.pptx", "out.pptx",
                                  "--image-gen", "mcp", "--mcp-server", "my-srv"])
        assert args.mcp_server == "my-srv"

    def test_mcp_config_default(self):
        parser = build_parser()
        args = parser.parse_args(["create", "o.md", "t.pptx", "out.pptx",
                                  "--image-gen", "mcp"])
        assert args.mcp_config == "mcp_servers.json"

    def test_mcp_config_custom(self):
        parser = build_parser()
        args = parser.parse_args(["create", "o.md", "t.pptx", "out.pptx",
                                  "--image-gen", "mcp", "--mcp-config", "custom.json"])
        assert args.mcp_config == "custom.json"


class TestAddSlideWithImagePrompt:
    """Test _add_slide handles IMAGE_PROMPT from LLM suggestions."""

    def _make_prs(self):
        return Presentation()

    @patch('aippt.images.generate_mcp_image')
    def test_generates_image_when_image_prompt_present(self, mock_gen):
        """When LLM output contains IMAGE_PROMPT, generate_mcp_image is called."""
        mock_gen.return_value = None

        prs = self._make_prs()
        content = [
            "CONTENT:",
            "- Point one",
            "NARRATIVE: Narrative text.",
            "LAYOUT: bullet",
            "IMAGE_PROMPT: A network topology diagram",
            "VISUALS: Tips",
            "TALKING_POINTS: Points",
        ]

        mock_manager = MagicMock()
        slide_data = {
            'title': 'Test Slide',
            'content': content,
            'slide_num': 1,
        }
        ctx = BuildContext(
            image_gen='mcp',
            mcp_manager=mock_manager,
            image_dir="/tmp/images",
        )
        build_slide(prs, slide_data, ctx)

        mock_gen.assert_called_once()

    def test_no_mcp_call_without_image_prompt(self):
        """When IMAGE_PROMPT is empty, no image generation occurs."""
        prs = self._make_prs()
        content = [
            "CONTENT:",
            "- Point one",
            "NARRATIVE: Narrative text.",
            "LAYOUT: bullet",
            "VISUALS: Tips",
            "TALKING_POINTS: Points",
        ]

        # Should not raise even with mcp_manager=None
        slide_data = {
            'title': 'Test Slide',
            'content': content,
            'slide_num': 1,
        }
        ctx = BuildContext(
            image_gen='mcp',
            mcp_manager=None,
        )
        build_slide(prs, slide_data, ctx)

    @patch('aippt.images.generate_mcp_image')
    def test_disclaimer_added_when_image_generated(self, mock_gen, tmp_path):
        """Slides with generated images get a disclaimer textbox."""
        img_path = str(tmp_path / "slide_01_gen.png")
        with open(img_path, 'wb') as f:
            f.write(_make_tiny_png())

        mock_gen.return_value = img_path

        prs = self._make_prs()
        content = [
            "CONTENT:",
            "- Point one",
            "NARRATIVE: Narrative.",
            "LAYOUT: bullet",
            "IMAGE_PROMPT: A diagram",
            "VISUALS: Tips",
            "TALKING_POINTS: Points",
        ]

        mock_manager = MagicMock()
        slide_data = {
            'title': 'Test Slide',
            'content': content,
            'slide_num': 1,
        }
        ctx = BuildContext(
            image_gen='mcp',
            mcp_manager=mock_manager,
            image_dir=str(tmp_path),
        )
        build_slide(prs, slide_data, ctx)

        slide = prs.slides[0]
        shape_texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("AI-Generated" in t for t in shape_texts)

    @patch('aippt.images.generate_mcp_image')
    def test_disclaimer_in_speaker_notes(self, mock_gen, tmp_path):
        """Speaker notes are prepended with disclaimer for generated images."""
        img_path = str(tmp_path / "slide_01_gen.png")
        with open(img_path, 'wb') as f:
            f.write(_make_tiny_png())

        mock_gen.return_value = img_path

        prs = self._make_prs()
        content = [
            "CONTENT:",
            "- Point one",
            "NARRATIVE: Narrative.",
            "LAYOUT: bullet",
            "IMAGE_PROMPT: A diagram",
            "VISUALS: Tips",
            "TALKING_POINTS: Points",
        ]

        mock_manager = MagicMock()
        slide_data = {
            'title': 'Test Slide',
            'content': content,
            'slide_num': 1,
        }
        ctx = BuildContext(
            image_gen='mcp',
            mcp_manager=mock_manager,
            image_dir=str(tmp_path),
        )
        build_slide(prs, slide_data, ctx)

        slide = prs.slides[0]
        notes = slide.notes_slide.notes_text_frame.text
        assert "[AI-GENERATED]" in notes

    def test_no_image_when_mcp_manager_none(self):
        """IMAGE_PROMPT present but mcp_manager=None — slide renders normally."""
        prs = self._make_prs()
        content = [
            "CONTENT:", "- Point one",
            "NARRATIVE: N.", "LAYOUT: bullet",
            "IMAGE_PROMPT: A diagram that cannot be generated",
            "VISUALS: Tips", "TALKING_POINTS: Points",
        ]

        slide_data = {
            'title': 'Test Slide',
            'content': content,
            'slide_num': 1,
        }
        ctx = BuildContext(
            image_gen='mcp',
            mcp_manager=None,
        )
        result = build_slide(prs, slide_data, ctx)

        assert len(prs.slides) == 1
        assert result == 'bullet'
        shape_texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
        assert not any("AI-Generated" in t for t in shape_texts)

    @patch('aippt.images.generate_mcp_image')
    def test_graceful_fallback_when_generation_fails(self, mock_gen):
        """Slide renders normally when MCP image generation returns None."""
        mock_gen.return_value = None

        prs = self._make_prs()
        content = [
            "CONTENT:",
            "- Point one",
            "NARRATIVE: Narrative.",
            "LAYOUT: bullet",
            "IMAGE_PROMPT: A diagram that will fail",
            "VISUALS: Tips",
            "TALKING_POINTS: Points",
        ]

        mock_manager = MagicMock()
        slide_data = {
            'title': 'Test Slide',
            'content': content,
            'slide_num': 1,
        }
        ctx = BuildContext(
            image_gen='mcp',
            mcp_manager=mock_manager,
            image_dir="/tmp/images",
        )
        result = build_slide(prs, slide_data, ctx)

        assert len(prs.slides) == 1
        assert result == 'bullet'
        shape_texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
        assert not any("AI-Generated" in t for t in shape_texts)


class TestImageGenMetadata:
    @patch('aippt.images.generate_mcp_image')
    def test_metadata_records_image_gen(self, mock_gen, tmp_path):
        """Image generation metadata is appended to slide notes."""
        img_path = str(tmp_path / "slide_01_gen.png")
        with open(img_path, 'wb') as f:
            f.write(_make_tiny_png())

        mock_gen.return_value = img_path

        prs = Presentation()
        content = [
            "CONTENT:", "- Point",
            "NARRATIVE: N.", "LAYOUT: bullet",
            "IMAGE_PROMPT: A diagram",
            "VISUALS: V", "TALKING_POINTS: T",
        ]

        mock_manager = MagicMock()
        slide_data = {
            'title': 'Test',
            'content': content,
            'slide_num': 1,
        }
        ctx = BuildContext(
            image_gen='mcp',
            mcp_manager=mock_manager,
            image_dir=str(tmp_path),
            model="test-model",
        )
        build_slide(prs, slide_data, ctx)

        from aippt.metadata import extract_metadata
        entries = extract_metadata(prs.slides[0])
        ops = [e['operation'] for e in entries]
        assert 'image-gen' in ops
