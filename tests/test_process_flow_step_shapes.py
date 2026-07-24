"""Regression test for MAT-78: addProcessFlow crashing on non-string steps.

LLM-generated pptxgenjs scripts sometimes pass {label, description}-style
objects to addProcessFlow's `steps` param (mirroring addCardGrid's
object-shaped `cards`), even though steps is documented as flat strings.
_addStepBox must normalize these instead of throwing TypeError: step.indexOf
is not a function.
"""

import os
import subprocess
import tempfile

from pptx import Presentation

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
NODE_PATH = subprocess.check_output(
    ["npm", "root", "-g"], cwd=REPO_ROOT, text=True
).strip()


def _run_script(script: str, out_name: str) -> str:
    helpers_path = os.path.join(REPO_ROOT, "lib", "pptxgenjs-helpers.mjs")
    theme_path = os.path.join(REPO_ROOT, "themes", "default.yaml")
    output_path = os.path.join(tempfile.mkdtemp(), out_name)
    full_script = script.format(helpers_path=helpers_path, theme_path=theme_path, output_path=output_path)

    script_path = os.path.join(tempfile.gettempdir(), f"test_process_flow_{out_name}.mjs")
    with open(script_path, "w") as f:
        f.write(full_script)

    env = {**os.environ, "NODE_PATH": NODE_PATH}
    result = subprocess.run(
        ["node", script_path],
        cwd=REPO_ROOT,
        env=env,
        capture_output=True,
        text=True,
        timeout=30,
    )
    assert result.returncode == 0, f"Generation failed:\n{result.stderr}"
    assert os.path.exists(output_path), f"Output not created: {output_path}"
    return output_path


class TestProcessFlowStringSteps:
    """Baseline: plain-string steps (the documented shape) keep working."""

    def test_string_steps_with_label_description_split(self, tmp_path):
        out = _run_script(
            """\
import {{ createDeck, addProcessFlow }} from '{helpers_path}';
const deck = createDeck('{theme_path}');
addProcessFlow(deck, 'Process', [
  'Step One\\nDoes the first thing',
  'Step Two\\nDoes the second thing',
  'Step Three\\nDoes the third thing',
], 1, '');
await deck.save('{output_path}');
""",
            "string_steps.pptx",
        )
        p = Presentation(out)
        assert len(p.slides) == 1
        texts = [s.text_frame.text for shape in p.slides[0].shapes if shape.has_text_frame for s in [shape]]
        assert any("Step One" in t for t in texts)
        assert any("Does the first thing" in t for t in texts)


class TestProcessFlowObjectSteps:
    """MAT-78: {label, description}-style objects must not crash the render."""

    def test_label_description_object_steps(self, tmp_path):
        out = _run_script(
            """\
import {{ createDeck, addProcessFlow }} from '{helpers_path}';
const deck = createDeck('{theme_path}');
addProcessFlow(deck, 'Process', [
  {{ label: 'Step One', description: 'Does the first thing' }},
  {{ label: 'Step Two', description: 'Does the second thing' }},
], 1, '');
await deck.save('{output_path}');
""",
            "object_steps.pptx",
        )
        p = Presentation(out)
        assert len(p.slides) == 1
        texts = [shape.text_frame.text for shape in p.slides[0].shapes if shape.has_text_frame]
        assert any("Step One" in t for t in texts)
        assert any("Does the first thing" in t for t in texts)

    def test_title_body_object_steps(self, tmp_path):
        """Alternate key names (title/body) are also normalized."""
        out = _run_script(
            """\
import {{ createDeck, addProcessFlow }} from '{helpers_path}';
const deck = createDeck('{theme_path}');
addProcessFlow(deck, 'Process', [
  {{ title: 'Alpha', body: 'First' }},
  {{ title: 'Beta', body: 'Second' }},
], 1, '');
await deck.save('{output_path}');
""",
            "title_body_steps.pptx",
        )
        p = Presentation(out)
        assert len(p.slides) == 1

    def test_five_steps_two_row_layout_with_object_steps(self, tmp_path):
        """5+ steps take the two-row branch (separate _addStepBox call sites)."""
        out = _run_script(
            """\
import {{ createDeck, addProcessFlow }} from '{helpers_path}';
const deck = createDeck('{theme_path}');
addProcessFlow(deck, 'Process', [
  {{ label: 'One', description: 'd1' }},
  {{ label: 'Two', description: 'd2' }},
  {{ label: 'Three', description: 'd3' }},
  {{ label: 'Four', description: 'd4' }},
  {{ label: 'Five', description: 'd5' }},
], 1, '');
await deck.save('{output_path}');
""",
            "five_steps.pptx",
        )
        p = Presentation(out)
        assert len(p.slides) == 1
