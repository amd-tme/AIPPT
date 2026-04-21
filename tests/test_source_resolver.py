"""Tests for aippt.source_resolver — deck name to source script resolution."""

import os
import pytest
from pathlib import Path
from pptx import Presentation

from aippt.catalog import catalog_deck, resolve_deck
from aippt.source_resolver import resolve_source


class TestResolveSource:
    @pytest.fixture
    def tracked_deck(self, tmp_path):
        """Create a cataloged deck with source tracking metadata."""
        db_path = str(tmp_path / "test.db")
        pptx_path = str(tmp_path / "test.pptx")
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(pptx_path)
        deck_id = catalog_deck(
            pptx_path, db_path=db_path,
            source_script_path="output/test.mjs",
            source_engine="pptxgenjs",
            source_theme="amd",
        )
        assert deck_id is not None, "catalog_deck returned None"
        deck = resolve_deck(str(deck_id), db_path=db_path)
        assert "source_script_path" in deck, (
            "resolve_deck() missing source columns — PRD 1 Task 8 not merged"
        )
        return db_path, deck_id

    @pytest.fixture
    def untracked_deck(self, tmp_path):
        """Create a cataloged deck WITHOUT source tracking."""
        db_path = str(tmp_path / "test.db")
        pptx_path = str(tmp_path / "nosource.pptx")
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(pptx_path)
        deck_id = catalog_deck(pptx_path, db_path=db_path)
        return db_path, deck_id

    def test_resolve_by_script_path(self, tmp_path):
        """Direct script path should be returned as-is with engine detected."""
        script_path = str(tmp_path / "deck.mjs")
        Path(script_path).write_text(
            "import { createDeck } from '../lib/pptxgenjs-helpers.mjs';\n"
        )
        result = resolve_source(script_path)
        assert result["script_path"] == script_path
        assert result["engine"] == "pptxgenjs"
        assert result["resolved_from"] == "path"

    def test_resolve_by_deck_id(self, tracked_deck):
        db_path, deck_id = tracked_deck
        result = resolve_source(str(deck_id), db_path=db_path)
        assert result["script_path"] == "output/test.mjs"
        assert result["engine"] == "pptxgenjs"
        assert result["resolved_from"] == "catalog"

    def test_resolve_by_deck_name(self, tracked_deck):
        db_path, _ = tracked_deck
        result = resolve_source("test", db_path=db_path)
        assert result["script_path"] == "output/test.mjs"
        assert result["resolved_from"] == "catalog"

    def test_resolve_untracked_deck_returns_error(self, untracked_deck):
        db_path, deck_id = untracked_deck
        result = resolve_source(str(deck_id), db_path=db_path)
        assert result["error"] is not None
        assert "No source script" in result["error"]

    def test_resolve_no_match_returns_error(self, tracked_deck):
        db_path, _ = tracked_deck
        result = resolve_source("nonexistent", db_path=db_path)
        assert result["error"] is not None
        assert "No deck found" in result["error"]

    def test_resolve_ambiguous_returns_choices(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        for name in ("security-review", "security-audit"):
            pptx_path = str(tmp_path / f"{name}.pptx")
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(pptx_path)
            catalog_deck(
                pptx_path, db_path=db_path,
                source_script_path=f"output/{name}.mjs",
                source_engine="pptxgenjs",
            )
        result = resolve_source("security", db_path=db_path)
        assert result["error"] is not None
        assert "Multiple decks" in result["error"]
        assert len(result["choices"]) == 2


from aippt.source_resolver import create_backup, restore_backup, has_backup


class TestBackupManagement:
    def test_create_backup(self, tmp_path):
        script = tmp_path / "deck.mjs"
        script.write_text("original content")
        bak_path = create_backup(str(script))
        assert bak_path == str(script) + ".bak"
        assert Path(bak_path).read_text() == "original content"

    def test_create_backup_skips_if_exists(self, tmp_path):
        script = tmp_path / "deck.mjs"
        script.write_text("v1")
        bak = tmp_path / "deck.mjs.bak"
        bak.write_text("original backup")
        result = create_backup(str(script))
        assert result is None  # skipped
        assert bak.read_text() == "original backup"  # not overwritten

    def test_has_backup(self, tmp_path):
        script = tmp_path / "deck.mjs"
        script.write_text("content")
        assert has_backup(str(script)) is False
        (tmp_path / "deck.mjs.bak").write_text("backup")
        assert has_backup(str(script)) is True

    def test_restore_backup(self, tmp_path):
        script = tmp_path / "deck.mjs"
        bak = tmp_path / "deck.mjs.bak"
        script.write_text("modified content")
        bak.write_text("original content")
        restored = restore_backup(str(script))
        assert restored is True
        assert script.read_text() == "original content"
        assert bak.exists()  # backup is preserved, not deleted

    def test_restore_backup_no_bak(self, tmp_path):
        script = tmp_path / "deck.mjs"
        script.write_text("content")
        restored = restore_backup(str(script))
        assert restored is False


import shutil
import subprocess
from aippt.source_resolver import run_script

node_installed = pytest.mark.skipif(
    shutil.which("node") is None,
    reason="Node.js not installed",
)


class TestRunScript:
    @node_installed
    def test_run_node_script(self, tmp_path):
        script = tmp_path / "hello.mjs"
        script.write_text("console.log('hello from node');")
        result = run_script(str(script), engine="pptxgenjs")
        assert result["success"] is True
        assert "hello from node" in result["stdout"]

    def test_run_python_script(self, tmp_path):
        script = tmp_path / "hello.py"
        script.write_text("print('hello from python')")
        result = run_script(str(script), engine="python-pptx")
        assert result["success"] is True
        assert "hello from python" in result["stdout"]

    @node_installed
    def test_run_failing_script(self, tmp_path):
        script = tmp_path / "bad.mjs"
        script.write_text("throw new Error('intentional failure');")
        result = run_script(str(script), engine="pptxgenjs")
        assert result["success"] is False
        assert "intentional failure" in result["stderr"]

    def test_run_detects_file_lock_error(self, tmp_path):
        """Simulate a file-in-use error message."""
        script = tmp_path / "lock.py"
        script.write_text(
            "import sys; sys.stderr.write('PermissionError: [Errno 13]'); sys.exit(1)"
        )
        result = run_script(str(script), engine="python-pptx")
        assert result["success"] is False
        assert result["file_locked"] is True
