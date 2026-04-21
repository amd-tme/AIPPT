"""End-to-end integration test for catalog-search-tag-export-remix flow."""
import csv
import os
import textwrap

import pytest
import yaml
from pptx import Presentation

from aippt.catalog import (
    catalog_deck,
    search_slides,
    add_tags,
    get_slide_tags,
    get_db,
    add_taxonomy_tags,
    list_taxonomy,
    import_taxonomy_csv,
    export_taxonomy_csv,
    get_taxonomy_names,
    remove_slide_tag,
)
from aippt.export import export_csv
from aippt.remix import generate_manifest, load_manifest, assemble_deck
from aippt.cli import build_parser, cmd_models, cmd_tags, cmd_tag, cmd_untag
from aippt.config import load_model_config, save_model_config, ConfigError
from aippt.writeback import write_notes_to_pptx, create_backup


@pytest.fixture
def sample_deck(tmp_path):
    """Create a minimal PPTX for testing."""
    prs = Presentation()
    layout = prs.slide_layouts[0]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Introduction"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Security Overview"

    slide3 = prs.slides.add_slide(layout)
    slide3.shapes.title.text = "Architecture Deep Dive"

    path = str(tmp_path / "test.pptx")
    prs.save(path)
    return path


@pytest.fixture
def db_path(tmp_path):
    return str(tmp_path / "test.db")


def test_full_catalog_search_tag_export_flow(tmp_path, sample_deck, db_path):
    """Test the complete workflow: catalog -> search -> tag -> export."""
    # Step 1: Catalog the deck
    deck_id = catalog_deck(sample_deck, db_path=db_path)
    assert deck_id > 0

    # Step 2: Search by title
    results = search_slides(db_path=db_path, title_contains="Security")
    assert len(results) == 1
    assert results[0]["title"] == "Security Overview"

    # Step 3: Tag a slide
    slide_id = results[0]["id"]
    add_tags(slide_id, ["security", "overview"], source="manual", db_path=db_path)
    tags = get_slide_tags(slide_id, db_path)
    assert "security" in tags
    assert "overview" in tags

    # Step 4: Search by tag
    results = search_slides(db_path=db_path, tags=["security"])
    assert len(results) == 1
    assert results[0]["title"] == "Security Overview"

    # Step 5: Export to CSV
    csv_path = str(tmp_path / "export.csv")
    count = export_csv(csv_path, db_path=db_path, export_all=True)
    assert count == 3  # All 3 slides exported
    assert os.path.exists(csv_path)

    with open(csv_path, encoding="utf-8") as f:
        reader = list(csv.DictReader(f))
        assert len(reader) == 3
        # Find the tagged slide in the CSV
        security_row = [r for r in reader if r["title"] == "Security Overview"][0]
        assert "security" in security_row["tags"]
        assert "overview" in security_row["tags"]


def test_search_and_manifest_generation(tmp_path, sample_deck, db_path):
    """Test search -> manifest generation -> manifest loading."""
    # Catalog
    catalog_deck(sample_deck, db_path=db_path)

    # Tag slides for searching
    results = search_slides(db_path=db_path, title_contains="Security")
    add_tags(results[0]["id"], ["security"], source="manual", db_path=db_path)

    results = search_slides(db_path=db_path, title_contains="Architecture")
    add_tags(results[0]["id"], ["architecture"], source="manual", db_path=db_path)

    # Search by multiple criteria
    security_results = search_slides(db_path=db_path, tags=["security"])
    assert len(security_results) == 1

    # Generate manifest
    manifest_yaml = generate_manifest(security_results, title="Security Deck")
    manifest_path = str(tmp_path / "manifest.yaml")
    with open(manifest_path, "w", encoding="utf-8") as f:
        f.write(manifest_yaml)

    # Load manifest back
    manifest = load_manifest(manifest_path)
    assert manifest["title"] == "Security Deck"
    assert len(manifest["slides"]) == 1
    assert manifest["slides"][0]["title"] == "Security Overview"


def test_remix_assembly(tmp_path, sample_deck, db_path):
    """Test assembling a remixed deck from a manifest."""
    # Catalog
    catalog_deck(sample_deck, db_path=db_path)

    # Create manifest picking slides 1 and 3
    manifest = {
        "title": "Cherry-picked Deck",
        "slides": [
            {"deck": "test.pptx", "deck_path": sample_deck, "position": 1, "title": "Introduction"},
            {"deck": "test.pptx", "deck_path": sample_deck, "position": 3, "title": "Architecture Deep Dive"},
        ],
    }
    manifest_path = str(tmp_path / "manifest.yaml")
    with open(manifest_path, "w") as f:
        yaml.dump(manifest, f)

    output_path = str(tmp_path / "remixed.pptx")
    count = assemble_deck(manifest_path, output_path, db_path=db_path)

    assert count == 2
    assert os.path.exists(output_path)

    # Verify the output deck has 2 slides
    result_prs = Presentation(output_path)
    assert len(result_prs.slides) == 2


def test_recatalog_detects_same_file(tmp_path, sample_deck, db_path):
    """Test that re-cataloging the same file is a no-op."""
    deck_id1 = catalog_deck(sample_deck, db_path=db_path)
    deck_id2 = catalog_deck(sample_deck, db_path=db_path)
    assert deck_id1 == deck_id2  # Same ID, no duplicate


def test_export_specific_deck(tmp_path, sample_deck, db_path):
    """Test exporting only a specific deck."""
    catalog_deck(sample_deck, db_path=db_path)

    csv_path = str(tmp_path / "export.csv")
    count = export_csv(csv_path, db_path=db_path, deck_path=sample_deck)
    assert count == 3

    with open(csv_path, encoding="utf-8") as f:
        reader = list(csv.DictReader(f))
        assert len(reader) == 3


# ---------------------------------------------------------------------------
# Model management integration tests
# ---------------------------------------------------------------------------

class TestModelsIntegration:
    def test_models_show(self, tmp_path, capsys, monkeypatch, models_yaml):
        """Test CLI 'models' command shows current configuration."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args(["models"])
        result = cmd_models(args)

        assert result == 0
        captured = capsys.readouterr()
        assert "enhance" in captured.out
        assert "claude-3.5-sonnet" in captured.out

    def test_models_show_errors_without_config(self, tmp_path, capsys, monkeypatch):
        """Test CLI 'models' errors when models.yaml is missing."""
        missing = str(tmp_path / "nonexistent.yaml")
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", missing)

        parser = build_parser()
        args = parser.parse_args(["models"])
        result = cmd_models(args)

        assert result == 1
        captured = capsys.readouterr()
        assert "Error" in captured.out

    def test_models_init_creates_config(self, tmp_path, capsys, monkeypatch):
        """Test CLI 'models init' copies example file."""
        import shutil
        import aippt.config as cfg_module
        dest = str(tmp_path / "models.yaml")
        # Use actual example file from project root
        example = cfg_module.EXAMPLE_CONFIG_PATH
        monkeypatch.setattr(cfg_module, "DEFAULT_CONFIG_PATH", dest)

        parser = build_parser()
        args = parser.parse_args(["models", "init"])
        result = cmd_models(args)

        assert result == 0
        assert os.path.exists(dest)
        # Should be a valid config
        config = load_model_config(dest)
        assert "registry" in config
        assert "defaults" in config
        captured = capsys.readouterr()
        assert "Created" in captured.out

    def test_models_init_fails_if_already_exists(self, tmp_path, capsys, monkeypatch, models_yaml):
        """Test CLI 'models init' fails if models.yaml already exists."""
        import aippt.config as cfg_module
        monkeypatch.setattr(cfg_module, "DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args(["models", "init"])
        result = cmd_models(args)

        assert result == 1
        captured = capsys.readouterr()
        assert "already exists" in captured.out

    def test_models_set(self, tmp_path, capsys, monkeypatch, models_yaml):
        """Test CLI 'models set' writes config file."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args(["models", "set", "enhance", "gpt-4.1"])
        result = cmd_models(args)

        assert result == 0
        config = load_model_config(models_yaml)
        assert config["defaults"]["enhance"] == "gpt-4.1"

    def test_models_set_rejects_unknown_model(self, tmp_path, capsys, monkeypatch, models_yaml):
        """Test that setting a model not in registry returns error."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args(["models", "set", "enhance", "nonexistent-model"])
        result = cmd_models(args)

        assert result == 1
        captured = capsys.readouterr()
        assert "nonexistent-model" in captured.out

    def test_models_reset_returns_error(self, tmp_path, capsys, monkeypatch, models_yaml):
        """Test CLI 'models reset' is deprecated and returns error."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args(["models", "reset"])
        result = cmd_models(args)

        assert result == 1
        captured = capsys.readouterr()
        assert "not supported" in captured.out or "no longer" in captured.out

    def test_models_list_available(self, tmp_path, capsys, monkeypatch, models_yaml):
        """Test CLI 'models list-available' shows registry."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args(["models", "list-available"])
        result = cmd_models(args)

        assert result == 0
        captured = capsys.readouterr()
        assert "gpt-4o" in captured.out
        assert "claude-3.5-sonnet" in captured.out
        assert "Provider" in captured.out

    def test_models_list_available_errors_without_config(self, tmp_path, capsys, monkeypatch):
        """Test 'models list-available' errors when models.yaml is missing."""
        missing = str(tmp_path / "nonexistent.yaml")
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", missing)

        parser = build_parser()
        args = parser.parse_args(["models", "list-available"])
        result = cmd_models(args)

        assert result == 1
        captured = capsys.readouterr()
        assert "Error" in captured.out

    def test_create_uses_config_default(self, tmp_path, monkeypatch, models_yaml):
        """Test that create --enhance without --model uses config default."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args([
            "create", "outline.md", "template.pptx", "output.pptx", "--enhance"
        ])
        # --model should be None when not specified
        assert args.model is None

        # Verify config resolution
        from aippt.config import get_model_default
        model = args.model or get_model_default("enhance", models_yaml)
        assert model == "claude-3.5-sonnet"

    def test_cli_model_overrides_config(self, tmp_path, monkeypatch, models_yaml):
        """Test --model flag takes priority over config."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args([
            "create", "outline.md", "template.pptx", "output.pptx",
            "--enhance", "--model", "claude-3.7-sonnet"
        ])
        assert args.model == "claude-3.7-sonnet"

        from aippt.config import get_model_default
        model = args.model or get_model_default("enhance", models_yaml)
        assert model == "claude-3.7-sonnet"

    def test_analyze_per_mode_defaults(self, tmp_path, monkeypatch, models_yaml):
        """Test that analyze uses per-mode model defaults."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        from aippt.config import get_model_default

        assert get_model_default("feedback", models_yaml) == "gpt-4o"
        assert get_model_default("notes", models_yaml) == "gpt-4o"
        assert get_model_default("tags", models_yaml) == "gpt-4o"

    def test_models_set_invalid_operation(self, tmp_path, capsys, monkeypatch, models_yaml):
        """Test that setting an invalid operation returns error."""
        monkeypatch.setattr("aippt.config.DEFAULT_CONFIG_PATH", models_yaml)

        parser = build_parser()
        args = parser.parse_args(["models", "set", "nonexistent", "some-model"])
        result = cmd_models(args)

        assert result == 1
        captured = capsys.readouterr()
        assert "Unknown operation" in captured.out


# ---------------------------------------------------------------------------
# Tag management integration tests
# ---------------------------------------------------------------------------


class TestTagManagementIntegration:

    def test_taxonomy_import_list_export_roundtrip(self, tmp_path, db_path, capsys):
        """Import CSV, list via CLI, export, compare."""
        # Create a taxonomy CSV
        csv_in = str(tmp_path / "taxonomy.csv")
        with open(csv_in, "w", encoding="utf-8", newline="") as f:
            writer = csv.DictWriter(f, fieldnames=["name", "category"])
            writer.writeheader()
            writer.writerow({"name": "security", "category": "topic"})
            writer.writerow({"name": "cloud", "category": "topic"})
            writer.writerow({"name": "executive", "category": "audience"})

        # Import via CLI
        parser = build_parser()
        args = parser.parse_args(["tags", "import", csv_in, "--db", db_path])
        result = cmd_tags(args)
        assert result == 0
        captured = capsys.readouterr()
        assert "3 new" in captured.out

        # List via CLI
        args = parser.parse_args(["tags", "--db", db_path])
        result = cmd_tags(args)
        assert result == 0
        captured = capsys.readouterr()
        assert "security" in captured.out
        assert "topic" in captured.out

        # Export via CLI
        csv_out = str(tmp_path / "exported.csv")
        args = parser.parse_args(["tags", "export", csv_out, "--db", db_path])
        result = cmd_tags(args)
        assert result == 0

        # Compare
        with open(csv_out, encoding="utf-8") as f:
            reader = list(csv.DictReader(f))
            assert len(reader) == 3
            names = {r["name"] for r in reader}
            assert names == {"security", "cloud", "executive"}

    def test_tag_untag_slide_cli(self, tmp_path, sample_deck, db_path, capsys):
        """CLI tag and untag commands."""
        catalog_deck(sample_deck, db_path=db_path)

        parser = build_parser()

        # Tag slide 1
        args = parser.parse_args(["tag", "1", "security,overview", "--db", db_path])
        result = cmd_tag(args)
        assert result == 0
        captured = capsys.readouterr()
        assert "Tagged slide 1" in captured.out

        tags = get_slide_tags(1, db_path)
        assert "security" in tags
        assert "overview" in tags

        # Untag one
        args = parser.parse_args(["untag", "1", "overview", "--db", db_path])
        result = cmd_untag(args)
        assert result == 0
        captured = capsys.readouterr()
        assert "Untagged slide 1" in captured.out

        tags = get_slide_tags(1, db_path)
        assert "security" in tags
        assert "overview" not in tags

        # Untag all
        args = parser.parse_args(["untag", "1", "--all", "--db", db_path])
        result = cmd_untag(args)
        assert result == 0

        tags = get_slide_tags(1, db_path)
        assert tags == []

    def test_tags_add_remove_cli(self, tmp_path, db_path, capsys):
        """CLI taxonomy add and remove commands."""
        parser = build_parser()

        # Add
        args = parser.parse_args(["tags", "add", "zero-trust", "--category", "security", "--db", db_path])
        result = cmd_tags(args)
        assert result == 0
        captured = capsys.readouterr()
        assert "Added" in captured.out

        tags = list_taxonomy(db_path)
        assert len(tags) == 1
        assert tags[0]["name"] == "zero-trust"

        # Remove
        args = parser.parse_args(["tags", "remove", "zero-trust", "--db", db_path])
        result = cmd_tags(args)
        assert result == 0
        captured = capsys.readouterr()
        assert "Removed" in captured.out

        tags = list_taxonomy(db_path)
        assert len(tags) == 0

    def test_tags_rename_cli(self, tmp_path, sample_deck, db_path, capsys):
        """CLI taxonomy rename updates tags and taxonomy."""
        catalog_deck(sample_deck, db_path=db_path)
        add_tags(1, ["cloud"], source="ai", db_path=db_path)
        add_taxonomy_tags([{"name": "cloud", "category": "topic"}], db_path)

        parser = build_parser()
        args = parser.parse_args(["tags", "rename", "cloud", "cloud-computing", "--db", db_path])
        result = cmd_tags(args)
        assert result == 0
        captured = capsys.readouterr()
        assert "Renamed" in captured.out
        assert "cloud-computing" in captured.out

        # Verify tag on slide was renamed
        tags = get_slide_tags(1, db_path)
        assert "cloud-computing" in tags
        assert "cloud" not in tags

        # Verify taxonomy was renamed
        tax = list_taxonomy(db_path)
        names = [t["name"] for t in tax]
        assert "cloud-computing" in names

    def test_analyze_uses_db_taxonomy(self, tmp_path, db_path):
        """Verify get_taxonomy_names returns DB taxonomy for analyze fallback."""
        # Empty taxonomy -> empty list
        names = get_taxonomy_names(db_path)
        assert names == []

        # Populate taxonomy
        add_taxonomy_tags([
            {"name": "security", "category": "topic"},
            {"name": "cloud", "category": "topic"},
        ], db_path)

        names = get_taxonomy_names(db_path)
        assert len(names) == 2
        assert "security" in names
        assert "cloud" in names


# ---------------------------------------------------------------------------
# Metadata API integration tests
# ---------------------------------------------------------------------------


@pytest.fixture
def metadata_deck(tmp_path):
    """Create a PPTX with explicit core_properties metadata."""
    from datetime import datetime, timezone

    prs = Presentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Meta Slide"

    prs.core_properties.author = "Integration Author"
    prs.core_properties.created = datetime(2024, 5, 10, 9, 0, 0, tzinfo=timezone.utc)
    prs.core_properties.modified = datetime(2025, 9, 25, 17, 0, 0, tzinfo=timezone.utc)

    path = str(tmp_path / "meta_deck.pptx")
    prs.save(path)
    return path


class TestMetadataAPIIntegration:
    """Integration tests for metadata fields surfaced through the REST API."""

    @pytest.fixture
    def client(self, tmp_path, metadata_deck):
        """Create a TestClient with the FastAPI app pointed at a fresh db."""
        from fastapi.testclient import TestClient
        from aippt.web.app import create_app

        db_path = str(tmp_path / "api_meta.db")
        uploads_dir = str(tmp_path / "uploads")
        catalog_deck(metadata_deck, db_path=db_path)
        app = create_app(db_path=db_path, uploads_dir=uploads_dir)
        return TestClient(app)

    def test_list_decks_includes_metadata_fields(self, client):
        """GET /api/decks must include author, created_date, modified_date."""
        response = client.get("/api/decks")
        assert response.status_code == 200
        decks = response.json()
        assert len(decks) >= 1
        deck = decks[0]
        assert "author" in deck
        assert "created_date" in deck
        assert "modified_date" in deck

    def test_list_decks_author_value(self, client):
        """GET /api/decks returns correct author value."""
        response = client.get("/api/decks")
        assert response.status_code == 200
        deck = response.json()[0]
        assert deck["author"] == "Integration Author"

    def test_list_decks_created_date_value(self, client):
        """GET /api/decks returns created_date containing expected date."""
        response = client.get("/api/decks")
        assert response.status_code == 200
        deck = response.json()[0]
        assert deck["created_date"] is not None
        assert "2024-05-10" in deck["created_date"]

    def test_list_decks_modified_date_value(self, client):
        """GET /api/decks returns modified_date containing expected date."""
        response = client.get("/api/decks")
        assert response.status_code == 200
        deck = response.json()[0]
        assert deck["modified_date"] is not None
        assert "2025-09-25" in deck["modified_date"]

    def test_get_slide_includes_metadata_fields(self, client):
        """GET /api/slides/{id} must include author, slide_created_date, updated_at."""
        response = client.get("/api/slides/1")
        assert response.status_code == 200
        slide = response.json()
        assert "author" in slide
        assert "slide_created_date" in slide
        assert "updated_at" in slide

    def test_get_slide_author_value(self, client):
        """GET /api/slides/{id} returns correct author value."""
        response = client.get("/api/slides/1")
        assert response.status_code == 200
        slide = response.json()
        assert slide["author"] == "Integration Author"

    def test_get_slide_created_date_value(self, client):
        """GET /api/slides/{id} returns slide_created_date containing expected date."""
        response = client.get("/api/slides/1")
        assert response.status_code == 200
        slide = response.json()
        assert slide["slide_created_date"] is not None
        assert "2024-05-10" in slide["slide_created_date"]


# ---------------------------------------------------------------------------
# Notes editing integration tests
# ---------------------------------------------------------------------------


class TestNotesEditingIntegration:
    """Integration tests for notes editing via the web API."""

    @pytest.fixture
    def notes_deck(self, tmp_path):
        """Create a PPTX with notes on the first slide."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Notes Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Initial notes from PPTX"
        path = str(tmp_path / "notes_deck.pptx")
        prs.save(path)
        return path

    @pytest.fixture
    def client(self, tmp_path, notes_deck):
        from fastapi.testclient import TestClient
        from aippt.web.app import create_app

        db_path = str(tmp_path / "notes.db")
        uploads_dir = str(tmp_path / "uploads")
        catalog_deck(notes_deck, db_path=db_path)
        app = create_app(db_path=db_path, uploads_dir=uploads_dir)
        return TestClient(app)

    @pytest.fixture
    def db_path(self, client):
        return client.app.state.db_path

    def test_save_creates_history_row(self, client, db_path):
        resp = client.post("/api/slides/1/notes/save", json={"notes": "Edited"})
        assert resp.status_code == 200

        conn = get_db(db_path)
        hist = conn.execute(
            "SELECT * FROM edit_history WHERE slide_id = 1"
        ).fetchone()
        assert hist is not None
        assert hist["old_value"] == "Initial notes from PPTX"
        assert hist["new_value"] == "Edited"
        conn.close()

    def test_two_saves_create_chain(self, client, db_path):
        client.post("/api/slides/1/notes/save", json={"notes": "Edit 1"})
        client.post("/api/slides/1/notes/save", json={"notes": "Edit 2"})

        conn = get_db(db_path)
        rows = conn.execute(
            "SELECT old_value, new_value FROM edit_history WHERE slide_id = 1 ORDER BY id"
        ).fetchall()
        assert len(rows) == 2
        assert rows[0]["old_value"] == "Initial notes from PPTX"
        assert rows[0]["new_value"] == "Edit 1"
        assert rows[1]["old_value"] == "Edit 1"
        assert rows[1]["new_value"] == "Edit 2"
        conn.close()

    def test_history_api_returns_correct_order(self, client, db_path):
        client.post("/api/slides/1/notes/save", json={"notes": "First edit"})
        client.post("/api/slides/1/notes/save", json={"notes": "Second edit"})

        resp = client.get("/api/slides/1/notes/history")
        assert resp.status_code == 200
        history = resp.json()["history"]
        assert len(history) == 2
        # Newest first
        assert history[0]["new_value"] == "Second edit"
        assert history[1]["new_value"] == "First edit"

    def test_updated_at_changes_on_save(self, client, db_path):
        conn = get_db(db_path)
        before = conn.execute(
            "SELECT updated_at FROM slides WHERE id = 1"
        ).fetchone()["updated_at"]
        conn.close()

        client.post("/api/slides/1/notes/save", json={"notes": "Timestamp test"})

        conn = get_db(db_path)
        after = conn.execute(
            "SELECT updated_at FROM slides WHERE id = 1"
        ).fetchone()["updated_at"]
        conn.close()
        assert after >= before


# ---------------------------------------------------------------------------
# Notes write-back integration tests
# ---------------------------------------------------------------------------


class TestNotesWritebackIntegration:
    """Integration tests for the full notes round-trip: catalog -> edit -> write back -> verify."""

    @pytest.fixture
    def notes_deck(self, tmp_path):
        """Create a 2-slide PPTX with notes on slide 1."""
        prs = Presentation()
        layout = prs.slide_layouts[0]

        s1 = prs.slides.add_slide(layout)
        s1.shapes.title.text = "First Slide"
        s1.notes_slide.notes_text_frame.text = "Original first"

        s2 = prs.slides.add_slide(layout)
        s2.shapes.title.text = "Second Slide"

        path = str(tmp_path / "roundtrip.pptx")
        prs.save(path)
        return path

    def test_full_round_trip(self, tmp_path, notes_deck):
        db_path = str(tmp_path / "rt.db")

        # 1. Catalog
        deck_id = catalog_deck(notes_deck, db_path=db_path)

        # 2. Edit notes via record_edit
        from aippt.catalog import record_edit
        record_edit(1, "notes", "Updated first slide notes", source="web", db_path=db_path)
        record_edit(2, "notes", "Brand new second slide notes", source="web", db_path=db_path)

        # 3. Write back
        result = write_notes_to_pptx(notes_deck, db_path=db_path, deck_id=deck_id)
        assert result.slides_written == 2
        assert result.slides_skipped == 0

        # 4. Verify PPTX
        prs = Presentation(notes_deck)
        assert prs.slides[0].notes_slide.notes_text_frame.text == "Updated first slide notes"
        assert prs.slides[1].notes_slide.notes_text_frame.text == "Brand new second slide notes"

    def test_write_to_output_path(self, tmp_path, notes_deck):
        db_path = str(tmp_path / "out.db")
        output = str(tmp_path / "output_copy.pptx")

        deck_id = catalog_deck(notes_deck, db_path=db_path)
        from aippt.catalog import record_edit
        record_edit(1, "notes", "Output path test", source="web", db_path=db_path)

        write_notes_to_pptx(notes_deck, db_path=db_path, deck_id=deck_id, output_path=output)

        # Output has new notes
        prs_out = Presentation(output)
        assert prs_out.slides[0].notes_slide.notes_text_frame.text == "Output path test"

        # Original unchanged
        prs_orig = Presentation(notes_deck)
        assert prs_orig.slides[0].notes_slide.notes_text_frame.text == "Original first"

    def test_round_trip_preserves_after_recatalog(self, tmp_path, notes_deck):
        db_path = str(tmp_path / "recat.db")

        # Catalog -> edit -> write back
        deck_id = catalog_deck(notes_deck, db_path=db_path)
        from aippt.catalog import record_edit
        record_edit(1, "notes", "Persisted notes", source="web", db_path=db_path)
        write_notes_to_pptx(notes_deck, db_path=db_path, deck_id=deck_id)

        # Re-catalog (simulates re-ingest)
        deck_id_2 = catalog_deck(notes_deck, db_path=db_path)

        # Notes should be preserved from the PPTX
        conn = get_db(db_path)
        row = conn.execute(
            "SELECT notes FROM slides WHERE deck_id = ? AND position = 1",
            (deck_id_2,),
        ).fetchone()
        conn.close()
        assert row["notes"] == "Persisted notes"


# ---------------------------------------------------------------------------
# LAYOUT: and IMAGE: directive integration tests
# ---------------------------------------------------------------------------


class TestDirectivesIntegration:
    """Integration tests for LAYOUT: and IMAGE: outline directives in create_deck()."""

    @pytest.fixture
    def minimal_template(self, tmp_path):
        """Create a minimal blank PPTX template and return its path."""
        prs = Presentation()
        path = str(tmp_path / "template.pptx")
        prs.save(path)
        return path

    def test_layout_directive_produces_correct_layout(self, tmp_path, minimal_template):
        """LAYOUT: two_column directive causes slide to use 'Two Content' layout."""
        from aippt.pipeline import run_pipeline, PipelineConfig

        outline = "# My Slide\nLAYOUT: two_column\n- Left point\n- Right point\n"
        output_path = str(tmp_path / "out.pptx")

        config = PipelineConfig(
            outline_text=outline,
            template_path=minimal_template,
            output_path=output_path,
        )
        run_pipeline(config)

        prs = Presentation(output_path)
        assert len(prs.slides) == 1
        assert prs.slides[0].slide_layout.name == "Two Content"

    def test_image_directive_embeds_image(self, tmp_path, minimal_template):
        """IMAGE: directive with a valid PNG results in a picture shape on the slide."""
        from PIL import Image as PILImage
        from aippt.pipeline import run_pipeline, PipelineConfig

        # Create a tiny red PNG next to the outline file
        img_path = tmp_path / "test.png"
        PILImage.new("RGB", (100, 100), "red").save(str(img_path))

        outline_path = tmp_path / "outline.md"
        outline_path.write_text(
            "# Image Slide\nIMAGE: test.png\n- Some bullet\n",
            encoding="utf-8",
        )
        output_path = str(tmp_path / "out.pptx")

        config = PipelineConfig(
            outline_text=outline_path.read_text(encoding="utf-8"),
            template_path=minimal_template,
            output_path=output_path,
            outline_path=str(outline_path),
        )
        run_pipeline(config)

        prs = Presentation(output_path)
        assert len(prs.slides) == 1

        shape_types = [shape.shape_type for shape in prs.slides[0].shapes]
        assert 13 in shape_types, "Expected a picture shape (type 13) on the slide"

    def test_layout_directive_without_enhance(self, tmp_path, minimal_template):
        """LAYOUT: numbered without --enhance applies numbered items to slide content."""
        from aippt.pipeline import run_pipeline, PipelineConfig

        outline = (
            "# Steps\n"
            "LAYOUT: numbered\n"
            "- Install dependencies\n"
            "- Run migrations\n"
            "- Start the service\n"
        )
        output_path = str(tmp_path / "out.pptx")

        config = PipelineConfig(
            outline_text=outline,
            template_path=minimal_template,
            output_path=output_path,
        )
        run_pipeline(config)

        prs = Presentation(output_path)
        assert len(prs.slides) == 1

        slide = prs.slides[0]
        # Collect all text from non-title placeholders and text boxes
        all_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_text.append(text)

        # At least one paragraph should begin with a numbered prefix like "1. "
        numbered = [t for t in all_text if t.startswith("1. ") or t.startswith("2. ") or t.startswith("3. ")]
        assert len(numbered) >= 1, f"Expected numbered items, got: {all_text}"

    def test_missing_image_gracefully_skipped(self, tmp_path, minimal_template):
        """IMAGE: with a nonexistent file does not crash; slide is created without a picture."""
        from aippt.pipeline import run_pipeline, PipelineConfig

        outline_path = tmp_path / "outline.md"
        outline_path.write_text(
            "# No Image Slide\nIMAGE: nonexistent.png\n- Bullet one\n",
            encoding="utf-8",
        )
        output_path = str(tmp_path / "out.pptx")

        # Should not raise
        config = PipelineConfig(
            outline_text=outline_path.read_text(encoding="utf-8"),
            template_path=minimal_template,
            output_path=output_path,
            outline_path=str(outline_path),
        )
        result = run_pipeline(config)

        assert result.slide_count == 1

        prs = Presentation(output_path)
        assert len(prs.slides) == 1

        # No picture shape (type 13) should be present
        shape_types = [shape.shape_type for shape in prs.slides[0].shapes]
        assert 13 not in shape_types, "Did not expect a picture shape when image is missing"


# ---------------------------------------------------------------------------
# Reverse round-trip integration tests
# ---------------------------------------------------------------------------


class TestReverseRoundTrip:
    """Test that reverse -> create round-trip preserves content without notes leakage."""

    def test_roundtrip_no_notes_leakage(self, tmp_path):
        """Create a deck, reverse it (with notes), create from reversed md — no notes in body."""
        from aippt.reverse import convert_pptx_to_outline
        from aippt.pipeline import run_pipeline, PipelineConfig
        from pptx import Presentation as PresentationClass

        # Step 1: Create a simple deck from outline
        outline = textwrap.dedent("""\
            # Test Slide
            - Bullet one
            - Bullet two
        """)
        outline_path = tmp_path / "original.md"
        outline_path.write_text(outline)

        template_pptx = tmp_path / "template.pptx"
        PresentationClass().save(str(template_pptx))

        first_pptx = str(tmp_path / "first.pptx")
        config = PipelineConfig(
            outline_text=outline,
            template_path=str(template_pptx),
            output_path=first_pptx,
            outline_path=str(outline_path),
        )
        run_pipeline(config)

        # Manually add speaker notes to the first slide
        prs = PresentationClass(first_pptx)
        slide = prs.slides[0]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "These are speaker notes that should NOT appear on slides"
        prs.save(first_pptx)

        # Step 2: Reverse the deck to markdown (notes included)
        reversed_md = str(tmp_path / "reversed.md")
        convert_pptx_to_outline(first_pptx, reversed_md, include_notes=True)

        # Verify reversed markdown has notes as HTML comments
        reversed_content = open(reversed_md).read()
        assert '<!-- notes' in reversed_content
        assert '*Notes:*' not in reversed_content

        # Step 3: Create a new deck from the reversed markdown
        second_pptx = str(tmp_path / "second.pptx")
        config2 = PipelineConfig(
            outline_text=reversed_content,
            template_path=str(template_pptx),
            output_path=second_pptx,
            outline_path=reversed_md,
        )
        run_pipeline(config2)

        # Step 4: Verify no notes content leaked into slide body
        prs2 = PresentationClass(second_pptx)
        for slide in prs2.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    assert 'speaker notes' not in text.lower(), \
                        f"Notes leaked into slide body: {text}"
                    assert 'should NOT appear' not in text

    def test_roundtrip_legacy_notes_stripped(self, tmp_path):
        """Legacy *Notes:* format in markdown should be stripped by create."""
        from aippt.pipeline import run_pipeline, PipelineConfig
        from pptx import Presentation as PresentationClass

        # Simulate old-format reversed markdown with *Notes:* sections
        legacy_md = textwrap.dedent("""\
            # Test Slide
            - Bullet one
            - Bullet two

            *Notes:*
            - Old format speaker notes
            - Should not appear on slide
        """)

        template_pptx = tmp_path / "template.pptx"
        PresentationClass().save(str(template_pptx))

        output_pptx = str(tmp_path / "output.pptx")
        md_path = tmp_path / "legacy.md"
        md_path.write_text(legacy_md)
        config = PipelineConfig(
            outline_text=legacy_md,
            template_path=str(template_pptx),
            output_path=output_pptx,
            outline_path=str(md_path),
        )
        run_pipeline(config)

        prs = PresentationClass(output_pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    assert 'Old format speaker notes' not in text
                    assert '*Notes:*' not in text
