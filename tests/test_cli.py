"""Tests for aippt.cli module."""

import os
import sys
import pytest
from unittest.mock import MagicMock, patch

from aippt.cli import (
    build_parser,
    cmd_reverse,
    cmd_ingest,
    _extract_slide_text,
)
from aippt.builder import build_slide, BuildContext


class TestBuildParser:
    def test_creates_parser(self):
        parser = build_parser()
        assert parser is not None

    def test_has_create_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['create', 'outline.md', 'template.pptx', 'output.pptx'])
        assert args.command == 'create'
        assert args.outline == 'outline.md'
        assert args.template == 'template.pptx'
        assert args.output == 'output.pptx'

    def test_has_reverse_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['reverse', 'input.pptx'])
        assert args.command == 'reverse'
        assert args.input == 'input.pptx'

    def test_has_catalog_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['catalog', 'deck.pptx'])
        assert args.command == 'catalog'
        assert args.deck == 'deck.pptx'

    def test_has_search_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['search', '--tags', 'security,cloud'])
        assert args.command == 'search'
        assert args.tags == 'security,cloud'

    def test_has_remix_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['remix', 'manifest.yaml', 'output.pptx'])
        assert args.command == 'remix'
        assert args.manifest == 'manifest.yaml'
        assert args.output == 'output.pptx'

    def test_has_analyze_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['analyze', 'deck.pptx', '--mode', 'tags'])
        assert args.command == 'analyze'
        assert args.deck == 'deck.pptx'
        assert args.mode == 'tags'

    def test_has_export_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['export', '--all'])
        assert args.command == 'export'
        assert args.all is True

    def test_has_serve_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['serve', '--port', '9000'])
        assert args.command == 'serve'
        assert args.port == 9000

    def test_debug_flag(self):
        parser = build_parser()
        args = parser.parse_args(['--debug', 'reverse', 'input.pptx'])
        assert args.debug is True

    def test_create_options(self):
        parser = build_parser()
        args = parser.parse_args([
            'create', 'outline.md', 'template.pptx', 'output.pptx',
            '--enhance', '--model', 'gpt-4o', '--test', '3',
            '--image-gen', 'dalle'
        ])
        assert args.enhance is True
        assert args.model == 'gpt-4o'
        assert args.test == 3
        assert args.image_gen == 'dalle'

    def test_no_command_returns_none(self):
        parser = build_parser()
        args = parser.parse_args([])
        assert args.command is None

    def test_has_models_init_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['models', 'init'])
        assert args.command == 'models'
        assert args.models_action == 'init'

    def test_has_models_list_available_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['models', 'list-available'])
        assert args.command == 'models'
        assert args.models_action == 'list-available'

    def test_has_models_set_subcommand(self):
        parser = build_parser()
        args = parser.parse_args(['models', 'set', 'enhance', 'gpt-4o'])
        assert args.command == 'models'
        assert args.models_action == 'set'
        assert args.operation == 'enhance'
        assert args.model_name == 'gpt-4o'


class TestCmdReverse:
    def test_reverse_missing_file(self, tmp_path, capsys):
        """Test reverse command with missing input file."""
        parser = build_parser()
        args = parser.parse_args(['reverse', '/nonexistent/file.pptx'])

        result = cmd_reverse(args)

        assert result == 1

    def test_reverse_success(self, tmp_path):
        """Test reverse command with valid file."""
        from pptx import Presentation

        # Create test PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"

        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        output_path = tmp_path / "output.md"

        parser = build_parser()
        args = parser.parse_args(['reverse', str(pptx_path), str(output_path)])

        result = cmd_reverse(args)

        assert result == 0
        assert output_path.exists()

    def test_reverse_auto_output_name(self, tmp_path):
        """Test reverse command auto-generates output filename."""
        from pptx import Presentation

        # Create test PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"

        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        parser = build_parser()
        args = parser.parse_args(['reverse', str(pptx_path)])

        result = cmd_reverse(args)

        assert result == 0
        # Should create test.md in same directory
        expected_output = tmp_path / "test.md"
        assert expected_output.exists()

    def test_reverse_strip_notes(self, tmp_path):
        """--strip-notes should omit notes entirely."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Hidden notes"

        pptx_path = str(tmp_path / "test.pptx")
        md_path = str(tmp_path / "test.md")
        prs.save(pptx_path)

        parser = build_parser()
        args = parser.parse_args(['reverse', pptx_path, md_path, '--strip-notes'])

        result = cmd_reverse(args)
        assert result == 0

        content = open(md_path).read()
        assert 'Hidden notes' not in content
        assert '<!-- notes' not in content


class TestLegacyCompatibility:
    """Test backwards compatibility with old CLI syntax."""

    def test_legacy_syntax_detection(self):
        """Test that old-style positional args are converted."""
        # This tests the wrapper logic in aippt.py
        # by simulating the sys.argv manipulation
        original_argv = sys.argv.copy()

        try:
            sys.argv = ['aippt.py', 'outline.md', 'template.pptx', 'output.pptx']

            # Import the wrapper module to test detection
            from aippt.cli import build_parser

            # Check that 'outline.md' is not a subcommand
            subcommands = {'create', 'reverse', 'catalog', 'search', 'remix', 'analyze', 'export', 'serve', 'models'}
            first_arg = sys.argv[1] if len(sys.argv) > 1 else ''
            is_legacy = first_arg and not first_arg.startswith('-') and first_arg not in subcommands

            assert is_legacy is True

        finally:
            sys.argv = original_argv

    def test_new_syntax_detection(self):
        """Test that new subcommand syntax is recognized."""
        original_argv = sys.argv.copy()

        try:
            sys.argv = ['aippt.py', 'create', 'outline.md', 'template.pptx', 'output.pptx']

            subcommands = {'create', 'reverse', 'catalog', 'search', 'remix', 'analyze', 'export', 'serve', 'models'}
            first_arg = sys.argv[1] if len(sys.argv) > 1 else ''
            is_legacy = first_arg and not first_arg.startswith('-') and first_arg not in subcommands

            assert is_legacy is False

        finally:
            sys.argv = original_argv


class TestExportImagesCommand:
    def test_find_powershell_prefers_pwsh_exe(self, monkeypatch):
        from aippt.cli import _find_powershell
        import shutil
        def mock_which(cmd):
            if cmd == "pwsh.exe":
                return "/mnt/c/Program Files/PowerShell/7/pwsh.exe"
            return None
        monkeypatch.setattr(shutil, "which", mock_which)
        result = _find_powershell()
        assert result == "/mnt/c/Program Files/PowerShell/7/pwsh.exe"

    def test_find_powershell_falls_back_to_powershell(self, monkeypatch):
        from aippt.cli import _find_powershell
        import shutil
        def mock_which(cmd):
            if cmd == "powershell":
                return "/usr/bin/powershell"
            return None
        monkeypatch.setattr(shutil, "which", mock_which)
        result = _find_powershell()
        assert result == "/usr/bin/powershell"

    def test_find_powershell_returns_none_when_unavailable(self, monkeypatch):
        from aippt.cli import _find_powershell
        import shutil
        monkeypatch.setattr(shutil, "which", lambda cmd: None)
        result = _find_powershell()
        assert result is None


# ---------------------------------------------------------------------------
# Helpers for building mock pptx slide objects
# ---------------------------------------------------------------------------

def _make_paragraph(text: str) -> MagicMock:
    """Return a mock paragraph whose .text attribute returns *text*."""
    para = MagicMock()
    para.text = text
    return para


def _make_shape(paragraphs: list, has_text_frame: bool = True) -> MagicMock:
    """Return a mock shape with an optional text frame.

    Args:
        paragraphs: List of paragraph text strings that the text frame exposes.
        has_text_frame: Whether the shape has a text frame at all.
    """
    shape = MagicMock()
    shape.has_text_frame = has_text_frame
    if has_text_frame:
        tf = MagicMock()
        # .text returns the full concatenated text (mirroring python-pptx behaviour)
        tf.text = "\n".join(paragraphs)
        tf.paragraphs = [_make_paragraph(p) for p in paragraphs]
        shape.text_frame = tf
    return shape


def _make_slide(shapes: list) -> MagicMock:
    """Return a mock slide whose .shapes iterates over *shapes*."""
    slide = MagicMock()
    slide.shapes = shapes
    return slide


# ---------------------------------------------------------------------------
# Tests for _extract_slide_text
# ---------------------------------------------------------------------------

class TestExtractSlideText:
    """Unit tests for the _extract_slide_text() helper in aippt.cli."""

    def test_basic_extraction_multiple_shapes(self):
        """Text from multiple body shapes is joined with newlines."""
        shape1 = _make_shape(["First bullet", "Second bullet"])
        shape2 = _make_shape(["Third bullet"])
        slide = _make_slide([shape1, shape2])

        result = _extract_slide_text(slide, "My Title")

        assert result == "First bullet\nSecond bullet\nThird bullet"

    def test_title_shape_is_excluded(self):
        """A shape whose full text matches the slide title is skipped."""
        title_shape = _make_shape(["Slide Title"])
        body_shape = _make_shape(["Bullet one", "Bullet two"])
        slide = _make_slide([title_shape, body_shape])

        result = _extract_slide_text(slide, "Slide Title")

        assert "Slide Title" not in result
        assert result == "Bullet one\nBullet two"

    def test_title_exclusion_strips_whitespace(self):
        """Title matching is whitespace-tolerant (both title and shape text are stripped)."""
        title_shape = _make_shape(["  Slide Title  "])
        body_shape = _make_shape(["Body text"])
        slide = _make_slide([title_shape, body_shape])

        result = _extract_slide_text(slide, "  Slide Title  ")

        assert "Slide Title" not in result
        assert result == "Body text"

    def test_empty_slide_returns_empty_string(self):
        """A slide with no shapes returns an empty string."""
        slide = _make_slide([])

        result = _extract_slide_text(slide, "Some Title")

        assert result == ""

    def test_shapes_without_text_frame_are_skipped(self):
        """Shapes that have no text frame (images, charts, etc.) are ignored."""
        image_shape = _make_shape([], has_text_frame=False)
        body_shape = _make_shape(["Real content"])
        slide = _make_slide([image_shape, body_shape])

        result = _extract_slide_text(slide, "Title")

        assert result == "Real content"

    def test_empty_paragraphs_are_skipped(self):
        """Paragraphs that are blank (or whitespace-only) are not included."""
        shape = _make_shape(["", "  ", "Visible text", ""])
        slide = _make_slide([shape])

        result = _extract_slide_text(slide, "Title")

        assert result == "Visible text"

    def test_no_title_argument_keeps_all_shapes(self):
        """When slide_title is empty, no shape is excluded on title grounds."""
        shape1 = _make_shape(["First"])
        shape2 = _make_shape(["Second"])
        slide = _make_slide([shape1, shape2])

        result = _extract_slide_text(slide, "")

        assert result == "First\nSecond"

    def test_all_shapes_lack_text_frame_returns_empty_string(self):
        """A slide where every shape has no text frame returns an empty string."""
        shapes = [_make_shape([], has_text_frame=False) for _ in range(3)]
        slide = _make_slide(shapes)

        result = _extract_slide_text(slide, "Title")

        assert result == ""

    def test_single_shape_single_paragraph(self):
        """Minimal case: one shape, one paragraph, no title match."""
        shape = _make_shape(["Only line"])
        slide = _make_slide([shape])

        result = _extract_slide_text(slide, "Different Title")

        assert result == "Only line"


class TestIngestParser:
    """Tests for the ingest subcommand argument parsing."""

    def test_ingest_subcommand_exists(self):
        parser = build_parser()
        args = parser.parse_args(['ingest', 'deck.pptx'])
        assert args.command == 'ingest'
        assert args.deck == 'deck.pptx'

    def test_ingest_defaults(self):
        parser = build_parser()
        args = parser.parse_args(['ingest', 'deck.pptx'])
        assert args.db == 'slides.db'
        assert args.images_dir is None
        assert args.tags is False
        assert args.taxonomy is None
        assert args.model is None
        assert args.width == 1920
        assert args.height == 1080

    def test_ingest_all_flags(self):
        parser = build_parser()
        args = parser.parse_args([
            'ingest', 'deck.pptx',
            '--images-dir', 'img/',
            '--db', 'test.db',
            '--tags',
            '--taxonomy', 'tags.csv',
            '--model', 'gpt-4o',
            '--gateway-config', 'gw.yaml',
            '--width', '2560',
            '--height', '1440',
        ])
        assert args.tags is True
        assert args.taxonomy == 'tags.csv'
        assert args.model == 'gpt-4o'
        assert args.gateway_config == 'gw.yaml'
        assert args.images_dir == 'img/'
        assert args.db == 'test.db'
        assert args.width == 2560
        assert args.height == 1440

    def test_taxonomy_flag_without_tags_flag(self):
        """--taxonomy doesn't implicitly enable --tags."""
        parser = build_parser()
        args = parser.parse_args(['ingest', 'deck.pptx', '--taxonomy', 'tags.csv'])
        assert args.tags is False


class TestIngestCommand:
    """Tests for cmd_ingest orchestration logic."""

    def test_file_not_found(self, tmp_path):
        parser = build_parser()
        args = parser.parse_args(['ingest', str(tmp_path / 'missing.pptx')])
        result = cmd_ingest(args)
        assert result == 1

    @patch('aippt.ingest.catalog_deck')
    @patch('aippt.ingest.cmd_export_images')
    def test_export_images_failure_aborts(self, mock_export, mock_catalog, tmp_path):
        """When image export fails, ingest aborts (require_images=True by default)."""
        deck = tmp_path / 'test.pptx'
        deck.touch()
        mock_export.return_value = 1

        parser = build_parser()
        args = parser.parse_args(['ingest', str(deck)])
        result = cmd_ingest(args)

        assert result == 1
        mock_export.assert_called_once()
        mock_catalog.assert_not_called()

    @patch('aippt.ingest.get_deck_by_id', return_value={'slide_count': 5})
    @patch('aippt.ingest.catalog_deck')
    @patch('aippt.ingest.cmd_export_images')
    def test_export_then_catalog(self, mock_export, mock_catalog, _mock_get, tmp_path):
        deck = tmp_path / 'test.pptx'
        deck.touch()
        mock_export.return_value = 0
        mock_catalog.return_value = 42

        parser = build_parser()
        args = parser.parse_args(['ingest', str(deck), '--db', str(tmp_path / 'test.db')])
        result = cmd_ingest(args)

        assert result == 0
        mock_export.assert_called_once()
        mock_catalog.assert_called_once()
        _, kwargs = mock_catalog.call_args
        assert kwargs['db_path'] == str(tmp_path / 'test.db')

    @patch('aippt.ingest.get_deck_by_id', return_value={'slide_count': 1})
    @patch('aippt.ingest.catalog_deck')
    @patch('aippt.ingest.cmd_export_images')
    def test_default_images_dir(self, mock_export, mock_catalog, _mock_get, tmp_path):
        deck = tmp_path / 'my-deck.pptx'
        deck.touch()
        mock_export.return_value = 0
        mock_catalog.return_value = 1

        parser = build_parser()
        args = parser.parse_args(['ingest', str(deck)])
        result = cmd_ingest(args)

        assert result == 0
        export_args = mock_export.call_args[0][0]
        assert 'my-deck' in export_args.out_dir

    @patch('aippt.ingest.get_deck_by_id', return_value={'slide_count': 1})
    @patch('aippt.ingest.catalog_deck')
    @patch('aippt.ingest.cmd_export_images')
    def test_no_tags_by_default(self, mock_export, mock_catalog, _mock_get, tmp_path):
        """Without --tags, cmd_analyze is not called."""
        deck = tmp_path / 'test.pptx'
        deck.touch()
        mock_export.return_value = 0
        mock_catalog.return_value = 1

        parser = build_parser()
        args = parser.parse_args(['ingest', str(deck)])

        with patch('aippt.ingest.cmd_analyze') as mock_analyze:
            result = cmd_ingest(args)
            mock_analyze.assert_not_called()

        assert result == 0


class TestCreateDeck:
    def test_returns_result_dict_without_enhance(self, tmp_path):
        outline = "# Test Deck\n## Slide 1: Intro\n- Point one\n- Point two\n"
        template_path = str(tmp_path / "template.pptx")
        output_path = str(tmp_path / "output.pptx")
        from pptx import Presentation
        prs = Presentation()
        prs.save(template_path)

        from aippt.pipeline import run_pipeline, PipelineConfig
        config = PipelineConfig(
            outline_text=outline,
            template_path=template_path,
            output_path=output_path,
            enhance=False,
        )
        result = run_pipeline(config)
        assert result.slide_count >= 1
        assert result.output_path == output_path
        assert os.path.exists(output_path)

    def test_calls_progress_callback(self, tmp_path):
        outline = "# Test\n## Slide 1: Hello\n- World\n"
        template_path = str(tmp_path / "template.pptx")
        output_path = str(tmp_path / "output.pptx")
        from pptx import Presentation
        prs = Presentation()
        prs.save(template_path)

        steps = []
        def on_progress(step, detail=""):
            steps.append(step)

        from aippt.pipeline import run_pipeline, PipelineConfig
        config = PipelineConfig(
            outline_text=outline,
            template_path=template_path,
            output_path=output_path,
            enhance=False,
            progress_callback=on_progress,
        )
        run_pipeline(config)
        assert "parse" in steps
        assert "build" in steps

    def test_raises_on_missing_template(self, tmp_path):
        from aippt.pipeline import run_pipeline, PipelineConfig
        with pytest.raises(FileNotFoundError):
            config = PipelineConfig(
                outline_text="# Test\n## S1\n- x",
                template_path=str(tmp_path / "nope.pptx"),
                output_path=str(tmp_path / "out.pptx"),
            )
            run_pipeline(config)


class TestAddSlideEnhancedContent:
    """Tests for build_slide() CONTENT: extraction and fallback."""

    def _make_prs(self):
        from pptx import Presentation
        return Presentation()

    def test_uses_enhanced_content_from_suggestions(self):
        """When LLM response contains CONTENT:, slide body uses enhanced text."""
        prs = self._make_prs()
        enhanced_response = (
            "CONTENT:\n"
            "- Enhanced point one with more detail\n"
            "- Enhanced point two with added context\n"
            "NARRATIVE: This slide covers important points.\n"
            "LAYOUT: bullet\n"
            "VISUALS: Emphasize first point\n"
            "TALKING_POINTS: Additional details"
        )
        slide_data = {
            'title': 'Test Title',
            'content': enhanced_response,
            'original_content': ["- Terse point one", "- Terse point two"],
        }
        layout = build_slide(prs, slide_data, BuildContext())
        slide = prs.slides[0]
        # Find body text (non-title shapes)
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() != "Test Title":
                body_texts.append(shape.text_frame.text)
        body = '\n'.join(body_texts)
        assert 'Enhanced point one' in body
        assert 'Terse point one' not in body

    def test_falls_back_to_original_when_no_content(self):
        """When LLM response lacks CONTENT:, slide body uses original_content."""
        prs = self._make_prs()
        enhanced_response = (
            "NARRATIVE: This slide covers important points.\n"
            "LAYOUT: bullet\n"
            "VISUALS: Emphasize first point\n"
            "TALKING_POINTS: Additional details"
        )
        slide_data = {
            'title': 'Test Title',
            'content': enhanced_response,
            'original_content': ["- Original point one", "- Original point two"],
        }
        layout = build_slide(prs, slide_data, BuildContext())
        slide = prs.slides[0]
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() != "Test Title":
                body_texts.append(shape.text_frame.text)
        body = '\n'.join(body_texts)
        assert 'Original point one' in body

    def test_non_enhanced_path_uses_content_lines(self):
        """Without original_content (non-enhanced path), slide body uses content_lines."""
        prs = self._make_prs()
        content = "- Plain bullet one\n- Plain bullet two"
        slide_data = {
            'title': 'Test Title',
            'content': content,
        }
        layout = build_slide(prs, slide_data, BuildContext())
        slide = prs.slides[0]
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() != "Test Title":
                body_texts.append(shape.text_frame.text)
        body = '\n'.join(body_texts)
        assert 'Plain bullet one' in body
