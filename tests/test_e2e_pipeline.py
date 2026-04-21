"""End-to-end pipeline tests: outline -> deck -> images -> catalog -> AI analysis.

These tests make REAL LLM API calls and optionally invoke PowerShell for image export.
They validate the full user workflow, not individual modules.

Run with:
    AMD_LLM_KEY=<key> python -m pytest tests/test_e2e_pipeline.py -m e2e -v -s

Requirements:
    - AMD_LLM_KEY environment variable set (for LLM calls via gateway)
    - models.yaml and gateway.yaml in project root (real config, not test fixtures)
    - Optional: pwsh.exe accessible from WSL (for image export step)
    - Optional: Microsoft PowerPoint installed on Windows side (for image export)
"""

import os
import re
import shutil
import struct
import zlib
from pathlib import Path

import pytest
from pptx import Presentation

# ---------------------------------------------------------------------------
# Skip logic
# ---------------------------------------------------------------------------

SKIP_NO_KEY = not os.environ.get("AMD_LLM_KEY")
SKIP_KEY_REASON = "AMD_LLM_KEY environment variable not set"

PROJECT_ROOT = Path(__file__).parent.parent
OUTLINES_DIR = Path(__file__).parent / "e2e_outlines"
MODELS_YAML = PROJECT_ROOT / "models.yaml"
GATEWAY_YAML = PROJECT_ROOT / "gateway.yaml"


def _pwsh_available():
    """Check if any PowerShell executable is accessible."""
    from aippt.cli import _find_powershell
    return _find_powershell() is not None


HAS_PWSH = _pwsh_available() if not SKIP_NO_KEY else False
SKIP_NO_PWSH_REASON = "PowerShell (pwsh.exe) not available from this environment"

# Deck definitions: (name, outline_file, expected_slide_count)
DECK_SPECS = [
    ("tech_overview", "tech_overview.md", 3),
    ("mini_presentation", "mini_presentation.md", 5),
    ("generic_info", "generic_info.md", 4),
]


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(scope="module")
def monkeypatch_module():
    """Module-scoped monkeypatch (pytest's monkeypatch is function-scoped)."""
    from _pytest.monkeypatch import MonkeyPatch
    mp = MonkeyPatch()
    yield mp
    mp.undo()


@pytest.fixture(scope="module")
def e2e_workspace(tmp_path_factory):
    """Create a workspace directory with subdirs for the entire E2E module.

    Structure:
        workspace/
            outlines/       -- copies of test outlines
            decks/          -- generated PPTX files
            images/         -- exported slide images (per deck)
            exports/        -- feedback and analysis exports
            slides.db       -- shared catalog database
    """
    workspace = tmp_path_factory.mktemp("e2e")
    for subdir in ("outlines", "decks", "images", "exports"):
        (workspace / subdir).mkdir()

    # Copy outline files into workspace
    for _, outline_file, _ in DECK_SPECS:
        src = OUTLINES_DIR / outline_file
        dst = workspace / "outlines" / outline_file
        shutil.copy2(str(src), str(dst))

    return workspace


@pytest.fixture(scope="module")
def db_path(e2e_workspace):
    """Path to the shared E2E database."""
    return str(e2e_workspace / "slides.db")


@pytest.fixture(autouse=True)
def patch_default_config_path(real_models_yaml):
    """Override the conftest autouse fixture that redirects to tmp_path.

    The conftest version (function-scoped, autouse) would overwrite our
    module-scoped real_models_yaml patch on every test.  This local override
    simply yields, letting the module-scoped patch remain in effect.
    """
    yield


@pytest.fixture(scope="module")
def real_models_yaml(monkeypatch_module):
    """Point config to the real models.yaml in the project root.

    Unlike unit tests (which patch to tmp_path), E2E tests use the actual
    models.yaml so we validate the real configuration.
    """
    import aippt.config as cfg_module
    monkeypatch_module.setattr(cfg_module, "DEFAULT_CONFIG_PATH", str(MODELS_YAML))
    return str(MODELS_YAML)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _get_default_template():
    """Get the path to the default python-pptx template."""
    import pptx
    return os.path.join(os.path.dirname(pptx.__file__), "templates", "default.pptx")


def _create_placeholder_images(images_dir, slide_count):
    """Create minimal valid PNG files as placeholders when PowerShell is unavailable.

    These are 100x75 pixel white PNGs -- enough for the LLM vision API to process
    (it will see a blank slide, which is fine for testing the pipeline).
    """
    def _minimal_png():
        width, height = 100, 75
        # IHDR
        ihdr_data = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
        ihdr = b"IHDR" + ihdr_data
        ihdr_chunk = struct.pack(">I", len(ihdr_data)) + ihdr + struct.pack(">I", zlib.crc32(ihdr) & 0xFFFFFFFF)
        # IDAT -- white pixels (RGB)
        raw_data = b""
        for _ in range(height):
            raw_data += b"\x00" + b"\xff" * (width * 3)
        compressed = zlib.compress(raw_data)
        idat = b"IDAT" + compressed
        idat_chunk = struct.pack(">I", len(compressed)) + idat + struct.pack(">I", zlib.crc32(idat) & 0xFFFFFFFF)
        # IEND
        iend = b"IEND"
        iend_chunk = struct.pack(">I", 0) + iend + struct.pack(">I", zlib.crc32(iend) & 0xFFFFFFFF)
        return b"\x89PNG\r\n\x1a\n" + ihdr_chunk + idat_chunk + iend_chunk

    os.makedirs(images_dir, exist_ok=True)
    png_data = _minimal_png()
    for i in range(1, slide_count + 1):
        with open(os.path.join(images_dir, f"Slide{i}.PNG"), "wb") as f:
            f.write(png_data)


def _count_h1_headers(md_path):
    """Count H1 (# ) headers in a markdown file."""
    with open(md_path, encoding="utf-8") as f:
        return sum(1 for line in f if line.startswith("# "))


# ---------------------------------------------------------------------------
# Capability matrix for feedback-to-code analysis
# ---------------------------------------------------------------------------

CAPABILITY_MATRIX = {
    "font": {"level": "full", "description": "Font size, family, style changes", "api": "paragraph.font.size = Pt(N)"},
    "color": {"level": "partial", "description": "Individual color changes (no theme swap)", "api": "font.color.rgb = RGBColor(...)"},
    "text": {"level": "full", "description": "Text content changes", "api": "shape.text_frame.text = '...'"},
    "bold": {"level": "full", "description": "Bold/italic/underline", "api": "run.font.bold = True"},
    "bullet": {"level": "full", "description": "Bullet level and formatting", "api": "paragraph.level = N"},
    "notes": {"level": "full", "description": "Speaker notes", "api": "slide.notes_slide.notes_text_frame.text"},
    "position": {"level": "full", "description": "Shape position and size", "api": "shape.left = Inches(N)"},
    "shape": {"level": "full", "description": "Add/remove shapes", "api": "slide.shapes.add_textbox(...)"},
    "background": {"level": "full", "description": "Slide background color", "api": "slide.background.fill.solid()"},
    "table": {"level": "full", "description": "Table cell text and formatting", "api": "table.cell(r,c).text = '...'"},
    "layout": {"level": "partial", "description": "Layout selection (template-dependent)", "api": "slide_layouts[N]"},
    "image": {"level": "partial", "description": "Image add/replace (not edit)", "api": "slide.shapes.add_picture(...)"},
    "split": {"level": "partial", "description": "Slide splitting/reordering (XML)", "api": "XML manipulation"},
    "diagram": {"level": "none", "description": "SmartArt / complex diagrams", "api": "Not supported"},
    "animation": {"level": "none", "description": "Animations and transitions", "api": "Not supported"},
    "chart": {"level": "limited", "description": "Chart data (limited styling)", "api": "chart.series[0].values = [...]"},
    "whitespace": {"level": "full", "description": "Margin and spacing adjustments", "api": "text_frame.margin_left = Inches(N)"},
}

CAPABILITY_KEYWORDS = {
    "font": ["font", "text size", "pt ", "point", "typeface", "serif", "sans-serif"],
    "color": ["color", "colour", "palette", "contrast", "rgb", "hex"],
    "text": ["text", "wording", "label", "title", "heading", "content", "rewrite", "rephrase"],
    "bold": ["bold", "italic", "underline", "emphasis", "weight"],
    "bullet": ["bullet", "indent", "list", "numbering", "hierarchy"],
    "notes": ["notes", "speaker", "talking point"],
    "position": ["position", "align", "center", "margin", "spacing", "move", "resize", "placement"],
    "shape": ["shape", "textbox", "box", "rectangle", "callout", "add a "],
    "background": ["background"],
    "table": ["table", "cell", "row", "column", "grid"],
    "layout": ["layout", "template", "slide design"],
    "image": ["image", "picture", "photo", "icon", "logo", "graphic", "visual"],
    "split": ["split", "divide", "separate", "two slides", "multiple slides", "break up"],
    "diagram": ["diagram", "smartart", "flowchart", "org chart", "process flow"],
    "animation": ["animation", "transition", "fade", "appear", "motion"],
    "chart": ["chart", "graph", "pie", "bar chart", "line chart", "data visualization"],
    "whitespace": ["whitespace", "white space", "spacing", "breathing room", "clutter", "dense", "crowded"],
}


def classify_feedback_item(text):
    """Classify a single feedback bullet point against the capability matrix.

    Returns: (category, level) or ("unknown", "unknown") if no match.
    """
    text_lower = text.lower()
    for category, keywords in CAPABILITY_KEYWORDS.items():
        for keyword in keywords:
            if keyword in text_lower:
                return category, CAPABILITY_MATRIX[category]["level"]
    return "unknown", "unknown"


def analyze_improvements_feedback(improvements_text):
    """Parse improvements markdown and classify each bullet point.

    Returns a list of dicts:
        [{"section": "...", "item": "...", "category": "...", "level": "..."}, ...]
    """
    results = []
    current_section = "General"

    for line in improvements_text.split("\n"):
        line = line.strip()
        if line.startswith("## "):
            current_section = line[3:].strip()
            continue
        if line.startswith("- ") or line.startswith("* ") or re.match(r"^\d+\.", line):
            item_text = re.sub(r"^[-*]\s*|\d+\.\s*", "", line).strip()
            if item_text:
                category, level = classify_feedback_item(item_text)
                results.append({
                    "section": current_section,
                    "item": item_text,
                    "category": category,
                    "level": level,
                })

    return results


# ---------------------------------------------------------------------------
# Pipeline tests
# ---------------------------------------------------------------------------

@pytest.mark.e2e
@pytest.mark.skipif(SKIP_NO_KEY, reason=SKIP_KEY_REASON)
class TestE2EPipeline:
    """Full pipeline test: outline -> deck -> images -> catalog -> analyze -> export.

    Uses module-scoped fixtures so state (database, files) persists across
    test methods within this class. Methods are ordered and must run sequentially.
    """

    # Store results across test methods (class-level state)
    _deck_paths = {}
    _image_dirs = {}
    _feedback_results = {}
    _improvements_results = {}

    # -- Step 1: Create decks from outlines ----------------------------------

    @pytest.mark.parametrize("deck_name, outline_file, expected_slides", DECK_SPECS)
    def test_01_create_deck(self, deck_name, outline_file, expected_slides,
                            e2e_workspace, real_models_yaml):
        """Create a PPTX deck from a markdown outline (no LLM enhancement)."""
        from aippt.cli import build_parser, cmd_create

        outline_path = str(e2e_workspace / "outlines" / outline_file)
        output_path = str(e2e_workspace / "decks" / f"{deck_name}.pptx")
        template_path = _get_default_template()

        # Verify outline has expected slide count
        assert _count_h1_headers(outline_path) == expected_slides, \
            f"{outline_file} should have {expected_slides} H1 headers"

        # Build args and run create
        parser = build_parser()
        args = parser.parse_args([
            "create", outline_path, template_path, output_path,
        ])
        result = cmd_create(args)
        assert result == 0, f"cmd_create failed for {deck_name}"

        # Verify output
        assert os.path.exists(output_path), f"Output PPTX not created: {output_path}"
        prs = Presentation(output_path)
        assert len(prs.slides) == expected_slides, \
            f"Expected {expected_slides} slides, got {len(prs.slides)}"

        # Store for later tests
        TestE2EPipeline._deck_paths[deck_name] = output_path
        print(f"  Created {deck_name}: {len(prs.slides)} slides")

    # -- Step 2: Export images -----------------------------------------------

    @pytest.mark.parametrize("deck_name, _, expected_slides", DECK_SPECS)
    def test_02_export_images(self, deck_name, _, expected_slides,
                              e2e_workspace, real_models_yaml):
        """Export slide images via PowerShell, or create placeholders."""
        deck_path = TestE2EPipeline._deck_paths.get(deck_name)
        assert deck_path, f"Deck not created yet: {deck_name} (run test_01 first)"

        images_dir = str(e2e_workspace / "images" / deck_name)

        if HAS_PWSH:
            from aippt.cli import build_parser, cmd_export_images
            parser = build_parser()
            args = parser.parse_args([
                "export-images", deck_path, images_dir,
                "--width", "1920", "--height", "1080",
            ])
            result = cmd_export_images(args)
            if result != 0:
                print(f"  PowerShell export failed (rc={result}), falling back to placeholders")
                _create_placeholder_images(images_dir, expected_slides)
        else:
            print(f"  PowerShell not available, creating placeholder images")
            _create_placeholder_images(images_dir, expected_slides)

        # Verify images exist
        for i in range(1, expected_slides + 1):
            found = False
            for ext in (".PNG", ".png", ".jpg", ".jpeg"):
                if os.path.exists(os.path.join(images_dir, f"Slide{i}{ext}")):
                    found = True
                    break
            assert found, f"No image found for Slide{i} in {images_dir}"

        TestE2EPipeline._image_dirs[deck_name] = images_dir
        print(f"  Images for {deck_name}: {expected_slides} slides")

    # -- Step 3: Catalog decks -----------------------------------------------

    @pytest.mark.parametrize("deck_name, _, expected_slides", DECK_SPECS)
    def test_03_catalog_deck(self, deck_name, _, expected_slides,
                             e2e_workspace, db_path, real_models_yaml):
        """Catalog each deck into the SQLite database."""
        from aippt.catalog import catalog_deck, get_db

        deck_path = TestE2EPipeline._deck_paths[deck_name]
        images_dir = TestE2EPipeline._image_dirs[deck_name]

        deck_id = catalog_deck(deck_path, db_path=db_path, images_dir=images_dir)
        assert deck_id > 0, f"catalog_deck returned invalid id: {deck_id}"

        # Verify slide count in DB
        conn = get_db(db_path)
        slides = conn.execute(
            "SELECT COUNT(*) as cnt FROM slides WHERE deck_id = ?", (deck_id,)
        ).fetchone()
        conn.close()
        assert slides["cnt"] == expected_slides, \
            f"Expected {expected_slides} slides in DB for {deck_name}, got {slides['cnt']}"

        print(f"  Cataloged {deck_name}: deck_id={deck_id}, {expected_slides} slides")

    # -- Step 4: Generate tags (real LLM) ------------------------------------

    @pytest.mark.parametrize("deck_name, _, expected_slides", DECK_SPECS)
    def test_04_analyze_tags(self, deck_name, _, expected_slides,
                             e2e_workspace, db_path, real_models_yaml):
        """Generate AI tags for each slide via real LLM call."""
        from aippt.catalog import get_db, get_slide_tags

        deck_path = TestE2EPipeline._deck_paths[deck_name]
        images_dir = TestE2EPipeline._image_dirs[deck_name]

        from aippt.cli import build_parser, cmd_analyze
        parser = build_parser()
        args = parser.parse_args([
            "analyze", deck_path,
            "--mode", "tags",
            "--images-dir", images_dir,
            "--db", db_path,
            "--gateway-config", str(GATEWAY_YAML),
        ])
        result = cmd_analyze(args)
        assert result == 0, f"cmd_analyze --mode tags failed for {deck_name}"

        # Verify tags were added
        conn = get_db(db_path)
        deck_row = conn.execute(
            "SELECT id FROM decks WHERE file_path = ?",
            (os.path.abspath(deck_path),)
        ).fetchone()
        slides = conn.execute(
            "SELECT id, position, title FROM slides WHERE deck_id = ? ORDER BY position",
            (deck_row["id"],)
        ).fetchall()
        conn.close()

        for slide in slides:
            tags = get_slide_tags(slide["id"], db_path)
            assert len(tags) >= 1, \
                f"Slide {slide['position']} ({slide['title']}) has no tags"
            print(f"  {deck_name} slide {slide['position']}: tags={tags}")

    # -- Step 5: Generate feedback (real LLM) --------------------------------

    @pytest.mark.parametrize("deck_name, _, expected_slides", DECK_SPECS)
    def test_05_analyze_feedback(self, deck_name, _, expected_slides,
                                 e2e_workspace, db_path, real_models_yaml, capsys):
        """Generate AI feedback for each slide via real LLM call."""
        deck_path = TestE2EPipeline._deck_paths[deck_name]
        images_dir = TestE2EPipeline._image_dirs[deck_name]

        from aippt.cli import build_parser, cmd_analyze
        parser = build_parser()
        args = parser.parse_args([
            "analyze", deck_path,
            "--mode", "feedback",
            "--images-dir", images_dir,
            "--db", db_path,
            "--gateway-config", str(GATEWAY_YAML),
        ])
        result = cmd_analyze(args)
        assert result == 0, f"cmd_analyze --mode feedback failed for {deck_name}"

        captured = capsys.readouterr()
        assert len(captured.out) > 50, \
            f"Feedback output too short for {deck_name}: {len(captured.out)} chars"

        TestE2EPipeline._feedback_results[deck_name] = captured.out
        print(f"  {deck_name} feedback: {len(captured.out)} chars captured")

    # -- Step 6: Generate notes (real LLM) -----------------------------------

    @pytest.mark.parametrize("deck_name, _, expected_slides", DECK_SPECS)
    def test_06_analyze_notes(self, deck_name, _, expected_slides,
                              e2e_workspace, db_path, real_models_yaml):
        """Generate AI speaker notes and save to PPTX."""
        deck_path = TestE2EPipeline._deck_paths[deck_name]
        images_dir = TestE2EPipeline._image_dirs[deck_name]

        from aippt.cli import build_parser, cmd_analyze
        parser = build_parser()
        args = parser.parse_args([
            "analyze", deck_path,
            "--mode", "notes",
            "--images-dir", images_dir,
            "--db", db_path,
            "--gateway-config", str(GATEWAY_YAML),
        ])
        result = cmd_analyze(args)
        assert result == 0, f"cmd_analyze --mode notes failed for {deck_name}"

        # Verify notes were written to the PPTX
        prs = Presentation(deck_path)
        notes_found = 0
        for i, slide in enumerate(prs.slides, 1):
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    notes_found += 1
                    preview = notes_text[:80].replace("\n", " ")
                    print(f"  {deck_name} slide {i} notes: {preview}...")

        assert notes_found >= 1, \
            f"No speaker notes found in {deck_name} PPTX after notes generation"
        print(f"  {deck_name}: {notes_found}/{len(prs.slides)} slides have notes")

    # -- Step 7: Generate improvements (real LLM) ----------------------------

    @pytest.mark.parametrize("deck_name, _, expected_slides", DECK_SPECS)
    def test_07_analyze_improvements(self, deck_name, _, expected_slides,
                                     e2e_workspace, db_path, real_models_yaml,
                                     capsys):
        """Generate structured improvement suggestions via real LLM call."""
        deck_path = TestE2EPipeline._deck_paths[deck_name]
        images_dir = TestE2EPipeline._image_dirs[deck_name]

        from aippt.cli import build_parser, cmd_analyze
        parser = build_parser()
        args = parser.parse_args([
            "analyze", deck_path,
            "--mode", "improvements",
            "--images-dir", images_dir,
            "--db", db_path,
            "--gateway-config", str(GATEWAY_YAML),
        ])
        result = cmd_analyze(args)
        assert result == 0, f"cmd_analyze --mode improvements failed for {deck_name}"

        captured = capsys.readouterr()
        assert len(captured.out) > 100, \
            f"Improvements output too short for {deck_name}: {len(captured.out)} chars"

        # Verify expected sections are present in the output
        output_lower = captured.out.lower()
        for section in ["visual design", "technical accuracy", "flow", "split"]:
            assert section in output_lower, \
                f"Missing '{section}' section in improvements for {deck_name}"

        TestE2EPipeline._improvements_results[deck_name] = captured.out
        print(f"  {deck_name} improvements: {len(captured.out)} chars captured")

    # -- Step 8: Export feedback to markdown ----------------------------------

    def test_08_export_feedback(self, e2e_workspace, real_models_yaml):
        """Export collected feedback results to markdown files."""
        exports_dir = e2e_workspace / "exports"

        for deck_name, feedback_text in TestE2EPipeline._feedback_results.items():
            output_path = exports_dir / f"{deck_name}-feedback.md"
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# Feedback: {deck_name}\n\n")
                f.write(f"Generated by E2E pipeline test\n\n")
                f.write(feedback_text)

            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 100
            print(f"  Exported feedback: {output_path}")

        assert len(TestE2EPipeline._feedback_results) == len(DECK_SPECS), \
            "Not all decks have feedback results"

    # -- Step 9: Capability analysis -----------------------------------------

    def test_09_capability_analysis(self, e2e_workspace, real_models_yaml):
        """Classify improvement feedback against python-pptx capability matrix."""
        exports_dir = e2e_workspace / "exports"

        for deck_name, improvements_text in TestE2EPipeline._improvements_results.items():
            classified = analyze_improvements_feedback(improvements_text)
            assert len(classified) > 0, \
                f"No feedback items parsed from improvements for {deck_name}"

            # Count by level
            counts = {"full": 0, "partial": 0, "limited": 0, "none": 0, "unknown": 0}
            for item in classified:
                counts[item["level"]] = counts.get(item["level"], 0) + 1

            total = len(classified)
            actionable = counts["full"] + counts["partial"] + counts["limited"]

            # Write analysis report
            output_path = exports_dir / f"{deck_name}-capability-analysis.md"
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(f"# Capability Analysis: {deck_name}\n\n")
                f.write(f"Total feedback items: {total}\n")
                f.write(f"Actionable (full/partial/limited): {actionable} ({actionable*100//total if total else 0}%)\n")
                f.write(f"Not actionable: {counts['none']}\n")
                f.write(f"Unclassified: {counts['unknown']}\n\n")

                f.write("## Detailed Classification\n\n")
                current_section = None
                for item in classified:
                    if item["section"] != current_section:
                        current_section = item["section"]
                        f.write(f"\n### {current_section}\n\n")

                    level_label = {"full": "YES", "partial": "PARTIAL", "limited": "LIMITED",
                                   "none": "NO", "unknown": "?"}
                    f.write(f"- [{level_label[item['level']]}] {item['item']}\n")
                    f.write(f"  Category: {item['category']}")
                    if item["category"] in CAPABILITY_MATRIX:
                        f.write(f" -- {CAPABILITY_MATRIX[item['category']]['description']}")
                        f.write(f" -- API: `{CAPABILITY_MATRIX[item['category']]['api']}`")
                    f.write("\n")

                f.write("\n## python-pptx Capability Reference\n\n")
                f.write("| Category | Support | API |\n")
                f.write("|----------|---------|-----|\n")
                for cat, info in CAPABILITY_MATRIX.items():
                    f.write(f"| {cat} | {info['level']} | `{info['api']}` |\n")

            assert os.path.exists(output_path)
            print(f"  {deck_name}: {actionable}/{total} items actionable via python-pptx")
            print(f"    Full: {counts['full']}, Partial: {counts['partial']}, "
                  f"None: {counts['none']}, Unknown: {counts['unknown']}")
            print(f"  Report: {output_path}")

    # -- Step 10: Print summary and export locations -------------------------

    def test_10_summary(self, e2e_workspace, db_path, real_models_yaml):
        """Print a summary of all E2E test artifacts."""
        print("\n" + "=" * 70)
        print("E2E PIPELINE TEST SUMMARY")
        print("=" * 70)

        print(f"\nWorkspace: {e2e_workspace}")
        print(f"Database: {e2e_workspace / 'slides.db'}")

        print("\nDecks created:")
        for name, path in TestE2EPipeline._deck_paths.items():
            prs = Presentation(path)
            print(f"  {name}: {len(prs.slides)} slides")

        print("\nImage directories:")
        for name, path in TestE2EPipeline._image_dirs.items():
            count = len([f for f in os.listdir(path) if f.lower().endswith((".png", ".jpg"))])
            print(f"  {name}: {count} images")

        print("\nExport files:")
        exports_dir = e2e_workspace / "exports"
        if exports_dir.exists():
            for f in sorted(exports_dir.iterdir()):
                size = f.stat().st_size
                print(f"  {f.name}: {size:,} bytes")

        # Verify database has expected content
        from aippt.catalog import get_db
        conn = get_db(db_path)
        deck_count = conn.execute("SELECT COUNT(*) FROM decks").fetchone()[0]
        slide_count = conn.execute("SELECT COUNT(*) FROM slides").fetchone()[0]
        tag_count = conn.execute("SELECT COUNT(DISTINCT tag_id) FROM slide_tags").fetchone()[0]
        conn.close()

        print(f"\nDatabase summary:")
        print(f"  Decks: {deck_count}")
        print(f"  Slides: {slide_count}")
        print(f"  Unique tags applied: {tag_count}")
        print("\n" + "=" * 70)
