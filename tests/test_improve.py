"""Tests for aippt.improve module."""

import pytest
from unittest.mock import MagicMock, patch
from pptx import Presentation
from pptx.util import Inches

from aippt.improve import (
    build_rewrite_prompt,
    parse_rewritten_content,
    parse_rewrite_response,
    parse_validation_response,
    is_generic_title,
    select_focus,
    has_converged,
    extract_slide_content,
    _try_reexport_images,
    FOCUS_GUIDANCE,
    REWRITE_SYSTEM_PROMPT,
    AUDIENCE_REWRITE_PROMPTS,
    VALIDATION_SYSTEM_PROMPT,
)


class TestBuildRewritePrompt:
    def test_includes_title(self):
        prompt = build_rewrite_prompt("My Title", "- bullet 1", "some feedback")
        assert "My Title" in prompt

    def test_includes_current_content(self):
        prompt = build_rewrite_prompt("Title", "- bullet 1\n- bullet 2", "feedback")
        assert "- bullet 1" in prompt
        assert "- bullet 2" in prompt

    def test_includes_feedback(self):
        prompt = build_rewrite_prompt("Title", "- content", "Add more detail about X")
        assert "Add more detail about X" in prompt

    def test_includes_constraints(self):
        prompt = build_rewrite_prompt("Title", "- content", "feedback")
        assert "Return ONLY improved bullet content" in prompt

    def test_two_column_prompt_includes_column_markers(self):
        prompt = build_rewrite_prompt(
            "Title", "## Left Column\n- A\n## Right Column\n- B",
            "feedback", is_two_column=True)
        assert "## Left Column" in prompt
        assert "## Right Column" in prompt

    def test_two_column_prompt_instructs_format(self):
        prompt = build_rewrite_prompt(
            "Title", "## Left Column\n- A\n## Right Column\n- B",
            "feedback", is_two_column=True)
        assert "Left Column" in prompt and "Right Column" in prompt


class TestParseRewrittenContent:
    def test_extracts_bullets(self):
        response = "- Improved bullet one\n- Improved bullet two\n- Third bullet"
        result = parse_rewritten_content(response)
        assert "- Improved bullet one" in result
        assert "- Improved bullet two" in result

    def test_strips_preamble_text(self):
        response = "Here are the improvements:\n\n- Actual bullet\n- Another bullet"
        result = parse_rewritten_content(response)
        assert "Here are the improvements" not in result
        assert "- Actual bullet" in result

    def test_handles_numbered_items(self):
        response = "1. First step\n2. Second step"
        result = parse_rewritten_content(response)
        assert "1. First step" in result

    def test_preserves_sub_bullets(self):
        response = "- Top level\n  - Sub bullet\n- Another top"
        result = parse_rewritten_content(response)
        assert "  - Sub bullet" in result


class TestExtractSlideContent:
    def test_extracts_title_and_body(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Add a textbox with content (simulating a body placeholder)
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(3))
        tf = tb.text_frame
        tf.paragraphs[0].text = "First bullet"
        tf.add_paragraph().text = "Second bullet"

        title, body = extract_slide_content(slide)
        assert "First bullet" in body
        assert "Second bullet" in body


    def test_extracts_from_real_content_placeholder(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        if slide.shapes.title:
            slide.shapes.title.text = "My Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text_frame.paragraphs[0].text = "Main point"
                ph.text_frame.add_paragraph().text = "Sub point"
        title, body = extract_slide_content(slide)
        assert title == "My Title"
        assert "Main point" in body
        assert "Sub point" in body

    def _make_two_column_slide(self):
        """Create a slide with Two Content layout (2 body placeholders)."""
        prs = Presentation()
        # slide_layouts[3] is "Two Content" in default python-pptx
        layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Two Column Title"
        content_phs = sorted(
            [(ph.placeholder_format.idx, ph)
             for ph in slide.placeholders if ph.placeholder_format.idx > 0],
            key=lambda x: x[0]
        )
        if len(content_phs) >= 2:
            content_phs[0][1].text_frame.paragraphs[0].text = "Left item A"
            content_phs[0][1].text_frame.add_paragraph().text = "Left item B"
            content_phs[1][1].text_frame.paragraphs[0].text = "Right item X"
            content_phs[1][1].text_frame.add_paragraph().text = "Right item Y"
        return prs, slide

    def test_extract_two_column_labels_sections(self):
        _, slide = self._make_two_column_slide()
        title, body = extract_slide_content(slide)
        assert "## Left Column" in body
        assert "Left item A" in body
        assert "## Right Column" in body
        assert "Right item X" in body

    def test_extract_single_column_no_section_labels(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Single bullet"
        _, body = extract_slide_content(slide)
        assert "## Left Column" not in body
        assert "Single bullet" in body

    def test_sub_bullet_gets_indented_prefix_from_level(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text_frame.paragraphs[0].text = "Top level"
                p = ph.text_frame.add_paragraph()
                p.text = "Indented"
                p.level = 1
        _, body = extract_slide_content(slide)
        assert "  - Indented" in body
        assert "- Top level" in body


class TestImproveSlide:
    """Test improve_slide with mocked LLM."""

    def _make_slide_with_content(self):
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        # Add title
        if slide.shapes.title:
            slide.shapes.title.text = "Test Title"
        # Add body content via textbox
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(3))
        tf = tb.text_frame
        tf.paragraphs[0].text = "Original bullet one"
        tf.add_paragraph().text = "Original bullet two"
        return prs, slide

    @patch('aippt.improve.analyze_slide')
    def test_status_no_content_when_empty_slide(self, mock_analyze):
        from aippt.improve import improve_slide
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        mock_client = MagicMock()
        result = improve_slide(slide, None, mock_client, dry_run=False)
        assert result['status'] == 'no_content'
        assert result['applied'] is False

    @patch('aippt.improve.analyze_slide')
    def test_status_dry_run(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved"
        _, slide = self._make_slide_with_content()
        result = improve_slide(slide, None, mock_client, dry_run=True)
        assert result['status'] == 'dry_run'
        assert result['applied'] is False

    @patch('aippt.improve.analyze_slide')
    def test_dry_run_does_not_modify(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "## Visual Design\nAdd more detail"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved bullet\n- Better bullet"

        _, slide = self._make_slide_with_content()
        result = improve_slide(slide, None, mock_client, dry_run=True)

        assert result['applied'] is False
        assert "Improved bullet" in result['improved']

    @patch('aippt.improve.analyze_slide')
    def test_returns_original_and_improved(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "## Visual Design\nFeedback here"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- New content\n- More content"

        _, slide = self._make_slide_with_content()
        result = improve_slide(slide, None, mock_client, dry_run=True)

        assert "Original bullet one" in result['original']
        assert "New content" in result['improved']
        assert result['title'] == "Test Title"

    def _make_real_content_slide(self):
        """Create slide with real body placeholder (layout[1])."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Real Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Original bullet"
                ph.text_frame.add_paragraph().text = "Second bullet"
        return prs, slide

    @patch('aippt.improve.analyze_slide')
    def test_apply_modifies_placeholder_text(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Brand new content"
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=False)
        assert result['status'] == 'applied'
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                assert "Brand new content" in ph.text_frame.text

    @patch('aippt.improve.analyze_slide')
    def test_two_column_improve_applies_to_both_placeholders(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback here"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = (
            "## Left Column\n- New left A\n- New left B\n"
            "## Right Column\n- New right X\n- New right Y"
        )
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content
        if slide.shapes.title:
            slide.shapes.title.text = "Two Col Test"
        content_phs = sorted(
            [(ph.placeholder_format.idx, ph)
             for ph in slide.placeholders if ph.placeholder_format.idx > 0],
            key=lambda x: x[0]
        )
        if len(content_phs) >= 2:
            content_phs[0][1].text_frame.paragraphs[0].text = "Old left"
            content_phs[1][1].text_frame.paragraphs[0].text = "Old right"

        result = improve_slide(slide, None, mock_client, dry_run=False)

        assert result['status'] == 'applied'
        left_text = content_phs[0][1].text_frame.text
        right_text = content_phs[1][1].text_frame.text
        assert "New left" in left_text
        assert "New right" in right_text


class TestImproveDeck:
    """Test improve_deck orchestration with mocked LLM."""

    @patch('aippt.improve.improve_slide')
    def test_filters_slides(self, mock_improve):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'Test', 'original': '', 'improved': '',
            'feedback': '', 'applied': False
        }
        mock_client = MagicMock()

        # Create a temp pptx with 3 slides
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp_path = f.name

        try:
            improve_deck(tmp_path, slides_filter=[1, 3], dry_run=True, client=mock_client)
            # Should only be called for slides 1 and 3
            assert mock_improve.call_count == 2
        finally:
            os.unlink(tmp_path)


class TestImproveDeckMultiPass:
    @patch('aippt.improve.improve_slide')
    def test_two_passes_calls_improve_twice_per_slide(self, mock_improve):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': 'o', 'improved': 'i',
            'feedback': 'f', 'applied': True, 'status': 'applied'
        }
        mock_client = MagicMock()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            results = improve_deck(tmp, passes=2, dry_run=True, client=mock_client)
            assert mock_improve.call_count == 2
            assert len(results) == 2
        finally:
            os.unlink(tmp)

    @patch('aippt.improve.improve_slide')
    def test_three_slides_two_passes_calls_six_times(self, mock_improve):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': '', 'improved': '',
            'feedback': '', 'applied': False, 'status': 'no_content'
        }
        mock_client = MagicMock()
        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            improve_deck(tmp, passes=2, dry_run=True, client=mock_client)
            assert mock_improve.call_count == 6
        finally:
            os.unlink(tmp)


class TestParseTwoColumnRewrittenContent:
    def test_splits_on_right_column_marker(self):
        from aippt.improve import parse_two_column_rewritten_content
        response = "## Left Column\n- A\n- B\n## Right Column\n- C\n- D"
        left, right = parse_two_column_rewritten_content(response)
        assert "- A" in left
        assert "- C" in right

    def test_fallback_when_no_markers(self):
        from aippt.improve import parse_two_column_rewritten_content
        response = "- A\n- B\n- C\n- D"
        left, right = parse_two_column_rewritten_content(response)
        assert left  # All content goes to left
        assert right == ""  # Empty right

    def test_strips_section_headers(self):
        from aippt.improve import parse_two_column_rewritten_content
        response = "## Left Column\n- bullet\n## Right Column\n- other"
        left, right = parse_two_column_rewritten_content(response)
        assert "## Left Column" not in left
        assert "## Right Column" not in right


class TestFocusGuidance:
    """Test FOCUS_GUIDANCE dict and focus parameter wiring."""

    def test_focus_guidance_has_all_choices(self):
        expected = {"accuracy", "detail", "brevity", "structure", "general"}
        assert set(FOCUS_GUIDANCE.keys()) == expected

    def test_general_focus_is_empty(self):
        assert FOCUS_GUIDANCE["general"] == ""

    def test_non_general_focus_values_non_empty(self):
        for key, value in FOCUS_GUIDANCE.items():
            if key != "general":
                assert value, f"FOCUS_GUIDANCE['{key}'] should not be empty"

    @patch('aippt.improve.analyze_slide')
    def test_focus_appended_to_system_prompt(self, mock_analyze):
        """When focus != 'general', guidance should be appended to system prompt."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved"

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Bullet one"

        improve_slide(slide, None, mock_client, dry_run=True, focus="brevity")

        call_kwargs = mock_client.generate_text.call_args[1]
        assert "Prioritize conciseness" in call_kwargs["system_prompt"]
        assert REWRITE_SYSTEM_PROMPT in call_kwargs["system_prompt"]

    @patch('aippt.improve.analyze_slide')
    def test_general_focus_uses_base_prompt(self, mock_analyze):
        """When focus == 'general', system prompt should be base + audience only."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved"

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Bullet one"

        improve_slide(slide, None, mock_client, dry_run=True, focus="general")

        call_kwargs = mock_client.generate_text.call_args[1]
        assert call_kwargs["system_prompt"].startswith(REWRITE_SYSTEM_PROMPT)
        # No focus text appended (general = empty)
        assert 'Prioritize conciseness' not in call_kwargs["system_prompt"]


class TestTryReexportImages:
    """Test _try_reexport_images graceful degradation."""

    @patch('aippt.improve.subprocess.run')
    @patch('os.path.exists')
    def test_returns_false_when_script_missing(self, mock_exists, mock_run):
        mock_exists.return_value = False
        result = _try_reexport_images('/tmp/test.pptx', '/tmp/images')
        assert result is False
        mock_run.assert_not_called()

    @patch('aippt.improve.subprocess.run')
    @patch('os.path.exists', return_value=True)
    def test_returns_false_when_no_powershell(self, mock_exists, mock_run):
        mock_run.side_effect = FileNotFoundError
        result = _try_reexport_images('/tmp/test.pptx', '/tmp/images')
        assert result is False

    @patch('aippt.improve.subprocess.run')
    @patch('os.path.exists', return_value=True)
    def test_returns_true_on_success(self, mock_exists, mock_run):
        # First call: PowerShell version check succeeds
        # Second call: actual export succeeds
        version_result = MagicMock(returncode=0)
        export_result = MagicMock(returncode=0)
        mock_run.side_effect = [version_result, export_result]
        result = _try_reexport_images('/tmp/test.pptx', '/tmp/images')
        assert result is True

    @patch('aippt.improve.subprocess.run')
    @patch('os.path.exists', return_value=True)
    def test_returns_false_on_export_failure(self, mock_exists, mock_run):
        version_result = MagicMock(returncode=0)
        export_result = MagicMock(returncode=1, stderr="PowerPoint not found")
        mock_run.side_effect = [version_result, export_result]
        result = _try_reexport_images('/tmp/test.pptx', '/tmp/images')
        assert result is False


class TestImproveDeckReexport:
    """Test that improve_deck attempts re-export on multi-pass."""

    @patch('aippt.improve._try_reexport_images')
    @patch('aippt.improve.improve_slide')
    def test_reexport_called_between_passes(self, mock_improve, mock_reexport):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': 'o', 'improved': 'i',
            'feedback': 'f', 'applied': True, 'status': 'applied'
        }
        mock_reexport.return_value = False
        mock_client = MagicMock()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            improve_deck(tmp, images_dir='/tmp/images', passes=2,
                         dry_run=False, client=mock_client)
            # Should attempt re-export after pass 1 (not after pass 2)
            assert mock_reexport.call_count == 1
        finally:
            os.unlink(tmp)

    @patch('aippt.improve._try_reexport_images')
    @patch('aippt.improve.improve_slide')
    def test_no_reexport_on_single_pass(self, mock_improve, mock_reexport):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': '', 'improved': '',
            'feedback': '', 'applied': True, 'status': 'applied'
        }
        mock_client = MagicMock()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            improve_deck(tmp, images_dir='/tmp/images', passes=1,
                         dry_run=False, client=mock_client)
            mock_reexport.assert_not_called()
        finally:
            os.unlink(tmp)


class TestImproveMetadata:
    """Verify improve appends structured metadata."""

    def _make_slide_with_content(self):
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = "Test Title"
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(3))
        tf = tb.text_frame
        tf.paragraphs[0].text = "Original bullet one"
        tf.add_paragraph().text = "Original bullet two"
        return prs, slide

    @patch('aippt.improve.analyze_slide')
    def test_improve_writes_metadata(self, mock_analyze):
        from aippt.improve import improve_slide
        from aippt.metadata import extract_metadata
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved bullet"
        mock_client.model = "claude-sonnet-4-6"
        _, slide = self._make_slide_with_content()
        result = improve_slide(slide, None, mock_client, dry_run=False, focus="brevity")
        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert entries[0]["operation"] == "improve"
        assert entries[0]["model"] == "claude-sonnet-4-6"
        assert entries[0]["focus"] == "brevity"
        assert "changes_summary" in entries[0]


class TestAudienceRewrite:
    def test_audience_rewrite_prompts_has_all_types(self):
        expected = {"engineers", "executives", "product", "mixed"}
        assert set(AUDIENCE_REWRITE_PROMPTS.keys()) == expected

    def test_each_prompt_is_nonempty_string(self):
        for key, value in AUDIENCE_REWRITE_PROMPTS.items():
            assert isinstance(value, str)
            assert len(value) > 20

    def test_engineers_mentions_technical(self):
        assert 'technical' in AUDIENCE_REWRITE_PROMPTS['engineers'].lower()

    def test_executives_mentions_business(self):
        assert 'business' in AUDIENCE_REWRITE_PROMPTS['executives'].lower()

    @patch('aippt.improve.analyze_slide')
    def test_audience_appended_to_rewrite_system_prompt(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved"

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Bullet one"

        improve_slide(slide, None, mock_client, dry_run=True, audience='executives')

        call_kwargs = mock_client.generate_text.call_args[1]
        assert 'business' in call_kwargs['system_prompt'].lower()
        assert REWRITE_SYSTEM_PROMPT in call_kwargs['system_prompt']

    @patch('aippt.improve.analyze_slide')
    def test_default_audience_is_mixed(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved"

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Bullet one"

        improve_slide(slide, None, mock_client, dry_run=True)

        call_kwargs = mock_client.generate_text.call_args[1]
        assert AUDIENCE_REWRITE_PROMPTS['mixed'] in call_kwargs['system_prompt']

    @patch('aippt.improve.analyze_slide')
    def test_audience_and_focus_both_appended(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved"

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Bullet one"

        improve_slide(slide, None, mock_client, dry_run=True,
                      focus='brevity', audience='executives')

        call_kwargs = mock_client.generate_text.call_args[1]
        system_prompt = call_kwargs['system_prompt']
        assert 'business' in system_prompt.lower()
        assert 'Prioritize conciseness' in system_prompt


class TestParseRewriteResponse:
    """Test parse_rewrite_response() title extraction from rewrite output."""

    def test_extracts_title_and_content(self):
        response = "TITLE: Engineers Waste 3 Hours/Day\nCONTENT:\n- bullet one\n- bullet two"
        new_title, content = parse_rewrite_response(response)
        assert new_title == "Engineers Waste 3 Hours/Day"
        assert "- bullet one" in content
        assert "- bullet two" in content

    def test_keep_returns_none_title(self):
        response = "TITLE: KEEP\nCONTENT:\n- bullet one\n- bullet two"
        new_title, content = parse_rewrite_response(response)
        assert new_title is None
        assert "- bullet one" in content

    def test_missing_title_line_returns_none(self):
        response = "- bullet one\n- bullet two\n- bullet three"
        new_title, content = parse_rewrite_response(response)
        assert new_title is None
        assert "- bullet one" in content

    def test_content_without_content_marker(self):
        """When TITLE is present but no CONTENT: marker, rest is content."""
        response = "TITLE: New Title\n- bullet one\n- bullet two"
        new_title, content = parse_rewrite_response(response)
        assert new_title == "New Title"
        assert "- bullet one" in content

    def test_title_with_extra_whitespace(self):
        response = "TITLE:   Spaced Title   \nCONTENT:\n- bullet"
        new_title, content = parse_rewrite_response(response)
        assert new_title == "Spaced Title"


class TestIsGenericTitle:
    """Test is_generic_title() heuristic for detecting generic/label titles."""

    def test_detects_common_generic_titles(self):
        generics = ["The Problem", "Architecture", "Results", "Overview",
                     "Summary", "Introduction", "Conclusion", "Agenda",
                     "Next Steps", "Questions", "Key Findings", "Key Takeaways"]
        for title in generics:
            assert is_generic_title(title), f"Expected '{title}' to be generic"

    def test_rejects_insight_driven_titles(self):
        specific = [
            "3 Hours Lost Per Engineer Per Day",
            "Production-Ready AI Stack",
            "290 Chunks, 2-Second Queries",
            "Engineers Waste 3 Hours/Day Searching for Answers",
            "Zero Trust Architecture Reduces Breach Risk by 80%",
        ]
        for title in specific:
            assert not is_generic_title(title), f"Expected '{title}' to NOT be generic"

    def test_case_insensitive(self):
        assert is_generic_title("the problem")
        assert is_generic_title("THE PROBLEM")
        assert is_generic_title("architecture")

    def test_empty_title_not_generic(self):
        assert not is_generic_title("")

    def test_background_is_generic(self):
        assert is_generic_title("Background")


class TestBuildRewritePromptTitle:
    """Test that build_rewrite_prompt requests TITLE when appropriate."""

    def test_requests_title_by_default(self):
        prompt = build_rewrite_prompt("The Problem", "- bullet", "feedback")
        assert "TITLE:" in prompt

    def test_omits_title_when_keep_titles(self):
        prompt = build_rewrite_prompt("The Problem", "- bullet", "feedback",
                                       keep_titles=True)
        assert "TITLE:" not in prompt

    def test_includes_title_format_instructions(self):
        prompt = build_rewrite_prompt("Architecture", "- bullet", "feedback")
        assert "KEEP" in prompt  # Should mention KEEP as an option


class TestTitleApplication:
    """Test title rewriting in improve_slide."""

    def _make_real_content_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "The Problem"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Original bullet"
                ph.text_frame.add_paragraph().text = "Second bullet"
        return prs, slide

    @patch('aippt.improve.analyze_slide')
    def test_applies_new_title(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Title is generic. Suggest: Engineers Waste 3 Hours/Day"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = (
            "TITLE: Engineers Waste 3 Hours/Day\n"
            "CONTENT:\n- Improved bullet"
        )
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=False)
        assert slide.shapes.title.text == "Engineers Waste 3 Hours/Day"
        assert result['title_rewritten'] is True
        assert result['original_title'] == "The Problem"
        assert result['new_title'] == "Engineers Waste 3 Hours/Day"

    @patch('aippt.improve.analyze_slide')
    def test_keeps_title_when_keep_response(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = (
            "TITLE: KEEP\nCONTENT:\n- Improved bullet"
        )
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=False)
        assert slide.shapes.title.text == "The Problem"
        assert result.get('title_rewritten', False) is False

    @patch('aippt.improve.analyze_slide')
    def test_skips_title_when_keep_titles_flag(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved bullet"
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=False,
                               keep_titles=True)
        assert slide.shapes.title.text == "The Problem"
        assert result.get('title_rewritten', False) is False

    @patch('aippt.improve.analyze_slide')
    def test_title_in_metadata(self, mock_analyze):
        from aippt.improve import improve_slide
        from aippt.metadata import extract_metadata
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = (
            "TITLE: Better Title\nCONTENT:\n- Improved bullet"
        )
        mock_client.model = "test-model"
        _, slide = self._make_real_content_slide()
        improve_slide(slide, None, mock_client, dry_run=False)
        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert entries[0]["title_rewritten"] is True
        assert entries[0]["original_title"] == "The Problem"
        assert entries[0]["new_title"] == "Better Title"

    @patch('aippt.improve.analyze_slide')
    def test_dry_run_reports_title_change(self, mock_analyze):
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = (
            "TITLE: New Title\nCONTENT:\n- Improved bullet"
        )
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True)
        # Title not actually changed in dry run
        assert slide.shapes.title.text == "The Problem"
        # But result should report the proposed change
        assert result['new_title'] == "New Title"
        assert result['original_title'] == "The Problem"


class TestKeepTitlesFlag:
    """Test --keep-titles CLI flag parsing and threading."""

    @patch('aippt.improve.improve_slide')
    def test_keep_titles_passed_to_improve_deck(self, mock_improve):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': '', 'improved': '',
            'feedback': '', 'applied': False, 'status': 'no_content'
        }
        mock_client = MagicMock()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            improve_deck(tmp, dry_run=True, client=mock_client, keep_titles=True)
            call_kwargs = mock_improve.call_args[1]
            assert call_kwargs.get('keep_titles') is True
        finally:
            os.unlink(tmp)


class TestSelectFocus:
    """Test select_focus() auto-detection from feedback text."""

    def test_vague_signals_select_detail(self):
        feedback = "The claims are too vague and lack specificity."
        assert select_focus(feedback) == "detail"

    def test_unclear_signals_select_detail(self):
        feedback = "Several points are unclear and unspecific."
        assert select_focus(feedback) == "detail"

    def test_verbose_signals_select_brevity(self):
        feedback = "Content is verbose and redundant."
        assert select_focus(feedback) == "brevity"

    def test_wordy_signals_select_brevity(self):
        feedback = "Too wordy. Remove repetitive bullet points."
        assert select_focus(feedback) == "brevity"

    def test_inaccurate_signals_select_accuracy(self):
        feedback = "The data is inaccurate and misleading."
        assert select_focus(feedback) == "accuracy"

    def test_unsupported_signals_select_accuracy(self):
        feedback = "Claims are unsupported by evidence."
        assert select_focus(feedback) == "accuracy"

    def test_disorganized_signals_select_structure(self):
        feedback = "Content is disorganized with poor flow."
        assert select_focus(feedback) == "structure"

    def test_no_hierarchy_signals_select_structure(self):
        feedback = "There is no hierarchy in the bullet points."
        assert select_focus(feedback) == "structure"

    def test_no_signal_returns_general(self):
        feedback = "Overall good slide with minor improvements needed."
        assert select_focus(feedback) == "general"

    def test_empty_feedback_returns_general(self):
        assert select_focus("") == "general"

    def test_case_insensitive(self):
        feedback = "The content is VAGUE and VERBOSE."
        # Both have 1 hit each — tie broken by dict order (detail first)
        result = select_focus(feedback)
        assert result in ("detail", "brevity")  # Either is valid (1 hit each)

    def test_mixed_signals_highest_count_wins(self):
        # 2 brevity signals ("verbose", "redundant") vs 1 accuracy ("inaccurate")
        feedback = "The content is verbose and redundant with an inaccurate claim."
        assert select_focus(feedback) == "brevity"


class TestParseValidationResponse:
    """Test parse_validation_response() extracts structured verdict."""

    def test_parses_pass_verdict(self):
        response = (
            "VERDICT: PASS\n"
            "ADDRESSED: clarity, specificity, examples\n"
            "UNADDRESSED: \n"
            "SUGGESTION: "
        )
        result = parse_validation_response(response)
        assert result['verdict'] == 'PASS'
        assert 'clarity' in result['addressed']
        assert result['unaddressed'] == ''

    def test_parses_fail_verdict(self):
        response = (
            "VERDICT: FAIL\n"
            "ADDRESSED: structure\n"
            "UNADDRESSED: specificity, concrete examples\n"
            "SUGGESTION: Add specific metrics and data points"
        )
        result = parse_validation_response(response)
        assert result['verdict'] == 'FAIL'
        assert 'specificity' in result['unaddressed']
        assert 'Add specific' in result['suggestion']

    def test_parses_partial_verdict(self):
        response = (
            "VERDICT: PARTIAL\n"
            "ADDRESSED: organization\n"
            "UNADDRESSED: accuracy of claims\n"
            "SUGGESTION: Verify the 60% figure"
        )
        result = parse_validation_response(response)
        assert result['verdict'] == 'PARTIAL'

    def test_missing_fields_default_empty(self):
        response = "VERDICT: PASS"
        result = parse_validation_response(response)
        assert result['verdict'] == 'PASS'
        assert result['addressed'] == ''
        assert result['unaddressed'] == ''
        assert result['suggestion'] == ''

    def test_unknown_verdict_treated_as_fail(self):
        response = "VERDICT: MAYBE\nADDRESSED: nothing"
        result = parse_validation_response(response)
        assert result['verdict'] == 'FAIL'

    def test_multiline_unaddressed_field(self):
        response = (
            "VERDICT: PARTIAL\n"
            "ADDRESSED: clarity\n"
            "UNADDRESSED: specificity\n"
            "- missing concrete metrics\n"
            "- no examples provided\n"
            "SUGGESTION: Add numbers"
        )
        result = parse_validation_response(response)
        assert result['verdict'] == 'PARTIAL'
        assert 'specificity' in result['unaddressed']
        assert 'missing concrete metrics' in result['unaddressed']
        assert 'no examples provided' in result['unaddressed']
        assert result['suggestion'] == 'Add numbers'


class TestHasConverged:
    """Test has_converged() similarity detection."""

    def test_identical_content_converged(self):
        assert has_converged("- bullet one\n- bullet two", "- bullet one\n- bullet two")

    def test_very_similar_content_converged(self):
        # >95% similarity
        original = "- bullet one\n- bullet two\n- bullet three"
        changed = "- bullet one\n- bullet two\n- bullet three."  # minor punctuation
        assert has_converged(original, changed)

    def test_substantially_different_not_converged(self):
        original = "- bullet one\n- bullet two"
        changed = "- completely new content\n- entirely different approach\n- new third point"
        assert not has_converged(original, changed)

    def test_empty_strings_converged(self):
        assert has_converged("", "")

    def test_one_empty_not_converged(self):
        assert not has_converged("- content", "")


class TestRetryLoop:
    """Test validation retry loop in improve_slide."""

    def _make_real_content_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Test Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                ph.text_frame.paragraphs[0].text = "Original bullet"
                ph.text_frame.add_paragraph().text = "Second bullet"
        return prs, slide

    @patch('aippt.improve.analyze_slide')
    def test_validation_pass_no_retry(self, mock_analyze):
        """When validation passes on first attempt, no retry occurs."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        # First call: rewrite. Second call: validation returns PASS
        mock_client.generate_text.side_effect = [
            "- Improved bullet",
            "VERDICT: PASS\nADDRESSED: all points\nUNADDRESSED: \nSUGGESTION: ",
        ]
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               max_retries=2, no_validate=False)
        assert result.get('validation', {}).get('passed') is True
        assert result.get('validation', {}).get('retries', 0) == 0

    @patch('aippt.improve.analyze_slide')
    def test_validation_fail_triggers_retry(self, mock_analyze):
        """When validation fails, retry occurs with adjusted prompt."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Content is vague."
        mock_client = MagicMock()
        # Calls: rewrite1, validate1(FAIL), rewrite2, validate2(PASS)
        mock_client.generate_text.side_effect = [
            "- Vague bullet",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: specificity\nSUGGESTION: Add numbers",
            "- Specific bullet with 42% metric",
            "VERDICT: PASS\nADDRESSED: specificity\nUNADDRESSED: \nSUGGESTION: ",
        ]
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               max_retries=2, no_validate=False)
        assert result.get('validation', {}).get('retries') == 1
        assert result.get('validation', {}).get('passed') is True

    @patch('aippt.improve.analyze_slide')
    def test_max_retries_stops_loop(self, mock_analyze):
        """When max retries reached, applies best attempt."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        # All validations fail
        mock_client.generate_text.side_effect = [
            "- Attempt 1",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: everything\nSUGGESTION: try harder",
            "- Attempt 2",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: still bad\nSUGGESTION: more effort",
            "- Attempt 3",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: gave up\nSUGGESTION: impossible",
        ]
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               max_retries=2, no_validate=False)
        assert result.get('validation', {}).get('passed') is False
        assert result.get('validation', {}).get('retries') == 2

    @patch('aippt.improve.analyze_slide')
    def test_no_validate_skips_validation(self, mock_analyze):
        """When no_validate=True, no validation LLM call is made."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved bullet"
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               no_validate=True)
        # Only 1 generate_text call (the rewrite), no validation call
        assert mock_client.generate_text.call_count == 1
        assert 'validation' not in result

    @patch('aippt.improve.analyze_slide')
    def test_convergence_stops_retry(self, mock_analyze):
        """When rewrite produces identical content, stop retrying."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        # Rewrite, validate(FAIL), retry produces same content
        mock_client.generate_text.side_effect = [
            "- Same bullet",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: detail\nSUGGESTION: add more",
            "- Same bullet",  # Converged — same output
        ]
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               max_retries=2, no_validate=False)
        # Should stop after convergence, not validate again
        assert result.get('validation', {}).get('retries') == 1
        assert result.get('validation', {}).get('converged') is True

    @patch('aippt.improve.analyze_slide')
    def test_partial_verdict_accepted_after_retry(self, mock_analyze):
        """PARTIAL verdict is accepted after at least one retry."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        # Calls: rewrite1, validate1(FAIL), rewrite2, validate2(PARTIAL)
        mock_client.generate_text.side_effect = [
            "- Attempt 1",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: detail\nSUGGESTION: add more",
            "- Attempt 2 with more detail",
            "VERDICT: PARTIAL\nADDRESSED: some\nUNADDRESSED: minor point\nSUGGESTION: ",
        ]
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               max_retries=2, no_validate=False)
        assert result.get('validation', {}).get('passed') is True
        assert result.get('validation', {}).get('retries') == 1

    @patch('aippt.improve.analyze_slide')
    def test_partial_on_first_attempt_not_accepted(self, mock_analyze):
        """PARTIAL verdict on first attempt (retries=0) triggers retry."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        # Calls: rewrite1, validate1(PARTIAL), rewrite2, validate2(PASS)
        mock_client.generate_text.side_effect = [
            "- Attempt 1",
            "VERDICT: PARTIAL\nADDRESSED: some\nUNADDRESSED: detail\nSUGGESTION: more",
            "- Attempt 2 improved",
            "VERDICT: PASS\nADDRESSED: all\nUNADDRESSED: \nSUGGESTION: ",
        ]
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               max_retries=2, no_validate=False)
        assert result.get('validation', {}).get('passed') is True
        assert result.get('validation', {}).get('retries') == 1

    @patch('aippt.improve.analyze_slide')
    def test_best_attempt_used_when_max_retries_exhausted(self, mock_analyze):
        """When all retries fail, the best attempt (PARTIAL) is used, not the last (FAIL)."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Feedback"
        mock_client = MagicMock()
        # Calls: rewrite1, validate1(PARTIAL), rewrite2, validate2(FAIL),
        #        rewrite3, validate3(FAIL)
        mock_client.generate_text.side_effect = [
            "- Attempt 1 decent",
            "VERDICT: PARTIAL\nADDRESSED: some\nUNADDRESSED: detail\nSUGGESTION: more",
            "- Attempt 2 worse",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: everything\nSUGGESTION: try harder",
            "- Attempt 3 also bad",
            "VERDICT: FAIL\nADDRESSED: \nUNADDRESSED: everything\nSUGGESTION: give up",
        ]
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               max_retries=2, no_validate=False)
        # PARTIAL was the best — but it was on first try (retries=0),
        # so it wasn't accepted. Best attempt content should be used.
        assert result['improved'] == "- Attempt 1 decent"

    @patch('aippt.improve.analyze_slide')
    def test_validation_metadata_recorded(self, mock_analyze):
        """Validation results appear in metadata."""
        from aippt.improve import improve_slide
        from aippt.metadata import extract_metadata
        mock_analyze.return_value = "Content is vague."
        mock_client = MagicMock()
        mock_client.model = "test-model"
        mock_client.generate_text.side_effect = [
            "- Improved bullet",
            "VERDICT: PASS\nADDRESSED: all\nUNADDRESSED: \nSUGGESTION: ",
        ]
        _, slide = self._make_real_content_slide()
        improve_slide(slide, None, mock_client, dry_run=False,
                      max_retries=2, no_validate=False)
        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert 'validation' in entries[0]
        assert entries[0]['validation']['passed'] is True

    @patch('aippt.improve.analyze_slide')
    def test_auto_focus_recorded_in_result(self, mock_analyze):
        """When focus is auto-selected, focus_source='auto' in result."""
        from aippt.improve import improve_slide
        mock_analyze.return_value = "Content is vague and unclear."
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "- Improved"
        _, slide = self._make_real_content_slide()
        result = improve_slide(slide, None, mock_client, dry_run=True,
                               no_validate=True)
        # Default focus="general" triggers auto-selection from feedback
        assert result.get('focus_source') == 'auto'
        assert result.get('focus_used') == 'detail'  # "vague" → detail


class TestNoValidateAndMaxRetriesCLI:
    """Test --no-validate and --max-retries CLI flag threading."""

    @patch('aippt.improve.improve_slide')
    def test_no_validate_passed_to_improve_deck(self, mock_improve):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': '', 'improved': '',
            'feedback': '', 'applied': False, 'status': 'no_content'
        }
        mock_client = MagicMock()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            improve_deck(tmp, dry_run=True, client=mock_client, no_validate=True)
            call_kwargs = mock_improve.call_args[1]
            assert call_kwargs.get('no_validate') is True
        finally:
            os.unlink(tmp)

    @patch('aippt.improve.improve_slide')
    def test_max_retries_passed_to_improve_deck(self, mock_improve):
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': '', 'improved': '',
            'feedback': '', 'applied': False, 'status': 'no_content'
        }
        mock_client = MagicMock()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            improve_deck(tmp, dry_run=True, client=mock_client, max_retries=3)
            call_kwargs = mock_improve.call_args[1]
            assert call_kwargs.get('max_retries') == 3
        finally:
            os.unlink(tmp)

    @patch('aippt.improve.improve_slide')
    def test_combined_audience_keep_titles_validation_threading(self, mock_improve):
        """All three flags (audience, keep_titles, no_validate) thread through to improve_slide."""
        from aippt.improve import improve_deck
        mock_improve.return_value = {
            'title': 'T', 'original': '', 'improved': '',
            'feedback': '', 'applied': False, 'status': 'no_content'
        }
        mock_client = MagicMock()
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        import tempfile, os
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as f:
            prs.save(f.name)
            tmp = f.name
        try:
            improve_deck(
                tmp, dry_run=True, client=mock_client,
                audience='executives',
                keep_titles=True,
                no_validate=True,
                max_retries=1,
                focus='brevity',
            )
            call_kwargs = mock_improve.call_args[1]
            assert call_kwargs.get('audience') == 'executives'
            assert call_kwargs.get('keep_titles') is True
            assert call_kwargs.get('no_validate') is True
            assert call_kwargs.get('max_retries') == 1
            assert call_kwargs.get('focus') == 'brevity'
        finally:
            os.unlink(tmp)
