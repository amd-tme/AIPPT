"""Regression test: addTwoColumn crashing on non-string column items.

Discovered via live Playwright testing against the deployed SLAI instance
(2026-07-27): enhanced-mode generation produced {text, subs}-style bullet
objects for a two_column slide, and addTwoColumn's renderColItems passed
them straight to _pushBulletItem without normalizing to a string first
(unlike addBulletSlide and the addImageTwoColumn sibling, which both do),
crashing with TypeError: text.match is not a function.
"""

import os
import subprocess
import tempfile

from pptx import Presentation

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
NODE_PATH = subprocess.check_output(
    ["npm", "root", "-g"], cwd=REPO_ROOT, text=True
).strip()


def _run_script(script: str, out_name: str) -> str:
    helpers_path = os.path.join(REPO_ROOT, "lib", "pptxgenjs-helpers.mjs")
    theme_path = os.path.join(REPO_ROOT, "themes", "default.yaml")
    output_path = os.path.join(tempfile.mkdtemp(), out_name)
    full_script = script.format(helpers_path=helpers_path, theme_path=theme_path, output_path=output_path)

    script_path = os.path.join(tempfile.gettempdir(), f"test_two_column_{out_name}.mjs")
    with open(script_path, "w") as f:
        f.write(full_script)

    env = {**os.environ, "NODE_PATH": NODE_PATH}
    result = subprocess.run(
        ["node", script_path],
        cwd=REPO_ROOT,
        env=env,
        capture_output=True,
        text=True,
        timeout=30,
    )
    assert result.returncode == 0, f"Generation failed:\n{result.stderr}"
    assert os.path.exists(output_path), f"Output not created: {output_path}"
    return output_path


def test_string_items_keep_working(tmp_path):
    out = _run_script(
        """\
import {{ createDeck, addTwoColumn }} from '{helpers_path}';
const deck = createDeck('{theme_path}');
addTwoColumn(deck, 'Compare', 'Left', 'Right',
  ['**Slow** — Building takes days', '**Manual** — Cycles repeatedly'],
  ['**Scattered** — Lives everywhere', '**Invisible** — Cannot find it'],
  1, '');
await deck.save('{output_path}');
""",
        "string_items.pptx",
    )
    p = Presentation(out)
    assert len(p.slides) == 1


def test_object_shaped_items_do_not_crash(tmp_path):
    """MAT-78-sibling bug: {text, subs}-style objects must not crash the render."""
    out = _run_script(
        """\
import {{ createDeck, addTwoColumn }} from '{helpers_path}';
const deck = createDeck('{theme_path}');
addTwoColumn(deck, 'Compare', 'Left', 'Right',
  [{{ text: '**Slow** — Building takes days', subs: ['detail one'] }}],
  [{{ text: '**Scattered** — Lives everywhere', subs: ['detail two'] }}],
  1, '');
await deck.save('{output_path}');
""",
        "object_items.pptx",
    )
    p = Presentation(out)
    assert len(p.slides) == 1
    texts = [shape.text_frame.text for shape in p.slides[0].shapes if shape.has_text_frame]
    assert any("Slow" in t for t in texts)
    assert any("Scattered" in t for t in texts)
