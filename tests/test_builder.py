"""Tests for aippt.builder module."""

import pytest
from unittest.mock import MagicMock, patch
from pptx import Presentation

from aippt.builder import BuildContext, build_slide


class TestBuildContext:
    def test_defaults(self):
        ctx = BuildContext()
        assert ctx.client is None
        assert ctx.image_gen == "none"
        assert ctx.image_dir is None
        assert ctx.model is None
        assert ctx.mcp_manager is None
        assert ctx.mcp_server == "txt2img"
        assert ctx.classification == "internal"
        assert ctx.audience == "mixed"
        assert ctx.audience_source == "default"


class TestBuildSlideContent:
    """Tests for build_slide() CONTENT: extraction and fallback.

    Mirrors existing TestAddSlideEnhancedContent from test_cli.py.
    """

    def _make_prs(self):
        return Presentation()

    def test_uses_enhanced_content_from_suggestions(self):
        """When LLM response contains CONTENT:, slide body uses enhanced text."""
        prs = self._make_prs()
        slide_data = {
            'title': 'Test Title',
            'content': (
                "CONTENT:\n"
                "- Enhanced point one with more detail\n"
                "- Enhanced point two with added context\n"
                "NARRATIVE: This slide covers important points.\n"
                "LAYOUT: bullet\n"
                "VISUALS: Emphasize first point\n"
                "TALKING_POINTS: Additional details"
            ),
            'original_content': ["- Terse point one", "- Terse point two"],
        }
        layout = build_slide(prs, slide_data, BuildContext())
        slide = prs.slides[0]
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() != "Test Title":
                body_texts.append(shape.text_frame.text)
        body = '\n'.join(body_texts)
        assert 'Enhanced point one' in body
        assert 'Terse point one' not in body

    def test_falls_back_to_original_when_no_content(self):
        """When LLM response lacks CONTENT:, slide body uses original_content."""
        prs = self._make_prs()
        slide_data = {
            'title': 'Test Title',
            'content': (
                "NARRATIVE: This slide covers important points.\n"
                "LAYOUT: bullet\n"
                "VISUALS: Emphasize first point\n"
                "TALKING_POINTS: Additional details"
            ),
            'original_content': ["- Original point one", "- Original point two"],
        }
        layout = build_slide(prs, slide_data, BuildContext())
        slide = prs.slides[0]
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() != "Test Title":
                body_texts.append(shape.text_frame.text)
        body = '\n'.join(body_texts)
        assert 'Original point one' in body

    def test_non_enhanced_path_uses_content_lines(self):
        """Without original_content (non-enhanced path), slide body uses content_lines."""
        prs = self._make_prs()
        slide_data = {
            'title': 'Test Title',
            'content': "- Plain bullet one\n- Plain bullet two",
        }
        layout = build_slide(prs, slide_data, BuildContext())
        slide = prs.slides[0]
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip() != "Test Title":
                body_texts.append(shape.text_frame.text)
        body = '\n'.join(body_texts)
        assert 'Plain bullet one' in body

    def test_returns_layout_type(self):
        """build_slide returns the layout type string."""
        prs = self._make_prs()
        slide_data = {
            'title': 'Test Title',
            'content': "LAYOUT: bullet\nCONTENT:\n- Bullet",
        }
        layout = build_slide(prs, slide_data, BuildContext())
        assert layout == 'bullet'
