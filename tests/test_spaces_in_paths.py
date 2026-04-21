"""Tests for filenames and paths containing spaces.

Real-world PowerPoint decks on Windows often have names like
"Q4 2025 Security Review.pptx".  These tests validate that spaces
flow correctly through catalog, search, remix, export, and CLI parsing.
"""

import csv
import os

import pytest
import yaml
from pptx import Presentation

from aippt.catalog import (
    catalog_deck,
    file_hash,
    get_db,
    get_deck_slides,
    list_decks,
    search_slides,
    add_tags,
)
from aippt.cli import build_parser, cmd_reverse
from aippt.export import export_csv
from aippt.remix import generate_manifest, load_manifest, assemble_deck


# -- Fixtures ----------------------------------------------------------------

SPACED_NAME = "Q4 2025 Security Review"


@pytest.fixture
def spaced_pptx(tmp_path):
    """Create a minimal PPTX whose filename contains spaces."""
    prs = Presentation()
    layout = prs.slide_layouts[0]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Introduction"

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Security Overview"

    path = tmp_path / f"{SPACED_NAME}.pptx"
    prs.save(str(path))
    return str(path)


@pytest.fixture
def db_path(tmp_path):
    return str(tmp_path / "test.db")


# -- Catalog -----------------------------------------------------------------

class TestCatalogSpacedFilenames:
    def test_catalog_deck_with_spaces(self, spaced_pptx, db_path):
        """Deck with spaces catalogs and stores the correct name."""
        deck_id = catalog_deck(spaced_pptx, db_path=db_path)
        assert deck_id > 0

        decks = list_decks(db_path)
        assert len(decks) == 1
        assert decks[0]["name"] == SPACED_NAME

    def test_catalog_stores_relative_path_with_spaces(self, spaced_pptx, db_path):
        """Database stores a relative path that preserves spaces."""
        catalog_deck(spaced_pptx, db_path=db_path)

        conn = get_db(db_path)
        row = conn.execute("SELECT file_path FROM decks").fetchone()
        conn.close()

        assert " " in row["file_path"]
        assert not os.path.isabs(row["file_path"])

    def test_catalog_slides_retrievable(self, spaced_pptx, db_path):
        """Slides from a spaced-name deck are searchable."""
        deck_id = catalog_deck(spaced_pptx, db_path=db_path)
        slides = get_deck_slides(deck_id, db_path)
        assert len(slides) == 2
        titles = [s["title"] for s in slides]
        assert "Introduction" in titles
        assert "Security Overview" in titles

    def test_catalog_with_spaced_images_dir(self, tmp_path, spaced_pptx, db_path):
        """Image directory with spaces is resolved correctly."""
        images_dir = tmp_path / "images" / SPACED_NAME
        images_dir.mkdir(parents=True)

        # Create a dummy image
        (images_dir / "Slide1.png").write_bytes(b"\x89PNG fake")

        deck_id = catalog_deck(spaced_pptx, db_path=db_path, images_dir=str(images_dir))

        slides = get_deck_slides(deck_id, db_path)
        slide1 = [s for s in slides if s["position"] == 1][0]
        assert slide1["image_path"] is not None
        assert " " in slide1["image_path"]

    def test_recatalog_spaced_file(self, spaced_pptx, db_path):
        """Re-cataloging a spaced-name file returns the same ID."""
        id1 = catalog_deck(spaced_pptx, db_path=db_path)
        id2 = catalog_deck(spaced_pptx, db_path=db_path)
        assert id1 == id2

    def test_file_hash_with_spaces(self, spaced_pptx):
        """file_hash works on paths with spaces."""
        h = file_hash(spaced_pptx)
        assert len(h) == 64  # SHA-256 hex digest


# -- Search & tags -----------------------------------------------------------

class TestSearchSpacedDecks:
    def test_search_returns_spaced_deck_path(self, spaced_pptx, db_path):
        """Search results include the full spaced deck path."""
        deck_id = catalog_deck(spaced_pptx, db_path=db_path)
        add_tags(1, ["security"], "manual", db_path)

        results = search_slides(db_path=db_path, tags=["security"])
        assert len(results) == 1
        assert " " in results[0]["deck_path"]


# -- Export CSV --------------------------------------------------------------

class TestExportSpacedDecks:
    def test_export_csv_with_spaced_deck(self, tmp_path, spaced_pptx, db_path):
        """CSV export captures the spaced deck name correctly."""
        catalog_deck(spaced_pptx, db_path=db_path)

        output = str(tmp_path / "out.csv")
        count = export_csv(output, db_path=db_path, export_all=True)

        assert count == 2
        with open(output, encoding="utf-8") as f:
            rows = list(csv.DictReader(f))
            assert all(r["deck_name"] == SPACED_NAME for r in rows)

    def test_csv_output_path_with_spaces(self, tmp_path, spaced_pptx, db_path):
        """CSV can be written to a path that contains spaces."""
        catalog_deck(spaced_pptx, db_path=db_path)

        spaced_dir = tmp_path / "My Reports"
        spaced_dir.mkdir()
        output = str(spaced_dir / "slide export.csv")

        count = export_csv(output, db_path=db_path, export_all=True)
        assert count == 2
        assert os.path.exists(output)


# -- Remix -------------------------------------------------------------------

class TestRemixSpacedPaths:
    def test_generate_manifest_preserves_spaced_basename(self, spaced_pptx):
        """Manifest YAML preserves the spaced filename."""
        slides = [
            {"deck_path": spaced_pptx, "position": 1, "title": "Introduction"},
        ]
        result = generate_manifest(slides, title="Test")
        parsed = yaml.safe_load(result)
        assert " " in parsed["slides"][0]["deck"]
        assert parsed["slides"][0]["deck"] == f"{SPACED_NAME}.pptx"
        assert parsed["slides"][0]["deck_path"] == spaced_pptx

    def test_manifest_yaml_roundtrip_with_spaces(self, tmp_path, spaced_pptx):
        """Manifest with spaced paths survives YAML write/read cycle."""
        slides = [
            {"deck_path": spaced_pptx, "position": 1, "title": "Introduction"},
        ]
        yaml_str = generate_manifest(slides, title="Spaced Test")

        manifest_path = str(tmp_path / "manifest.yaml")
        with open(manifest_path, "w", encoding="utf-8") as f:
            f.write(yaml_str)

        loaded = load_manifest(manifest_path)
        assert loaded["slides"][0]["deck_path"] == spaced_pptx
        assert " " in loaded["slides"][0]["deck"]

    def test_assemble_deck_with_spaced_source(self, tmp_path, spaced_pptx, db_path):
        """assemble_deck works when the source deck path has spaces."""
        get_db(db_path).close()

        manifest = {
            "title": "From Spaced Deck",
            "slides": [
                {
                    "deck": f"{SPACED_NAME}.pptx",
                    "deck_path": spaced_pptx,
                    "position": 1,
                    "title": "Introduction",
                },
                {
                    "deck": f"{SPACED_NAME}.pptx",
                    "deck_path": spaced_pptx,
                    "position": 2,
                    "title": "Security Overview",
                },
            ],
        }
        manifest_path = str(tmp_path / "manifest.yaml")
        with open(manifest_path, "w") as f:
            yaml.dump(manifest, f)

        output_path = str(tmp_path / "assembled output.pptx")
        count = assemble_deck(manifest_path, output_path, db_path=db_path)

        assert count == 2
        assert os.path.exists(output_path)

        result_prs = Presentation(output_path)
        assert len(result_prs.slides) == 2


# -- CLI arg parsing ---------------------------------------------------------

class TestCliSpacedArgs:
    def test_create_with_spaced_paths(self):
        """Parser accepts file arguments containing spaces."""
        parser = build_parser()
        args = parser.parse_args([
            "create",
            "my outline notes.md",
            "Company Template 2025.pptx",
            "Q4 2025 Output.pptx",
        ])
        assert args.outline == "my outline notes.md"
        assert args.template == "Company Template 2025.pptx"
        assert args.output == "Q4 2025 Output.pptx"

    def test_catalog_with_spaced_path(self):
        parser = build_parser()
        args = parser.parse_args(["catalog", "My Big Deck.pptx"])
        assert args.deck == "My Big Deck.pptx"

    def test_reverse_auto_output_with_spaces(self, tmp_path):
        """Reverse command auto-generates .md from a spaced .pptx name."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"

        pptx_path = tmp_path / f"{SPACED_NAME}.pptx"
        prs.save(str(pptx_path))

        parser = build_parser()
        args = parser.parse_args(["reverse", str(pptx_path)])

        result = cmd_reverse(args)
        assert result == 0

        expected_md = tmp_path / f"{SPACED_NAME}.md"
        assert expected_md.exists()

    def test_reverse_spaced_output_path(self, tmp_path):
        """Reverse command writes to an explicitly spaced output path."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"

        pptx_path = tmp_path / "test.pptx"
        prs.save(str(pptx_path))

        output_path = tmp_path / "my output notes.md"

        parser = build_parser()
        args = parser.parse_args(["reverse", str(pptx_path), str(output_path)])

        result = cmd_reverse(args)
        assert result == 0
        assert output_path.exists()
