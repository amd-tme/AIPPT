"""Tests for source tracking: schema migration, catalog, ingest, CLI, and detection."""
import os

import pytest
from unittest.mock import patch, MagicMock

from aippt.catalog import (
    get_db,
    catalog_deck,
    get_deck_by_id,
    resolve_deck,
    detect_source_engine,
    detect_source_theme,
)
from aippt.ingest import ingest_deck


class TestDetectSourceEngine:
    """Tests for detect_source_engine()."""

    def test_detects_pptxgenjs_single_quotes(self, tmp_path):
        script = tmp_path / "deck.js"
        script.write_text("const pptx = require('pptxgenjs');\n")
        assert detect_source_engine(str(script)) == "pptxgenjs"

    def test_detects_pptxgenjs_double_quotes(self, tmp_path):
        script = tmp_path / "deck.js"
        script.write_text('const pptx = require("pptxgenjs");\n')
        assert detect_source_engine(str(script)) == "pptxgenjs"

    def test_detects_python_pptx_from_import(self, tmp_path):
        script = tmp_path / "deck.py"
        script.write_text("from pptx import Presentation\n")
        assert detect_source_engine(str(script)) == "python-pptx"

    def test_detects_python_pptx_import_pptx(self, tmp_path):
        script = tmp_path / "deck.py"
        script.write_text("import pptx\n")
        assert detect_source_engine(str(script)) == "python-pptx"

    def test_returns_none_for_unknown(self, tmp_path):
        script = tmp_path / "deck.txt"
        script.write_text("just some text\n")
        assert detect_source_engine(str(script)) is None

    def test_returns_none_for_missing_file(self, tmp_path):
        assert detect_source_engine(str(tmp_path / "missing.js")) is None

    def test_only_scans_first_50_lines(self, tmp_path):
        script = tmp_path / "deck.js"
        lines = ["// filler line\n"] * 55
        lines.append("const pptx = require('pptxgenjs');\n")
        script.write_text("".join(lines))
        assert detect_source_engine(str(script)) is None

    def test_first_match_wins(self, tmp_path):
        script = tmp_path / "deck.js"
        script.write_text("const pptx = require('pptxgenjs');\nfrom pptx import Presentation\n")
        assert detect_source_engine(str(script)) == "pptxgenjs"


class TestDetectSourceTheme:
    """Tests for detect_source_theme()."""

    def test_detects_amd_theme(self, tmp_path):
        script = tmp_path / "deck.js"
        script.write_text("const theme = loadYaml('themes/amd.yaml');\n")
        assert detect_source_theme(str(script)) == "amd"

    def test_detects_default_theme(self, tmp_path):
        script = tmp_path / "deck.js"
        script.write_text("const theme = loadYaml('themes/default.yaml');\n")
        assert detect_source_theme(str(script)) == "default"

    def test_detects_yml_extension(self, tmp_path):
        script = tmp_path / "deck.js"
        script.write_text("const theme = loadYaml('themes/custom.yml');\n")
        assert detect_source_theme(str(script)) == "custom"

    def test_returns_none_when_no_theme(self, tmp_path):
        script = tmp_path / "deck.js"
        script.write_text("const pptx = require('pptxgenjs');\n")
        assert detect_source_theme(str(script)) is None

    def test_returns_none_for_missing_file(self, tmp_path):
        assert detect_source_theme(str(tmp_path / "missing.js")) is None


class TestSchemaMigrationSourceCols:
    """Tests that source tracking columns exist after migration."""

    def test_fresh_db_has_source_columns(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn = get_db(db_path)
        cols = {row[1] for row in conn.execute("PRAGMA table_info(decks)").fetchall()}
        conn.close()
        assert "source_script_path" in cols
        assert "source_engine" in cols
        assert "source_theme" in cols
        assert "outline_path" in cols
        assert "source_generated_at" in cols

    def test_migration_idempotent(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        conn1 = get_db(db_path)
        conn1.close()
        # Open again — migration runs a second time
        conn2 = get_db(db_path)
        cols = {row[1] for row in conn2.execute("PRAGMA table_info(decks)").fetchall()}
        conn2.close()
        assert "source_script_path" in cols


class TestCatalogDeckSourceTracking:
    """Tests for catalog_deck() source tracking fields."""

    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """Create a minimal PPTX for testing."""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        path = tmp_path / "deck.pptx"
        prs.save(str(path))
        return str(path)

    def test_catalog_with_source_fields(self, tmp_path, sample_pptx):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(
            sample_pptx, db_path=db_path,
            source_script_path="output/deck.js",
            source_engine="pptxgenjs",
            source_theme="amd",
            outline_path="outlines/deck.md",
        )
        deck = get_deck_by_id(deck_id, db_path)
        assert deck["source_script_path"] == "output/deck.js"
        assert deck["source_engine"] == "pptxgenjs"
        assert deck["source_theme"] == "amd"
        assert deck["outline_path"] == "outlines/deck.md"
        assert deck["source_generated_at"] is not None

    def test_catalog_without_source_fields(self, tmp_path, sample_pptx):
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(sample_pptx, db_path=db_path)
        deck = get_deck_by_id(deck_id, db_path)
        assert deck["source_script_path"] is None
        assert deck["source_engine"] is None
        assert deck["source_theme"] is None
        assert deck["outline_path"] is None
        assert deck["source_generated_at"] is None

    def test_recatalog_preserves_source_fields(self, tmp_path, sample_pptx):
        """Re-cataloging without source fields should preserve existing ones."""
        db_path = str(tmp_path / "test.db")
        # First catalog with source
        deck_id = catalog_deck(
            sample_pptx, db_path=db_path,
            source_script_path="output/deck.js",
            source_engine="pptxgenjs",
        )

        # Modify the PPTX so hash changes (triggers re-catalog path)
        from pptx import Presentation
        prs = Presentation(sample_pptx)
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(sample_pptx)

        # Re-catalog without passing source fields
        deck_id2 = catalog_deck(sample_pptx, db_path=db_path)
        assert deck_id == deck_id2  # Same deck, re-cataloged

        deck = get_deck_by_id(deck_id, db_path)
        # Source fields should be preserved
        assert deck["source_script_path"] == "output/deck.js"
        assert deck["source_engine"] == "pptxgenjs"

    def test_recatalog_updates_source_when_provided(self, tmp_path, sample_pptx):
        """Re-cataloging with new source fields should update them."""
        db_path = str(tmp_path / "test.db")
        deck_id = catalog_deck(
            sample_pptx, db_path=db_path,
            source_script_path="output/old.js",
            source_engine="pptxgenjs",
        )

        from pptx import Presentation
        prs = Presentation(sample_pptx)
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(sample_pptx)

        deck_id2 = catalog_deck(
            sample_pptx, db_path=db_path,
            source_script_path="output/new.py",
            source_engine="python-pptx",
        )

        deck = get_deck_by_id(deck_id, db_path)
        assert deck["source_script_path"] == "output/new.py"
        assert deck["source_engine"] == "python-pptx"

    def test_resolve_deck_includes_source_fields(self, tmp_path, sample_pptx):
        db_path = str(tmp_path / "test.db")
        catalog_deck(
            sample_pptx, db_path=db_path,
            source_script_path="output/deck.js",
            source_engine="pptxgenjs",
        )
        result = resolve_deck("deck", db_path=db_path)
        assert result["source_script_path"] == "output/deck.js"
        assert result["source_engine"] == "pptxgenjs"


class TestIngestSourceTracking:
    """Tests for ingest_deck() source tracking."""

    @patch("aippt.ingest.get_deck_by_id", return_value={"slide_count": 3})
    @patch("aippt.ingest.catalog_deck", return_value=1)
    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_ingest_with_source(self, mock_export, mock_catalog, mock_get, tmp_path):
        deck = tmp_path / "test.pptx"
        deck.touch()

        script = tmp_path / "test.js"
        script.write_text("const pptx = require('pptxgenjs');\n")

        result = ingest_deck(
            str(deck),
            db_path=str(tmp_path / "test.db"),
            source_script_path=str(script),
        )

        assert result["source_tracked"] is True
        # Verify catalog_deck was called with source fields
        _, kwargs = mock_catalog.call_args
        assert kwargs["source_script_path"] == str(script)
        assert kwargs["source_engine"] == "pptxgenjs"

    @patch("aippt.ingest.get_deck_by_id", return_value={"slide_count": 3})
    @patch("aippt.ingest.catalog_deck", return_value=1)
    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_ingest_without_source(self, mock_export, mock_catalog, mock_get, tmp_path):
        deck = tmp_path / "test.pptx"
        deck.touch()

        result = ingest_deck(str(deck), db_path=str(tmp_path / "test.db"))

        assert result["source_tracked"] is False

    @patch("aippt.ingest.get_deck_by_id", return_value={"slide_count": 3})
    @patch("aippt.ingest.catalog_deck", return_value=1)
    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_ingest_theme_override(self, mock_export, mock_catalog, mock_get, tmp_path):
        deck = tmp_path / "test.pptx"
        deck.touch()

        script = tmp_path / "test.js"
        script.write_text("const pptx = require('pptxgenjs');\nconst t = loadYaml('themes/default.yaml');\n")

        result = ingest_deck(
            str(deck),
            db_path=str(tmp_path / "test.db"),
            source_script_path=str(script),
            source_theme="amd",  # Override auto-detected 'default'
        )

        _, kwargs = mock_catalog.call_args
        assert kwargs["source_theme"] == "amd"

    @patch("aippt.ingest.get_deck_by_id", return_value={"slide_count": 3})
    @patch("aippt.ingest.catalog_deck", return_value=1)
    @patch("aippt.ingest.cmd_export_images", return_value=0)
    def test_ingest_auto_detects_theme(self, mock_export, mock_catalog, mock_get, tmp_path):
        deck = tmp_path / "test.pptx"
        deck.touch()

        script = tmp_path / "test.js"
        script.write_text("const pptx = require('pptxgenjs');\nconst t = loadYaml('themes/amd.yaml');\n")

        result = ingest_deck(
            str(deck),
            db_path=str(tmp_path / "test.db"),
            source_script_path=str(script),
        )

        _, kwargs = mock_catalog.call_args
        assert kwargs["source_theme"] == "amd"
