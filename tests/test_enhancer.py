"""Tests for aippt.enhancer module."""

import json
import pytest
from unittest.mock import MagicMock, patch

from aippt.enhancer import (
    enhance_with_llm,
    format_slide_notes,
    parse_deck_plan,
    plan_deck,
    _repair_truncated_json,
    SYSTEM_PROMPT,
    AUDIENCE_PROMPTS,
    PLANNING_SYSTEM_PROMPT,
    VALID_LAYOUTS,
)


class TestFormatSlideNotes:
    def test_formats_all_sections(self):
        suggestions = {
            'NARRATIVE': 'This is the narrative.',
            'LAYOUT': 'bullet',
            'VISUALS': 'Use icons for each point.',
            'TALKING_POINTS': 'Emphasize the first point.',
        }
        result = format_slide_notes(suggestions)

        assert 'NARRATIVE:' in result
        assert 'This is the narrative.' in result
        assert 'LAYOUT:' in result
        assert 'bullet' in result
        assert 'VISUALS:' in result
        assert 'Use icons for each point.' in result
        assert 'TALKING POINTS:' in result

    def test_handles_missing_keys(self):
        suggestions = {
            'NARRATIVE': 'Just narrative',
        }
        result = format_slide_notes(suggestions)

        assert 'NARRATIVE:' in result
        assert 'Just narrative' in result
        assert 'LAYOUT:' in result  # Should still have section header
        assert 'VISUALS:' in result

    def test_handles_empty_dict(self):
        suggestions = {}
        result = format_slide_notes(suggestions)

        assert 'NARRATIVE:' in result
        assert 'LAYOUT:' in result
        assert 'VISUALS:' in result
        assert 'TALKING POINTS:' in result

    def test_result_is_stripped(self):
        suggestions = {'NARRATIVE': 'test'}
        result = format_slide_notes(suggestions)

        assert not result.startswith('\n')
        assert not result.endswith('\n')


class TestEnhanceWithLlm:
    @patch('aippt.enhancer.LLMClient')
    def test_calls_generate_text(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test narrative\nLAYOUT: bullet"

        slide = {
            'title': 'Test Slide',
            'content': ['- Point 1', '- Point 2']
        }

        result = enhance_with_llm(slide, mock_client)

        mock_client.generate_text.assert_called_once()
        assert 'NARRATIVE: Test narrative' in result

    @patch('aippt.enhancer.LLMClient')
    def test_includes_slide_title_in_prompt(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"

        slide = {
            'title': 'Important Title',
            'content': ['content']
        }

        enhance_with_llm(slide, mock_client)

        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert 'Important Title' in call_kwargs['prompt']

    @patch('aippt.enhancer.LLMClient')
    def test_uses_correct_system_prompt(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"

        slide = {
            'title': 'Title',
            'content': ['content']
        }

        enhance_with_llm(slide, mock_client)

        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert call_kwargs['system_prompt'].startswith(SYSTEM_PROMPT)

    @patch('aippt.enhancer.LLMClient')
    def test_returns_raw_llm_response(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "LLM response"

        slide = {
            'title': 'Title',
            'content': ['Original line 1', 'Original line 2']
        }

        result = enhance_with_llm(slide, mock_client)

        assert result == "LLM response"
        assert 'ORIGINAL CONTENT:' not in result


    @patch('aippt.enhancer.LLMClient')
    def test_prompt_discourages_diagram_when_image_gen_none(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"

        slide = {
            'title': 'Test Slide',
            'content': ['- Point 1', '- Point 2']
        }

        enhance_with_llm(slide, mock_client, image_gen='none')

        call_kwargs = mock_client.generate_text.call_args.kwargs
        prompt = call_kwargs['prompt'].lower()
        assert 'disabled' in prompt or 'do not' in prompt
        assert 'diagram' in prompt

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_mentions_two_column_use_cases(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"

        slide = {
            'title': 'Test Slide',
            'content': ['- Point 1', '- Point 2']
        }

        enhance_with_llm(slide, mock_client)

        call_kwargs = mock_client.generate_text.call_args.kwargs
        prompt = call_kwargs['prompt'].lower()
        assert 'before/after' in prompt or 'input/output' in prompt
        assert 'problem/solution' in prompt or 'pairs or contrasts' in prompt


    @patch('aippt.enhancer.LLMClient')
    def test_prompt_has_variety_rule(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"
        slide = {'title': 'T', 'content': ['- p1']}
        enhance_with_llm(slide, mock_client)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'VARIETY' in prompt or 'variety' in prompt.lower()

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_shows_numbered_format_example(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"
        slide = {'title': 'Steps', 'content': ['- Install', '- Configure']}
        enhance_with_llm(slide, mock_client)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert '1.' in prompt and '2.' in prompt

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_requests_content_section(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "CONTENT:\n- Enhanced\nNARRATIVE: test"
        slide = {'title': 'Test', 'content': ['- Point 1', '- Point 2']}
        enhance_with_llm(slide, mock_client)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'CONTENT:' in prompt
        assert 'Enhanced bullet' in prompt or 'Enhanced bullet 1' in prompt

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_instructs_bullet_preservation(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"
        slide = {'title': 'Test', 'content': ['- A', '- B']}
        enhance_with_llm(slide, mock_client)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'Preserve' in prompt
        assert 'number' in prompt.lower() or 'order' in prompt.lower()


class TestSystemPrompt:
    def test_system_prompt_is_defined(self):
        assert SYSTEM_PROMPT
        assert 'professional' in SYSTEM_PROMPT.lower()
        assert 'PowerPoint' in SYSTEM_PROMPT or 'slides' in SYSTEM_PROMPT.lower()

    def test_system_prompt_guides_two_column_usage(self):
        assert '2-4 two_column' in SYSTEM_PROMPT

    def test_system_prompt_marks_bullet_as_default(self):
        assert 'DEFAULT' in SYSTEM_PROMPT


class TestMCPImagePrompt:
    """Tests for IMAGE_PROMPT guidance when image_gen='mcp'."""

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_includes_image_prompt_guidance_when_mcp(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"
        slide = {'title': 'Architecture', 'content': ['- Service A', '- Service B']}

        enhance_with_llm(slide, mock_client, image_gen='mcp')

        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'IMAGE_PROMPT:' in prompt
        assert 'IMAGE_PROMPT' in prompt

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_excludes_image_prompt_when_not_mcp(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"
        slide = {'title': 'Test', 'content': ['- Point']}

        enhance_with_llm(slide, mock_client, image_gen='none')

        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'IMAGE_PROMPT:' not in prompt

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_excludes_image_prompt_when_dalle(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"
        slide = {'title': 'Test', 'content': ['- Point']}

        enhance_with_llm(slide, mock_client, image_gen='dalle')

        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'IMAGE_PROMPT:' not in prompt

    @patch('aippt.enhancer.LLMClient')
    def test_prompt_mentions_image_prompt_format(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"
        slide = {'title': 'Test', 'content': ['- Point']}

        enhance_with_llm(slide, mock_client, image_gen='mcp')

        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'diagram' in prompt.lower() or 'illustration' in prompt.lower()


class TestEnhanceMetadata:
    """Verify enhance writes metadata to slide notes."""

    def test_add_slide_writes_enhance_metadata(self):
        import json
        from pptx import Presentation
        from aippt.metadata import extract_metadata, append_metadata

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Simulate what _add_slide does after enhancement
        append_metadata(
            slide, "enhance",
            model="claude-sonnet-4-6",
            layout_selected="bullet",
            original_content_hash="abc123",
            directives={"LAYOUT": None, "IMAGE": None},
        )
        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert entries[0]["operation"] == "enhance"
        assert entries[0]["model"] == "claude-sonnet-4-6"
        assert entries[0]["layout_selected"] == "bullet"


class TestAudiencePrompts:
    def test_audience_prompts_has_all_types(self):
        expected = {"engineers", "executives", "product", "mixed"}
        assert set(AUDIENCE_PROMPTS.keys()) == expected

    def test_each_prompt_is_nonempty_string(self):
        for key, value in AUDIENCE_PROMPTS.items():
            assert isinstance(value, str), f"AUDIENCE_PROMPTS['{key}'] is not a string"
            assert len(value) > 20, f"AUDIENCE_PROMPTS['{key}'] is too short"

    def test_engineers_mentions_technical(self):
        assert 'technical' in AUDIENCE_PROMPTS['engineers'].lower()

    def test_executives_mentions_business(self):
        assert 'business' in AUDIENCE_PROMPTS['executives'].lower()

    @patch('aippt.enhancer.LLMClient')
    def test_audience_appended_to_system_prompt(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"
        slide = {'title': 'Test', 'content': ['- Point']}

        enhance_with_llm(slide, mock_client, audience='engineers')

        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert 'technical' in call_kwargs['system_prompt'].lower()
        assert SYSTEM_PROMPT in call_kwargs['system_prompt']

    @patch('aippt.enhancer.LLMClient')
    def test_mixed_audience_is_default(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"
        slide = {'title': 'Test', 'content': ['- Point']}

        enhance_with_llm(slide, mock_client)

        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert AUDIENCE_PROMPTS['mixed'] in call_kwargs['system_prompt']

    @patch('aippt.enhancer.LLMClient')
    def test_unknown_audience_falls_back_to_mixed(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"
        slide = {'title': 'Test', 'content': ['- Point']}

        enhance_with_llm(slide, mock_client, audience='aliens')

        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert AUDIENCE_PROMPTS['mixed'] in call_kwargs['system_prompt']


class TestParseDeckPlan:
    def test_parses_valid_json(self):
        raw = '''{
          "narrative_arc": "problem-solution",
          "arc_assessment": "Strong opening with clear call to action.",
          "slides": [
            {
              "index": 0,
              "title": "The Problem",
              "role": "hook",
              "suggested_layout": "basic",
              "transition_to_next": "Now let's look at the impact...",
              "context_hint": "Open with a compelling problem statement"
            },
            {
              "index": 1,
              "title": "Impact",
              "role": "context",
              "suggested_layout": "two_column",
              "transition_to_next": "Here is our solution...",
              "context_hint": "Use before/after parallel structure"
            }
          ]
        }'''
        plan = parse_deck_plan(raw)
        assert plan['narrative_arc'] == 'problem-solution'
        assert plan['arc_assessment'] == 'Strong opening with clear call to action.'
        assert len(plan['slides']) == 2
        assert plan['slides'][0]['role'] == 'hook'
        assert plan['slides'][0]['suggested_layout'] == 'basic'
        assert plan['slides'][1]['transition_to_next'] == 'Here is our solution...'

    def test_extracts_json_from_markdown_code_block(self):
        raw = '''Here is the deck plan:
```json
{
  "narrative_arc": "chronological",
  "arc_assessment": "Good flow.",
  "slides": [
    {"index": 0, "title": "T", "role": "hook", "suggested_layout": "bullet",
     "transition_to_next": "next", "context_hint": "hint"}
  ]
}
```
'''
        plan = parse_deck_plan(raw)
        assert plan['narrative_arc'] == 'chronological'
        assert len(plan['slides']) == 1

    def test_returns_empty_plan_on_invalid_json(self):
        plan = parse_deck_plan("not json at all")
        assert plan['narrative_arc'] == 'unknown'
        assert plan['slides'] == []

    def test_returns_empty_plan_on_missing_slides(self):
        raw = '{"narrative_arc": "problem-solution"}'
        plan = parse_deck_plan(raw)
        assert plan['slides'] == []

    def test_normalizes_layout_to_valid_type(self):
        raw = '''{
          "narrative_arc": "x",
          "arc_assessment": "x",
          "slides": [
            {"index": 0, "title": "T", "role": "hook",
             "suggested_layout": "BULLET",
             "transition_to_next": "", "context_hint": ""}
          ]
        }'''
        plan = parse_deck_plan(raw)
        assert plan['slides'][0]['suggested_layout'] == 'bullet'

    def test_unknown_layout_falls_back_to_bullet(self):
        raw = '''{
          "narrative_arc": "x",
          "arc_assessment": "x",
          "slides": [
            {"index": 0, "title": "T", "role": "hook",
             "suggested_layout": "fancy_grid",
             "transition_to_next": "", "context_hint": ""}
          ]
        }'''
        plan = parse_deck_plan(raw)
        assert plan['slides'][0]['suggested_layout'] == 'bullet'


class TestRepairTruncatedJson:
    def test_repairs_missing_closing_brackets(self):
        truncated = '{"narrative_arc": "x", "slides": [{"index": 0, "title": "T"'
        repaired = _repair_truncated_json(truncated)
        result = json.loads(repaired)
        assert result['narrative_arc'] == 'x'

    def test_repairs_truncated_mid_array(self):
        truncated = ('{"narrative_arc": "x", "arc_assessment": "", "slides": ['
                      '{"index": 0, "title": "A", "role": "hook", "suggested_layout": "basic"}, '
                      '{"index": 1, "title": "B", "role": "context"')
        repaired = _repair_truncated_json(truncated)
        result = json.loads(repaired)
        assert len(result['slides']) == 2

    def test_noop_on_valid_json(self):
        valid = '{"key": "value"}'
        assert _repair_truncated_json(valid) == valid

    def test_strips_trailing_comma(self):
        truncated = '{"slides": [{"a": 1},'
        repaired = _repair_truncated_json(truncated)
        result = json.loads(repaired)
        assert len(result['slides']) == 1


class TestParseDeckPlanTruncated:
    def test_repairs_truncated_json_response(self):
        # Simulate a response truncated mid-slide by token limit
        truncated = ('{"narrative_arc": "problem-solution", "arc_assessment": "Good.", '
                      '"slides": [{"index": 0, "title": "Intro", "role": "hook", '
                      '"suggested_layout": "basic", "transition_to_next": "", "context_hint": ""}, '
                      '{"index": 1, "title": "Details", "role": "context", '
                      '"suggested_layout": "bullet", "transition_to_next": "next"')
        plan = parse_deck_plan(truncated)
        assert plan['narrative_arc'] == 'problem-solution'
        assert len(plan['slides']) >= 1  # At least the complete slide parsed

    def test_repairs_truncated_code_block(self):
        truncated = ('```json\n'
                      '{"narrative_arc": "x", "arc_assessment": "", '
                      '"slides": [{"index": 0, "title": "T", "role": "hook", '
                      '"suggested_layout": "basic", "transition_to_next": "", "context_hint": ""}\n'
                      '```')
        plan = parse_deck_plan(truncated)
        assert plan['narrative_arc'] == 'x'
        assert len(plan['slides']) == 1


class TestPlanDeck:
    def test_calls_llm_with_all_slide_titles(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = json.dumps({
            "narrative_arc": "problem-solution",
            "arc_assessment": "Good.",
            "slides": [
                {"index": 0, "title": "Intro", "role": "hook",
                 "suggested_layout": "basic", "transition_to_next": "",
                 "context_hint": ""},
                {"index": 1, "title": "Details", "role": "context",
                 "suggested_layout": "bullet", "transition_to_next": "",
                 "context_hint": ""},
            ]
        })
        slides = [
            {'title': 'Intro', 'content': ['- Point A', '- Point B']},
            {'title': 'Details', 'content': ['- Detail 1']},
        ]
        plan = plan_deck(slides, mock_client)
        assert plan['narrative_arc'] == 'problem-solution'
        assert len(plan['slides']) == 2

        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert 'Intro' in call_kwargs['prompt']
        assert 'Details' in call_kwargs['prompt']

    def test_uses_planning_system_prompt(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': 'T', 'content': ['- p']}]
        plan_deck(slides, mock_client)
        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert PLANNING_SYSTEM_PROMPT in call_kwargs['system_prompt']

    def test_includes_audience_in_prompt_when_provided(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': 'T', 'content': ['- p']}]
        plan_deck(slides, mock_client, audience='executives')
        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert 'executives' in call_kwargs['prompt'].lower()

    def test_includes_content_summary_max_3_bullets(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': 'T', 'content': [
            '- Bullet 1', '- Bullet 2', '- Bullet 3',
            '- Bullet 4', '- Bullet 5'
        ]}]
        plan_deck(slides, mock_client)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'Bullet 1' in prompt
        assert 'Bullet 3' in prompt
        assert 'Bullet 4' not in prompt

    def test_max_tokens_scales_with_slide_count(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': f'Slide {i}', 'content': ['- p']} for i in range(14)]
        plan_deck(slides, mock_client)
        call_kwargs = mock_client.generate_text.call_args.kwargs
        # 14 slides * 250 = 3500, which exceeds the 2000 minimum
        assert call_kwargs['max_tokens'] >= 3500

    def test_max_tokens_minimum_for_small_decks(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': 'T', 'content': ['- p']}]
        plan_deck(slides, mock_client)
        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert call_kwargs['max_tokens'] == 2000

    def test_max_tokens_capped_at_4000(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': f'S{i}', 'content': ['- p']} for i in range(20)]
        plan_deck(slides, mock_client)
        call_kwargs = mock_client.generate_text.call_args.kwargs
        assert call_kwargs['max_tokens'] == 4000

    def test_handles_llm_failure_gracefully(self):
        mock_client = MagicMock()
        mock_client.generate_text.side_effect = Exception("API error")
        slides = [
            {'title': 'Slide 1', 'content': ['- p']},
            {'title': 'Slide 2', 'content': ['- q']},
        ]
        plan = plan_deck(slides, mock_client)
        assert plan['narrative_arc'] == 'unknown'
        assert len(plan['slides']) == 0


class TestDeckContextInjection:
    @patch('aippt.enhancer.LLMClient')
    def test_deck_context_adds_role_to_prompt(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"
        slide = {'title': 'Test', 'content': ['- Point']}
        deck_context = {
            'role': 'hook',
            'suggested_layout': 'basic',
            'transition_to_next': 'Move to details...',
            'context_hint': 'Open with impact',
        }
        enhance_with_llm(slide, mock_client, deck_context=deck_context)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'hook' in prompt
        assert 'basic' in prompt.lower()
        assert 'Move to details' in prompt
        assert 'Open with impact' in prompt

    @patch('aippt.enhancer.LLMClient')
    def test_deck_context_none_produces_same_prompt_as_before(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"
        slide = {'title': 'Test', 'content': ['- Point']}

        enhance_with_llm(slide, mock_client, deck_context=None)
        prompt_without = mock_client.generate_text.call_args.kwargs['prompt']

        enhance_with_llm(slide, mock_client)
        prompt_default = mock_client.generate_text.call_args.kwargs['prompt']

        assert prompt_without == prompt_default
        assert 'Role in narrative' not in prompt_without

    @patch('aippt.enhancer.LLMClient')
    def test_deck_context_mentions_override_guidance(self, mock_client_cls):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "response"
        slide = {'title': 'Test', 'content': ['- Point']}
        deck_context = {
            'role': 'evidence',
            'suggested_layout': 'two_column',
            'transition_to_next': '',
            'context_hint': '',
        }
        enhance_with_llm(slide, mock_client, deck_context=deck_context)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'override' in prompt.lower() or 'recommendation' in prompt.lower()


class TestLayoutVariety:
    def test_no_three_consecutive_same_layouts(self):
        """plan_deck prompt should instruct no more than 2 consecutive same layouts."""
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': f'Slide {i}', 'content': ['- p']} for i in range(10)]
        plan_deck(slides, mock_client)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        # The prompt must mention the consecutive layout constraint
        assert '2 consecutive' in prompt.lower() or 'two consecutive' in prompt.lower()

    def test_prompt_lists_available_layouts(self):
        mock_client = MagicMock()
        mock_client.generate_text.return_value = '{"narrative_arc":"x","arc_assessment":"","slides":[]}'
        slides = [{'title': 'T', 'content': ['- p']}]
        plan_deck(slides, mock_client)
        prompt = mock_client.generate_text.call_args.kwargs['prompt']
        assert 'bullet' in prompt
        assert 'two_column' in prompt
        assert 'numbered' in prompt
        assert 'basic' in prompt


class TestNoPlanFlag:
    """Verify --no-plan suppresses deck planning."""

    @patch('aippt.enhancer.plan_deck')
    @patch('aippt.llm.LLMClient')
    @patch('aippt.config.get_model_default', return_value='test-model')
    def test_no_plan_true_skips_planning(self, mock_get_model, mock_llm_cls, mock_plan_deck):
        """When no_plan=True, plan_deck() must not be called."""
        from aippt.pipeline import run_pipeline, PipelineConfig
        import tempfile, os
        from pptx import Presentation

        # Mock LLMClient so enhance=True doesn't crash without credentials
        mock_client = MagicMock()
        mock_client.generate_text.return_value = "NARRATIVE: Test\nLAYOUT: bullet"
        mock_client.model_config = MagicMock()
        mock_client.model_config.provider = "openai"
        mock_llm_cls.return_value = mock_client

        prs = Presentation()
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tf:
            prs.save(tf.name)
            template = tf.name
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as of:
            output = of.name
        try:
            config = PipelineConfig(
                outline_text="## Test\n- Point\n",
                template_path=template,
                output_path=output,
                enhance=True,
                no_plan=True,
            )
            run_pipeline(config)
            mock_plan_deck.assert_not_called()
        finally:
            os.unlink(template)
            os.unlink(output)


class TestShowPlanOutput:
    """Verify --show-plan prints the deck plan to stdout."""

    @patch('aippt.llm.LLMClient')
    @patch('aippt.config.get_model_default', return_value='test-model')
    def test_show_plan_prints_to_stdout(self, mock_get_model, mock_llm_cls, capsys):
        """When show_plan=True, deck plan is printed to stdout."""
        from aippt.pipeline import run_pipeline, PipelineConfig
        import tempfile, os
        from pptx import Presentation as PptxPresentation

        mock_client = MagicMock()
        plan_json = json.dumps({
            "narrative_arc": "problem-solution",
            "arc_assessment": "Strong opening with clear call to action.",
            "slides": [
                {"index": 0, "title": "The Problem", "role": "hook",
                 "suggested_layout": "basic",
                 "transition_to_next": "Now the solution...",
                 "context_hint": "Open with impact"},
                {"index": 1, "title": "Our Solution", "role": "solution",
                 "suggested_layout": "bullet",
                 "transition_to_next": "",
                 "context_hint": "Present the fix"},
            ]
        })
        # First call = plan_deck, subsequent calls = per-slide enhance
        mock_client.generate_text.side_effect = [
            plan_json,
            "NARRATIVE: Test\nLAYOUT: basic",
            "NARRATIVE: Test\nLAYOUT: bullet",
        ]
        mock_client.model_config = MagicMock()
        mock_client.model_config.provider = "openai"
        mock_llm_cls.return_value = mock_client

        prs = PptxPresentation()
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tf:
            prs.save(tf.name)
            template = tf.name
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as of:
            output = of.name
        try:
            config = PipelineConfig(
                outline_text="## The Problem\n- Issue one\n## Our Solution\n- Fix one\n",
                template_path=template,
                output_path=output,
                enhance=True,
                show_plan=True,
            )
            run_pipeline(config)
            captured = capsys.readouterr()
            assert "Deck Narrative Plan" in captured.out
            assert "problem-solution" in captured.out
            assert "hook" in captured.out
            assert "The Problem" in captured.out
            assert "Our Solution" in captured.out
            assert "Open with impact" in captured.out
            assert "Now the solution" in captured.out
        finally:
            os.unlink(template)
            os.unlink(output)


class TestFrontmatterToEnhanceIntegration:
    """Integration: frontmatter audience flows through to enhance + deck planning."""

    @patch('aippt.llm.LLMClient')
    @patch('aippt.config.get_model_default', return_value='test-model')
    def test_frontmatter_audience_reaches_enhance_and_plan(
        self, mock_get_model, mock_llm_cls
    ):
        """Frontmatter audience=executives flows into both plan_deck and enhance_with_llm."""
        from aippt.pipeline import run_pipeline, PipelineConfig
        import tempfile, os
        from pptx import Presentation as PptxPresentation

        mock_client = MagicMock()
        plan_json = json.dumps({
            "narrative_arc": "persuasive",
            "arc_assessment": "Good executive flow.",
            "slides": [
                {"index": 0, "title": "Revenue Impact", "role": "hook",
                 "suggested_layout": "basic",
                 "transition_to_next": "", "context_hint": "Lead with ROI"},
            ]
        })
        mock_client.generate_text.side_effect = [
            plan_json,  # plan_deck call
            "NARRATIVE: Business narrative\nLAYOUT: basic",  # enhance call
        ]
        mock_client.model_config = MagicMock()
        mock_client.model_config.provider = "openai"
        mock_llm_cls.return_value = mock_client

        prs = PptxPresentation()
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tf:
            prs.save(tf.name)
            template = tf.name
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as of:
            output = of.name
        try:
            outline_with_frontmatter = (
                "---\naudience: executives\n---\n\n"
                "## Revenue Impact\n- Q1 revenue up 40%\n- Costs down 20%\n- Margin expanded\n"
            )
            config = PipelineConfig(
                outline_text=outline_with_frontmatter,
                template_path=template,
                output_path=output,
                enhance=True,
            )
            run_pipeline(config)

            # Verify generate_text was called at least twice (plan + enhance)
            assert mock_client.generate_text.call_count >= 2

            # First call = plan_deck — should include "executives" in the prompt
            plan_call = mock_client.generate_text.call_args_list[0]
            assert 'executives' in plan_call.kwargs['prompt'].lower()

            # Second call = enhance_with_llm — system prompt should include
            # executive audience guidance
            enhance_call = mock_client.generate_text.call_args_list[1]
            assert 'business' in enhance_call.kwargs['system_prompt'].lower()
        finally:
            os.unlink(template)
            os.unlink(output)


class TestEnhanceMetadataWithDeckPlan:
    """Verify deck plan fields appear in enhance metadata."""

    def test_metadata_includes_deck_plan_fields(self):
        from pptx import Presentation
        from aippt.metadata import extract_metadata, append_metadata

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        append_metadata(
            slide, "enhance",
            model="claude-sonnet-4-6",
            layout_selected="two_column",
            deck_plan_role="evidence",
            deck_plan_layout="two_column",
            deck_plan_context="Show data comparison",
            narrative_arc="problem-solution",
        )
        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert entries[0]["deck_plan_role"] == "evidence"
        assert entries[0]["deck_plan_layout"] == "two_column"
        assert entries[0]["deck_plan_context"] == "Show data comparison"
        assert entries[0]["narrative_arc"] == "problem-solution"

    def test_metadata_omits_deck_plan_when_no_plan(self):
        from pptx import Presentation
        from aippt.metadata import extract_metadata, append_metadata

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        append_metadata(
            slide, "enhance",
            model="claude-sonnet-4-6",
            layout_selected="bullet",
        )
        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert "deck_plan_role" not in entries[0]
        assert "narrative_arc" not in entries[0]
