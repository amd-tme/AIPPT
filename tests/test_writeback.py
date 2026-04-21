"""Tests for PPTX notes write-back."""
import os
import pytest
from pptx import Presentation
from aippt.catalog import catalog_deck, get_db, record_edit
from aippt.writeback import write_notes_to_pptx, create_backup, WritebackResult
from aippt.cli import build_parser, cmd_write_notes


@pytest.fixture
def deck_with_notes(tmp_path):
    """Create a 3-slide PPTX where slide 1 has notes, slides 2-3 don't."""
    prs = Presentation()
    layout = prs.slide_layouts[0]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Slide One"
    s1.notes_slide.notes_text_frame.text = "Original notes 1"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Slide Two"
    s3 = prs.slides.add_slide(layout)
    s3.shapes.title.text = "Slide Three"
    path = str(tmp_path / "deck.pptx")
    prs.save(path)
    return path

@pytest.fixture
def db_path(tmp_path):
    return str(tmp_path / "test.db")

@pytest.fixture
def cataloged(deck_with_notes, db_path):
    deck_id = catalog_deck(deck_with_notes, db_path=db_path)
    return deck_id, deck_with_notes, db_path


class TestWriteNotesToPptx:
    def test_writes_edited_notes(self, cataloged):
        deck_id, deck_path, db_path = cataloged
        record_edit(1, "notes", "Edited notes for slide 1", source="web", db_path=db_path)
        result = write_notes_to_pptx(deck_path, db_path=db_path, deck_id=deck_id)
        assert isinstance(result, WritebackResult)
        assert result.deck_id == deck_id
        assert result.slides_written == 1
        assert result.slides_total == 3
        prs = Presentation(deck_path)
        assert prs.slides[0].notes_slide.notes_text_frame.text == "Edited notes for slide 1"

    def test_skips_empty_notes(self, cataloged):
        deck_id, deck_path, db_path = cataloged
        result = write_notes_to_pptx(deck_path, db_path=db_path, deck_id=deck_id)
        assert result.slides_written == 1  # only slide 1 has notes from original
        assert result.slides_skipped == 2

    def test_creates_notes_frame_when_missing(self, cataloged):
        deck_id, deck_path, db_path = cataloged
        record_edit(2, "notes", "New notes for slide 2", source="web", db_path=db_path)
        write_notes_to_pptx(deck_path, db_path=db_path, deck_id=deck_id)
        prs = Presentation(deck_path)
        assert prs.slides[1].notes_slide.notes_text_frame.text == "New notes for slide 2"

    def test_output_path_leaves_original_untouched(self, cataloged, tmp_path):
        deck_id, deck_path, db_path = cataloged
        record_edit(1, "notes", "Modified notes", source="web", db_path=db_path)
        output = str(tmp_path / "copy.pptx")
        result = write_notes_to_pptx(deck_path, db_path=db_path, deck_id=deck_id, output_path=output)
        prs_out = Presentation(output)
        assert prs_out.slides[0].notes_slide.notes_text_frame.text == "Modified notes"
        prs_orig = Presentation(deck_path)
        assert prs_orig.slides[0].notes_slide.notes_text_frame.text == "Original notes 1"

    def test_lookup_by_file_path(self, cataloged):
        deck_id, deck_path, db_path = cataloged
        record_edit(1, "notes", "Via path lookup", source="web", db_path=db_path)
        result = write_notes_to_pptx(deck_path, db_path=db_path)
        assert result.deck_id == deck_id
        assert result.slides_written == 1


class TestWritebackErrors:
    def test_file_not_found(self, db_path):
        with pytest.raises(FileNotFoundError):
            write_notes_to_pptx("/nonexistent/deck.pptx", db_path=db_path)

    def test_deck_not_in_db(self, deck_with_notes, db_path):
        with pytest.raises(ValueError, match="not found in database"):
            write_notes_to_pptx(deck_with_notes, db_path=db_path)

    def test_slide_count_mismatch(self, cataloged, tmp_path):
        deck_id, deck_path, db_path = cataloged
        prs = Presentation(deck_path)
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(deck_path)
        with pytest.raises(ValueError, match="Slide count mismatch"):
            write_notes_to_pptx(deck_path, db_path=db_path, deck_id=deck_id)


class TestCreateBackup:
    def test_creates_bak_file(self, deck_with_notes):
        backup_path = create_backup(deck_with_notes)
        assert os.path.exists(backup_path)
        assert backup_path.endswith(".pptx.bak")
        assert os.path.dirname(backup_path) == os.path.dirname(deck_with_notes)

    def test_backup_is_valid_pptx(self, deck_with_notes):
        backup_path = create_backup(deck_with_notes)
        prs = Presentation(backup_path)
        assert len(prs.slides) == 3

    def test_file_not_found(self):
        with pytest.raises(FileNotFoundError):
            create_backup("/nonexistent/deck.pptx")


class TestCmdWriteNotes:
    def test_writes_notes_and_creates_backup(self, cataloged, capsys):
        deck_id, deck_path, db_path = cataloged
        record_edit(1, "notes", "CLI edited", source="web", db_path=db_path)

        parser = build_parser()
        args = parser.parse_args(["write-notes", deck_path, "--db", db_path])
        result = cmd_write_notes(args)

        assert result == 0
        captured = capsys.readouterr()
        assert "Wrote notes to 1" in captured.out
        assert ".pptx.bak" in captured.out

        # Verify PPTX was modified
        prs = Presentation(deck_path)
        assert prs.slides[0].notes_slide.notes_text_frame.text == "CLI edited"

    def test_error_on_missing_file(self, capsys):
        parser = build_parser()
        args = parser.parse_args(["write-notes", "/nonexistent.pptx"])
        result = cmd_write_notes(args)

        assert result == 1
        captured = capsys.readouterr()
        assert "not found" in captured.err.lower() or "not found" in captured.out.lower()
