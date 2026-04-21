"""Tests for the 'metadata' CLI subcommand."""

import argparse
import json

import pytest
from pptx import Presentation

from aippt.cli import cmd_metadata


def _make_args(**kwargs):
    defaults = {
        "command": "metadata",
        "deck": None,
        "slide": None,
        "debug": False,
    }
    defaults.update(kwargs)
    return argparse.Namespace(**defaults)


class TestMetadataCmd:
    def test_no_metadata_in_deck(self, tmp_path, capsys):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = str(tmp_path / "test.pptx")
        prs.save(path)
        args = _make_args(deck=path)
        result = cmd_metadata(args)
        assert result == 0
        output = json.loads(capsys.readouterr().out)
        assert output == []

    def test_metadata_present(self, tmp_path, capsys):
        from aippt.metadata import append_metadata
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "Slide One"
        append_metadata(slide, "enhance", model="claude-sonnet-4-6")
        path = str(tmp_path / "test.pptx")
        prs.save(path)
        args = _make_args(deck=path)
        result = cmd_metadata(args)
        assert result == 0
        output = json.loads(capsys.readouterr().out)
        assert len(output) == 1
        assert output[0]["slide"] == 1
        assert output[0]["title"] == "Slide One"
        assert len(output[0]["operations"]) == 1

    def test_slide_filter(self, tmp_path, capsys):
        from aippt.metadata import append_metadata
        prs = Presentation()
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s2 = prs.slides.add_slide(prs.slide_layouts[0])
        append_metadata(s1, "enhance", model="m1")
        append_metadata(s2, "improve", model="m2")
        path = str(tmp_path / "test.pptx")
        prs.save(path)
        args = _make_args(deck=path, slide=2)
        result = cmd_metadata(args)
        assert result == 0
        output = json.loads(capsys.readouterr().out)
        assert len(output) == 1
        assert output[0]["slide"] == 2

    def test_file_not_found(self, capsys):
        args = _make_args(deck="/nonexistent.pptx")
        result = cmd_metadata(args)
        assert result == 1
