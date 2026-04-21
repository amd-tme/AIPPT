"""Tests for web app relative path resolution."""
import os
import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from aippt.catalog import catalog_deck, get_db
from aippt.web.app import create_app


@pytest.fixture
def project_dir(tmp_path):
    """Create a project directory structure with uploads and images."""
    uploads = tmp_path / "uploads"
    uploads.mkdir()
    images = tmp_path / "images" / "test"
    images.mkdir(parents=True)

    # Create a minimal PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Slide"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Test notes"
    deck_path = uploads / "test.pptx"
    prs.save(str(deck_path))

    # Create a fake slide image
    (images / "Slide1.PNG").write_bytes(b"\x89PNG fake image data")

    return tmp_path


@pytest.fixture
def client_relative(project_dir):
    """Create a TestClient with relative paths in the database."""
    db_path = str(project_dir / "slides.db")
    uploads_dir = str(project_dir / "uploads")
    deck_path = str(project_dir / "uploads" / "test.pptx")
    images_dir = str(project_dir / "images" / "test")

    catalog_deck(
        deck_path,
        db_path=db_path,
        images_dir=images_dir,
        base_dir=str(project_dir),
    )

    app = create_app(
        db_path=db_path,
        uploads_dir=uploads_dir,
        project_root=str(project_dir),
    )
    return TestClient(app)


class TestSlideImageRelativePath:
    """Verify /slide-image/{id} resolves relative image_path."""

    def test_serves_image_from_relative_path(self, client_relative):
        """Slide image endpoint resolves relative path and serves the file."""
        resp = client_relative.get("/slide-image/1")
        assert resp.status_code == 200
        assert b"PNG" in resp.content

    def test_relative_path_stored_in_db(self, client_relative):
        """Verify the DB actually has a relative path (not absolute)."""
        db_path = client_relative.app.state.db_path
        conn = get_db(db_path)
        row = conn.execute("SELECT image_path FROM slides WHERE id = 1").fetchone()
        conn.close()
        assert row["image_path"] is not None
        assert not os.path.isabs(row["image_path"])

    def test_missing_image_returns_404(self, project_dir):
        """Slide with no image returns 404."""
        db_path = str(project_dir / "no_images.db")
        uploads_dir = str(project_dir / "uploads")
        deck_path = str(project_dir / "uploads" / "test.pptx")

        # Catalog without images_dir
        catalog_deck(deck_path, db_path=db_path, base_dir=str(project_dir))

        app = create_app(
            db_path=db_path,
            uploads_dir=uploads_dir,
            project_root=str(project_dir),
        )
        client = TestClient(app)
        resp = client.get("/slide-image/1")
        assert resp.status_code == 404


class TestDownloadRelativePath:
    """Verify /api/decks/{id}/download resolves relative file_path."""

    def test_download_resolves_relative_path(self, client_relative):
        """Download endpoint resolves relative file_path and serves PPTX."""
        resp = client_relative.get("/api/decks/1/download")
        assert resp.status_code == 200
        content_type = resp.headers.get("content-type", "")
        assert "openxmlformats" in content_type or "application/" in content_type

    def test_download_file_path_is_relative(self, client_relative):
        """Verify the DB has a relative file_path."""
        db_path = client_relative.app.state.db_path
        conn = get_db(db_path)
        row = conn.execute("SELECT file_path FROM decks WHERE id = 1").fetchone()
        conn.close()
        assert not os.path.isabs(row["file_path"])


class TestAbsolutePathBackcompat:
    """Verify absolute paths still work during transition."""

    def test_absolute_image_path_still_works(self, project_dir):
        """If DB has absolute paths (pre-migration), images still serve."""
        db_path = str(project_dir / "abs.db")
        uploads_dir = str(project_dir / "uploads")

        # Manually insert with absolute paths
        conn = get_db(db_path)
        abs_deck_path = str(project_dir / "uploads" / "test.pptx")
        abs_image_path = str(project_dir / "images" / "test" / "Slide1.PNG")

        conn.execute(
            "INSERT INTO decks (name, file_path, file_hash, slide_count) VALUES (?, ?, ?, ?)",
            ("test", abs_deck_path, "abc123", 1),
        )
        conn.execute(
            "INSERT INTO slides (deck_id, position, title, content_text, content_hash, image_path) "
            "VALUES (?, ?, ?, ?, ?, ?)",
            (1, 1, "Slide 1", "content", "hash1", abs_image_path),
        )
        conn.commit()
        conn.close()

        app = create_app(
            db_path=db_path,
            uploads_dir=uploads_dir,
            project_root=str(project_dir),
        )
        client = TestClient(app)

        resp = client.get("/slide-image/1")
        assert resp.status_code == 200
        assert b"PNG" in resp.content
