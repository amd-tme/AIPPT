"""Tests for the remix module."""
import pytest
import yaml
from pptx import Presentation

from aippt.remix import generate_manifest, load_manifest, copy_slide, assemble_deck
from aippt.catalog import get_db


class TestGenerateManifest:
    def test_basic(self):
        slides = [
            {"deck_path": "/path/deck.pptx", "position": 1, "title": "Intro"},
            {"deck_path": "/path/deck.pptx", "position": 5, "title": "Summary"},
        ]
        result = generate_manifest(slides, title="Test Deck")
        parsed = yaml.safe_load(result)
        assert parsed["title"] == "Test Deck"
        assert len(parsed["slides"]) == 2
        assert parsed["slides"][0]["position"] == 1
        assert parsed["slides"][1]["position"] == 5

    def test_default_title_and_template(self):
        slides = [{"deck_path": "/path/a.pptx", "position": 1, "title": "Slide 1"}]
        result = generate_manifest(slides)
        parsed = yaml.safe_load(result)
        assert parsed["title"] == "Remixed Presentation"
        assert parsed["template"] == "template.pptx"

    def test_empty_slides(self):
        result = generate_manifest([], title="Empty")
        parsed = yaml.safe_load(result)
        assert parsed["slides"] == []

    def test_preserves_deck_basename(self):
        slides = [{"deck_path": "/long/path/to/my-deck.pptx", "position": 3, "title": "T"}]
        result = generate_manifest(slides)
        parsed = yaml.safe_load(result)
        assert parsed["slides"][0]["deck"] == "my-deck.pptx"
        assert parsed["slides"][0]["deck_path"] == "/long/path/to/my-deck.pptx"


class TestLoadManifest:
    def test_valid(self, tmp_path):
        f = tmp_path / "manifest.yaml"
        f.write_text(yaml.dump({
            "title": "Test",
            "slides": [{"deck": "a.pptx", "position": 1, "title": "Slide 1"}],
        }))
        result = load_manifest(str(f))
        assert result["title"] == "Test"
        assert len(result["slides"]) == 1

    def test_missing_title(self, tmp_path):
        f = tmp_path / "manifest.yaml"
        f.write_text(yaml.dump({"slides": []}))
        with pytest.raises(ValueError, match="missing required key"):
            load_manifest(str(f))

    def test_missing_slides(self, tmp_path):
        f = tmp_path / "manifest.yaml"
        f.write_text(yaml.dump({"title": "Test"}))
        with pytest.raises(ValueError, match="missing required key"):
            load_manifest(str(f))

    def test_preserves_extra_fields(self, tmp_path):
        f = tmp_path / "manifest.yaml"
        f.write_text(yaml.dump({
            "title": "Test",
            "template": "custom.pptx",
            "slides": [{"deck": "a.pptx", "position": 1, "title": "S1"}],
        }))
        result = load_manifest(str(f))
        assert result["template"] == "custom.pptx"


class TestCopySlide:
    @pytest.fixture
    def source_pptx(self, tmp_path):
        """Create a minimal source PPTX with two slides."""
        prs = Presentation()
        layout = prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Slide One"

        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Slide Two"

        path = str(tmp_path / "source.pptx")
        prs.save(path)
        return path

    def test_copy_slide_creates_new_slide(self, source_pptx):
        source_prs = Presentation(source_pptx)
        target_prs = Presentation()

        initial_count = len(target_prs.slides)
        copy_slide(source_prs, 0, target_prs)
        assert len(target_prs.slides) == initial_count + 1

    def test_copy_multiple_slides(self, source_pptx):
        source_prs = Presentation(source_pptx)
        target_prs = Presentation()

        copy_slide(source_prs, 0, target_prs)
        copy_slide(source_prs, 1, target_prs)
        assert len(target_prs.slides) == 2


class TestAssembleDeck:
    @pytest.fixture
    def source_deck(self, tmp_path):
        """Create a minimal source PPTX."""
        prs = Presentation()
        layout = prs.slide_layouts[0]

        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Intro"

        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Details"

        slide3 = prs.slides.add_slide(layout)
        slide3.shapes.title.text = "Summary"

        path = str(tmp_path / "source.pptx")
        prs.save(path)
        return path

    def test_assemble_from_manifest(self, tmp_path, source_deck):
        db_path = str(tmp_path / "test.db")
        get_db(db_path).close()

        manifest = {
            "title": "Assembled Deck",
            "slides": [
                {"deck": "source.pptx", "deck_path": source_deck, "position": 1, "title": "Intro"},
                {"deck": "source.pptx", "deck_path": source_deck, "position": 3, "title": "Summary"},
            ],
        }
        manifest_path = str(tmp_path / "manifest.yaml")
        with open(manifest_path, "w") as f:
            yaml.dump(manifest, f)

        output_path = str(tmp_path / "output.pptx")
        count = assemble_deck(manifest_path, output_path, db_path=db_path)

        assert count == 2
        result_prs = Presentation(output_path)
        assert len(result_prs.slides) == 2

    def test_assemble_skips_missing_deck(self, tmp_path):
        db_path = str(tmp_path / "test.db")
        get_db(db_path).close()

        # Create a source deck so the first slide works as template
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Real Slide"
        real_path = str(tmp_path / "real.pptx")
        prs.save(real_path)

        manifest = {
            "title": "Test",
            "slides": [
                {"deck": "real.pptx", "deck_path": real_path, "position": 1, "title": "Real Slide"},
                {"deck": "missing.pptx", "deck_path": "/nonexistent/missing.pptx", "position": 1, "title": "Missing"},
            ],
        }
        manifest_path = str(tmp_path / "manifest.yaml")
        with open(manifest_path, "w") as f:
            yaml.dump(manifest, f)

        output_path = str(tmp_path / "output.pptx")
        count = assemble_deck(manifest_path, output_path, db_path=db_path)
        assert count == 1  # Only the real slide copied

    def test_assemble_skips_invalid_position(self, tmp_path, source_deck):
        db_path = str(tmp_path / "test.db")
        get_db(db_path).close()

        manifest = {
            "title": "Test",
            "slides": [
                {"deck": "source.pptx", "deck_path": source_deck, "position": 1, "title": "Intro"},
                {"deck": "source.pptx", "deck_path": source_deck, "position": 99, "title": "Invalid"},
            ],
        }
        manifest_path = str(tmp_path / "manifest.yaml")
        with open(manifest_path, "w") as f:
            yaml.dump(manifest, f)

        output_path = str(tmp_path / "output.pptx")
        count = assemble_deck(manifest_path, output_path, db_path=db_path)
        assert count == 1  # Only the valid slide copied
