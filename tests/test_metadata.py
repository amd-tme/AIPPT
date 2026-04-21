"""Tests for the metadata module."""

import json

import pytest
from pptx import Presentation


class TestFormatNotesWithMetadata:
    def test_notes_with_single_entry(self):
        from aippt.metadata import format_notes_with_metadata
        entries = [{"operation": "enhance", "timestamp": "2026-03-08T14:30:00Z", "model": "claude-sonnet-4-6"}]
        result = format_notes_with_metadata("Speaker notes here.", entries)
        assert result.startswith("Speaker notes here.")
        assert "\n\n---\n[AIPPT-META]\n" in result
        assert "[/AIPPT-META]" in result
        parsed = json.loads(result.split("[AIPPT-META]\n")[1].split("\n[/AIPPT-META]")[0])
        assert parsed == entries

    def test_empty_notes_with_metadata(self):
        from aippt.metadata import format_notes_with_metadata
        entries = [{"operation": "enhance"}]
        result = format_notes_with_metadata("", entries)
        assert result.startswith("[AIPPT-META]\n")
        assert not result.startswith("\n\n---")

    def test_empty_entries_returns_notes_only(self):
        from aippt.metadata import format_notes_with_metadata
        result = format_notes_with_metadata("Just notes.", [])
        assert result == "Just notes."

    def test_multiple_entries(self):
        from aippt.metadata import format_notes_with_metadata
        entries = [
            {"operation": "enhance", "model": "claude-sonnet-4-6"},
            {"operation": "improve", "model": "claude-sonnet-4-6"},
        ]
        result = format_notes_with_metadata("Notes.", entries)
        parsed = json.loads(result.split("[AIPPT-META]\n")[1].split("\n[/AIPPT-META]")[0])
        assert len(parsed) == 2


class TestExtractMetadata:
    def test_extract_from_notes_with_metadata(self):
        from aippt.metadata import extract_metadata
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        entries = [{"operation": "enhance", "model": "claude-sonnet-4-6"}]
        slide.notes_slide.notes_text_frame.text = (
            "Human notes\n\n---\n[AIPPT-META]\n"
            + json.dumps(entries, indent=2)
            + "\n[/AIPPT-META]"
        )
        result = extract_metadata(slide)
        assert result == entries

    def test_extract_no_metadata(self):
        from aippt.metadata import extract_metadata
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Just human notes"
        result = extract_metadata(slide)
        assert result == []

    def test_extract_malformed_json(self):
        from aippt.metadata import extract_metadata
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = (
            "Notes\n\n---\n[AIPPT-META]\n{bad json\n[/AIPPT-META]"
        )
        result = extract_metadata(slide)
        assert result == []

    def test_extract_empty_notes(self):
        from aippt.metadata import extract_metadata
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        result = extract_metadata(slide)
        assert result == []


class TestExtractNotesText:
    def test_notes_before_separator(self):
        from aippt.metadata import extract_notes_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = (
            "Human notes here\n\n---\n[AIPPT-META]\n[]\n[/AIPPT-META]"
        )
        result = extract_notes_text(slide)
        assert result == "Human notes here"

    def test_notes_no_metadata(self):
        from aippt.metadata import extract_notes_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Just regular notes"
        result = extract_notes_text(slide)
        assert result == "Just regular notes"

    def test_empty_notes(self):
        from aippt.metadata import extract_notes_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        result = extract_notes_text(slide)
        assert result == ""


class TestAppendMetadata:
    def test_append_to_empty_notes(self):
        from aippt.metadata import append_metadata, extract_metadata, extract_notes_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        append_metadata(slide, "enhance", model="claude-sonnet-4-6", layout_selected="bullet")
        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert entries[0]["operation"] == "enhance"
        assert entries[0]["model"] == "claude-sonnet-4-6"
        assert entries[0]["layout_selected"] == "bullet"
        assert "timestamp" in entries[0]

    def test_append_to_existing_notes(self):
        from aippt.metadata import append_metadata, extract_metadata, extract_notes_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Existing human notes"
        append_metadata(slide, "enhance", model="claude-sonnet-4-6")
        assert extract_notes_text(slide) == "Existing human notes"
        assert len(extract_metadata(slide)) == 1

    def test_append_to_existing_metadata(self):
        from aippt.metadata import append_metadata, extract_metadata
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Notes"
        append_metadata(slide, "enhance", model="m1")
        append_metadata(slide, "improve", model="m2", focus="brevity")
        entries = extract_metadata(slide)
        assert len(entries) == 2
        assert entries[0]["operation"] == "enhance"
        assert entries[1]["operation"] == "improve"
        assert entries[1]["focus"] == "brevity"

    def test_preserves_human_notes_across_appends(self):
        from aippt.metadata import append_metadata, extract_notes_text
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "My important notes"
        append_metadata(slide, "enhance", model="m1")
        append_metadata(slide, "improve", model="m2")
        assert extract_notes_text(slide) == "My important notes"
