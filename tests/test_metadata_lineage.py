"""Tests for metadata lineage tracking helpers."""
import json

import pytest
from pptx import Presentation

from aippt.metadata import (
    create_lineage_entry,
    append_history_entry,
    get_slide_lineage,
    append_metadata,
    extract_metadata,
    extract_notes_text,
    format_notes_with_metadata,
    HISTORY_CAP,
)


class TestCreateLineageEntry:
    def test_basic_lineage(self):
        entry = create_lineage_entry("outline -> pptxgenjs")
        assert entry["source"] == "outline -> pptxgenjs"
        assert "created" in entry
        assert len(entry["history"]) == 1
        assert "Created from outline -> pptxgenjs" in entry["history"][0]

    def test_lineage_with_layout_and_theme(self):
        entry = create_lineage_entry(
            "outline -> pptxgenjs", layout="bullet", theme="amd"
        )
        assert entry["layout"] == "bullet"
        assert entry["theme"] == "amd"
        assert "(bullet layout)" in entry["history"][0]

    def test_lineage_custom_created_date(self):
        entry = create_lineage_entry("outline -> pptxgenjs", created="2026-03-11")
        assert entry["created"] == "2026-03-11"

    def test_no_layout_key_when_none(self):
        entry = create_lineage_entry("outline -> pptxgenjs")
        assert "layout" not in entry

    def test_no_theme_key_when_none(self):
        entry = create_lineage_entry("outline -> pptxgenjs")
        assert "theme" not in entry


class TestAppendHistoryEntry:
    def test_append_to_existing_lineage(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Set up initial lineage
        lineage = create_lineage_entry("outline -> pptxgenjs", layout="bullet")
        append_metadata(slide, "create", **lineage)

        append_history_entry(slide, "Changed to two_column layout", "/edit-deck")

        entries = extract_metadata(slide)
        # The lineage entry should be updated in-place
        lineage_entry = None
        for e in entries:
            if "history" in e:
                lineage_entry = e
                break
        assert lineage_entry is not None
        assert len(lineage_entry["history"]) == 2
        assert "[/edit-deck]" in lineage_entry["history"][1]
        assert "two_column" in lineage_entry["history"][1]

    def test_append_without_source_tag(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        lineage = create_lineage_entry("outline -> pptxgenjs")
        append_metadata(slide, "create", **lineage)

        append_history_entry(slide, "Fixed text overflow")

        entries = extract_metadata(slide)
        for e in entries:
            if "history" in e:
                assert "[" not in e["history"][-1] or "Created" in e["history"][-1]

    def test_append_creates_entry_when_no_lineage(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Human notes"

        append_history_entry(slide, "Edited content", "/edit-deck")

        entries = extract_metadata(slide)
        assert len(entries) == 1
        assert entries[0]["operation"] == "edit"
        assert len(entries[0]["history"]) == 1
        assert "[/edit-deck]" in entries[0]["history"][0]

    def test_history_cap(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        lineage = create_lineage_entry("outline -> pptxgenjs")
        append_metadata(slide, "create", **lineage)

        for i in range(HISTORY_CAP + 5):
            append_history_entry(slide, f"Edit {i}")

        entries = extract_metadata(slide)
        for e in entries:
            if "history" in e:
                assert len(e["history"]) == HISTORY_CAP
                # Oldest entries (including the "Created" one) should be trimmed
                assert "Edit 5" in e["history"][0]
                break
        else:
            pytest.fail("No entry with history found")

    def test_preserves_human_notes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Important notes"
        lineage = create_lineage_entry("outline -> pptxgenjs")
        append_metadata(slide, "create", **lineage)

        append_history_entry(slide, "Changed layout")

        assert extract_notes_text(slide) == "Important notes"


class TestGetSlideLineage:
    def test_full_lineage(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        lineage = create_lineage_entry(
            "outline -> pptxgenjs", layout="bullet", theme="amd", created="2026-03-11"
        )
        append_metadata(slide, "create", **lineage)

        result = get_slide_lineage(slide)
        assert result["source"] == "outline -> pptxgenjs"
        assert result["created"] == "2026-03-11"
        assert result["layout"] == "bullet"
        assert result["theme"] == "amd"
        assert len(result["history"]) == 1

    def test_no_lineage(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "Just notes"

        result = get_slide_lineage(slide)
        assert result["source"] is None
        assert result["created"] is None
        assert result["layout"] is None
        assert result["theme"] is None
        assert result["history"] == []

    def test_empty_notes(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        result = get_slide_lineage(slide)
        assert result["source"] is None
        assert result["history"] == []

    def test_layout_updated_by_later_entry(self):
        """get_slide_lineage should return the latest layout."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        lineage = create_lineage_entry("outline -> pptxgenjs", layout="bullet")
        append_metadata(slide, "create", **lineage)
        # Simulate an edit that changes layout
        append_metadata(slide, "edit", layout="two_column")

        result = get_slide_lineage(slide)
        assert result["layout"] == "two_column"

    def test_backward_compatible_with_old_metadata(self):
        """Slides with pre-lineage metadata (no source/history) should work."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        append_metadata(slide, "enhance", model="claude-sonnet-4-6", layout_selected="bullet")

        result = get_slide_lineage(slide)
        assert result["source"] is None
        assert result["history"] == []
