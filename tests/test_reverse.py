"""Tests for aippt.reverse module."""

import os
import pytest
from unittest.mock import MagicMock, patch
from pptx import Presentation

from aippt.reverse import (
    extract_text_from_shape,
    convert_pptx_to_outline,
    _resolve_slide_image,
    _enhance_slide_with_llm,
    _should_skip_shape,
    _extract_slide_title,
    REVERSE_SYSTEM_PROMPT,
    REVERSE_TEXT_ONLY_SYSTEM_PROMPT,
)


class TestExtractTextFromShape:
    def test_extracts_simple_text(self):
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        para = MagicMock()
        para.text = "Simple text content"
        para.level = 0
        mock_shape.text_frame.paragraphs = [para]
        del mock_shape.shapes  # not a group
        mock_shape.has_table = False

        result = extract_text_from_shape(mock_shape)
        assert result == "- Simple text content"

    def test_extracts_text_with_whitespace(self):
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        para = MagicMock()
        para.text = "  Text with spaces  "
        para.level = 0
        mock_shape.text_frame.paragraphs = [para]
        del mock_shape.shapes
        mock_shape.has_table = False

        result = extract_text_from_shape(mock_shape)
        assert result == "- Text with spaces"

    def test_handles_group_shapes(self):
        mock_subshape1 = MagicMock()
        mock_subshape1.has_text_frame = True
        p1 = MagicMock(); p1.text = "First"; p1.level = 0
        mock_subshape1.text_frame.paragraphs = [p1]
        del mock_subshape1.shapes
        mock_subshape1.has_table = False

        mock_subshape2 = MagicMock()
        mock_subshape2.has_text_frame = True
        p2 = MagicMock(); p2.text = "Second"; p2.level = 0
        mock_subshape2.text_frame.paragraphs = [p2]
        del mock_subshape2.shapes
        mock_subshape2.has_table = False

        mock_group = MagicMock()
        mock_group.has_text_frame = False
        mock_group.shapes = [mock_subshape1, mock_subshape2]
        mock_group.has_table = False

        result = extract_text_from_shape(mock_group)
        assert "- First" in result
        assert "- Second" in result

    def test_handles_table_shapes(self):
        mock_cell1 = MagicMock()
        mock_cell1.text = "Header 1"
        mock_cell2 = MagicMock()
        mock_cell2.text = "Header 2"

        mock_cell3 = MagicMock()
        mock_cell3.text = "Data 1"
        mock_cell4 = MagicMock()
        mock_cell4.text = "Data 2"

        mock_row1 = MagicMock()
        mock_row1.cells = [mock_cell1, mock_cell2]
        mock_row2 = MagicMock()
        mock_row2.cells = [mock_cell3, mock_cell4]

        mock_table = MagicMock()
        mock_table.rows = [mock_row1, mock_row2]
        mock_table.columns = [MagicMock(), MagicMock()]  # 2 columns

        mock_shape = MagicMock()
        mock_shape.has_text_frame = False
        mock_shape.has_table = True
        mock_shape.table = mock_table
        del mock_shape.shapes

        result = extract_text_from_shape(mock_shape)
        assert "| Header 1 | Header 2 |" in result
        assert "| --- | --- |" in result
        assert "| Data 1 | Data 2 |" in result

    def test_returns_empty_for_empty_shape(self):
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        para = MagicMock()
        para.text = ""
        mock_shape.text_frame.paragraphs = [para]
        del mock_shape.shapes
        mock_shape.has_table = False

        result = extract_text_from_shape(mock_shape)
        assert result == ""


class TestBulletHierarchy:
    def test_preserves_paragraph_levels(self):
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        p0 = MagicMock(); p0.text = "Top level"; p0.level = 0
        p1 = MagicMock(); p1.text = "Indented"; p1.level = 1
        p2 = MagicMock(); p2.text = "Deep indent"; p2.level = 2
        mock_shape.text_frame.paragraphs = [p0, p1, p2]
        del mock_shape.shapes
        mock_shape.has_table = False

        result = extract_text_from_shape(mock_shape)
        lines = result.split("\n")
        assert lines[0] == "- Top level"
        assert lines[1] == "  - Indented"
        assert lines[2] == "    - Deep indent"

    def test_none_level_treated_as_zero(self):
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        p = MagicMock(); p.text = "No level set"; p.level = None
        mock_shape.text_frame.paragraphs = [p]
        del mock_shape.shapes
        mock_shape.has_table = False

        result = extract_text_from_shape(mock_shape)
        assert result == "- No level set"

    def test_skips_empty_paragraphs(self):
        mock_shape = MagicMock()
        mock_shape.has_text_frame = True
        p0 = MagicMock(); p0.text = "Content"; p0.level = 0
        p1 = MagicMock(); p1.text = ""; p1.level = 0
        p2 = MagicMock(); p2.text = "More content"; p2.level = 0
        mock_shape.text_frame.paragraphs = [p0, p1, p2]
        del mock_shape.shapes
        mock_shape.has_table = False

        result = extract_text_from_shape(mock_shape)
        lines = result.split("\n")
        assert len(lines) == 2
        assert lines[0] == "- Content"
        assert lines[1] == "- More content"


class TestTableFormatting:
    def test_single_row_table(self):
        """Single row table still gets header separator."""
        mock_cell1 = MagicMock(); mock_cell1.text = "Only"
        mock_cell2 = MagicMock(); mock_cell2.text = "Row"

        mock_row = MagicMock()
        mock_row.cells = [mock_cell1, mock_cell2]

        mock_table = MagicMock()
        mock_table.rows = [mock_row]
        mock_table.columns = [MagicMock(), MagicMock()]

        mock_shape = MagicMock()
        mock_shape.has_text_frame = False
        mock_shape.has_table = True
        mock_shape.table = mock_table
        del mock_shape.shapes

        result = extract_text_from_shape(mock_shape)
        assert "| Only | Row |" in result
        assert "| --- | --- |" in result


class TestConvertPptxToOutline:
    @pytest.fixture
    def sample_pptx(self, tmp_path):
        """Create a minimal PPTX for testing."""
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[0]

        # Slide 1
        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Introduction"

        # Slide 2
        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Main Content"

        path = tmp_path / "test.pptx"
        prs.save(str(path))
        return str(path)

    def test_creates_markdown_file(self, sample_pptx, tmp_path):
        output = str(tmp_path / "output.md")

        result = convert_pptx_to_outline(sample_pptx, output)

        assert result is True
        assert os.path.exists(output)

    def test_includes_slide_titles(self, sample_pptx, tmp_path):
        output = str(tmp_path / "output.md")

        convert_pptx_to_outline(sample_pptx, output)

        with open(output) as f:
            content = f.read()

        assert "# Introduction" in content
        assert "# Main Content" in content

    def test_handles_missing_file(self, tmp_path):
        output = str(tmp_path / "output.md")

        result = convert_pptx_to_outline("/nonexistent/file.pptx", output)

        assert result is False

    def test_exclude_notes_option(self, tmp_path):
        """Test that notes can be excluded."""
        from pptx import Presentation

        # Create PPTX with notes
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes here"

        pptx_path = tmp_path / "with_notes.pptx"
        prs.save(str(pptx_path))

        # Convert without notes
        output = str(tmp_path / "output.md")
        convert_pptx_to_outline(str(pptx_path), output, include_notes=False)

        with open(output) as f:
            content = f.read()

        assert "Speaker notes here" not in content

    def test_include_notes_by_default(self, tmp_path):
        """Test that notes are included by default."""
        from pptx import Presentation

        # Create PPTX with notes
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes here"

        pptx_path = tmp_path / "with_notes.pptx"
        prs.save(str(pptx_path))

        # Convert with notes (default)
        output = str(tmp_path / "output.md")
        convert_pptx_to_outline(str(pptx_path), output, include_notes=True)

        with open(output) as f:
            content = f.read()

        assert "Speaker notes here" in content


class TestResolveSlideImage:
    def test_finds_png_uppercase(self, tmp_path):
        img = tmp_path / "Slide1.PNG"
        img.write_bytes(b"fake")
        result = _resolve_slide_image(1, str(tmp_path))
        assert result == str(img)

    def test_finds_png_lowercase(self, tmp_path):
        img = tmp_path / "Slide2.png"
        img.write_bytes(b"fake")
        result = _resolve_slide_image(2, str(tmp_path))
        assert result == str(img)

    def test_finds_jpg(self, tmp_path):
        img = tmp_path / "Slide3.jpg"
        img.write_bytes(b"fake")
        result = _resolve_slide_image(3, str(tmp_path))
        assert result == str(img)

    def test_returns_none_when_missing(self, tmp_path):
        result = _resolve_slide_image(99, str(tmp_path))
        assert result is None

    def test_returns_none_with_no_images_dir(self):
        result = _resolve_slide_image(1, None)
        assert result is None

    def test_returns_none_with_nonexistent_dir(self):
        result = _resolve_slide_image(1, "/nonexistent/path")
        assert result is None


class TestEnhanceSlideWithLlm:
    def test_uses_image_path_when_available(self, tmp_path):
        img = tmp_path / "Slide1.PNG"
        img.write_bytes(b"fake image data")
        client = MagicMock()
        client.generate_text_with_image.return_value = "## Title\n- Point"

        result = _enhance_slide_with_llm(client, 1, "Title", "raw text", str(img))

        assert result == "## Title\n- Point"
        client.generate_text_with_image.assert_called_once()
        client.generate_text.assert_not_called()

    def test_falls_back_to_text_when_no_image(self, tmp_path):
        client = MagicMock()
        client.generate_text.return_value = "## Title\n- Point"

        result = _enhance_slide_with_llm(client, 1, "Title", "raw text", None)

        assert result == "## Title\n- Point"
        client.generate_text.assert_called_once()
        client.generate_text_with_image.assert_not_called()

    def test_falls_back_to_text_when_image_call_fails(self, tmp_path):
        img = tmp_path / "Slide1.PNG"
        img.write_bytes(b"fake image data")
        client = MagicMock()
        client.generate_text_with_image.side_effect = Exception("vision error")
        client.generate_text.return_value = "## Title\n- Fallback"

        result = _enhance_slide_with_llm(client, 1, "Title", "raw text", str(img))

        assert result == "## Title\n- Fallback"
        client.generate_text.assert_called_once()

    def test_returns_none_when_both_fail(self, tmp_path):
        client = MagicMock()
        client.generate_text.side_effect = Exception("API Error")
        client.generate_text_with_image.side_effect = Exception("API Error")

        result = _enhance_slide_with_llm(client, 1, "Title", "raw text", None)

        assert result is None

    def test_passes_correct_system_prompt_for_image(self, tmp_path):
        img = tmp_path / "Slide1.PNG"
        img.write_bytes(b"fake")
        client = MagicMock()
        client.generate_text_with_image.return_value = "## T"

        _enhance_slide_with_llm(client, 1, "T", "text", str(img))

        call_kwargs = client.generate_text_with_image.call_args
        assert call_kwargs.kwargs["system_prompt"] == REVERSE_SYSTEM_PROMPT

    def test_passes_correct_system_prompt_for_text_only(self):
        client = MagicMock()
        client.generate_text.return_value = "## T"

        _enhance_slide_with_llm(client, 1, "T", "text", None)

        call_kwargs = client.generate_text.call_args
        assert call_kwargs.kwargs["system_prompt"] == REVERSE_TEXT_ONLY_SYSTEM_PROMPT


class TestEnhancedReverse:
    @pytest.fixture
    def sample_pptx_path(self, tmp_path):
        from pptx import Presentation
        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(layout)
        slide1.shapes.title.text = "Architecture Overview"
        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Network Topology"
        path = tmp_path / "test.pptx"
        prs.save(str(path))
        return str(path)

    @pytest.fixture
    def mock_llm_client(self):
        client = MagicMock()
        client.generate_text_with_image.return_value = (
            "## Architecture Overview\n\n"
            "- System uses microservices architecture\n"
            "- Three main components: API, Worker, Database\n"
        )
        client.generate_text.return_value = (
            "## Architecture Overview\n\n"
            "- System uses microservices architecture\n"
            "- Three main components: API, Worker, Database\n"
        )
        return client

    def test_enhance_uses_llm_output(self, sample_pptx_path, tmp_path, mock_llm_client):
        output = str(tmp_path / "output.md")
        images_dir = str(tmp_path / "images")
        os.makedirs(images_dir)
        # Create non-trivial fake image (>5KB to avoid placeholder detection)
        with open(os.path.join(images_dir, "Slide1.PNG"), "wb") as f:
            f.write(b"\x89PNG" + b"\x00" * 10000)
        result = convert_pptx_to_outline(
            sample_pptx_path, output, include_notes=True,
            enhance=True, llm_client=mock_llm_client, images_dir=images_dir,
        )
        assert result is True
        mock_llm_client.generate_text_with_image.assert_called()
        with open(output) as f:
            content = f.read()
        assert "microservices architecture" in content

    def test_enhance_falls_back_without_images(self, sample_pptx_path, tmp_path, mock_llm_client):
        output = str(tmp_path / "output.md")
        result = convert_pptx_to_outline(
            sample_pptx_path, output, include_notes=True,
            enhance=True, llm_client=mock_llm_client, images_dir=None,
        )
        assert result is True
        mock_llm_client.generate_text.assert_called()

    def test_enhance_graceful_degradation_on_llm_failure(self, sample_pptx_path, tmp_path):
        output = str(tmp_path / "output.md")
        failing_client = MagicMock()
        failing_client.generate_text.side_effect = Exception("API Error")
        failing_client.generate_text_with_image.side_effect = Exception("API Error")
        result = convert_pptx_to_outline(
            sample_pptx_path, output, include_notes=True,
            enhance=True, llm_client=failing_client, images_dir=None,
        )
        assert result is True
        with open(output) as f:
            content = f.read()
        assert "Architecture Overview" in content

    def test_enhance_without_client_falls_back(self, sample_pptx_path, tmp_path):
        output = str(tmp_path / "output.md")
        result = convert_pptx_to_outline(
            sample_pptx_path, output, include_notes=True,
            enhance=True, llm_client=None, images_dir=None,
        )
        assert result is True
        with open(output) as f:
            content = f.read()
        assert "Architecture Overview" in content

    def test_non_enhanced_unchanged(self, sample_pptx_path, tmp_path):
        output = str(tmp_path / "output.md")
        result = convert_pptx_to_outline(sample_pptx_path, output)
        assert result is True
        with open(output) as f:
            content = f.read()
        assert "Architecture Overview" in content


class TestShapeFiltering:
    def test_skips_title_shape(self):
        title = MagicMock()
        assert _should_skip_shape(title, title_shape=title) is True

    def test_skips_slide_number_placeholder(self):
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MagicMock()
        shape.is_placeholder = True
        shape.placeholder_format.type = PP_PLACEHOLDER.SLIDE_NUMBER
        shape.shape_type = 1
        shape.text = "42"
        assert _should_skip_shape(shape) is True

    def test_skips_footer_placeholder(self):
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MagicMock()
        shape.is_placeholder = True
        shape.placeholder_format.type = PP_PLACEHOLDER.FOOTER
        shape.shape_type = 1
        shape.text = "Confidential"
        assert _should_skip_shape(shape) is True

    def test_skips_date_placeholder(self):
        from pptx.enum.shapes import PP_PLACEHOLDER
        shape = MagicMock()
        shape.is_placeholder = True
        shape.placeholder_format.type = PP_PLACEHOLDER.DATE
        shape.shape_type = 1
        shape.text = "2026-03-03"
        assert _should_skip_shape(shape) is True

    def test_skips_connector_line(self):
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shape = MagicMock()
        shape.is_placeholder = False
        shape.shape_type = MSO_SHAPE_TYPE.LINE
        shape.text = ""
        assert _should_skip_shape(shape) is True

    def test_skips_freeform(self):
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shape = MagicMock()
        shape.is_placeholder = False
        shape.shape_type = MSO_SHAPE_TYPE.FREEFORM
        shape.text = ""
        assert _should_skip_shape(shape) is True

    def test_skips_short_numeric_text(self):
        shape = MagicMock()
        shape.is_placeholder = False
        shape.shape_type = 1
        shape.text = "2"
        assert _should_skip_shape(shape) is True

    def test_keeps_normal_content_shape(self):
        shape = MagicMock()
        shape.is_placeholder = False
        shape.shape_type = 1
        shape.text = "Real content here"
        assert _should_skip_shape(shape) is False

    def test_keeps_non_numeric_short_text(self):
        shape = MagicMock()
        shape.is_placeholder = False
        shape.shape_type = 1
        shape.text = "OK"
        assert _should_skip_shape(shape) is False


class TestTitleExtraction:
    def test_standard_title(self):
        """Uses title placeholder when available."""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "My Title"
        assert _extract_slide_title(slide) == "My Title"

    def test_multi_paragraph_title(self):
        """Joins multi-paragraph titles with em dash."""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tf = slide.shapes.title.text_frame
        tf.paragraphs[0].text = "Main Title"
        tf.add_paragraph().text = "Subtitle Line"
        assert _extract_slide_title(slide) == "Main Title — Subtitle Line"

    def test_fallback_to_untitled(self):
        """Returns 'Untitled Slide' when no title source found."""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.paragraphs[0].text = ""
        assert _extract_slide_title(slide) == "Untitled Slide"

    def test_whitespace_only_title_triggers_fallback(self):
        """Title with only whitespace triggers fallback."""
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "   "
        result = _extract_slide_title(slide)
        # Should not be the whitespace — either subtitle, first shape, or Untitled
        assert result.strip() != ""


class TestNotesFormat:
    """Tests that reverse outputs notes as HTML comments, not bullet lists."""

    def test_notes_as_html_comments(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes here"

        pptx_path = str(tmp_path / "test.pptx")
        md_path = str(tmp_path / "test.md")
        prs.save(pptx_path)

        convert_pptx_to_outline(pptx_path, md_path, include_notes=True)

        content = open(md_path).read()
        assert '<!-- notes' in content
        assert 'Speaker notes here' in content
        assert '-->' in content
        assert '*Notes:*' not in content

    def test_notes_html_comment_multiline(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Line one\nLine two\nLine three"

        pptx_path = str(tmp_path / "test.pptx")
        md_path = str(tmp_path / "test.md")
        prs.save(pptx_path)

        convert_pptx_to_outline(pptx_path, md_path, include_notes=True)

        content = open(md_path).read()
        assert '<!-- notes' in content
        assert 'Line one' in content
        assert 'Line two' in content
        assert 'Line three' in content

    def test_strips_analysis_artifacts(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = (
            "Good speaker notes\n"
            "[Note: analysis based on slide text only — no image was available]\n"
            "More notes"
        )

        pptx_path = str(tmp_path / "test.pptx")
        md_path = str(tmp_path / "test.md")
        prs.save(pptx_path)

        convert_pptx_to_outline(pptx_path, md_path, include_notes=True)

        content = open(md_path).read()
        assert 'Good speaker notes' in content
        assert 'More notes' in content
        assert '[Note: analysis based on slide text only' not in content

    def test_no_notes_flag_omits_notes(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes"

        pptx_path = str(tmp_path / "test.pptx")
        md_path = str(tmp_path / "test.md")
        prs.save(pptx_path)

        convert_pptx_to_outline(pptx_path, md_path, include_notes=False)

        content = open(md_path).read()
        assert '<!-- notes' not in content
        assert 'Speaker notes' not in content

    def test_notes_with_arrow_notation_escaped(self, tmp_path):
        """Notes containing --> should be escaped so they don't break the HTML comment."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "A --> B flow"

        pptx_path = str(tmp_path / "test.pptx")
        md_path = str(tmp_path / "test.md")
        prs.save(pptx_path)

        convert_pptx_to_outline(pptx_path, md_path, include_notes=True)

        content = open(md_path).read()
        # The --> should be escaped so HTML comment isn't broken
        assert '<!-- notes' in content
        assert 'A -- > B flow' in content
        assert content.count('-->') == 1  # Only the closing -->
