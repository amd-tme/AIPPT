"""Tests for AI-generated image disclaimer."""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from aippt.layouts import add_disclaimer_textbox


DISCLAIMER_TEXT = "AI-Generated Image -- Not Approved for External Use"


class TestAddDisclaimerTextbox:
    def _make_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        return slide

    def test_adds_textbox_shape(self):
        slide = self._make_slide()
        shape_count_before = len(slide.shapes)
        add_disclaimer_textbox(slide)
        assert len(slide.shapes) == shape_count_before + 1

    def test_textbox_contains_disclaimer_text(self):
        slide = self._make_slide()
        add_disclaimer_textbox(slide)
        tb = slide.shapes[-1]
        assert DISCLAIMER_TEXT in tb.text_frame.text

    def test_custom_disclaimer_text(self):
        slide = self._make_slide()
        custom = "DRAFT -- Internal Only"
        add_disclaimer_textbox(slide, text=custom)
        tb = slide.shapes[-1]
        assert custom in tb.text_frame.text

    def test_textbox_is_near_bottom(self):
        slide = self._make_slide()
        add_disclaimer_textbox(slide)
        tb = slide.shapes[-1]
        assert tb.top >= Inches(6.5)

    def test_textbox_uses_small_font(self):
        slide = self._make_slide()
        add_disclaimer_textbox(slide)
        tb = slide.shapes[-1]
        run = tb.text_frame.paragraphs[0].runs[0]
        assert run.font.size <= Pt(10)
