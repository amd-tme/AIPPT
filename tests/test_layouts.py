"""Tests for aippt.layouts module."""

import io
import pytest
from unittest.mock import MagicMock, patch
from pptx import Presentation
from pptx.util import Pt, Inches

from aippt.layouts import (
    get_layout_index,
    select_slide_layout,
    parse_layout_suggestion,
    split_content_for_columns,
    remove_all_slides,
    _apply_bullets_to_text_frame,
    _auto_number_content,
    apply_layout_content,
    SLIDE_LAYOUTS,
    KNOWN_LAYOUT_TYPES,
)


class TestSlideLayouts:
    def test_slide_layouts_dict_has_required_keys(self):
        assert 'title_and_content' in SLIDE_LAYOUTS
        assert 'two_content' in SLIDE_LAYOUTS
        assert 'title_only' in SLIDE_LAYOUTS

    def test_layout_has_required_fields(self):
        for name, layout in SLIDE_LAYOUTS.items():
            assert 'name' in layout
            assert 'use_case' in layout
            assert 'placeholders' in layout


class TestParseLayoutSuggestion:
    def test_detects_diagram_keyword(self):
        result = parse_layout_suggestion("diagram")
        assert result['type'] == 'diagram'

    def test_detects_diagram_with_description(self):
        result = parse_layout_suggestion("diagram — flow chart showing process")
        assert result['type'] == 'diagram'

    def test_detects_two_column_keyword(self):
        result = parse_layout_suggestion("two_column")
        assert result['type'] == 'two_column'

    def test_detects_two_column_with_description(self):
        result = parse_layout_suggestion("two_column: comparison layout")
        assert result['type'] == 'two_column'

    def test_detects_bullet_keyword(self):
        result = parse_layout_suggestion("bullet")
        assert result['type'] == 'bullet'

    def test_detects_bullet_with_description(self):
        result = parse_layout_suggestion("bullet — a visual list of key points")
        assert result['type'] == 'bullet'

    def test_descriptive_text_defaults_to_basic(self):
        """Descriptive text like 'a visual list' should NOT trigger diagram."""
        result = parse_layout_suggestion("Use a visual list of key points")
        assert result['type'] == 'basic'

    def test_defaults_to_basic(self):
        result = parse_layout_suggestion("Just some regular text")
        assert result['type'] == 'basic'

    def test_basic_keyword(self):
        result = parse_layout_suggestion("basic")
        assert result['type'] == 'basic'

    def test_empty_input_defaults_to_basic(self):
        result = parse_layout_suggestion("")
        assert result['type'] == 'basic'

    def test_returns_required_keys(self):
        result = parse_layout_suggestion("anything")
        assert 'type' in result
        assert 'structure' in result
        assert 'elements' in result


class TestSplitContentForColumns:
    def test_splits_at_midpoint(self):
        content = "Line 1\nLine 2\nLine 3\nLine 4"
        left, right = split_content_for_columns(content)
        assert "Line 1" in left or "Line 2" in left
        assert "Line 3" in right or "Line 4" in right

    def test_splits_at_empty_line_near_middle(self):
        content = "Line 1\nLine 2\n\nLine 4\nLine 5"
        left, right = split_content_for_columns(content)
        assert left == "Line 1\nLine 2"
        assert right == "Line 4\nLine 5"

    def test_handles_single_line(self):
        content = "Single line"
        left, right = split_content_for_columns(content)
        # With a single line, it goes to the right side (midpoint is 0)
        assert right == "Single line"
        assert left == ""

    def test_handles_empty_content(self):
        content = ""
        left, right = split_content_for_columns(content)
        assert left == ""
        assert right == ""


class TestGetLayoutIndex:
    def test_finds_existing_layout(self):
        mock_prs = MagicMock()
        mock_layout_1 = MagicMock()
        mock_layout_1.name = "Title Slide"
        mock_layout_2 = MagicMock()
        mock_layout_2.name = "Title and Content"
        mock_prs.slide_layouts = [mock_layout_1, mock_layout_2]

        result = get_layout_index(mock_prs, "Title and Content")
        assert result == 1

    def test_returns_default_when_not_found(self):
        mock_prs = MagicMock()
        mock_layout = MagicMock()
        mock_layout.name = "Other Layout"
        mock_prs.slide_layouts = [mock_layout]

        result = get_layout_index(mock_prs, "Nonexistent Layout")
        assert result == 3  # Default fallback


class TestSelectSlideLayout:
    def test_selects_title_only_for_diagram(self):
        mock_prs = MagicMock()
        mock_layout_title_only = MagicMock()
        mock_layout_title_only.name = "Title Only"
        mock_layout_content = MagicMock()
        mock_layout_content.name = "Title and Content"
        mock_prs.slide_layouts = [mock_layout_title_only, mock_layout_content]

        result = select_slide_layout(mock_prs, "diagram")
        assert result.name == "Title Only"

    def test_selects_two_content_for_two_column(self):
        mock_prs = MagicMock()
        mock_layout_two = MagicMock()
        mock_layout_two.name = "Two Content"
        mock_layout_default = MagicMock()
        mock_layout_default.name = "Title and Content"
        mock_prs.slide_layouts = [mock_layout_two, mock_layout_default]

        result = select_slide_layout(mock_prs, "two_column")
        assert result.name == "Two Content"

    def test_falls_back_to_first_layout(self):
        mock_prs = MagicMock()
        mock_layout = MagicMock()
        mock_layout.name = "Completely Different"
        mock_prs.slide_layouts = [mock_layout]

        result = select_slide_layout(mock_prs, "bullet")
        assert result == mock_layout


class TestRemoveAllSlides:
    """Tests for remove_all_slides using real Presentation objects."""

    def _make_prs_with_slides(self, n):
        """Create a real Presentation with n blank slides."""
        prs = Presentation()
        layout = prs.slide_layouts[0]
        for _ in range(n):
            prs.slides.add_slide(layout)
        return prs

    def test_removes_slides_from_template_with_slides(self):
        prs = self._make_prs_with_slides(3)
        assert len(prs.slides) == 3
        remove_all_slides(prs)
        assert len(prs.slides) == 0

    def test_empty_presentation_is_noop(self):
        prs = Presentation()
        assert len(prs.slides) == 0
        remove_all_slides(prs)
        assert len(prs.slides) == 0

    def test_layouts_preserved_after_removal(self):
        prs = self._make_prs_with_slides(2)
        layout_count_before = len(prs.slide_layouts)
        remove_all_slides(prs)
        assert len(prs.slide_layouts) == layout_count_before

    def test_saveable_after_removal(self):
        prs = self._make_prs_with_slides(3)
        remove_all_slides(prs)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        reloaded = Presentation(buf)
        assert len(reloaded.slides) == 0


class TestApplyBulletsFormatting:
    """Tests for _apply_bullets_to_text_frame font sizing."""

    def _make_text_frame(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
        return tb.text_frame

    def test_level_0_font_size_sparse(self):
        """Single line (< 3 top-level) uses Pt(26) for sparse slide scaling."""
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "Plain text line")
        run = tf.paragraphs[0].runs[0]
        assert run.font.size == Pt(26)

    def test_level_1_bullet_font_size_sparse(self):
        """Single bullet (< 3 top-level) uses Pt(26) for sparse slide scaling."""
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Bullet item")
        run = tf.paragraphs[0].runs[0]
        assert run.font.size == Pt(26)

    def test_level_2_subbullet_font_size(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "  - Sub bullet")
        run = tf.paragraphs[0].runs[0]
        assert run.font.size == Pt(18)

    def test_mixed_levels_get_correct_sizes_sparse(self):
        """2 top-level bullets + 1 sub: sparse (< 3), so Pt(26) for top-level."""
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Top bullet\n  - Sub bullet\n- Another top")
        assert tf.paragraphs[0].runs[0].font.size == Pt(26)
        assert tf.paragraphs[1].runs[0].font.size == Pt(18)
        assert tf.paragraphs[2].runs[0].font.size == Pt(26)

    def test_mixed_levels_dense_uses_pt22(self):
        """3+ top-level bullets: dense, so Pt(22) for top-level."""
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Bullet one\n- Bullet two\n  - Sub bullet\n- Bullet three")
        assert tf.paragraphs[0].runs[0].font.size == Pt(22)
        assert tf.paragraphs[1].runs[0].font.size == Pt(22)
        assert tf.paragraphs[2].runs[0].font.size == Pt(18)
        assert tf.paragraphs[3].runs[0].font.size == Pt(22)

    def test_bold_lead_in_with_colon(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Content hashing: SHA-256 for dedup")
        p = tf.paragraphs[0]
        assert len(p.runs) == 2
        assert p.runs[0].text == "Content hashing: "
        assert p.runs[0].font.bold is True
        assert p.runs[0].font.size == Pt(26)  # sparse (1 bullet)
        assert p.runs[1].text == "SHA-256 for dedup"
        assert p.runs[1].font.bold is not True

    def test_bold_lead_in_with_em_dash(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Graceful degradation \u2014 failures don't crash")
        p = tf.paragraphs[0]
        assert len(p.runs) == 2
        assert p.runs[0].text == "Graceful degradation \u2014 "
        assert p.runs[0].font.bold is True

    def test_no_bold_when_lead_in_too_long(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- This is a much longer phrase that has: a colon somewhere")
        p = tf.paragraphs[0]
        # More than 4 words before colon - no bold
        assert len(p.runs) == 1
        assert p.runs[0].font.bold is not True

    def test_no_bold_when_no_separator(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Just a plain bullet point")
        p = tf.paragraphs[0]
        assert len(p.runs) == 1
        assert p.runs[0].font.bold is not True


class TestAutoNumberContent:
    """Tests for _auto_number_content helper."""

    def test_numbers_bullet_items(self):
        content = "- First item\n- Second item\n- Third item"
        result = _auto_number_content(content)
        assert result == "1. First item\n2. Second item\n3. Third item"

    def test_does_not_number_sub_bullets(self):
        content = "- First item\n  - Sub-bullet\n- Second item"
        result = _auto_number_content(content)
        assert "1. First item" in result
        assert "  - Sub-bullet" in result
        assert "2. Second item" in result

    def test_does_not_double_number(self):
        content = "1. Already numbered\n2. Also numbered"
        result = _auto_number_content(content)
        assert result == "1. Already numbered\n2. Also numbered"

    def test_mixed_numbered_and_bullets(self):
        content = "1. Already numbered\n- Not numbered\n- Another"
        result = _auto_number_content(content)
        assert "1. Already numbered" in result
        assert "2. Not numbered" in result
        assert "3. Another" in result

    def test_plain_text_lines_get_numbered(self):
        content = "First item\nSecond item"
        result = _auto_number_content(content)
        assert result == "1. First item\n2. Second item"

    def test_preserves_empty_lines(self):
        content = "- First\n\n- Second"
        result = _auto_number_content(content)
        lines = result.split('\n')
        assert lines[0] == "1. First"
        assert lines[1] == ""
        assert lines[2] == "2. Second"


class TestNumberedLayout:
    """Tests for numbered list layout support."""

    def test_numbered_in_known_layout_types(self):
        assert 'numbered' in KNOWN_LAYOUT_TYPES

    def test_parse_numbered_layout(self):
        result = parse_layout_suggestion("numbered")
        assert result['type'] == 'numbered'

    def test_apply_bullets_numbered_items(self):
        """Numbered items (1., 2.) should render with level 0 and correct font size."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
        tf = tb.text_frame
        _apply_bullets_to_text_frame(tf, "1. First step\n2. Second step\n3. Third step")
        assert tf.paragraphs[0].runs[0].text == "1. First step"
        assert tf.paragraphs[0].level == 0
        assert tf.paragraphs[0].runs[0].font.size == Pt(22)
        assert tf.paragraphs[1].runs[0].text == "2. Second step"
        assert tf.paragraphs[2].runs[0].text == "3. Third step"

    def test_apply_layout_content_numbers_bullets_for_numbered_type(self):
        """apply_layout_content with layout_type='numbered' should auto-number bullets."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        content = "- First item\n- Second item\n  - Sub-bullet\n- Third item"
        apply_layout_content(slide, content, 'numbered')
        # Find body placeholder text
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                texts = [p.text for p in ph.text_frame.paragraphs if p.text]
                # Should have numbered top-level items, sub-bullet unchanged
                assert any("1. First item" in t for t in texts)
                assert any("2. Second item" in t for t in texts)
                assert any("Sub-bullet" in t for t in texts)
                assert any("3. Third item" in t for t in texts)


class TestTwoColumnWithHeaders:
    def test_column_header_applied_as_bold_first_para(self):
        from pptx.util import Inches, Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
        tf = tb.text_frame
        _apply_bullets_to_text_frame(tf, "- Item one\n- Item two", header="My Header")
        # First paragraph should be bold header
        assert tf.paragraphs[0].runs[0].text == "My Header"
        assert tf.paragraphs[0].runs[0].font.bold is True
        assert tf.paragraphs[0].runs[0].font.size == Pt(26)  # sparse (2 bullets)
        # Second paragraph should be actual content
        assert tf.paragraphs[1].runs[0].text == "Item one"

    def test_no_header_works_normally(self):
        from pptx.util import Inches, Pt
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
        tf = tb.text_frame
        _apply_bullets_to_text_frame(tf, "- Item one\n- Item two")
        assert tf.paragraphs[0].runs[0].text == "Item one"


class TestPlaceholderImage:
    """Tests for diagram placeholder shape."""

    def test_adds_placeholder_shape(self):
        from aippt.layouts import apply_placeholder_image
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        shape_count_before = len(slide.shapes)
        apply_placeholder_image(slide, "Architecture flow diagram")
        assert len(slide.shapes) == shape_count_before + 1

    def test_placeholder_contains_description(self):
        from aippt.layouts import apply_placeholder_image
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        apply_placeholder_image(slide, "Network topology diagram")
        # Find the last shape (newly added)
        new_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame') and 'Network topology diagram' in s.text_frame.text]
        assert len(new_shapes) == 1

    def test_placeholder_has_gray_fill(self):
        from aippt.layouts import apply_placeholder_image
        from pptx.dml.color import RGBColor
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        apply_placeholder_image(slide, "Test diagram")
        # Last non-placeholder shape should be our rectangle
        shape = [s for s in slide.shapes if s.shape_type == 1][0]  # AUTO_SHAPE = 1
        assert shape.fill.fore_color.rgb == RGBColor(0xE0, 0xE0, 0xE0)


class TestAutoFontScaling:
    """Tests for auto font scaling on sparse slides (< 3 top-level bullets)."""

    def _make_text_frame(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
        return tb.text_frame

    def test_one_bullet_uses_pt26(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Single bullet")
        assert tf.paragraphs[0].runs[0].font.size == Pt(26)

    def test_two_bullets_uses_pt26(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- First bullet\n- Second bullet")
        assert tf.paragraphs[0].runs[0].font.size == Pt(26)
        assert tf.paragraphs[1].runs[0].font.size == Pt(26)

    def test_three_bullets_uses_pt22(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- One\n- Two\n- Three")
        assert tf.paragraphs[0].runs[0].font.size == Pt(22)
        assert tf.paragraphs[1].runs[0].font.size == Pt(22)
        assert tf.paragraphs[2].runs[0].font.size == Pt(22)

    def test_four_bullets_uses_pt22(self):
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- A\n- B\n- C\n- D")
        for p in tf.paragraphs:
            assert p.runs[0].font.size == Pt(22)

    def test_sub_bullets_not_counted_as_top_level(self):
        """2 top-level + 2 sub-bullets = still sparse (< 3 top-level)."""
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- Top one\n  - Sub A\n  - Sub B\n- Top two")
        # Top-level should be Pt(26)
        assert tf.paragraphs[0].runs[0].font.size == Pt(26)
        # Sub-bullets should always be Pt(18)
        assert tf.paragraphs[1].runs[0].font.size == Pt(18)
        assert tf.paragraphs[2].runs[0].font.size == Pt(18)
        # Second top-level
        assert tf.paragraphs[3].runs[0].font.size == Pt(26)

    def test_sub_bullets_always_pt18(self):
        """Sub-bullets stay Pt(18) regardless of density."""
        tf = self._make_text_frame()
        _apply_bullets_to_text_frame(tf, "- A\n- B\n- C\n  - Sub")
        assert tf.paragraphs[3].runs[0].font.size == Pt(18)


class TestAuthorImageInsertion:
    """Tests for image_path parameter on apply_layout_content() and _apply_author_image()."""

    from aippt.layouts import _apply_author_image  # noqa: F401 (imported for use in tests)

    def test_apply_layout_content_image_diagram_full_image(self, tmp_path):
        """IMAGE: + LAYOUT: diagram uses full-image behavior (text to notes)."""
        from PIL import Image

        img_file = tmp_path / "test_image.png"
        Image.new('RGB', (100, 100), color='blue').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        content = "- Bullet one\n- Bullet two"

        apply_layout_content(slide, content, 'diagram', image_path=str(img_file))

        # Verify a picture shape was added
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) >= 1

        # Verify content text appears in speaker notes (full-image mode)
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Bullet one" in notes_text or "Bullet two" in notes_text

    def test_apply_layout_content_image_two_column_full_image(self, tmp_path):
        """IMAGE: + LAYOUT: two_column uses full-image behavior (text to notes)."""
        from PIL import Image

        img_file = tmp_path / "test_image.png"
        Image.new('RGB', (100, 100), color='blue').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        content = "- Left\n- Right"

        apply_layout_content(slide, content, 'two_column', image_path=str(img_file))

        # Verify content text appears in speaker notes (full-image mode)
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Left" in notes_text or "Right" in notes_text

    def test_apply_layout_content_image_bullet_co_display(self, tmp_path):
        """IMAGE: + LAYOUT: bullet uses co-display (text stays on slide)."""
        from PIL import Image

        img_file = tmp_path / "test_image.png"
        Image.new('RGB', (100, 100), color='blue').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        content = "- Bullet one\n- Bullet two"

        apply_layout_content(slide, content, 'bullet', image_path=str(img_file))

        # Verify a picture shape was added
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) >= 1

        # Content should NOT be in notes (co-display keeps text on slide)
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Bullet one" not in notes_text

    def test_apply_layout_content_image_no_layout_co_display(self, tmp_path):
        """IMAGE: without explicit LAYOUT uses co-display (text stays on slide)."""
        from PIL import Image

        img_file = tmp_path / "test_image.png"
        Image.new('RGB', (100, 100), color='blue').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        content = "- Bullet one\n- Bullet two"

        apply_layout_content(slide, content, 'basic', image_path=str(img_file))

        # Verify a picture shape was added
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) >= 1

        # Content should NOT be in notes (co-display keeps text on slide)
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Bullet one" not in notes_text

    def test_apply_layout_content_without_image_path(self):
        """apply_layout_content without image_path puts content in body placeholder, not notes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        content = "- Alpha\n- Beta\n- Gamma"

        apply_layout_content(slide, content, 'bullet')

        # No picture shapes should be present
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) == 0

        # Content should appear in the body placeholder, not notes
        body_texts = []
        for ph in slide.placeholders:
            if ph.placeholder_format.idx > 0:
                for p in ph.text_frame.paragraphs:
                    if p.text:
                        body_texts.append(p.text)

        assert any("Alpha" in t for t in body_texts)

        # Notes should remain empty (or not contain the slide content)
        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Alpha" not in notes_text

    def test_author_image_moves_content_to_notes(self, tmp_path):
        """_apply_author_image() places the content text in speaker notes."""
        from PIL import Image
        from aippt.layouts import _apply_author_image

        img_file = tmp_path / "author.png"
        Image.new('RGB', (100, 100), color='blue').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        _apply_author_image(slide, str(img_file), "Bullet content here")

        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Bullet content here" in notes_text

    def test_author_image_appends_to_existing_notes(self, tmp_path):
        """_apply_author_image() appends new content after any existing notes text."""
        from PIL import Image
        from aippt.layouts import _apply_author_image

        img_file = tmp_path / "author2.png"
        Image.new('RGB', (100, 100), color='blue').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Pre-populate notes
        slide.notes_slide.notes_text_frame.text = "Existing presenter note."

        _apply_author_image(slide, str(img_file), "New bullet content")

        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Existing presenter note." in notes_text
        assert "New bullet content" in notes_text


class TestImageWithText:
    """Tests for _apply_image_with_text() co-display function."""

    def test_apply_image_with_text_adds_picture(self, tmp_path):
        """_apply_image_with_text() inserts image as a shape when no PICTURE placeholder."""
        from PIL import Image
        from aippt.layouts import _apply_image_with_text

        img_file = tmp_path / "test.png"
        Image.new('RGB', (100, 100), color='red').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        _apply_image_with_text(slide, str(img_file), "- Bullet A\n- Bullet B")

        # Should have added a picture shape (fallback since default layout has no PICTURE ph)
        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) >= 1

    def test_apply_image_with_text_does_not_write_notes(self, tmp_path):
        """_apply_image_with_text() does NOT move content to speaker notes."""
        from PIL import Image
        from aippt.layouts import _apply_image_with_text

        img_file = tmp_path / "test.png"
        Image.new('RGB', (100, 100), color='red').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        _apply_image_with_text(slide, str(img_file), "- Content here")

        notes_text = slide.notes_slide.notes_text_frame.text
        assert "Content here" not in notes_text

    def test_apply_image_with_text_empty_content(self, tmp_path):
        """_apply_image_with_text() handles empty content gracefully."""
        from PIL import Image
        from aippt.layouts import _apply_image_with_text

        img_file = tmp_path / "test.png"
        Image.new('RGB', (100, 100), color='red').save(str(img_file))

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Should not raise
        _apply_image_with_text(slide, str(img_file), "")

        picture_shapes = [s for s in slide.shapes if s.shape_type == 13]
        assert len(picture_shapes) >= 1


class TestSelectSlideLayoutImageText:
    """Tests for image_text layout type in select_slide_layout()."""

    def test_selects_screenshot_and_caption_layout(self):
        mock_prs = MagicMock()
        mock_layout_screenshot = MagicMock()
        mock_layout_screenshot.name = "Screenshot and caption"
        mock_layout_default = MagicMock()
        mock_layout_default.name = "Title and Content"
        mock_prs.slide_layouts = [mock_layout_default, mock_layout_screenshot]

        result = select_slide_layout(mock_prs, "image_text")
        assert result.name == "Screenshot and caption"

    def test_falls_back_when_screenshot_layout_missing(self):
        mock_prs = MagicMock()
        mock_layout = MagicMock()
        mock_layout.name = "Title and Content"
        mock_prs.slide_layouts = [mock_layout]

        result = select_slide_layout(mock_prs, "image_text")
        # Falls back to first layout when target not found
        assert result == mock_layout
