"""
pptx_helpers — reusable helper library for python-pptx slide generation.

Usage from output/ scripts:
    import sys; sys.path.insert(0, 'lib')
    from pptx_helpers import load_template, save_deck, set_placeholder, add_bullets, ...
"""

import os
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Template management
# ---------------------------------------------------------------------------

def load_template(template_path):
    """Load a PPTX template and remove all sample slides.

    Returns a Presentation object ready for adding new slides.
    """
    prs = Presentation(template_path)
    # Remove sample slides (templates often ship with example slides)
    sample_count = len(prs.slides)
    for i in range(sample_count - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    return prs


def save_deck(prs, output_path):
    """Save the presentation and print file size."""
    prs.save(output_path)
    size_mb = os.path.getsize(output_path) / (1024 * 1024)
    print(f"Deck saved: {output_path} ({size_mb:.1f} MB, {len(prs.slides)} slides)")


# ---------------------------------------------------------------------------
# Placeholder helpers
# ---------------------------------------------------------------------------

def get_placeholder(slide, idx):
    """Get a placeholder by its format index."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def set_placeholder(slide, idx, text):
    """Find placeholder by idx and set its text."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text_frame.clear()
            ph.text_frame.text = text
            return ph
    return None


# ---------------------------------------------------------------------------
# Bullet formatting internals
# ---------------------------------------------------------------------------

def suppress_bullet(paragraph):
    """Remove bullet dot from a paragraph (for headers, numbered items)."""
    pPr = paragraph._p.get_or_add_pPr()
    for child in list(pPr):
        if etree.QName(child.tag).localname.startswith('bu'):
            pPr.remove(child)
    pPr.append(pPr.makeelement(qn('a:buNone'), {}))
    pPr.set('marL', '0')
    pPr.set('indent', '0')


def _add_bullet_with_lead_in(paragraph, lead_in, rest):
    """Add a bullet with bold lead-in text followed by normal text."""
    run1 = paragraph.add_run()
    run1.text = lead_in
    run1.font.bold = True
    run1.font.size = Pt(18)
    run2 = paragraph.add_run()
    run2.text = rest
    run2.font.size = Pt(18)


def _format_bullet(paragraph, text):
    """Format a single bullet paragraph with optional bold lead-in."""
    # Pass 1: explicit separator after bold
    m = re.match(r'^\*\*(.+?)\*\*\s*[—\-:]\s*(.*)', text)
    if m:
        _add_bullet_with_lead_in(paragraph, m.group(1) + " \u2014 ", m.group(2))
        return
    # Pass 2: bold prefix no separator
    m = re.match(r'^\*\*(.+?)\*\*(.*)', text)
    if m:
        _add_bullet_with_lead_in(paragraph, m.group(1), m.group(2))
        return
    paragraph.text = text
    paragraph.font.size = Pt(18)


# ---------------------------------------------------------------------------
# Bullet formatting — public API
# ---------------------------------------------------------------------------

def add_bullets(slide, idx, items):
    """Add bulleted items to a placeholder with bold lead-in detection.

    items: list of str or (text, level) tuples.
    Strings get level 0. Tuples specify (text, indent_level).
    Bold lead-in patterns (**Bold** — rest) are auto-detected.
    """
    ph = get_placeholder(slide, idx)
    if not ph:
        return
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        # Check for bold lead-in with explicit separator
        match_sep = re.match(r'^\*\*(.+?)\*\*\s*[—\-:]\s*(.*)', text)
        match_bold = re.match(r'^\*\*(.+?)\*\*(.*)', text)
        if match_sep:
            bold_text = match_sep.group(1)
            rest_text = match_sep.group(2).strip()
            if rest_text:
                run1 = p.add_run()
                run1.text = bold_text + " \u2014 "
                run1.font.bold = True
                run1.font.size = Pt(18 if level == 0 else 16)
                run2 = p.add_run()
                run2.text = rest_text
                run2.font.size = Pt(18 if level == 0 else 16)
            else:
                run1 = p.add_run()
                run1.text = bold_text
                run1.font.bold = True
                run1.font.size = Pt(18 if level == 0 else 16)
        elif match_bold:
            bold_text = match_bold.group(1)
            rest_text = match_bold.group(2)
            run1 = p.add_run()
            run1.text = bold_text
            run1.font.bold = True
            run1.font.size = Pt(18 if level == 0 else 16)
            if rest_text:
                run2 = p.add_run()
                run2.text = rest_text
                run2.font.size = Pt(18 if level == 0 else 16)
        else:
            p.text = text
            p.font.size = Pt(18 if level == 0 else 16)

        p.level = level


def add_bullets_with_sub(slide, idx, items):
    """Add bulleted items with sub-bullet support and bold lead-in detection.

    items: list of str or dict. Dicts have 'text' and optional 'subs' list.
    """
    ph = get_placeholder(slide, idx)
    if not ph:
        return
    tf = ph.text_frame
    tf.clear()
    first = True
    for item in items:
        if isinstance(item, dict):
            text = item['text']
            subs = item.get('subs', [])
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            _format_bullet(p, text)
            for sub in subs:
                p = tf.add_paragraph()
                p.level = 1
                p.text = sub
                p.font.size = Pt(16)
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = 0
            _format_bullet(p, item)


def add_numbered_bullets(slide, idx, items):
    """Add numbered items with bold lead-in handling, suppressing template bullets.

    items: list of str or dict. Dicts have 'text' and optional 'subs' list.
    Sub-items render with monospace font at indent level 1.
    """
    ph = get_placeholder(slide, idx)
    if not ph:
        return
    tf = ph.text_frame
    tf.clear()
    num = 0
    first_para = True
    for item in items:
        if isinstance(item, dict):
            text = item['text']
            subs = item.get('subs', [])
            num += 1
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.level = 0
            suppress_bullet(p)
            m = re.match(r'^\*\*(.+?)\*\*\s*[:\-\u2014]?\s*(.*)', text)
            if m:
                bold_text = m.group(1).rstrip(':')
                run1 = p.add_run()
                run1.text = f"{num}. {bold_text}: "
                run1.font.bold = True
                run1.font.size = Pt(18)
                run2 = p.add_run()
                run2.text = m.group(2)
                run2.font.size = Pt(18)
            else:
                p.text = f"{num}. {text}"
                p.font.size = Pt(18)
            for sub in subs:
                sp = tf.add_paragraph()
                sp.level = 1
                sp.font.size = Pt(14)
                run = sp.add_run()
                run.text = sub
                run.font.name = "Consolas"
                run.font.size = Pt(14)
        else:
            num += 1
            p = tf.paragraphs[0] if first_para else tf.add_paragraph()
            first_para = False
            p.level = 0
            suppress_bullet(p)
            m = re.match(r'^\*\*(.+?)\*\*\s*[:\-\u2014]?\s*(.*)', item)
            if m:
                bold_text = m.group(1).rstrip(':')
                run1 = p.add_run()
                run1.text = f"{num}. {bold_text}: "
                run1.font.bold = True
                run1.font.size = Pt(18)
                run2 = p.add_run()
                run2.text = m.group(2)
                run2.font.size = Pt(18)
            else:
                p.text = f"{num}. {item}"
                p.font.size = Pt(18)


# ---------------------------------------------------------------------------
# Two-column helpers
# ---------------------------------------------------------------------------

def add_two_column_with_header(slide, idx, header, items):
    """Add a column header + bullets to a two-column placeholder.

    header: teal accent-colored header text (centered, bold, 20pt).
    items: list of str or dict (same format as add_bullets_with_sub).
    """
    ph = get_placeholder(slide, idx)
    if not ph:
        return
    tf = ph.text_frame
    tf.clear()
    # Header — teal accent color
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = header
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0xC2, 0xDE)  # teal accent
    p.alignment = PP_ALIGN.CENTER
    suppress_bullet(p)
    # Bullets
    for item in items:
        if isinstance(item, dict):
            text = item['text']
            subs = item.get('subs', [])
            p = tf.add_paragraph()
            p.level = 0
            _format_bullet(p, text)
            for sub in subs:
                sp = tf.add_paragraph()
                sp.level = 1
                sp.text = sub
                sp.font.size = Pt(16)
        else:
            p = tf.add_paragraph()
            p.level = 0
            _format_bullet(p, item)


def add_column_divider(slide, prs):
    """Add a vertical divider line between two-column placeholders."""
    slide_width = prs.slide_width
    center_x = slide_width // 2
    top = Inches(1.8)
    bottom = Inches(6.8)
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        center_x, top, center_x, bottom
    )
    line.line.color.rgb = RGBColor(0x63, 0x64, 0x66)  # surface gray
    line.line.width = Pt(1)


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------

def set_notes(slide, text):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text
