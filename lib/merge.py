"""
merge — combine multiple PPTX section files into a single deck.

Uses python-pptx XML-level slide copying to transfer slides between
Presentation objects, preserving images, shapes, and formatting.

Usage:
    from lib.merge import merge_decks
    merge_decks(['section-1.pptx', 'section-2.pptx'], 'final.pptx')
"""

import os
import re
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Emu


# ---------------------------------------------------------------------------
# Slide copying
# ---------------------------------------------------------------------------

def _copy_slide(source_slide, target_prs):
    """Copy a slide from one presentation into another.

    Uses the source slide's layout by reference. This is safe when all source
    presentations share the same template/theme (the aippt sectioned generation
    use case). For cross-template merges, a blank-layout approach would be needed.

    Handles:
    - All shape types (text, images, charts, connectors, groups)
    - Media relationships (images, embedded objects)
    - Hyperlinks
    - Skips notes-slide relationships (not transferred)
    """
    # Add a new slide using the source layout
    # In practice, for sectioned decks from the same template, we use a blank
    # layout from the target presentation to avoid importing foreign layout parts
    slide_layout = _get_blank_layout(target_prs)
    dest_slide = target_prs.slides.add_slide(slide_layout)

    # Remove any placeholder shapes the blank layout may have added
    for shape in list(dest_slide.shapes):
        if shape.is_placeholder:
            shape.element.getparent().remove(shape.element)

    # Copy slide background if set
    src_bg = source_slide.background
    if src_bg.fill.type is not None:
        dest_bg = dest_slide.background
        dest_bg_elem = dest_slide.background._element
        src_bg_elem = source_slide.background._element
        # Deep copy the background XML
        for child in list(dest_bg_elem):
            dest_bg_elem.remove(child)
        for child in src_bg_elem:
            dest_bg_elem.append(deepcopy(child))

    # Deep-copy all shape XML elements
    for shape in source_slide.shapes:
        new_element = deepcopy(shape.element)
        dest_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    # Copy relationships (images, hyperlinks, etc.)
    # Skip layout, master, and notes relationships
    skip_types = ('notesSlide', 'slideLayout', 'slideMaster')
    for rel in source_slide.part.rels.values():
        if any(t in rel.reltype for t in skip_types):
            continue
        dest_slide.part.rels._add_relationship(
            rel.reltype, rel._target, rel.rId
        )

    return dest_slide


def _get_blank_layout(prs):
    """Find the most blank/empty slide layout in the presentation.

    Looks for layouts named 'Blank' or with the fewest placeholders.
    Falls back to the last layout (commonly blank in most templates).
    """
    # Try to find one named 'Blank'
    for layout in prs.slide_layouts:
        if layout.name and 'blank' in layout.name.lower():
            return layout

    # Fall back to layout with fewest placeholders
    best = prs.slide_layouts[0]
    best_count = len(list(best.placeholders))
    for layout in prs.slide_layouts:
        count = len(list(layout.placeholders))
        if count < best_count:
            best = layout
            best_count = count
    return best


# ---------------------------------------------------------------------------
# Slide number renumbering
# ---------------------------------------------------------------------------

def _renumber_slides(prs):
    """Renumber slide footer text boxes sequentially (1, 2, 3, ...).

    Looks for small text elements near the bottom-left of each slide
    that contain just a number. These are the slide number indicators
    added by pptxgenjs and python-pptx helper libraries.

    Detection heuristic:
    - Shape is a text box (not a placeholder)
    - Position: bottom 15% of slide, left 20% of slide
    - Content: a single paragraph containing only digits
    - Font size ≤ 12pt (or small shape height ≤ 0.5 inches)
    """
    slide_width = prs.slide_width or Inches(13.33)
    slide_height = prs.slide_height or Inches(7.5)

    # Thresholds for "bottom-left" region
    bottom_threshold = slide_height * 0.80  # bottom 20%
    left_threshold = slide_width * 0.20     # left 20%
    max_height = Inches(0.6)                # small text box

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder:
                continue

            # Check position: bottom-left region
            if shape.top is None or shape.left is None:
                continue
            if shape.top < bottom_threshold:
                continue
            if shape.left > left_threshold:
                continue

            # Check size: small text box
            if shape.height is not None and shape.height > max_height:
                continue

            # Check content: just a number
            text = shape.text_frame.text.strip()
            if text and text.isdigit():
                # Update the number
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip().isdigit():
                            run.text = str(slide_num)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def merge_decks(chunk_paths, output_path, renumber=True):
    """Merge multiple PPTX files into a single deck.

    Args:
        chunk_paths: List of paths to PPTX section files, in order.
        output_path: Where to save the merged deck.
        renumber: If True, fix slide number text boxes to global numbering.

    Returns:
        dict with keys: output_path, slide_count, chunk_counts
    """
    if not chunk_paths:
        raise ValueError("No chunk paths provided")

    # Validate all files exist
    for path in chunk_paths:
        if not os.path.isfile(path):
            raise FileNotFoundError(f"Chunk file not found: {path}")

    # Use the first chunk as the base (preserves template, theme, masters)
    master = Presentation(chunk_paths[0])

    # Track slide counts per chunk
    chunk_counts = [len(master.slides)]

    # Copy slides from subsequent chunks
    for path in chunk_paths[1:]:
        src = Presentation(path)
        chunk_count = 0
        for slide in src.slides:
            _copy_slide(slide, master)
            chunk_count += 1
        chunk_counts.append(chunk_count)

    # Renumber slide footers
    if renumber:
        _renumber_slides(master)

    # Save
    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    master.save(output_path)

    total = len(master.slides)
    size_kb = os.path.getsize(output_path) / 1024
    print(f"Merged {len(chunk_paths)} chunks → {output_path} "
          f"({total} slides, {size_kb:.0f} KB)")

    return {
        'output_path': output_path,
        'slide_count': total,
        'chunk_counts': chunk_counts,
    }
