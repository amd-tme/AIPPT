"""Slide layout selection and content application.

This module handles the mapping of content types to PowerPoint layouts
and the application of content to slides using various layout strategies.
"""

import logging
import re
from typing import Dict, Optional, Tuple, Union

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Layout definitions
# ---------------------------------------------------------------------------

SLIDE_LAYOUTS = {
    'title_and_content': {
        'name': 'Title and Content',
        'use_case': 'default',
        'placeholders': ['title', 'content']
    },
    'two_content': {
        'name': 'Two Content',
        'use_case': 'comparison',
        'placeholders': ['title', 'left_content', 'right_content']
    },
    'title_only': {
        'name': 'Title Only',
        'use_case': 'custom',
        'placeholders': ['title']
    },
    'code': {
        'name': 'Developer Code Layout',
        'use_case': 'code',
        'placeholders': ['title', 'code']
    }
}


# ---------------------------------------------------------------------------
# Layout selection functions
# ---------------------------------------------------------------------------

def get_layout_index(prs, layout_name: str) -> int:
    """Get the index of a layout by its name.

    Args:
        prs: A python-pptx Presentation object
        layout_name: The name of the layout to find

    Returns:
        The index of the layout, or 3 (Title and Content) if not found
    """
    for idx, layout in enumerate(prs.slide_layouts):
        if layout.name == layout_name:
            return idx
    return 3


def select_slide_layout(prs, layout_type: str):
    """Select appropriate slide layout based on content type.

    Args:
        prs: A python-pptx Presentation object
        layout_type: One of 'diagram', 'two_column', 'bullet', 'basic'

    Returns:
        A SlideLayout object from the presentation
    """
    layout_map = {
        'diagram': 'Title Only',
        'two_column': 'Two Content',
        'bullet': 'Title and Content',
        'numbered': 'Title and Content',
        'basic': 'Title and Content',
        'image_text': 'Screenshot and caption',
    }

    layout_name = layout_map.get(layout_type, 'Title and Content')

    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout

    return prs.slide_layouts[0]


KNOWN_LAYOUT_TYPES = {'bullet', 'two_column', 'diagram', 'basic', 'numbered'}


def parse_layout_suggestion(layout_text: str) -> Dict[str, any]:
    """Parse LLM's layout suggestion into actionable structure.

    Extracts the first token from the layout text and matches it against
    known layout types. This prevents false matches from descriptive text
    like "bullet — a visual list" triggering the diagram matcher.

    Args:
        layout_text: Raw text from LLM containing layout suggestions

    Returns:
        Dictionary with 'type', 'structure', and 'elements' keys
    """
    layout_info = {
        'type': 'basic',
        'structure': [],
        'elements': []
    }

    # Extract the first word/token and match against known types
    first_token = layout_text.strip().split()[0].lower().rstrip(':,') if layout_text.strip() else ''

    if first_token in KNOWN_LAYOUT_TYPES:
        layout_info['type'] = first_token
    else:
        logger.debug(f"Unrecognized layout token '{first_token}', defaulting to basic")

    return layout_info


# ---------------------------------------------------------------------------
# Content splitting utilities
# ---------------------------------------------------------------------------

def split_content_for_columns(content: str) -> Tuple[str, str]:
    """Split content into two columns at a natural break point.

    Args:
        content: Multi-line content string

    Returns:
        Tuple of (left_content, right_content)
    """
    # Strip leading/trailing blank lines so they don't skew the split
    lines = [l for l in content.split('\n') if l.strip()]
    if not lines:
        return '', ''

    mid_point = len(lines) // 2

    left_content = '\n'.join(lines[:mid_point]).strip()
    right_content = '\n'.join(lines[mid_point:]).strip()

    return left_content, right_content


# ---------------------------------------------------------------------------
# Layout application functions
# ---------------------------------------------------------------------------

def _auto_number_content(content: str) -> str:
    """Prepend sequential numbers to top-level content lines.

    For use with the 'numbered' layout type. Lines that already start with
    a number prefix (e.g., '1. ') are left unchanged. Sub-bullets (indented
    lines) are not numbered.

    Args:
        content: Multi-line content string

    Returns:
        Content with top-level lines numbered sequentially
    """
    lines = content.split('\n')
    result = []
    counter = 1
    for line in lines:
        stripped = line.strip()
        if not stripped:
            result.append(line)
            continue
        # Sub-bullets: indented lines starting with bullet markers
        if line.startswith('  ') and stripped.startswith(('-', '•', '*')):
            result.append(line)
            continue
        # Already numbered
        if re.match(r'^\d+\.\s', stripped):
            result.append(line)
            counter += 1
            continue
        # Top-level line: strip bullet marker and add number
        text = stripped.lstrip('-•* ') if stripped.startswith(('-', '•', '*')) else stripped
        result.append(f"{counter}. {text}")
        counter += 1
    return '\n'.join(result)


def _apply_author_image(slide, image_path: str, content: str):
    """Insert an author-provided image and move text content to speaker notes.

    The image occupies the full content area. Any bullet text is preserved
    in the slide's speaker notes so the presenter still has it.

    Args:
        slide: A python-pptx Slide object
        image_path: Absolute path to the image file
        content: Text content to move to speaker notes
    """
    slide.shapes.add_picture(
        image_path,
        Inches(1), Inches(1.5),
        Inches(8), Inches(4),
    )

    # Preserve content in speaker notes
    if content and content.strip():
        notes_slide = slide.notes_slide
        existing = notes_slide.notes_text_frame.text
        if existing and existing.strip():
            notes_slide.notes_text_frame.text = existing + "\n\n" + content.strip()
        else:
            notes_slide.notes_text_frame.text = content.strip()


def _apply_image_with_text(slide, image_path: str, content: str):
    """Insert image into picture placeholder and text into body placeholder.

    Used with template layouts that have both a PICTURE and BODY placeholder
    (e.g., 'Screenshot and caption'). Falls back to full-image behavior if
    the expected placeholders are not found.

    Args:
        slide: A python-pptx Slide object
        image_path: Absolute path to the image file
        content: Text content to display alongside the image
    """
    pic_ph = None
    text_ph = None
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx == 0:  # skip title
            continue
        if shape.placeholder_format.type == 18:  # PP_PLACEHOLDER.PICTURE
            pic_ph = shape
        elif text_ph is None:
            text_ph = shape

    if pic_ph:
        pic_ph.insert_picture(open(image_path, 'rb'))
    else:
        # Fallback: add picture as a free shape
        slide.shapes.add_picture(
            image_path,
            Inches(0.6), Inches(1.6),
            Inches(7.9), Inches(5.3),
        )

    if text_ph and content and content.strip():
        _apply_bullets_to_text_frame(text_ph.text_frame, content, narrow=True)


def apply_layout_content(
    slide,
    content: str,
    layout_type: str,
    suggestions: Optional[Dict] = None,
    image_dir: Optional[str] = None,
    slide_num: Optional[int] = None,
    client=None,
    image_gen: str = 'none',
    image_path: Optional[str] = None,
):
    """Apply content to slide based on layout type.

    Args:
        slide: A python-pptx Slide object
        content: The text content to apply
        layout_type: One of 'bullet', 'two_column', 'diagram', 'basic'
        suggestions: Dict with 'VISUALS' and other LLM suggestions
        image_dir: Directory for storing generated images
        slide_num: Current slide number (for image naming)
        client: LLMClient instance for image generation
        image_gen: Image generation mode ('none', 'claude', 'dalle')
        image_path: Optional path to an author-provided image (from IMAGE: directive)
    """
    try:
        # Author-provided image takes priority over normal layout rendering
        if image_path:
            if layout_type in ('diagram', 'two_column'):
                # Full-image: text moves to speaker notes
                _apply_author_image(slide, image_path, content)
            else:
                # Co-display: image and text side-by-side
                _apply_image_with_text(slide, image_path, content)
            return

        if layout_type == 'numbered':
            content = _auto_number_content(content)
            apply_bullet_layout(slide, content, suggestions)
        elif layout_type == 'bullet':
            apply_bullet_layout(slide, content, suggestions)
        elif layout_type == 'two_column':
            apply_two_column_layout(slide, content, suggestions)
        elif layout_type == 'diagram':
            if not all([suggestions, image_dir, slide_num]):
                logger.warning(
                    "Missing required parameters for diagram layout, "
                    "falling back to basic layout"
                )
                apply_basic_layout(slide, content)
            else:
                # Import here to avoid circular dependency
                from aippt.images import apply_diagram_layout
                apply_diagram_layout(
                    slide=slide,
                    content=content,
                    suggestions=suggestions,
                    image_dir=image_dir,
                    slide_num=slide_num,
                    client=client,
                    image_gen=image_gen
                )
        else:
            apply_basic_layout(slide, content)
    except Exception as e:
        logger.error(f"Error applying layout: {e}")
        apply_basic_layout(slide, content)


def apply_basic_layout(slide, content: str):
    """Apply basic bullet point layout to slide.

    Args:
        slide: A python-pptx Slide object
        content: Multi-line text content with bullet points
    """
    try:
        # Try to use the first non-title content placeholder by index.
        # Template placeholders may be typed as BODY (2) or OBJECT (7);
        # matching on idx > 0 handles both.
        content_placeholder = None
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx > 0:
                content_placeholder = shape
                break

        if content_placeholder:
            tf = content_placeholder.text_frame
        else:
            content_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.5),
                Inches(8), Inches(5)
            )
            tf = content_box.text_frame

        _apply_bullets_to_text_frame(tf, content)

    except Exception as e:
        logger.error(f"Error in apply_basic_layout: {e}")
        raise


def apply_bullet_layout(slide, content: str, suggestions: Optional[Dict] = None):
    """Apply bullet point layout to slide.

    This is similar to basic layout but may incorporate visual suggestions.

    Args:
        slide: A python-pptx Slide object
        content: Multi-line text content with bullet points
        suggestions: Optional dict with layout suggestions (unused currently)
    """
    apply_basic_layout(slide, content)


def _detect_lead_in(text: str):
    """Detect a bold lead-in pattern (1-4 words followed by : or \u2014).

    Returns:
        Tuple of (lead_in, rest) if pattern found, or (None, text) if not.
    """
    match = re.match(r'^(\S+(?:\s+\S+){0,3}(?::\s| \u2014 ))(.*)', text)
    if match:
        return match.group(1), match.group(2)
    return None, text


def _apply_bullets_to_text_frame(tf, content: str, header: str = None,
                                 narrow: bool = False):
    """Apply bullet-processed content lines to a text frame.

    Processes markdown-style bullet markers and indentation into
    proper paragraph levels, matching the behavior of apply_basic_layout.
    Sets font size: Pt(26) for level 0 and level 1 when fewer than 3 top-level
    bullets are present (sparse slide), otherwise Pt(22). Sub-bullets stay Pt(18).

    When ``narrow`` is True (e.g. the text sits beside an image in a
    picture+text layout), font sizes are reduced so content fits the
    smaller area: Pt(18) base / Pt(14) sub-bullets.

    Args:
        tf: A python-pptx TextFrame object
        content: Multi-line text content with bullet points
        header: Optional column header rendered as a bold first paragraph
        narrow: If True, use smaller fonts suitable for narrow text areas
    """
    # Count top-level lines to determine font size
    top_level_count = 0
    for line in content.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue
        # Sub-bullets start with 2+ spaces then a bullet marker
        if line.startswith('  ') and stripped.startswith(('-', '•', '*')):
            continue
        top_level_count += 1

    if narrow:
        base_font_size = Pt(18)
        sub_font_size = Pt(14)
    else:
        base_font_size = Pt(26) if top_level_count < 3 else Pt(22)
        sub_font_size = Pt(18)

    if header:
        p = tf.paragraphs[0]
        p.level = 0
        run = p.add_run()
        run.text = header
        run.font.bold = True
        run.font.size = base_font_size
        first_para = False
    else:
        first_para = True
    for line in content.split('\n'):
        if line.strip():
            if first_para:
                # Reuse the existing first paragraph to avoid a leading blank line
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            stripped = line.strip()
            # Numbered items (1., 2., ...) are level 0, keep the number prefix.
            num_match = re.match(r'^(\d+\.\s)', stripped)
            if num_match:
                p.level = 0
                run = p.add_run()
                run.text = stripped
                run.font.size = base_font_size
            # Check indentation on the original line to distinguish level 2
            # sub-bullets (leading spaces before the dash) from level 1 bullets.
            elif line.startswith('  ') and stripped.startswith(('-', '•', '*')):
                p.level = 2
                run = p.add_run()
                run.text = stripped.lstrip('-•* ')
                run.font.size = sub_font_size
            elif stripped.startswith(('-', '•', '*')):
                p.level = 1
                text = stripped.lstrip('-•* ')
                lead_in, rest = _detect_lead_in(text)
                if lead_in:
                    run_bold = p.add_run()
                    run_bold.text = lead_in
                    run_bold.font.size = base_font_size
                    run_bold.font.bold = True
                    run_rest = p.add_run()
                    run_rest.text = rest
                    run_rest.font.size = base_font_size
                else:
                    run = p.add_run()
                    run.text = text
                    run.font.size = base_font_size
            else:
                lead_in, rest = _detect_lead_in(stripped)
                if lead_in:
                    run_bold = p.add_run()
                    run_bold.text = lead_in
                    run_bold.font.size = base_font_size
                    run_bold.font.bold = True
                    run_rest = p.add_run()
                    run_rest.text = rest
                    run_rest.font.size = base_font_size
                else:
                    run = p.add_run()
                    run.text = stripped
                    run.font.size = base_font_size


def apply_two_column_layout(slide, content: str, suggestions: Optional[Dict] = None):
    """Apply two-column layout to slide.

    Uses template placeholders when the "Two Content" layout provides them,
    falling back to manual textboxes. Processes bullet markers into proper
    paragraph levels.

    Args:
        slide: A python-pptx Slide object
        content: Content to split between columns
        suggestions: Optional dict with layout suggestions (unused currently)
    """
    left_content, right_content = split_content_for_columns(content)

    # Try to use template content placeholders (skip title at idx 0)
    content_phs = []
    for shape in slide.placeholders:
        idx = shape.placeholder_format.idx
        if idx > 0:
            content_phs.append((idx, shape))
    content_phs.sort(key=lambda x: x[0])

    if len(content_phs) >= 2:
        left_tf = content_phs[0][1].text_frame
        right_tf = content_phs[1][1].text_frame
    else:
        # Fallback to manual textboxes
        logger.debug("Two Content placeholders not found, using textboxes")
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(4.25), Inches(5)
        )
        left_tf = left_box.text_frame
        right_box = slide.shapes.add_textbox(
            Inches(5), Inches(1.5),
            Inches(4.25), Inches(5)
        )
        right_tf = right_box.text_frame

    from aippt.parser import parse_column_headers
    left_header, right_header = None, None
    if suggestions:
        left_header, right_header = parse_column_headers(suggestions.get('LAYOUT', ''))

    _apply_bullets_to_text_frame(left_tf, left_content, header=left_header)
    _apply_bullets_to_text_frame(right_tf, right_content, header=right_header)


def apply_comparison_layout(slide, content: str, suggestions: Optional[Dict] = None):
    """Apply comparison layout with divider to slide.

    Args:
        slide: A python-pptx Slide object
        content: Content with ' vs. ' or ' versus ' separator
        suggestions: Optional dict with layout suggestions (unused currently)
    """
    try:
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(4.25), Inches(5)
        )
        right_box = slide.shapes.add_textbox(
            Inches(5), Inches(1.5),
            Inches(4.25), Inches(5)
        )

        # Add dividing line
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.9), Inches(1.5),
            Inches(0.1), Inches(5)
        )

        # Split content
        if ' vs. ' in content:
            parts = content.split(' vs. ')
        elif ' versus ' in content:
            parts = content.split(' versus ')
        else:
            parts = split_content_for_columns(content)
            parts = [parts[0], parts[1]]

        left_content = parts[0] if len(parts) > 0 else ''
        right_content = parts[1] if len(parts) > 1 else ''

        # Add content to each side
        for text_frame, col_content in [
            (left_box.text_frame, left_content),
            (right_box.text_frame, right_content)
        ]:
            for line in col_content.split('\n'):
                if line.strip():
                    p = text_frame.add_paragraph()
                    stripped = line.strip()
                    if stripped.startswith(('-', '•', '*')):
                        p.level = 1
                        p.text = stripped.lstrip('-•* ')
                    else:
                        p.text = stripped

    except Exception as e:
        logger.error(f"Error in apply_comparison_layout: {e}")
        raise


# ---------------------------------------------------------------------------
# Placeholder image (diagram fallback)
# ---------------------------------------------------------------------------

def apply_placeholder_image(slide, description: str):
    """Insert a placeholder rectangle with description text for future image generation.

    Creates a centered light gray rectangle in the content area with
    descriptive text indicating what image should be generated.

    Args:
        slide: A python-pptx Slide object
        description: Description of the intended diagram/image
    """
    left = Inches(1.5)
    top = Inches(2.0)
    width = Inches(7.0)
    height = Inches(4.0)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )

    # Light gray fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Thin border
    shape.line.color.rgb = RGBColor(0xB0, 0xB0, 0xB0)
    shape.line.width = Pt(1)

    # Add description text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"[Image: {description}]"
    run.font.size = Pt(16)
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    run.font.italic = True


# ---------------------------------------------------------------------------
# Disclaimer textbox helper
# ---------------------------------------------------------------------------

def add_disclaimer_textbox(
    slide,
    text: str = "AI-Generated Image -- Not Approved for External Use",
) -> None:
    """Add a disclaimer text box at the bottom of the slide.

    Used when a slide contains an AI-generated image to indicate
    the content is not approved for external distribution.

    Args:
        slide: A python-pptx Slide object
        text: Disclaimer text to display
    """
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8),
        Inches(9), Inches(0.4),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    run.font.italic = True


# ---------------------------------------------------------------------------
# Template inspection
# ---------------------------------------------------------------------------

def inspect_template(template_path: str) -> Dict[str, any]:
    """Inspect a PowerPoint template to determine available layouts.

    Args:
        template_path: Path to the .pptx template file

    Returns:
        Dictionary with 'layouts' list and 'slide_size' info
    """
    from pptx import Presentation

    prs = Presentation(template_path)
    template_info = {
        'layouts': [],
        'slide_size': {
            'width': prs.slide_width,
            'height': prs.slide_height
        }
    }

    for layout in prs.slide_layouts:
        layout_info = {
            'name': layout.name,
            'placeholders': []
        }

        for placeholder in layout.placeholders:
            ph_info = {
                'type': (
                    placeholder.placeholder_format.type
                    if hasattr(placeholder, 'placeholder_format')
                    else 'unknown'
                ),
                'name': (
                    placeholder.name
                    if hasattr(placeholder, 'name')
                    else 'unnamed'
                ),
                'idx': (
                    placeholder.placeholder_format.idx
                    if hasattr(placeholder, 'placeholder_format')
                    else -1
                )
            }
            layout_info['placeholders'].append(ph_info)

        template_info['layouts'].append(layout_info)

    return template_info


# ---------------------------------------------------------------------------
# Layout creation helpers (for custom slides)
# ---------------------------------------------------------------------------

def remove_all_slides(prs):
    """Remove all slides from a presentation, keeping layouts and masters.

    Args:
        prs: A python-pptx Presentation object
    """
    while len(prs.slides) > 0:
        slide = prs.slides[0]
        rId = None
        # Find the relationship ID for this slide
        for rel in prs.part.rels.values():
            if rel.target_part is slide.part:
                rId = rel.rId
                break
        if rId:
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        else:
            # Fallback: just remove from the slide list XML
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def create_basic_layout(slide, visuals: str = ""):
    """Create a basic slide layout with title and content boxes."""
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    return title, content


def create_two_column_layout(slide, visuals: str = ""):
    """Create a two-column layout."""
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5))
    right = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.25), Inches(5))
    return title, left, right


def create_bullet_list_layout(slide, visuals: str = ""):
    """Create a layout optimized for bullet points."""
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    return title, content
