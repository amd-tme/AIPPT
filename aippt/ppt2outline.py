"""PowerPoint to markdown outline conversion."""
import logging
import os

from pptx import Presentation

logger = logging.getLogger(__name__)

REVERSE_SYSTEM_PROMPT = """You are an expert at converting presentation slides into structured markdown outlines.
Given a slide image and its extracted text content, produce a clean markdown outline that:
- Uses the slide title as an H2 heading (## Title)
- Converts bullet points into a hierarchical list with proper indentation
- Describes diagrams, charts, and visual elements as concise bullet points
- Omits decorative text, watermarks, and slide furniture (page numbers, dates, footers)
- Preserves technical accuracy — do not invent content not present on the slide
- Keep descriptions concise: 1-2 sentences per visual element

Return ONLY the markdown outline. Do not include commentary or explanation."""

REVERSE_TEXT_ONLY_SYSTEM_PROMPT = """You are an expert at converting presentation slide content into structured markdown outlines.
Given a slide's extracted text content, produce a clean markdown outline that:
- Uses the slide title as an H2 heading (## Title)
- Converts bullet points into a hierarchical list with proper indentation
- Cleans up duplicated or garbled text from shape extraction
- Omits decorative text, watermarks, and slide furniture (page numbers, dates, footers)
- Preserves technical accuracy — do not invent content not present on the slide

Return ONLY the markdown outline. Do not include commentary or explanation."""


def _should_skip_shape(shape, title_shape=None) -> bool:
    """Return True if shape is decorative and should be excluded from extraction.

    Filters: title shape (already handled), slide number/date/footer placeholders,
    connector/freeform shapes, and shapes with short numeric-only text (callout labels).
    """
    # Skip title shape (handled separately)
    if title_shape is not None and shape == title_shape:
        return True

    # Skip placeholder types: slide number, date, footer
    if shape.is_placeholder:
        from pptx.enum.shapes import PP_PLACEHOLDER
        ph_type = shape.placeholder_format.type
        if ph_type in (
            PP_PLACEHOLDER.SLIDE_NUMBER,
            PP_PLACEHOLDER.DATE,
            PP_PLACEHOLDER.FOOTER,
        ):
            return True

    # Skip connectors and freeforms
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
        return True

    # Skip shapes with only short numeric text (callout numbers like "1", "2")
    if hasattr(shape, "text") and shape.text.strip():
        text = shape.text.strip()
        if len(text) <= 3 and text.isdigit():
            return True

    return False


def _extract_slide_title(slide) -> str:
    """Extract slide title with fallback chain.

    Resolution order:
    1. Standard title placeholder (paragraph-aware for multi-line titles)
    2. Subtitle placeholder
    3. First short text shape (<=80 chars)
    4. "Untitled Slide" (last resort)
    """
    # 1. Standard title placeholder
    if slide.shapes.title and slide.shapes.title.text.strip():
        tf = slide.shapes.title.text_frame
        parts = [p.text.strip() for p in tf.paragraphs if p.text.strip()]
        if len(parts) > 1:
            return " — ".join(parts)
        return parts[0] if parts else "Untitled Slide"

    # 2. Subtitle placeholder
    from pptx.enum.shapes import PP_PLACEHOLDER
    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
            if shape.text.strip():
                return shape.text.strip()

    # 3. First short text shape
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            first_line = shape.text_frame.paragraphs[0].text.strip()
            if first_line and len(first_line) <= 80:
                return first_line

    return "Untitled Slide"


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape, preserving bullet hierarchy and table structure.

    Returns pre-formatted lines:
    - Text frames: each paragraph as indented bullet using paragraph.level
    - Groups: recursive extraction of sub-shapes
    - Tables: markdown table with header separator
    """
    lines = []

    # Text frame: extract with paragraph-level indentation
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                level = paragraph.level or 0
                indent = "  " * level
                lines.append(f"{indent}- {text}")

    # Group shapes: recursively extract
    elif hasattr(shape, "shapes"):
        for subshape in shape.shapes:
            subtext = extract_text_from_shape(subshape)
            if subtext:
                lines.append(subtext)

    # Table: render as markdown table
    elif hasattr(shape, "has_table") and shape.has_table:
        table = shape.table
        rows_text = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_text.append("| " + " | ".join(cells) + " |")
        if rows_text:
            col_count = len(table.columns)
            separator = "| " + " | ".join(["---"] * col_count) + " |"
            lines = [rows_text[0], separator] + rows_text[1:]

    return "\n".join(lines)


def _resolve_slide_image(slide_number: int, images_dir):
    """Find a slide image file in images_dir. Returns path or None."""
    if not images_dir or not os.path.isdir(images_dir):
        return None
    for ext in (".PNG", ".png", ".jpg", ".jpeg"):
        candidate = os.path.join(images_dir, f"Slide{slide_number}{ext}")
        if os.path.exists(candidate):
            return candidate
    return None


def _enhance_slide_with_llm(llm_client, slide_number, title, mechanical_text, image_path):
    """Send slide to LLM and return markdown string or None on failure.

    Tries image+text first (if image available), falls back to text-only.
    If both fail, returns None.
    """
    if image_path and os.path.exists(image_path):
        prompt = (
            f"Convert this slide to a markdown outline.\n\n"
            f"Slide title: {title}\n\n"
            f"Extracted text (may be noisy):\n{mechanical_text}"
        )
        try:
            return llm_client.generate_text_with_image(
                prompt=prompt,
                image_path=image_path,
                system_prompt=REVERSE_SYSTEM_PROMPT,
                max_tokens=1000,
                temperature=0.2,
            )
        except Exception as exc:
            logger.warning(
                "LLM image call failed for slide %d: %s — falling back to text-only",
                slide_number,
                exc,
            )

    # Text-only LLM path
    prompt = (
        f"Convert this slide content to a clean markdown outline.\n\n"
        f"Slide title: {title}\n\n"
        f"Extracted text:\n{mechanical_text}"
    )
    try:
        return llm_client.generate_text(
            prompt=prompt,
            system_prompt=REVERSE_TEXT_ONLY_SYSTEM_PROMPT,
            max_tokens=1000,
            temperature=0.2,
        )
    except Exception as exc:
        logger.warning(
            "LLM text call failed for slide %d: %s — falling back to mechanical extraction",
            slide_number,
            exc,
        )
        return None


def _format_notes_as_comment(notes_text: str) -> str:
    """Format speaker notes as an HTML comment block, stripping analysis artifacts."""
    lines = []
    for line in notes_text.split('\n'):
        if line.strip().startswith('[Note: analysis based on slide text only'):
            continue
        if line.strip():
            lines.append(line.replace('-->', '-- >'))
    if not lines:
        return ''
    return '<!-- notes\n' + '\n'.join(lines) + '\n-->\n'


def convert_pptx_to_outline(
    pptx_file: str,
    output_file: str,
    include_notes: bool = True,
    enhance: bool = False,
    llm_client=None,
    images_dir=None,
) -> bool:
    """Convert PowerPoint file to markdown outline.

    Args:
        pptx_file: Path to the input PowerPoint file
        output_file: Path for the output markdown file
        include_notes: Whether to include slide notes in output
        enhance: Whether to use LLM enhancement for better output
        llm_client: LLMClient instance (required when enhance=True)
        images_dir: Directory containing slide images for visual analysis

    Returns:
        True if conversion succeeded, False otherwise
    """
    try:
        presentation = Presentation(pptx_file)

        # Read sections from PowerPoint
        from aippt.sections import read_sections
        ppt_sections = read_sections(presentation)

        # Build slide_id → section_name map
        slide_to_section = {}
        for section in ppt_sections:
            for slide_id in section.slide_ids:
                slide_to_section[slide_id] = section.name

        current_section = None

        with open(output_file, 'w', encoding='utf-8') as f:
            for i, slide in enumerate(presentation.slides, 1):
                # Check if we're entering a new section
                section_name = slide_to_section.get(slide.slide_id)
                if section_name and section_name != current_section:
                    current_section = section_name
                    if section_name != "Default Section":
                        f.write(f"# {section_name}\n\n")

                # Extract slide title (with fallback chain)
                title = _extract_slide_title(slide)

                # Attempt LLM enhancement if requested and client is available
                if enhance and llm_client is not None:
                    logger.info("Enhancing slide %d: %s", i, title)

                    # Build mechanical text for context
                    mechanical_parts = []
                    for shape in slide.shapes:
                        if _should_skip_shape(shape, title_shape=slide.shapes.title):
                            continue
                        text = extract_text_from_shape(shape)
                        if text:
                            mechanical_parts.append(text)
                    mechanical_text = "\n".join(mechanical_parts)

                    image_path = _resolve_slide_image(i, images_dir)
                    llm_output = _enhance_slide_with_llm(
                        llm_client, i, title, mechanical_text, image_path
                    )

                    if llm_output:
                        f.write(llm_output.rstrip() + "\n")
                        # Include slide notes if requested
                        if include_notes and slide.has_notes_slide:
                            notes_text = slide.notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                comment = _format_notes_as_comment(notes_text)
                                if comment:
                                    f.write("\n" + comment)
                        f.write("\n\n")
                        continue

                # Mechanical extraction (default path, or LLM fallback)
                # Write slide title as H2 if we have sections, H1 otherwise
                header_level = "##" if ppt_sections else "#"
                f.write(f"{header_level} {title}\n\n")

                # Extract and write slide content
                content = []
                for shape in slide.shapes:
                    if _should_skip_shape(shape, title_shape=slide.shapes.title):
                        continue
                    text = extract_text_from_shape(shape)
                    if text:
                        content.append(text)

                # Write content (lines are pre-formatted with bullets and indentation)
                for text in content:
                    for line in text.split('\n'):
                        if line.strip():
                            f.write(f"{line}\n")

                # Include slide notes if requested
                if include_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        comment = _format_notes_as_comment(notes_text)
                        if comment:
                            f.write("\n" + comment)

                # Add space between slides
                f.write("\n\n")

        logger.info(f"Successfully converted {pptx_file} to {output_file}")
        return True

    except Exception as e:
        logger.error(f"Error converting {pptx_file}: {str(e)}")
        return False
