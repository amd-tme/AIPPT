"""Slide remix -- search, manifest, and deck assembly."""
import copy
import logging
import os
from typing import Dict, List, Optional

import yaml
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from aippt.catalog import search_slides, check_newer_versions, get_slide_section

logger = logging.getLogger(__name__)


def generate_manifest(
    slides: List[Dict],
    title: str = "Remixed Presentation",
    template: str = "template.pptx",
    db_path: str = "slides.db",
) -> str:
    """Generate a YAML manifest from search results.

    Args:
        slides: List of slide dictionaries from search_slides()
        title: Title for the remixed presentation
        template: Path to template PPTX
        db_path: Path to database for section lookups

    Returns:
        YAML string
    """
    manifest = {
        "title": title,
        "template": template,
        "slides": [],
    }
    for s in slides:
        section = get_slide_section(s["id"], db_path) if "id" in s else None
        slide_entry = {
            "deck": os.path.basename(s["deck_path"]),
            "deck_path": s["deck_path"],
            "position": s["position"],
            "title": s["title"],
        }
        if section:
            slide_entry["section"] = section
        manifest["slides"].append(slide_entry)
    return yaml.dump(manifest, default_flow_style=False, sort_keys=False)


def load_manifest(manifest_path: str) -> Dict:
    """Load and validate a remix manifest.

    Args:
        manifest_path: Path to YAML manifest file

    Returns:
        Parsed manifest dictionary

    Raises:
        ValueError: If required keys are missing
    """
    with open(manifest_path, encoding="utf-8") as f:
        manifest = yaml.safe_load(f)

    required = ["title", "slides"]
    for key in required:
        if key not in manifest:
            raise ValueError(f"Manifest missing required key: {key}")

    return manifest


def copy_slide(source_prs, source_index, target_prs):
    """Copy a slide from source presentation to target.

    Uses python-pptx XML manipulation to copy slide content.
    This is the standard approach since python-pptx doesn't have a
    built-in slide copy method.

    Args:
        source_prs: Source Presentation object
        source_index: 0-based index of the slide to copy
        target_prs: Target Presentation object

    Returns:
        The newly created slide
    """
    source_slide = source_prs.slides[source_index]

    # Find or use the first layout in target as base
    slide_layout = target_prs.slide_layouts[0]

    # Try to find matching layout by name
    source_layout_name = source_slide.slide_layout.name
    for layout in target_prs.slide_layouts:
        if layout.name == source_layout_name:
            slide_layout = layout
            break

    new_slide = target_prs.slides.add_slide(slide_layout)

    # Remove default placeholder shapes that came with the layout
    for shape in list(new_slide.placeholders):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copy shapes from source to target
    for shape in source_slide.shapes:
        sp = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(sp)

    # Copy notes if they exist
    if source_slide.has_notes_slide:
        notes_slide = new_slide.notes_slide
        notes_slide.notes_text_frame.text = source_slide.notes_slide.notes_text_frame.text

    return new_slide


from aippt.layouts import remove_all_slides as _remove_all_slides


def assemble_deck(
    manifest_path: str,
    output_path: str,
    db_path: str = "slides.db",
):
    """Assemble a new deck from a manifest file.

    Args:
        manifest_path: Path to YAML manifest
        output_path: Output .pptx path
        db_path: Catalog database path for version checking
    """
    manifest = load_manifest(manifest_path)

    # Use template if specified, otherwise use first source deck as base
    template = manifest.get("template")
    if template and os.path.exists(template):
        target_prs = Presentation(template)
    else:
        # Use the first source deck's template
        first_deck = manifest["slides"][0].get("deck_path", manifest["slides"][0]["deck"])
        target_prs = Presentation(first_deck)

    # Remove existing slides from the presentation
    _remove_all_slides(target_prs)

    # Check for newer versions
    slide_dicts = []
    for entry in manifest["slides"]:
        slide_dicts.append({
            "title": entry["title"],
            "content_hash": entry.get("content_hash", ""),
            "deck_name": entry.get("deck", ""),
            "updated_at": entry.get("updated_at", ""),
        })

    warnings = check_newer_versions(slide_dicts, db_path)
    for w in warnings:
        logger.warning(
            f"Newer version of '{w['slide_title']}' found in {w['newer_deck']} "
            f"(updated {w['newer_updated']})"
        )

    # Cache opened presentations
    prs_cache = {}
    copied_count = 0

    for entry in manifest["slides"]:
        deck_path = entry.get("deck_path", entry["deck"])
        position = entry["position"]

        if deck_path not in prs_cache:
            if not os.path.exists(deck_path):
                logger.error(f"Deck not found: {deck_path}")
                continue
            prs_cache[deck_path] = Presentation(deck_path)

        source_prs = prs_cache[deck_path]
        slide_index = position - 1  # 0-based

        if slide_index >= len(source_prs.slides):
            logger.error(f"Slide {position} not found in {deck_path} (only {len(source_prs.slides)} slides)")
            continue

        try:
            copy_slide(source_prs, slide_index, target_prs)
            copied_count += 1
            logger.info(f"Copied slide {position} from {os.path.basename(deck_path)}: {entry['title']}")
        except Exception as e:
            logger.error(f"Error copying slide {position} from {deck_path}: {e}")
            continue

    # Reconstruct sections from manifest
    sections_to_write = []
    current_section = None
    section_slide_ids = []

    for idx, entry in enumerate(manifest["slides"]):
        section_name = entry.get("section")

        # Check if we're starting a new section
        if section_name and section_name != current_section:
            # Save previous section if exists
            if current_section and section_slide_ids:
                from aippt.sections import Section
                sections_to_write.append(Section(name=current_section, slide_ids=section_slide_ids))
                section_slide_ids = []

            current_section = section_name

        # Add slide ID if in a section
        if current_section and idx < len(target_prs.slides):
            section_slide_ids.append(target_prs.slides[idx].slide_id)

    # Save final section
    if current_section and section_slide_ids:
        from aippt.sections import Section
        sections_to_write.append(Section(name=current_section, slide_ids=section_slide_ids))

    # Write sections to presentation
    if sections_to_write:
        try:
            from aippt.sections import write_sections
            write_sections(target_prs, sections_to_write)
            logger.info(f"Applied {len(sections_to_write)} sections to remixed presentation")
        except Exception as e:
            logger.error(f"Error applying sections to remixed presentation: {e}")

    target_prs.save(output_path)
    logger.info(f"Assembled {copied_count} slides into {output_path}")
    return copied_count
