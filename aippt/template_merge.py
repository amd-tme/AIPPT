"""Corporate template merge — inject corporate master/layouts into generated decks.

Opens a corporate template (e.g. templates/corp.pptx) as the base presentation,
copies each slide from the generated deck into it, and re-assigns slides to
matching corporate layouts based on [AIPPT-META] layout_selected metadata.
"""

import json
import logging
import os
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches

from aippt.metadata import extract_metadata, append_history_entry

logger = logging.getLogger(__name__)

CORP_LAYOUT_MAP = {
    "title":           "Title Slide - No Image",
    "bullet":          "Title and Content",
    "two_column":      "Two Content",
    "code":            "Developer Code Layout",
    "section_divider": "Divider slide",
    "closing":         "Closing logo slide",
}
FALLBACK_LAYOUT = "Blank"


def _find_layout_by_name(prs, name: str):
    """Find a slide layout by exact name. Returns None if not found."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return None


def _get_layout_for_slide(slide, layout_map: dict, fallback_layout_name: str, target_prs) -> tuple:
    """Determine the target layout for a source slide.

    Reads [AIPPT-META] layout_selected, maps through layout_map,
    falls back to fallback_layout_name.

    Returns: (layout_object, source_layout_type, target_layout_name)
    """
    entries = extract_metadata(slide)
    layout_type = None
    for entry in entries:
        if "layout_selected" in entry:
            layout_type = entry["layout_selected"]
            break

    if layout_type and layout_type in layout_map:
        target_name = layout_map[layout_type]
    else:
        if layout_type:
            logger.warning(f"Unmapped layout type '{layout_type}', using fallback '{fallback_layout_name}'")
        target_name = fallback_layout_name

    layout = _find_layout_by_name(target_prs, target_name)
    if layout is None:
        raise ValueError(
            f"Layout '{target_name}' not found in template. "
            f"Available layouts: {[l.name for l in target_prs.slide_layouts]}"
        )

    return layout, layout_type, target_name


def _copy_slide_to_template(source_slide, target_prs, target_layout):
    """Copy a slide into the target presentation using the given layout.

    Standalone copy function (not a wrapper around lib/merge._copy_slide)
    because the two serve different purposes: lib/merge copies between
    same-template decks, this copies cross-template using named layouts.

    Only non-placeholder shapes are copied from the source slide.
    Placeholders from the target layout are removed before copying.
    """
    dest_slide = target_prs.slides.add_slide(target_layout)

    for shape in list(dest_slide.shapes):
        if shape.is_placeholder:
            shape.element.getparent().remove(shape.element)

    src_bg = source_slide.background
    if src_bg.fill.type is not None:
        dest_bg_elem = dest_slide.background._element
        src_bg_elem = source_slide.background._element
        for child in list(dest_bg_elem):
            dest_bg_elem.remove(child)
        for child in src_bg_elem:
            dest_bg_elem.append(deepcopy(child))

    for shape in source_slide.shapes:
        if shape.is_placeholder:
            continue
        new_element = deepcopy(shape.element)
        dest_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

    skip_types = ('notesSlide', 'slideLayout', 'slideMaster')
    for rel in source_slide.part.rels.values():
        if any(t in rel.reltype for t in skip_types):
            continue
        dest_slide.part.rels._add_relationship(
            rel.reltype, rel._target, rel.rId
        )

    return dest_slide


def _transfer_notes(source_slide, dest_slide, target_layout_name: str) -> None:
    """Copy speaker notes from source to destination slide.

    Preserves [AIPPT-META] blocks. Appends a history entry recording
    the template merge operation.
    """
    try:
        source_notes = source_slide.notes_slide.notes_text_frame.text
    except Exception:
        source_notes = ""

    if source_notes:
        dest_slide.notes_slide.notes_text_frame.text = source_notes
        append_history_entry(
            dest_slide,
            f"Merged into corporate template ({target_layout_name})",
            source_tag="/template-merge",
        )


def _remove_template_slides(prs) -> None:
    """Remove all pre-built slides from the template presentation.

    Drops both the sldIdLst entries and the corresponding part relationships
    so that saved output has no duplicate entry warnings.
    """
    sldIdLst = prs.slides._sldIdLst
    ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    slide_rids = [
        sldId.get(f'{{{ns}}}id')
        for sldId in list(sldIdLst)
    ]
    for sldId in list(sldIdLst):
        sldIdLst.remove(sldId)
    for rId in slide_rids:
        if rId and rId in prs.part.rels:
            del prs.part.rels._rels[rId]


def merge_with_template(
    generated_pptx: str,
    template_pptx: str,
    output_pptx: str,
    layout_map: dict[str, str] | None = None,
    fallback_layout: str = "Blank",
) -> dict:
    """Merge a generated deck into a corporate template.

    Args:
        generated_pptx: Path to the AI-generated PPTX.
        template_pptx: Path to the corporate template PPTX.
        output_pptx: Where to save the merged result.
        layout_map: Override the default CORP_LAYOUT_MAP.
        fallback_layout: Layout name for unmapped slides (default: "Blank").

    Returns:
        dict with keys: output_path, slide_count, layout_assignments
    """
    if not os.path.isfile(generated_pptx):
        raise FileNotFoundError(f"Generated deck not found: {generated_pptx}")
    if not os.path.isfile(template_pptx):
        raise FileNotFoundError(f"Template not found: {template_pptx}")

    effective_map = layout_map if layout_map is not None else CORP_LAYOUT_MAP
    source_prs = Presentation(generated_pptx)

    if len(source_prs.slides) == 0:
        target_prs = Presentation(template_pptx)
        # Remove pre-built template slides so the output is truly empty
        _remove_template_slides(target_prs)
        os.makedirs(os.path.dirname(output_pptx) or '.', exist_ok=True)
        target_prs.save(output_pptx)
        return {"output_path": output_pptx, "slide_count": 0, "layout_assignments": []}

    target_prs = Presentation(template_pptx)

    # Remove any slides that come pre-built in the corporate template.
    # The template is used only for its master/layouts — we replace all slides
    # with the generated content.
    _remove_template_slides(target_prs)

    layout_assignments = []

    for i, slide in enumerate(source_prs.slides, 1):
        layout, source_type, target_name = _get_layout_for_slide(
            slide, effective_map, fallback_layout, target_prs
        )

        dest_slide = _copy_slide_to_template(slide, target_prs, layout)
        _transfer_notes(slide, dest_slide, target_name)

        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip()
                break

        layout_assignments.append({
            "slide_num": i,
            "title": title,
            "source_layout": source_type,
            "target_layout": target_name,
        })

        logger.info(f"Slide {i}: '{title}' -> {target_name} (from {source_type or 'no metadata'})")

    os.makedirs(os.path.dirname(output_pptx) or '.', exist_ok=True)
    target_prs.save(output_pptx)

    size_kb = os.path.getsize(output_pptx) / 1024
    logger.info(f"Template merge complete: {output_pptx} ({len(layout_assignments)} slides, {size_kb:.0f} KB)")

    return {
        "output_path": output_pptx,
        "slide_count": len(layout_assignments),
        "layout_assignments": layout_assignments,
    }
