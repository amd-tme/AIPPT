"""API routes for the web UI."""
import asyncio
import json
import logging
import os
import re
import sys
import tempfile
import uuid
from pathlib import Path
from typing import List

logger = logging.getLogger(__name__)

from fastapi import APIRouter, Request, UploadFile, File, Form, WebSocket, WebSocketDisconnect
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse, StreamingResponse

from aippt.catalog import (
    get_db,
    display_name,
    get_deck_by_file_hash,
    get_deck_origin,
    search_slides,
    add_tags,
    get_slide_tags,
    get_all_tags,
    remove_slide_tag,
    list_taxonomy,
    add_taxonomy_tags,
    remove_taxonomy_tag,
    import_taxonomy_csv,
    export_taxonomy_csv,
    rename_tag,
    catalog_deck,
    get_deck_by_id,
    delete_deck,
    record_edit,
)
from aippt.export import export_csv
from aippt import graph
from aippt.ingest import ingest_deck
from aippt.thumbnails import store_client_images
from aippt.web.asset_sync import persist_file, persist_tree, materialize_file
import aippt.pipeline as _pipeline_module
from aippt.pipeline import PipelineConfig
from aippt.config import get_template_default, TemplateConfigError, base_path_prefix

router = APIRouter()
STATIC_DIR = Path(__file__).parent / "static"


def _base_prefix() -> str:
    """Mount prefix (no trailing slash) for origin-resolved absolute URLs.

    Thin wrapper over :func:`aippt.config.base_path_prefix` — kept as a
    module-local name for the preview session/WS URL builders below.
    """
    return base_path_prefix()


@router.get("/healthz")
async def healthz():
    """Health check for Kubernetes liveness/readiness probes."""
    return {"status": "ok"}


@router.get("/", response_class=HTMLResponse)
async def dashboard(request: Request):
    """Main dashboard -- serves the single-page app.

    Injects ``<base href="${BASE_PATH}">`` so the SPA works whether it is
    served at the apex (``/``) or under a path prefix (``/aippt/`` when the
    slai-app-platform ingress routes by path under the apex). All asset and
    fetch references inside index.html use *relative* URLs so the document
    base resolves them correctly under either mount point.
    """
    html = (STATIC_DIR / "index.html").read_text(encoding="utf-8")
    base_path = os.environ.get("BASE_PATH", "/")
    if not base_path.endswith("/"):
        base_path += "/"
    base_tag = f'<base href="{base_path}">'
    return html.replace("<head>", f"<head>\n    {base_tag}", 1)


@router.get("/api/config")
async def get_config(request: Request):
    """API: Return frontend configuration flags.

    Public, unauthenticated. The SPA reads this once at boot to learn the
    upload size cap so it can pre-check ``file.size`` before the multipart
    POST, and to pick up the ``view_only`` toggle. Keep the payload to
    deployment-public values only (no secrets, no per-user state).
    """
    return {
        "view_only": getattr(request.app.state, "view_only", False),
        "max_upload_bytes": getattr(request.app.state, "max_upload_bytes", 0),
    }


_FILE_HASH_RE = re.compile(r"^[a-f0-9]{64}$")


@router.get("/api/decks/by-hash/{sha256}")
async def get_deck_by_hash_route(request: Request, sha256: str):
    """API: Return existing deck metadata for the given file SHA-256, or 404.

    Lets the SPA detect duplicate uploads *before* any bytes leave the
    browser. The hash matches ``catalog.file_hash`` (SHA-256 of the raw PPTX
    bytes). Returns 400 for malformed input so spurious DB lookups don't
    accumulate.
    """
    if not _FILE_HASH_RE.match(sha256):
        return JSONResponse(
            {"error": "Hash must be 64 lowercase hex chars (SHA-256)."},
            status_code=400,
        )
    db_path = request.app.state.db_path
    deck = get_deck_by_file_hash(sha256, db_path=db_path)
    if not deck:
        return JSONResponse(
            {"error": "No deck with that file hash."}, status_code=404,
        )
    return {
        "id": deck["id"],
        "name": deck["name"],
        "display_name": display_name(deck["name"]),
        "slide_count": deck["slide_count"],
        "file_hash": deck["file_hash"],
        "cataloged_at": deck["cataloged_at"],
        "updated_at": deck["updated_at"],
    }


@router.get("/api/decks")
async def list_decks(request: Request):
    """API: List all cataloged decks including derived origin block."""
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    decks = conn.execute(
        "SELECT id, name, slide_count, cataloged_at, updated_at, author, "
        "created_date, modified_date, subject, description, "
        "outline_path, source_script_path, source_engine, source_theme, "
        "source_generated_at FROM decks ORDER BY updated_at DESC"
    ).fetchall()
    # First slide id per deck for thumbnail
    first_slides = {
        row["deck_id"]: row["slide_id"]
        for row in conn.execute(
            "SELECT deck_id, MIN(id) AS slide_id FROM slides GROUP BY deck_id"
        ).fetchall()
    }
    first_slide_updated = {}
    if first_slides:
        placeholders = ",".join("?" * len(first_slides))
        rows = conn.execute(
            f"SELECT id, image_path, updated_at FROM slides WHERE id IN ({placeholders})",
            list(first_slides.values()),
        ).fetchall()
        for row in rows:
            first_slide_updated[row["id"]] = (row["image_path"], row["updated_at"])
    conn.close()
    result = []
    for d in decks:
        deck = dict(d)
        deck["display_name"] = display_name(deck["name"])
        sid = first_slides.get(deck["id"])
        if sid:
            img_path, img_updated = first_slide_updated.get(sid, (None, None))
            if img_path:
                v = img_updated or ""
                deck["thumbnail_url"] = f"slide-image/{sid}?v={v}"
            else:
                deck["thumbnail_url"] = None
        else:
            deck["thumbnail_url"] = None
        # Derive origin block inline (avoid per-row DB query)
        outline_path = deck.pop("outline_path", None)
        script_path = deck.pop("source_script_path", None)
        engine = deck.pop("source_engine", None)
        theme = deck.pop("source_theme", None)
        generated_at = deck.pop("source_generated_at", None)
        if script_path:
            kind = "script"
        elif outline_path:
            kind = "outline"
        else:
            kind = "upload"
        deck["origin"] = {
            "kind": kind,
            "outline_path": outline_path,
            "source_script_path": script_path,
            "engine": engine,
            "theme": theme,
            "generated_at": generated_at,
        }
        result.append(deck)
    return result


@router.get("/api/decks/{deck_id}")
async def get_deck_metadata(deck_id: int, request: Request):
    """API: Return metadata for a single deck including the ``origin`` block.

    The ``origin`` block includes a derived ``kind`` field
    (``"outline"`` | ``"script"`` | ``"upload"``) so the SPA can decide
    whether to show the Regenerate button.
    """
    db_path = request.app.state.db_path
    deck = get_deck_by_id(deck_id, db_path)
    if deck is None:
        return JSONResponse({"error": "Deck not found"}, status_code=404)
    origin = get_deck_origin(deck_id, db_path)
    return {
        "id": deck["id"],
        "name": deck["name"],
        "display_name": display_name(deck["name"]),
        "slide_count": deck["slide_count"],
        "author": deck.get("author", ""),
        "cataloged_at": deck.get("cataloged_at"),
        "updated_at": deck.get("updated_at"),
        "subject": deck.get("subject", ""),
        "description": deck.get("description", ""),
        "origin": origin,
    }


@router.post("/api/decks/{deck_id}/regenerate")
async def regenerate_deck(deck_id: int, request: Request):
    """API: Rerun the pipeline against the deck's recorded source and replace in place.

    Uses SSE progress streaming (same pattern as ``/api/decks/create``).

    * 403 — view-only mode or missing Bearer token
    * 404 — deck not found
    * 409 — deck has no recorded source (upload-only deck)
    * 410 — source file missing on disk
    * 200 — streaming SSE response
    """
    import asyncio
    import json
    import queue as _queue
    import shutil as _shutil

    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": VIEW_ONLY_MSG}, status_code=403)

    ms_token = _extract_bearer_token(request)
    if not ms_token:
        return JSONResponse(
            {"error": "Microsoft sign-in required for regeneration. "
                      "Sign in with the Microsoft button and retry."},
            status_code=403,
        )
    ntid, _ntid_err = _ntid_or_400(request)
    if _ntid_err is not None:
        return _ntid_err

    db_path = request.app.state.db_path
    uploads_dir = request.app.state.uploads_dir
    gateway_config = request.app.state.gateway_config
    project_root = getattr(request.app.state, "project_root", os.getcwd())

    deck = get_deck_by_id(deck_id, db_path)
    if deck is None:
        return JSONResponse({"error": "Deck not found"}, status_code=404)

    origin = get_deck_origin(deck_id, db_path)
    if origin["kind"] == "upload":
        return JSONResponse(
            {"error": "deck has no recorded source; regeneration requires an outline or script"},
            status_code=409,
        )

    # Resolve the source path: stable per-deck location first, then fallback.
    # materialize_file fetches from object storage on a cold pod (no-op in fs).
    source_path = None
    stable_outline = os.path.join(_sources_dir(uploads_dir, deck_id), "outline.md")
    if materialize_file(request.app.state, stable_outline):
        source_path = stable_outline
    elif origin.get("outline_path") and materialize_file(request.app.state, origin["outline_path"]):
        source_path = origin["outline_path"]

    if source_path is None:
        return JSONResponse(
            {"error": "Source file is missing on disk; regeneration not possible. "
                      "Re-upload the outline to regenerate."},
            status_code=410,
        )

    try:
        template_path = get_template_default()
    except TemplateConfigError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)

    if not os.path.exists(template_path):
        return JSONResponse(
            {"error": f"Template not found: {template_path}"},
            status_code=404,
        )

    # Read the source outline
    with open(source_path, encoding="utf-8") as f:
        md_text = f.read()

    # Determine output path (reuse existing deck file_path if present)
    existing_file_path = deck.get("file_path", "")
    if existing_file_path:
        existing_file_path = _resolve_db_path(existing_file_path, project_root)
    if existing_file_path and os.path.exists(os.path.dirname(existing_file_path)):
        output_path = existing_file_path
    else:
        # Reconstruct output path under uploads
        _short_id = uuid.uuid4().hex[:8]
        _base_name = re.sub(r'[^\w\s-]', '', display_name(deck["name"])).strip()[:80]
        output_path = os.path.join(uploads_dir, f"{_short_id}_{_base_name}.pptx")

    engine = origin.get("engine") or "python-pptx"
    theme = origin.get("theme")

    # Log the regeneration action
    logger.info(
        "regenerate_deck: deck_id=%s ntid=%s source=%s engine=%s",
        deck_id, ntid, source_path, engine,
    )

    event_q: _queue.Queue = _queue.Queue()

    def create_progress(step, detail=""):
        status = "running"
        if any(detail.startswith(w) for w in ("Parsed", "All", "Built")):
            status = "done"
        event_q.put(("progress", {"step": step, "status": status, "detail": detail}))

    def ingest_progress(step, detail=""):
        ingest_map = {
            "export_images": "running", "export_images_done": "running",
            "export_images_skipped": "running", "catalog": "running",
            "catalog_done": "running", "complete": "done",
        }
        if step in ingest_map:
            event_q.put(("progress", {"step": "ingest", "status": ingest_map[step], "detail": detail}))

    async def _event_generator():
        loop = asyncio.get_running_loop()

        def _worker():
            pipeline_config = PipelineConfig(
                outline_text=md_text,
                template_path=template_path,
                output_path=output_path,
                enhance=False,  # regenerate uses stored outline as-is
                gateway_config=gateway_config,
                progress_callback=create_progress,
                outline_path=source_path,
                source_engine=engine,
                source_kind="outline",
                source_theme=theme,
            )
            pipeline_result = _pipeline_module.run_pipeline(pipeline_config)
            # The regenerated deck overwrites the existing file in place.
            persist_file(request.app.state, output_path)

            event_q.put(("progress", {"step": "ingest", "status": "running", "detail": "Re-cataloging generated deck..."}))

            # Re-catalog in place: delete old slide rows, re-insert from new PPTX.
            # Preserve the existing deck_id by updating the existing row.
            import datetime as _dt
            now_iso = _dt.datetime.utcnow().isoformat()

            from aippt.catalog import file_hash as _file_hash
            new_hash = _file_hash(output_path)

            _conn = get_db(db_path)
            try:
                # Preserve each old slide's thumbnail state keyed by its content
                # hash. A regenerated slide may reuse a prior canvas-captured
                # image only if its content is unchanged (hash matches); slides
                # whose text changed drop to a placeholder rather than showing a
                # stale thumbnail. Keyed on content_hash, not position, so a
                # reordered/inserted slide doesn't inherit a neighbor's image.
                _old_thumbs = {
                    r["image_content_hash"]: r["image_path"]
                    for r in _conn.execute(
                        "SELECT image_path, image_content_hash FROM slides "
                        "WHERE deck_id = ? AND image_path IS NOT NULL "
                        "AND image_content_hash IS NOT NULL",
                        (deck_id,),
                    ).fetchall()
                }
                # Remove old slides (CASCADE would also work but be explicit)
                _conn.execute("DELETE FROM slides WHERE deck_id = ?", (deck_id,))
                # Re-catalog the new PPTX into the existing deck row
                from pptx import Presentation as _Prs
                prs_new = _Prs(output_path)
                from aippt.catalog import _resolve_slide_title, content_hash as _chash
                from aippt.reverse import extract_text_from_shape
                for i, slide in enumerate(prs_new.slides, 1):
                    title, title_fallback = _resolve_slide_title(slide)
                    texts = []
                    for shape in slide.shapes:
                        if shape == slide.shapes.title:
                            continue
                        if title_fallback and shape == title_fallback:
                            continue
                        txt = extract_text_from_shape(shape)
                        if txt:
                            texts.append(txt)
                    content_text = "\n".join(texts)
                    chash = _chash(title, content_text)
                    notes = ""
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                    # Canvas-only: reuse a prior thumbnail only when the slide's
                    # content is byte-for-byte the same (content_hash match).
                    # Changed slides get NULL image_path → placeholder card,
                    # never a stale image. A fresh capture repopulates on the
                    # next Live Preview → Save.
                    image_path = _old_thumbs.get(chash)
                    image_content_hash = chash if image_path else None
                    _conn.execute(
                        """INSERT INTO slides (deck_id, position, title, content_text,
                           content_hash, notes, image_path, image_content_hash)
                           VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
                        (deck_id, i, title, content_text, chash, notes,
                         image_path, image_content_hash),
                    )
                # Update the deck row (keep same id)
                _conn.execute(
                    """UPDATE decks SET file_hash = ?, slide_count = ?,
                       outline_path = ?, source_engine = ?, source_theme = ?,
                       source_generated_at = ?, updated_at = datetime('now')
                       WHERE id = ?""",
                    (new_hash, len(prs_new.slides),
                     source_path, engine, theme, now_iso, deck_id),
                )
                _conn.commit()
            finally:
                _conn.close()

            # Push any re-rendered slide images to durable storage.
            persist_tree(request.app.state, _per_deck_images_dir(request, output_path))

            return {
                "deck_id": deck_id,
                "deck_name": deck["name"],
                "slide_count": pipeline_result.slide_count,
                "output_path": output_path,
            }

        future = loop.run_in_executor(None, _worker)

        def _format_sse(event_name, payload):
            return f"event: {event_name}\ndata: {json.dumps(payload)}\n\n"

        while not future.done():
            await asyncio.sleep(0)
            while True:
                try:
                    event_name, payload = event_q.get_nowait()
                    yield _format_sse(event_name, payload)
                except _queue.Empty:
                    break

        while True:
            try:
                event_name, payload = event_q.get_nowait()
                yield _format_sse(event_name, payload)
            except _queue.Empty:
                break

        try:
            result = await future
        except Exception as exc:
            yield _format_sse("error", {"detail": str(exc)})
            return

        yield _format_sse("complete", {
            "deck_id": result["deck_id"],
            "deck_name": result["deck_name"],
            "display_name": display_name(result["deck_name"]),
            "slide_count": result["slide_count"],
            "output_path": result["output_path"],
        })

    return StreamingResponse(_event_generator(), media_type="text/event-stream")


@router.get("/api/decks/{deck_id}/slides")
async def deck_slides(deck_id: int, request: Request):
    """API: List slides for a deck."""
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    slides = conn.execute(
        "SELECT id, position, title, notes, image_path, content_hash, author, slide_created_date, updated_at, layout_type FROM slides WHERE deck_id = ? ORDER BY position",
        (deck_id,),
    ).fetchall()
    result = []
    for s in slides:
        tags = get_slide_tags(s["id"], db_path)
        d = dict(s)
        d["tags"] = tags
        result.append(d)
    conn.close()
    return result


@router.get("/api/search")
async def search(request: Request, tags: str = "", title: str = ""):
    """API: Search slides by tags and/or title."""
    db_path = request.app.state.db_path
    tag_list = [t.strip() for t in tags.split(",") if t.strip()] or None
    title_contains = title if title else None
    results = search_slides(db_path=db_path, tags=tag_list, title_contains=title_contains)
    for r in results:
        if "deck_name" in r:
            r["display_deck_name"] = display_name(r["deck_name"])
    return results


@router.get("/api/tags")
async def list_all_tags(request: Request):
    """API: List all tags in use across all slides."""
    result = get_all_tags(db_path=request.app.state.db_path)
    return {"tags": result}


@router.get("/api/slides/{slide_id}")
async def get_slide(slide_id: int, request: Request):
    """API: Get a single slide by ID with tags."""
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    row = conn.execute(
        """SELECT s.id, s.deck_id, s.position, s.title, s.notes, s.content_hash, s.image_path,
                  s.author, s.slide_created_date, s.updated_at, s.layout_type,
                  d.name as deck_name
           FROM slides s JOIN decks d ON s.deck_id = d.id
           WHERE s.id = ?""",
        (slide_id,),
    ).fetchone()
    conn.close()
    if not row:
        return JSONResponse({"error": "Slide not found"}, status_code=404)
    slide = dict(row)
    slide["tags"] = get_slide_tags(slide_id, db_path)
    return slide


@router.get("/api/slides/{slide_id}/tags")
async def get_slide_tags_endpoint(slide_id: int, request: Request):
    """API: Get tags for a slide."""
    db_path = request.app.state.db_path
    tags = get_slide_tags(slide_id, db_path)
    return {"tags": tags}


@router.post("/api/slides/{slide_id}/tags")
async def tag_slide(slide_id: int, request: Request):
    """API: Add tags to a slide."""
    db_path = request.app.state.db_path
    body = await request.json()
    tag_names = body.get("tags", [])
    source = body.get("source", "manual")
    add_tags(slide_id, tag_names, source=source, db_path=db_path)
    return {"status": "ok", "tags": get_slide_tags(slide_id, db_path)}


@router.get("/api/export")
async def export(request: Request):
    """API: Export all cataloged slides to CSV and return file."""
    db_path = request.app.state.db_path
    tmp = tempfile.NamedTemporaryFile(suffix=".csv", delete=False)
    tmp.close()
    export_csv(tmp.name, db_path=db_path, export_all=True)
    return FileResponse(tmp.name, filename="slides.csv", media_type="text/csv")


@router.get("/api/models")
async def get_models():
    """API: Get current model configuration."""
    from aippt.config import load_model_config, ConfigError
    try:
        config = load_model_config()
        # Return only the serializable parts (registry values are ModelConfig dataclasses)
        return {
            "defaults": config["defaults"],
            "source": config["source"],
        }
    except ConfigError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)


@router.put("/api/models")
async def update_models(request: Request):
    """API: Update one or more model defaults."""
    from aippt.config import load_model_config, save_model_config, VALID_OPERATIONS, ConfigError

    body = await request.json()
    defaults = body.get("defaults", {})

    try:
        config = load_model_config()
        registry = config["registry"]
        for op, model in defaults.items():
            if op not in VALID_OPERATIONS:
                return JSONResponse({"error": f"Unknown operation '{op}'"}, status_code=400)
            if not isinstance(model, str) or not model:
                return JSONResponse({"error": f"Model for '{op}' must be a non-empty string"}, status_code=400)
            if model not in registry:
                return JSONResponse(
                    {"error": f"Model '{model}' is not in the registry. Add it to models.yaml first."},
                    status_code=400,
                )
            config["defaults"][op] = model

        save_model_config(config["defaults"])
        return {"defaults": config["defaults"], "source": config["source"]}
    except ConfigError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)


@router.post("/api/models/reset")
async def reset_models():
    """API: Reset model configuration to built-in defaults (deprecated)."""
    return JSONResponse(
        {"error": "Reset is no longer supported. Edit models.yaml directly."},
        status_code=410,
    )


@router.get("/api/models/available")
async def available_models():
    """API: List all models from the registry with capabilities."""
    from aippt.config import get_model_registry, ConfigError
    try:
        registry = get_model_registry()
    except ConfigError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)
    result = []
    for name, cfg in registry.items():
        result.append({
            "name": name,
            "provider": cfg.provider,
            "supports_vision": cfg.supports_vision,
            "supports_images": cfg.supports_images,
            "max_input_tokens": cfg.max_input_tokens,
        })
    return result


# ---------------------------------------------------------------------------
# Template configuration endpoints
# ---------------------------------------------------------------------------


@router.get("/api/templates")
async def get_templates():
    """API: Get current template configuration."""
    from aippt.config import load_template_config, TemplateConfigError
    try:
        config = load_template_config()
        return {
            "default_template": config["default_template"],
            "source": config["source"],
        }
    except TemplateConfigError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)


@router.put("/api/templates")
async def update_templates(request: Request):
    """API: Update the default template path."""
    from aippt.config import set_template_default, load_template_config, TemplateConfigError

    body = await request.json()
    template_path = body.get("default_template", "").strip()

    if not template_path:
        return JSONResponse({"error": "default_template must be a non-empty string"}, status_code=400)

    try:
        set_template_default(template_path)
        config = load_template_config()
        return {
            "default_template": config["default_template"],
            "source": config["source"],
        }
    except (TemplateConfigError, ValueError) as exc:
        return JSONResponse({"error": str(exc)}, status_code=400)


# ---------------------------------------------------------------------------
# Taxonomy endpoints
# ---------------------------------------------------------------------------


@router.get("/api/taxonomy")
async def get_taxonomy(request: Request):
    """API: List all taxonomy tags grouped by category."""
    db_path = request.app.state.db_path
    tags = list_taxonomy(db_path)
    # Group by category for the UI
    categories = {}
    for t in tags:
        cat = t["category"] or "(uncategorized)"
        categories.setdefault(cat, []).append(t["name"])
    return {"tags": tags, "by_category": categories}


@router.post("/api/taxonomy")
async def add_taxonomy(request: Request):
    """API: Add a tag to the taxonomy."""
    db_path = request.app.state.db_path
    body = await request.json()
    name = body.get("name", "").strip()
    category = body.get("category", "").strip()
    if not name:
        return JSONResponse({"error": "name is required"}, status_code=400)
    new, updated = add_taxonomy_tags([{"name": name, "category": category}], db_path)
    return {"status": "ok", "new": new, "updated": updated}


@router.delete("/api/taxonomy/{tag_name}")
async def delete_taxonomy(tag_name: str, request: Request):
    """API: Remove a tag from the taxonomy."""
    db_path = request.app.state.db_path
    removed = remove_taxonomy_tag(tag_name, db_path)
    if not removed:
        return JSONResponse({"error": "tag not found"}, status_code=404)
    return {"status": "ok"}


@router.put("/api/taxonomy/{tag_name}")
async def rename_taxonomy(tag_name: str, request: Request):
    """API: Rename a taxonomy tag (also updates slide associations)."""
    db_path = request.app.state.db_path
    body = await request.json()
    new_name = body.get("new_name", "").strip()
    if not new_name:
        return JSONResponse({"error": "new_name is required"}, status_code=400)
    assoc_count = rename_tag(tag_name, new_name, db_path)
    return {"status": "ok", "updated_associations": assoc_count}


@router.post("/api/taxonomy/import")
async def import_taxonomy(request: Request, file: UploadFile = File(...)):
    """API: Import taxonomy from uploaded CSV."""
    db_path = request.app.state.db_path
    tmp = tempfile.NamedTemporaryFile(suffix=".csv", delete=False, mode="wb")
    content = await file.read()
    tmp.write(content)
    tmp.close()
    try:
        new, updated = import_taxonomy_csv(tmp.name, db_path)
        return {"status": "ok", "new": new, "updated": updated, "total": new + updated}
    finally:
        os.unlink(tmp.name)


@router.get("/api/taxonomy/export")
async def export_taxonomy(request: Request):
    """API: Export taxonomy as CSV download."""
    db_path = request.app.state.db_path
    tmp = tempfile.NamedTemporaryFile(suffix=".csv", delete=False)
    tmp.close()
    count = export_taxonomy_csv(tmp.name, db_path)
    return FileResponse(tmp.name, filename="taxonomy.csv", media_type="text/csv")


# ---------------------------------------------------------------------------
# AI analysis endpoints
# ---------------------------------------------------------------------------


def _resolve_user_ntid(request: Request, body: dict) -> str:
    """Resolve user NTID from request body, then env var fallback."""
    ntid = body.get("user_ntid", "").strip()
    if not ntid:
        ntid = os.environ.get("AIPPT_USER_NTID", "").strip()
    return ntid


def _make_llm_client(request: Request, operation: str, model_override: str = None, user_ntid: str = None):
    """Create an LLM client using gateway config and model registry.

    Args:
        request: FastAPI request (for app.state access)
        operation: Config operation name ('feedback', 'notes', 'improvements')
        model_override: Optional model name to use instead of configured default
        user_ntid: Optional NTID for the gateway user header

    Returns:
        Configured LLMClient

    Raises:
        ValueError: If model config is missing or gateway cannot be loaded
    """
    from aippt.llm import LLMClient, load_gateway_config
    from aippt.config import get_model_default, ConfigError

    try:
        model = model_override or get_model_default(operation)
    except ConfigError as exc:
        raise ValueError(str(exc))

    gateway = None
    gateway_config_path = getattr(request.app.state, "gateway_config", None)
    if gateway_config_path:
        if os.path.exists(gateway_config_path):
            gateway = load_gateway_config(gateway_config_path)

    if gateway and gateway.user_header and not user_ntid and not gateway.user_value:
        raise ValueError(
            "User NTID is required for LLM gateway requests. "
            "Set your NTID in the web UI or the AIPPT_USER_NTID environment variable."
        )

    try:
        return LLMClient(model=model, gateway=gateway, user_ntid=user_ntid)
    except (ConfigError, ValueError) as exc:
        raise ValueError(str(exc))


def _resolve_db_path(path: str, project_root: str) -> str:
    """Resolve a path from the database (may be relative or absolute)."""
    if os.path.isabs(path):
        return path
    return os.path.normpath(os.path.join(project_root, path))


def _get_slide_image_path(slide_id: int, db_path: str, project_root: str = None, state=None):
    """Return the image path for a slide, or None if not found.

    In object-storage mode, fetches the PNG from storage into the local cache
    first (cold-pod read-through) when *state* is provided.
    """
    conn = get_db(db_path)
    row = conn.execute(
        "SELECT image_path FROM slides WHERE id = ?", (slide_id,)
    ).fetchone()
    conn.close()
    if not row or not row["image_path"]:
        return None
    path = _resolve_db_path(row["image_path"], project_root or os.getcwd())
    if state is not None and materialize_file(state, path):
        return path
    return path if os.path.exists(path) else None


VIEW_ONLY_MSG = "LLM features are disabled in view-only mode"


@router.post("/api/slides/{slide_id}/analyze")
async def analyze_slide_endpoint(slide_id: int, request: Request):
    """API: Run feedback analysis on a single slide."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": VIEW_ONLY_MSG}, status_code=403)
    from aippt.analyze import analyze_slide

    db_path = request.app.state.db_path
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass

    image_path = _get_slide_image_path(slide_id, db_path, request.app.state.project_root, state=request.app.state)
    if not image_path:
        return JSONResponse({"error": "No image available for this slide"}, status_code=404)

    conn = get_db(db_path)
    row = conn.execute("SELECT title FROM slides WHERE id = ?", (slide_id,)).fetchone()
    conn.close()
    if not row:
        return JSONResponse({"error": "Slide not found"}, status_code=404)

    user_ntid = _resolve_user_ntid(request, body)
    try:
        client = _make_llm_client(request, "feedback", body.get("model"), user_ntid=user_ntid)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)

    try:
        result = analyze_slide(client=client, image_path=image_path, mode="feedback", title=row["title"] or "")
    except Exception as exc:
        return JSONResponse({"error": f"Analysis failed: {exc}"}, status_code=500)

    return {"result": result, "model": client.model}


@router.post("/api/slides/{slide_id}/notes")
async def suggest_notes_endpoint(slide_id: int, request: Request):
    """API: Generate speaker notes for a single slide."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": VIEW_ONLY_MSG}, status_code=403)
    from aippt.analyze import analyze_slide

    db_path = request.app.state.db_path
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass

    image_path = _get_slide_image_path(slide_id, db_path, request.app.state.project_root, state=request.app.state)
    if not image_path:
        return JSONResponse({"error": "No image available for this slide"}, status_code=404)

    conn = get_db(db_path)
    row = conn.execute("SELECT title FROM slides WHERE id = ?", (slide_id,)).fetchone()
    conn.close()
    if not row:
        return JSONResponse({"error": "Slide not found"}, status_code=404)

    user_ntid = _resolve_user_ntid(request, body)
    try:
        client = _make_llm_client(request, "notes", body.get("model"), user_ntid=user_ntid)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)

    try:
        result = analyze_slide(client=client, image_path=image_path, mode="notes", title=row["title"] or "")
    except Exception as exc:
        return JSONResponse({"error": f"Notes generation failed: {exc}"}, status_code=500)

    return {"result": result, "model": client.model}


@router.post("/api/slides/{slide_id}/notes/save")
async def save_notes_endpoint(slide_id: int, request: Request):
    """API: Save notes to the slide record with edit-history tracking."""
    db_path = request.app.state.db_path
    body = await request.json()
    notes = body.get("notes", "").strip()
    source = body.get("source", "web")
    if not notes:
        return JSONResponse({"error": "notes is required"}, status_code=400)

    try:
        changed = record_edit(slide_id, "notes", notes, source=source, db_path=db_path)
    except ValueError:
        return JSONResponse({"error": "Slide not found"}, status_code=404)

    return {"status": "ok", "changed": changed}


@router.get("/api/slides/{slide_id}/notes/history")
async def notes_history_endpoint(slide_id: int, request: Request):
    """API: Return edit history for a slide's notes field, newest first."""
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    row = conn.execute("SELECT id FROM slides WHERE id = ?", (slide_id,)).fetchone()
    if not row:
        conn.close()
        return JSONResponse({"error": "Slide not found"}, status_code=404)

    rows = conn.execute(
        "SELECT old_value, new_value, source, created_at FROM edit_history "
        "WHERE slide_id = ? AND field = 'notes' ORDER BY id DESC",
        (slide_id,),
    ).fetchall()
    conn.close()

    history = [
        {
            "old_value": r["old_value"],
            "new_value": r["new_value"],
            "source": r["source"],
            "created_at": r["created_at"],
        }
        for r in rows
    ]
    return {"history": history}


@router.post("/api/decks/{deck_id}/write-notes")
async def write_notes_to_deck_endpoint(deck_id: int, request: Request):
    """API: Write DB notes back to the original PPTX file (with backup)."""
    from aippt.writeback import write_notes_to_pptx, create_backup

    db_path = request.app.state.db_path
    project_root = request.app.state.project_root
    deck = get_deck_by_id(deck_id, db_path)
    if deck is None:
        return JSONResponse({"error": "Deck not found"}, status_code=404)

    file_path = deck.get("file_path", "")
    if file_path:
        file_path = _resolve_db_path(file_path, project_root)
    if not file_path or not materialize_file(request.app.state, file_path):
        return JSONResponse({"error": "Source file not found"}, status_code=404)

    try:
        backup_path = create_backup(file_path)
    except FileNotFoundError:
        return JSONResponse({"error": "Source file not found"}, status_code=404)

    try:
        result = write_notes_to_pptx(
            file_path, db_path=db_path, deck_id=deck_id
        )
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=409)

    # The original deck was modified in place; push the new version to storage.
    persist_file(request.app.state, file_path)

    return {
        "status": "ok",
        "slides_written": result.slides_written,
        "slides_skipped": result.slides_skipped,
        "slides_total": result.slides_total,
        "backup_path": backup_path,
        "warnings": result.warnings,
    }


@router.post("/api/slides/{slide_id}/improvements")
async def improvements_endpoint(slide_id: int, request: Request):
    """API: Run improvement analysis on a single slide."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": VIEW_ONLY_MSG}, status_code=403)
    from aippt.analyze import analyze_slide

    db_path = request.app.state.db_path
    body = {}
    try:
        body = await request.json()
    except Exception:
        pass

    image_path = _get_slide_image_path(slide_id, db_path, request.app.state.project_root, state=request.app.state)
    if not image_path:
        return JSONResponse({"error": "No image available for this slide"}, status_code=404)

    conn = get_db(db_path)
    row = conn.execute("SELECT title FROM slides WHERE id = ?", (slide_id,)).fetchone()
    conn.close()
    if not row:
        return JSONResponse({"error": "Slide not found"}, status_code=404)

    user_ntid = _resolve_user_ntid(request, body)
    try:
        client = _make_llm_client(request, "feedback", body.get("model"), user_ntid=user_ntid)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=503)

    try:
        result = analyze_slide(client=client, image_path=image_path, mode="improvements", title=row["title"] or "")
    except Exception as exc:
        return JSONResponse({"error": f"Improvements analysis failed: {exc}"}, status_code=500)

    return {"result": result, "model": client.model}


# ---------------------------------------------------------------------------
# Per-slide tag removal
# ---------------------------------------------------------------------------


@router.delete("/api/slides/{slide_id}/tags/{tag_name}")
async def remove_tag_from_slide(slide_id: int, tag_name: str, request: Request):
    """API: Remove a specific tag from a slide."""
    db_path = request.app.state.db_path
    removed = remove_slide_tag(slide_id, tag_name, db_path)
    if not removed:
        return JSONResponse({"error": "tag not found on slide"}, status_code=404)
    return {"status": "ok", "tags": get_slide_tags(slide_id, db_path)}


@router.get("/slide-image/{slide_id}")
async def serve_slide_image(slide_id: int, request: Request):
    """Serve a slide image by slide ID."""
    db_path = request.app.state.db_path
    project_root = request.app.state.project_root
    conn = get_db(db_path)
    row = conn.execute(
        "SELECT image_path FROM slides WHERE id = ?", (slide_id,)
    ).fetchone()
    conn.close()

    if not row or not row["image_path"]:
        return HTMLResponse("No image", status_code=404)

    image_path = _resolve_db_path(row["image_path"], project_root)
    if materialize_file(request.app.state, image_path):
        return FileResponse(
            image_path,
            headers={"Cache-Control": "no-cache"},
        )
    return HTMLResponse("Image not found", status_code=404)


# ---------------------------------------------------------------------------
# Microsoft Graph device-code auth endpoints
#
# These are unauthenticated and ignore view-only mode -- they ARE the auth
# path. Tokens are returned to the browser, which holds them in localStorage
# and forwards them as Authorization: Bearer on subsequent calls.
# ---------------------------------------------------------------------------


def _sources_dir(uploads_dir: str, deck_id: int) -> str:
    """Return the stable per-deck source storage directory.

    ``uploads/sources/<deck_id>/`` is the canonical location for the originating
    outline (.md) and generated scripts (.mjs/.py).  Created on first use.

    Args:
        uploads_dir: The app's uploads root (``app.state.uploads_dir``).
        deck_id: The deck's integer DB ID.

    Returns:
        Absolute path to the per-deck sources directory (created if absent).
    """
    path = os.path.join(uploads_dir, "sources", str(deck_id))
    os.makedirs(path, exist_ok=True)
    return path


def _stage_writable_script(script_path: str, deck_id: int, uploads_dir: str, project_root: str) -> str:
    """Copy a preview script to a writable per-deck location and return the copy's path.

    Chat "edit" patches (and their reverts) write the modified source back to
    ``source_script_path`` in place. Scripts discovered from the bundled
    ``examples/`` tree live on the container's read-only ``/app`` filesystem, so
    an in-place write raises ``OSError: [Errno 30]``. Staging a copy onto the
    writable data volume makes the patch/revert path work.

    The copy preserves the script's path *relative to ``project_root``* under a
    per-deck staging base, and copies the repo ``lib/`` alongside, so relative
    imports of any depth resolve exactly as they do in the repo::

        <uploads>/preview-scripts/<deck_id>/examples/<name>/<script>.mjs  # copy
        <uploads>/preview-scripts/<deck_id>/lib/                         # repo lib/

    e.g. from a staged ``examples/demo/demo.mjs`` the script's own
    ``../../lib/...`` import resolves to the copied ``lib/``. Theme files (read
    relative to the render cwd, which stays ``project_root``) are unaffected.
    Idempotent: re-staging refreshes the copy in place.

    Returns the staged script path, or the original ``script_path`` unchanged if
    staging fails for any reason (the caller then falls back to in-place —
    correct for local dev where the tree is already writable).
    """
    import shutil

    try:
        src = Path(script_path).resolve()
        if not src.is_file():
            return script_path
        base = Path(uploads_dir) / "preview-scripts" / str(deck_id)

        # Preserve the script's path relative to project_root so its relative
        # imports (../../lib, etc.) resolve after the copy. If the script sits
        # outside project_root, fall back to just its basename.
        root = Path(project_root).resolve()
        try:
            rel = src.relative_to(root)
        except ValueError:
            rel = Path(src.name)
        dst = base / rel
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(src, dst)

        # Copy the repo lib/ so relative "../../lib/..." imports resolve.
        lib_src = root / "lib"
        if lib_src.is_dir():
            # dirs_exist_ok keeps re-staging idempotent (Python 3.8+).
            shutil.copytree(lib_src, base / "lib", dirs_exist_ok=True)
        return str(dst)
    except OSError as exc:
        logger.warning(
            "Could not stage writable copy of %s for deck %s (%s); using in-place path",
            script_path, deck_id, exc,
        )
        return script_path


def _stage_render_workspace(uploads_dir: str, short_id: str, base_name: str, project_root: str) -> tuple:
    """Build a per-job workspace for a freshly generated pptxgenjs script.

    The generated script imports the shared helpers with a hardcoded relative
    path (``../lib/pptxgenjs-helpers.mjs``). ESM resolves that against the
    *script file's own URL*, not the process cwd — so the ``lib/`` it needs must
    sit one directory above the script. In the container the script is written
    to the writable ``/app/data`` volume while the repo ``lib/`` is baked at
    ``/app/lib`` (read-only ``/app``), so ``../lib`` misses. Mirror the preview
    path's staging (see :func:`_stage_writable_script`) by laying out, per job::

        <uploads>/<short_id>/script/<base_name>.mjs   # generated script
        <uploads>/<short_id>/lib/                      # fresh copy of repo lib/
        <uploads>/<short_id>/out/                      # render output

    so ``../lib`` from ``script/`` resolves to the copied ``lib/``. Staging is
    per-job (unique ``short_id``) and re-copies ``lib/`` from ``project_root``
    every run, so it is never stale after a redeploy and never races a
    concurrent render truncating a shared copy. Theme/asset files are read
    relative to the render cwd (which stays ``project_root``) and are unaffected.

    Returns ``(script_path, render_out_dir)``.
    """
    import shutil

    job = Path(uploads_dir) / short_id
    script_dir = job / "script"
    out_dir = job / "out"
    script_dir.mkdir(parents=True, exist_ok=True)
    out_dir.mkdir(parents=True, exist_ok=True)

    lib_src = Path(project_root) / "lib"
    if lib_src.is_dir():
        # dirs_exist_ok keeps re-staging idempotent (Python 3.8+).
        shutil.copytree(lib_src, job / "lib", dirs_exist_ok=True)

    return str(script_dir / f"{base_name}.mjs"), str(out_dir)


def _per_deck_images_dir(request: Request, deck_path: str) -> str:
    """Resolve the per-deck images directory under app.state.images_dir.

    Container deployments override the parent images dir via
    ``serve --images-dir /app/data/images`` so PNGs land on the data
    volume; without this helper, ingest_deck falls back to a cwd-relative
    ``images/`` path that lives in ephemeral container storage and is
    lost on pod restart.
    """
    base = getattr(request.app.state, "images_dir", None) or "images"
    deck_name = os.path.splitext(os.path.basename(deck_path))[0]
    return os.path.join(base, deck_name)


def _require_images_for_render() -> bool:
    """Whether the render pipeline must succeed for the upload to succeed.

    On Linux the only image path is Microsoft Graph; a render failure means
    no images at all, which makes the catalog useless. Surface it as an
    HTTP error instead of silently completing.
    """
    return sys.platform.startswith("linux")


def _graph_error_status(exc) -> int:
    """Map a graph.GraphError status_code → an HTTP response status.

    4xx Graph errors pass through (the caller's token or permissions are at
    fault). 5xx and anything weird become 502 Bad Gateway — we couldn't
    complete the request because the upstream service couldn't.
    """
    code = getattr(exc, "status_code", 0) or 0
    if 400 <= code < 500:
        return code
    return 502


def _extract_bearer_token(request: Request) -> str:
    """Return the bearer token from the Authorization header, or '' if absent
    or not a Bearer-scheme header.

    Strict: only ``Authorization: Bearer <token>`` is accepted (case-insensitive
    on 'Bearer'). Anything else — Basic, Digest, a raw token without a scheme,
    or 'Bearer ' with no token — returns '' so the caller treats it as
    unauthenticated.

    Tolerating raw tokens here would let any non-empty Authorization value
    (e.g. ``Authorization: Basic dXNlcjpwYXNz``) slip past the gate.
    """
    raw = request.headers.get("Authorization", "").strip()
    if not raw:
        return ""
    parts = raw.split(None, 1)
    if len(parts) != 2 or parts[0].lower() != "bearer":
        return ""
    token = parts[1].strip()
    return token  # empty after strip → '' → treated as unauthenticated


class InvalidNtid(ValueError):
    """Raised when X-AIPPT-NTID is present but fails the allowlist check."""


# NTIDs are interpolated into the SharePoint path
# /sites/.../root:/<root>/<NTID>/<job>.pptx — any character outside this set
# either splits the path ('/', '\\', ':'), escapes the folder ('..'), or
# causes Graph to 4xx deep inside the pipeline (whitespace, '?', '*', '#').
_NTID_RE = re.compile(r"^[A-Za-z0-9._-]+$")


def _extract_ntid_header(request: Request) -> str:
    """Return the trimmed X-AIPPT-NTID header, or '' if absent.

    Used by the upload endpoints to scope per-user SharePoint render-staging
    folders. We deliberately do NOT fall back to the server-side USER env
    var here — that would route every signed-in user's renders through one
    shared folder. The CLI layer applies env fallbacks for non-web callers.

    Raises InvalidNtid if the header is present but doesn't match
    ``[A-Za-z0-9._-]+``. Validating at the edge keeps malformed input out of
    the Graph request URL.
    """
    raw = request.headers.get("X-AIPPT-NTID", "").strip()
    if not raw:
        return ""
    if not _NTID_RE.match(raw):
        raise InvalidNtid(
            "Invalid X-AIPPT-NTID: must match [A-Za-z0-9._-]+ "
            "(no path separators, whitespace, or URL metacharacters)."
        )
    return raw


def _ntid_or_400(request: Request):
    """Wrap ``_extract_ntid_header`` for the edge endpoints.

    Returns ``(ntid, None)`` on success or ``("", JSONResponse(400))`` if the
    header is malformed. The JSONResponse is shaped like the other 4xx errors
    (``{"error": "..."}``) so the existing client error handling works.
    """
    try:
        return _extract_ntid_header(request), None
    except InvalidNtid as exc:
        return "", JSONResponse({"error": str(exc)}, status_code=400)


def _bearer_identity_unverified(request: Request) -> str:
    """Best-effort identity extraction from the Bearer JWT, **without**
    signature verification. For audit logging only — never use to gate.

    Returns a short string like ``"upn:melliott@amd.com"`` or
    ``"oid:1234-..."`` for the audit log when an admin action fires. Returns
    ``"unparseable"`` if the Bearer isn't a JWT (some Graph tokens aren't),
    or ``"absent"`` if no Bearer is present at all.

    The point of logging this alongside the (trusted) X-AIPPT-NTID is
    impersonation detection: if a user with Bearer token for ``jdoe`` calls
    a DELETE with ``X-AIPPT-NTID: melliott``, we get a paper trail to
    investigate from. We do not block on the discrepancy because in v1 the
    NTID header is the gate (per design), but the audit log makes the
    weakness recoverable.
    """
    import base64
    import json as _json

    token = _extract_bearer_token(request)
    if not token:
        return "absent"

    parts = token.split(".")
    if len(parts) != 3:
        return "unparseable"

    try:
        padded = parts[1] + "=" * (-len(parts[1]) % 4)
        payload = _json.loads(base64.urlsafe_b64decode(padded))
    except Exception:
        return "unparseable"

    if not isinstance(payload, dict):
        return "unparseable"

    for key in ("upn", "preferred_username", "unique_name", "oid", "sub"):
        val = payload.get(key)
        if isinstance(val, str) and val:
            return f"{key}:{val}"
    return "unparseable"


def _suggested_ntid_from_bearer(request: Request) -> str:
    """Best-effort NTID suggestion derived from the Bearer JWT, **unverified**.

    Pulls ``preferred_username`` / ``upn`` / ``unique_name`` from the token
    payload, takes the local-part (before ``@``), lowercases it, and returns
    it only if it matches ``_NTID_RE``. Returns ``""`` when there is no Bearer,
    the token isn't a parseable JWT, no usable claim is present, or the
    derived value is malformed.

    This is a UX hint only — the SPA uses it to pre-fill the NTID field so it
    can't drift from the signed-in identity (the live 403s traced back to a
    hand-typed typo, ``melliot`` vs ``melliott``). It is never used to gate;
    the admin check still trusts the explicit ``X-AIPPT-NTID`` header. Adding
    signature verification is the deferred v2 AAD-groups PRD.
    """
    import base64
    import json as _json

    token = _extract_bearer_token(request)
    if not token:
        return ""

    parts = token.split(".")
    if len(parts) != 3:
        return ""

    try:
        padded = parts[1] + "=" * (-len(parts[1]) % 4)
        payload = _json.loads(base64.urlsafe_b64decode(padded))
    except Exception:
        return ""

    if not isinstance(payload, dict):
        return ""

    for key in ("preferred_username", "upn", "unique_name"):
        val = payload.get(key)
        if not isinstance(val, str) or not val:
            continue
        local = val.split("@", 1)[0].strip().lower()
        if local and _NTID_RE.match(local):
            return local
    return ""


def _is_admin(request: Request) -> bool:
    """Admin-tier v1 gate. True iff: Bearer present, X-AIPPT-NTID is well-
    formed, and the NTID is in ``app.state.admin_ntids`` (sourced from
    gateway.yaml ``admin_ntids:``).

    The membership test is **case-insensitive**: the header is lowercased
    here before comparison and ``load_admin_ntids`` lowercases the allowlist
    at load, so case drift between config and client can't cause a silent
    403. Only the membership test is lowercased — ``_extract_ntid_header``
    still returns the original-case value for SharePoint paths and audit
    logging, which should reflect what the client actually sent.

    This is the v1 design — see ``aippt.config.load_admin_ntids`` for the
    threat model. Callers should also short-circuit on view-only mode before
    asking; admin actions all mutate state, so view-only deployments have
    no admins by definition.
    """
    if not _extract_bearer_token(request):
        return False
    try:
        ntid = _extract_ntid_header(request)
    except InvalidNtid:
        return False
    if not ntid:
        return False
    admins = getattr(request.app.state, "admin_ntids", set()) or set()
    return ntid.lower() in admins


def _require_admin(request: Request, action: str):
    """Convert ``_is_admin`` into a 403 JSONResponse on failure and emit an
    audit log line on success or denial. Returns ``None`` on success (caller
    may proceed) or a ``JSONResponse`` on failure (caller must return it).

    The audit line records ``X-AIPPT-NTID`` (claimed) and the Bearer-derived
    identity (unverified) so impersonation attempts are recoverable from
    ``/api/logs`` even though the gate itself trusts the header.
    """
    claimed_ntid = request.headers.get("X-AIPPT-NTID", "").strip() or "<absent>"
    bearer_id = _bearer_identity_unverified(request)
    if _is_admin(request):
        logger.info(
            "admin_action action=%s ntid=%s bearer_identity_unverified=%s",
            action, claimed_ntid, bearer_id,
        )
        return None
    logger.warning(
        "admin_denied action=%s ntid=%s bearer_identity_unverified=%s",
        action, claimed_ntid, bearer_id,
    )
    return JSONResponse(
        {"error": "Admin access required for this endpoint."},
        status_code=403,
    )


async def _read_json_body_or_400(request: Request):
    """Parse the request body as JSON; return ({}, None) when empty,
    (body, None) on success, or ({}, JSONResponse(400)) for malformed input.

    Without this guard, FastAPI hands ``json.loads`` an arbitrary byte
    stream — a client sending ``Content-Type: application/x-www-form-urlencoded``
    triggers ``JSONDecodeError`` and bubbles a 500 to the user. Auth
    endpoints in particular should never 500 on malformed input; the JS
    client treats 5xx very differently from 4xx.
    """
    raw = await request.body()
    if not raw:
        return {}, None
    try:
        return json.loads(raw), None
    except (json.JSONDecodeError, ValueError):
        return {}, JSONResponse(
            {"error": "Request body must be valid JSON"}, status_code=400,
        )


def _user_auth_error_status(exc) -> int:
    """Status mapping for /poll and /refresh.

    These endpoints' GraphErrors come in two flavors:
      - 4xx from AAD (expired_token, access_denied, invalid_grant) → the user
        needs to re-auth. The browser keys off HTTP 401 to abort the flow,
        so we collapse all 4xx into 401.
      - 5xx from AAD (service outage) → re-auth won't help. Surface as 502
        so the UI shows a transient error instead of bouncing the user back
        to sign-in repeatedly.
    """
    code = getattr(exc, "status_code", 0) or 0
    if 500 <= code:
        return 502
    return 401


@router.post('/api/auth/microsoft/start')
async def auth_microsoft_start(request: Request):
    """Start a Microsoft device-code flow. Unauthenticated."""
    try:
        return graph.start_device_code()
    except graph.GraphError as exc:
        # Pass 4xx through (server-config bugs like invalid_client are easier
        # to debug when the real status reaches the client); collapse 5xx to
        # 502 so transient AAD outages don't masquerade as our bugs.
        return JSONResponse(
            {"error": f"Microsoft auth start failed: {exc.message}",
             "code": exc.error_code},
            status_code=_graph_error_status(exc),
        )


@router.post('/api/auth/microsoft/poll')
async def auth_microsoft_poll(request: Request):
    """Poll the device-code endpoint once. Unauthenticated.

    Body: ``{"device_code": "..."}``.
    Returns either ``{"status": "pending"}`` or the token-bearing dict.
    """
    body, err = await _read_json_body_or_400(request)
    if err is not None:
        return err
    device_code = (body or {}).get("device_code", "").strip()
    if not device_code:
        return JSONResponse(
            {"error": "device_code is required"}, status_code=400,
        )
    try:
        return graph.poll_device_code(device_code)
    except graph.GraphError as exc:
        return JSONResponse(
            {"error": exc.message, "code": exc.error_code},
            status_code=_user_auth_error_status(exc),
        )


@router.get('/api/auth/whoami')
async def auth_whoami(request: Request):
    """API: Identity + capabilities the SPA can use to gate admin UI.

    Returns the caller's signed-in state (Bearer present), the claimed
    NTID from ``X-AIPPT-NTID`` (whatever the SPA sent, validated against
    the allowlist regex; empty if absent or malformed), and whether the
    caller is recognized as an admin under the v1 NTID-allowlist rules.

    This is the SPA's hint endpoint for showing/hiding admin controls.
    The actual gate runs server-side on each admin endpoint -- a malicious
    client that hides whoami's response and still calls DELETE gets 403.

    ``suggested_ntid`` is the local-part of the Bearer token's identity
    claim (lowercased, unverified) — a UX hint the SPA uses to pre-fill the
    NTID field so it can't drift from who the user signed in as. It never
    affects the gate; ``is_admin`` still keys off the explicit header.
    """
    signed_in = bool(_extract_bearer_token(request))
    try:
        ntid = _extract_ntid_header(request)
    except InvalidNtid:
        ntid = ""
    return {
        "signed_in": signed_in,
        "ntid": ntid,
        "is_admin": _is_admin(request),
        "suggested_ntid": _suggested_ntid_from_bearer(request),
    }


@router.post('/api/auth/microsoft/refresh')
async def auth_microsoft_refresh(request: Request):
    """Exchange a refresh token for a fresh access token. Unauthenticated."""
    body, err = await _read_json_body_or_400(request)
    if err is not None:
        return err
    refresh_token = (body or {}).get("refresh_token", "").strip()
    if not refresh_token:
        return JSONResponse(
            {"error": "refresh_token is required"}, status_code=400,
        )
    try:
        return graph.refresh_access_token(refresh_token)
    except graph.GraphError as exc:
        return JSONResponse(
            {"error": exc.message, "code": exc.error_code},
            status_code=_user_auth_error_status(exc),
        )


# ---------------------------------------------------------------------------
# Deck upload and download endpoints
# ---------------------------------------------------------------------------


@router.post('/api/decks/upload')
async def upload_deck(
    request: Request,
    file: UploadFile = File(...),
    generate_tags: bool = Form(False),
):
    """API: Upload a .pptx file, export images, catalog, and optionally generate tags.

    Requires an ``Authorization: Bearer <ms-token>`` header so the Linux
    Graph render path has a token. Returns 403 in view-only mode and 401
    when no token was provided.
    """
    db_path = request.app.state.db_path
    uploads_dir = request.app.state.uploads_dir
    gateway_config = request.app.state.gateway_config

    # Validate file extension
    if not file.filename or not file.filename.lower().endswith('.pptx'):
        return JSONResponse({'error': 'Only .pptx files are supported'}, status_code=400)

    # View-only deployments cannot ingest at all (LLM access required for
    # downstream tagging, and the render path needs a per-user MS token).
    if getattr(request.app.state, "view_only", False):
        return JSONResponse(
            {"error": "Deck ingest is disabled in view-only mode"},
            status_code=403,
        )

    ms_token = _extract_bearer_token(request)
    if not ms_token:
        return JSONResponse(
            {"error": "Microsoft sign-in required for ingest. "
                      "Sign in with the Microsoft button and retry."},
            status_code=401,
        )
    ntid, _ntid_err = _ntid_or_400(request)
    if _ntid_err is not None:
        return _ntid_err

    # Build a collision-safe filename and save to uploads_dir
    unique_prefix = uuid.uuid4().hex
    safe_name = f'{unique_prefix}_{file.filename}'
    dest_path = os.path.join(uploads_dir, safe_name)

    content = await file.read()
    max_bytes = getattr(request.app.state, "max_upload_bytes", 0) or 0
    if max_bytes > 0 and len(content) > max_bytes:
        # Backstop for chunked uploads that bypass the Content-Length middleware.
        return JSONResponse(
            {"error": f"Upload exceeds maximum size of {max_bytes} bytes "
                      f"({max_bytes // (1024 * 1024)} MB).",
             "max_bytes": max_bytes,
             "observed_bytes": len(content)},
            status_code=413,
        )
    with open(dest_path, 'wb') as fh:
        fh.write(content)
    persist_file(request.app.state, dest_path)

    # Run ingest pipeline: export images → catalog → optional tags
    per_deck_images = _per_deck_images_dir(request, dest_path)
    try:
        result = ingest_deck(
            deck_path=dest_path,
            db_path=db_path,
            images_dir=per_deck_images,
            generate_tags=generate_tags,
            gateway_config=gateway_config,
            require_images=_require_images_for_render(),
            ms_token=ms_token,
            ntid=ntid,
        )
    except graph.GraphError as exc:
        return JSONResponse(
            {'error': f'Microsoft Graph error: {exc.message}',
             'code': exc.error_code},
            status_code=_graph_error_status(exc),
        )
    except RuntimeError as exc:
        # ingest_deck raises RuntimeError when require_images=True and the
        # render path failed for a non-Graph reason (pdftoppm missing, etc.)
        return JSONResponse(
            {'error': str(exc)}, status_code=502,
        )
    except Exception as exc:
        return JSONResponse({'error': f'Failed to ingest deck: {exc}'}, status_code=500)

    # Push rendered slide images to durable storage (no-op in fs mode).
    persist_tree(request.app.state, per_deck_images)

    # Build descriptive message
    parts = [f"Deck '{result['deck_name']}' uploaded"]
    parts.append(f"{result['slide_count']} slide{'s' if result['slide_count'] != 1 else ''}")
    if result['images_exported']:
        parts.append("images exported")
    if result['tags_generated']:
        parts.append("tags generated")

    return {
        'id': result['deck_id'],
        'name': result['deck_name'],
        'display_name': display_name(result['deck_name']),
        'slide_count': result['slide_count'],
        'images_exported': result['images_exported'],
        'tags_generated': result['tags_generated'],
        'message': " — ".join(parts),
    }


def _save_image_uploads(md_text, outline_save_path, image_data):
    """Save uploaded images to paths matching IMAGE: directives.

    Parses IMAGE: directives from the outline to determine subdirectory
    structure (e.g., ``IMAGE: memes/photo.jpg`` -> save as ``memes/photo.jpg``
    relative to the outline). Falls back to saving with the original filename
    if no matching directive is found.

    Args:
        md_text: Raw markdown outline text.
        outline_save_path: Path where the outline was saved on disk.
        image_data: List of (filename, bytes) tuples from uploaded images.
    """
    outline_dir = os.path.dirname(outline_save_path)

    # Build map: basename -> relative path from IMAGE: directives
    image_map = {}
    for match in re.finditer(r'^IMAGE:\s*(.+)$', md_text, re.MULTILINE):
        path = match.group(1).strip()
        basename = os.path.basename(path)
        image_map[basename] = path

    for filename, data in image_data:
        # Look up full relative path from IMAGE: directives
        rel_path = image_map.get(filename, filename)
        save_path = os.path.join(outline_dir, rel_path)
        os.makedirs(os.path.dirname(save_path), exist_ok=True)
        with open(save_path, 'wb') as f:
            f.write(data)


@router.post('/api/decks/create')
async def create_deck_stream(
    request: Request,
    outline_text: str = Form(None),
    outline_file: UploadFile = File(None),
    enhance: bool = Form(True),
    model: str = Form(None),
    audience: str = Form(None),
    title: str = Form(None),
    image_files: List[UploadFile] = File(default=[]),
    engine: str = Form("pptxgenjs"),
):
    """API: Create a deck from markdown outline, streaming progress as SSE."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": VIEW_ONLY_MSG}, status_code=403)

    ms_token = _extract_bearer_token(request)
    if not ms_token:
        return JSONResponse(
            {"error": "Microsoft sign-in required for ingest. "
                      "Sign in with the Microsoft button and retry."},
            status_code=401,
        )
    ntid, _ntid_err = _ntid_or_400(request)
    if _ntid_err is not None:
        return _ntid_err

    import asyncio
    import json
    import queue as _queue

    # --- Validate inputs before entering SSE mode ---
    md_text = None
    if outline_text and outline_text.strip():
        md_text = outline_text
    elif outline_file and outline_file.filename:
        content = await outline_file.read()
        md_text = content.decode("utf-8")

    if not md_text or not md_text.strip():
        return JSONResponse(
            {"error": "Provide outline text or upload a .md file"},
            status_code=400,
        )

    # Template is only needed for the python-pptx engine
    template_path = None
    if engine == "python-pptx":
        try:
            template_path = get_template_default()
        except TemplateConfigError as exc:
            return JSONResponse({"error": str(exc)}, status_code=503)
        if not os.path.exists(template_path):
            return JSONResponse(
                {"error": f"Template not found: {template_path}. Update the path in Settings."},
                status_code=404,
            )

    db_path = request.app.state.db_path
    uploads_dir = request.app.state.uploads_dir
    gateway_config = request.app.state.gateway_config

    # Derive filename from outline title (first H1 or H2)
    _title_match = re.search(r'^#{1,2}\s+(.+)', md_text, re.MULTILINE)
    _base_name = _title_match.group(1).strip() if _title_match else "generated"
    # Sanitize for filesystem: keep alphanumeric, spaces, hyphens, underscores
    _base_name = re.sub(r'[^\w\s-]', '', _base_name).strip()
    _base_name = re.sub(r'\s+', ' ', _base_name)[:80]  # cap length
    _short_id = uuid.uuid4().hex[:8]
    output_path = os.path.join(uploads_dir, f"{_short_id}_{_base_name}.pptx")

    # Save outline to disk for IMAGE: directive resolution
    outline_filename = outline_file.filename if (outline_file and outline_file.filename) else "outline.md"
    outline_save_path = os.path.join(uploads_dir, f"{_short_id}_{outline_filename}")
    with open(outline_save_path, 'w', encoding='utf-8') as f:
        f.write(md_text)

    # Save uploaded images relative to the outline
    image_data = []
    for img in (image_files or []):
        if img.filename:
            data = await img.read()
            if data:
                image_data.append((img.filename, data))
    if image_data:
        _save_image_uploads(md_text, outline_save_path, image_data)

    event_q: _queue.Queue = _queue.Queue()

    def create_progress(step, detail=""):
        # Map completion markers to done status
        status = "running"
        if any(detail.startswith(w) for w in ("Parsed", "All", "Built")):
            status = "done"
        event_q.put(("progress", {"step": step, "status": status, "detail": detail}))

    def ingest_progress(step, detail=""):
        ingest_map = {
            "export_images": "running", "export_images_done": "running",
            "export_images_skipped": "running", "catalog": "running",
            "catalog_done": "running", "complete": "done",
        }
        if step in ingest_map:
            event_q.put(("progress", {"step": "ingest", "status": ingest_map[step], "detail": detail}))

    async def _event_generator():
        loop = asyncio.get_running_loop()

        def _worker():
            if engine == "pptxgenjs":
                return _worker_pptxgenjs()
            return _worker_python_pptx()

        def _worker_pptxgenjs():
            from aippt.pptxgenjs_gen import generate_script, ScriptGenerationError
            from aippt.preview import Renderer

            # Stage a per-job workspace so the generated script's ``../lib``
            # import resolves to a fresh copy of the repo lib/ on the writable
            # data volume (the baked /app/lib is unreachable via ../lib from
            # /app/data). See _stage_render_workspace.
            project_root = str(Path(__file__).resolve().parents[2])
            script_path, render_out_dir = _stage_render_workspace(
                uploads_dir, _short_id, _base_name, project_root
            )
            try:
                llm = _make_llm_client(request, operation="create", model_override=model, user_ntid=ntid)
            except ValueError as exc:
                event_q.put(("error", {"message": str(exc)}))
                return None

            try:
                generate_script(
                    outline_text=md_text,
                    output_script_path=script_path,
                    llm_client=llm,
                    theme="amd",
                    progress_callback=create_progress,
                )
            except ScriptGenerationError as exc:
                event_q.put(("error", {"message": str(exc)}))
                return None

            persist_file(request.app.state, script_path)
            create_progress("render", "Executing script…")

            renderer = Renderer(project_root=project_root)
            render_result = renderer.render(script_path, render_out_dir)
            if not render_result.success:
                msg = render_result.stderr_tail or "Script render failed (no .pptx produced)"
                event_q.put(("error", {"message": msg}))
                return None

            pptx_path = render_result.pptx_path
            persist_file(request.app.state, pptx_path)
            event_q.put(("progress", {"step": "ingest", "status": "running", "detail": "Cataloging generated deck…"}))

            per_deck_images = _per_deck_images_dir(request, pptx_path)
            ingest_result = ingest_deck(
                deck_path=pptx_path,
                db_path=db_path,
                images_dir=per_deck_images,
                gateway_config=gateway_config,
                require_images=_require_images_for_render(),
                progress_callback=ingest_progress,
                source_script_path=script_path,
                outline_path=outline_save_path,
                ms_token=ms_token,
                ntid=ntid,
            )
            persist_tree(request.app.state, per_deck_images)

            deck_id = ingest_result["deck_id"]
            try:
                import shutil as _shutil
                stable_dir = _sources_dir(uploads_dir, deck_id)
                stable_outline = os.path.join(stable_dir, "outline.md")
                stable_script = os.path.join(stable_dir, f"{_base_name}.mjs")
                _shutil.copy2(outline_save_path, stable_outline)
                _shutil.copy2(script_path, stable_script)
                persist_file(request.app.state, stable_outline)
                persist_file(request.app.state, stable_script)
                _conn = get_db(db_path)
                _conn.execute(
                    "UPDATE decks SET outline_path = ?, source_engine = ?, source_script_path = ?, "
                    "source_generated_at = datetime('now'), updated_at = datetime('now') "
                    "WHERE id = ?",
                    (stable_outline, "pptxgenjs", stable_script, deck_id),
                )
                _conn.commit()
                _conn.close()
            except Exception as _exc:
                logger.warning("Failed to persist pptxgenjs origin for deck %s: %s", deck_id, _exc)

            return {"output_path": pptx_path, "slide_count": ingest_result.get("slide_count", 0),
                    "title": _base_name, **ingest_result}

        def _worker_python_pptx():
            pipeline_config = PipelineConfig(
                outline_text=md_text,
                template_path=template_path,
                output_path=output_path,
                enhance=enhance,
                model=model,
                audience=audience,
                gateway_config=gateway_config,
                progress_callback=create_progress,
                outline_path=outline_save_path,
                source_engine="python-pptx",
                source_kind="outline",
            )
            pipeline_result = _pipeline_module.run_pipeline(pipeline_config)
            result = {
                "output_path": pipeline_result.output_path,
                "slide_count": pipeline_result.slide_count,
                "title": pipeline_result.title,
            }
            # Persist the generated deck + originating outline to durable storage.
            persist_file(request.app.state, output_path)
            persist_file(request.app.state, outline_save_path)
            event_q.put(("progress", {"step": "ingest", "status": "running", "detail": "Cataloging generated deck..."}))
            per_deck_images = _per_deck_images_dir(request, output_path)
            ingest_result = ingest_deck(
                deck_path=output_path,
                db_path=db_path,
                images_dir=per_deck_images,
                gateway_config=gateway_config,
                require_images=_require_images_for_render(),
                progress_callback=ingest_progress,
                ms_token=ms_token,
                ntid=ntid,
                # outline_path will be set below after we know deck_id
            )
            persist_tree(request.app.state, per_deck_images)
            # Copy outline to stable per-deck source location
            deck_id = ingest_result["deck_id"]
            try:
                import shutil as _shutil
                stable_dir = _sources_dir(uploads_dir, deck_id)
                stable_outline = os.path.join(stable_dir, "outline.md")
                _shutil.copy2(outline_save_path, stable_outline)
                persist_file(request.app.state, stable_outline)
                # Update the deck row with the stable outline path
                _conn = get_db(db_path)
                _conn.execute(
                    "UPDATE decks SET outline_path = ?, source_engine = ?, "
                    "source_generated_at = datetime('now'), updated_at = datetime('now') "
                    "WHERE id = ?",
                    (stable_outline, "python-pptx", deck_id),
                )
                _conn.commit()
                _conn.close()
            except Exception as _exc:
                logger.warning("Failed to persist outline origin for deck %s: %s", deck_id, _exc)
            return {**result, **ingest_result}

        future = loop.run_in_executor(None, _worker)

        def _format_sse(event_name, payload):
            return f"event: {event_name}\ndata: {json.dumps(payload)}\n\n"

        while not future.done():
            await asyncio.sleep(0)
            while True:
                try:
                    event_name, payload = event_q.get_nowait()
                    yield _format_sse(event_name, payload)
                except _queue.Empty:
                    break

        while True:
            try:
                event_name, payload = event_q.get_nowait()
                yield _format_sse(event_name, payload)
            except _queue.Empty:
                break

        try:
            result = await future
        except graph.GraphError as exc:
            yield _format_sse("error", {
                "detail": f"Microsoft Graph error: {exc.message}",
                "code": exc.error_code,
                "status": _graph_error_status(exc),
            })
            return
        except Exception as exc:
            yield _format_sse("error", {"detail": str(exc)})
            return

        if result is None:
            return  # worker already emitted an error event

        yield _format_sse("complete", {
            "deck_id": result["deck_id"],
            "deck_name": result.get("deck_name", ""),
            "display_name": display_name(result.get("deck_name", "")),
            "slide_count": result.get("slide_count", 0),
            "output_path": result.get("output_path", ""),
        })

    return StreamingResponse(_event_generator(), media_type="text/event-stream")


@router.post('/api/decks/upload-stream')
async def upload_deck_stream(
    request: Request,
    file: UploadFile = File(...),
    generate_tags: bool = Form(False),
):
    """API: Upload a .pptx file and stream progress as Server-Sent Events.

    The response is a ``text/event-stream`` that emits ``progress`` events for
    each ingest step (export_images, catalog, tags) and a final ``complete`` or
    ``error`` event.  Clients that do not support SSE should use the regular
    ``/api/decks/upload`` endpoint instead.
    """
    import asyncio
    import json
    import queue as _queue

    # NOTE: ingest_deck is imported at module top; do NOT re-import locally,
    # or test patches on `aippt.web.routes.ingest_deck` won't take effect for
    # the SSE handler.

    # Validate file extension before entering SSE mode so we can return a
    # plain JSON 400 response (SSE has already started once we yield the first
    # byte, so validation must happen first).
    if not file.filename or not file.filename.lower().endswith('.pptx'):
        return JSONResponse({'error': 'Only .pptx files are supported'}, status_code=400)

    # Same gates as /api/decks/upload — view-only blocks ingest entirely, and
    # the Linux Graph render path requires a per-user MS token from the browser.
    if getattr(request.app.state, "view_only", False):
        return JSONResponse(
            {"error": "Deck ingest is disabled in view-only mode"},
            status_code=403,
        )

    ms_token = _extract_bearer_token(request)
    if not ms_token:
        return JSONResponse(
            {"error": "Microsoft sign-in required for ingest. "
                      "Sign in with the Microsoft button and retry."},
            status_code=401,
        )
    ntid, _ntid_err = _ntid_or_400(request)
    if _ntid_err is not None:
        return _ntid_err

    db_path = request.app.state.db_path
    uploads_dir = request.app.state.uploads_dir
    gateway_config = request.app.state.gateway_config

    # Save uploaded file to uploads_dir
    unique_prefix = uuid.uuid4().hex
    safe_name = f'{unique_prefix}_{file.filename}'
    dest_path = os.path.join(uploads_dir, safe_name)

    content = await file.read()
    max_bytes = getattr(request.app.state, "max_upload_bytes", 0) or 0
    if max_bytes > 0 and len(content) > max_bytes:
        # Backstop for chunked uploads that bypass the Content-Length middleware.
        return JSONResponse(
            {"error": f"Upload exceeds maximum size of {max_bytes} bytes "
                      f"({max_bytes // (1024 * 1024)} MB).",
             "max_bytes": max_bytes,
             "observed_bytes": len(content)},
            status_code=413,
        )
    with open(dest_path, 'wb') as fh:
        fh.write(content)

    # Step-name → (step field, status field) mapping
    _STEP_MAP = {
        'export_images':         ('export_images', 'running'),
        'export_images_done':    ('export_images', 'done'),
        'export_images_skipped': ('export_images', 'skipped'),
        'catalog':               ('catalog',        'running'),
        'catalog_done':          ('catalog',        'done'),
        'tags':                  ('tags',           'running'),
        'tags_done':             ('tags',           'done'),
        'tags_partial':          ('tags',           'done'),
        'tags_error':            ('tags',           'error'),
    }

    # Queue is used to pass events from the worker thread to the generator.
    event_q: _queue.Queue = _queue.Queue()

    def progress_callback(step: str, detail: str = ''):
        """Called by ingest_deck at each pipeline stage."""
        if step in _STEP_MAP:
            mapped_step, status = _STEP_MAP[step]
            event_q.put(('progress', {'step': mapped_step, 'status': status, 'detail': detail}))
        # 'complete' is handled after the future resolves

    async def _event_generator():
        """Async generator that drives the ingest in a thread and yields SSE lines."""
        loop = asyncio.get_running_loop()

        # Run the blocking ingest in a thread pool executor so we don't block
        # the event loop.  Events land on event_q via progress_callback.
        future = loop.run_in_executor(
            None,
            lambda: ingest_deck(
                deck_path=dest_path,
                db_path=db_path,
                images_dir=_per_deck_images_dir(request, dest_path),
                generate_tags=generate_tags,
                gateway_config=gateway_config,
                require_images=_require_images_for_render(),
                progress_callback=progress_callback,
                ms_token=ms_token,
                ntid=ntid,
            ),
        )

        def _format_sse(event_name: str, payload: dict) -> str:
            return f'event: {event_name}\ndata: {json.dumps(payload)}\n\n'

        # Drain queued progress events while the future runs.
        while not future.done():
            # Give the event loop a chance to advance the future.
            await asyncio.sleep(0)
            while True:
                try:
                    event_name, payload = event_q.get_nowait()
                    yield _format_sse(event_name, payload)
                except _queue.Empty:
                    break

        # Drain any remaining events that arrived before the loop exited.
        while True:
            try:
                event_name, payload = event_q.get_nowait()
                yield _format_sse(event_name, payload)
            except _queue.Empty:
                break

        # Retrieve the result (or propagate the exception).
        try:
            result = await future
        except graph.GraphError as exc:
            yield _format_sse('error', {
                'detail': f'Microsoft Graph error: {exc.message}',
                'code': exc.error_code,
                'status': _graph_error_status(exc),
            })
            return
        except Exception as exc:  # noqa: BLE001
            yield _format_sse('error', {'detail': str(exc)})
            return

        yield _format_sse('complete', {
            'deck_id':        result['deck_id'],
            'deck_name':      result['deck_name'],
            'display_name':   display_name(result['deck_name']),
            'slide_count':    result['slide_count'],
            'images_exported': result['images_exported'],
            'tags_generated':  result['tags_generated'],
        })

    return StreamingResponse(_event_generator(), media_type='text/event-stream')


@router.get('/api/decks/{deck_id}/download')
async def download_deck(deck_id: int, request: Request):
    """API: Download a .pptx file with DB notes applied (temp copy, original untouched)."""
    from aippt.writeback import write_notes_to_pptx

    db_path = request.app.state.db_path
    project_root = request.app.state.project_root

    deck = get_deck_by_id(deck_id, db_path)
    if deck is None:
        return JSONResponse({'error': 'Deck not found'}, status_code=404)

    file_path = deck.get('file_path', '')
    if file_path:
        file_path = _resolve_db_path(file_path, project_root)
    if not file_path or not materialize_file(request.app.state, file_path):
        return JSONResponse({'error': 'Source file not found'}, status_code=404)

    # Create temp copy with DB notes applied
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    tmp.close()
    try:
        write_notes_to_pptx(
            file_path, db_path=db_path, deck_id=deck_id, output_path=tmp.name
        )
    except (FileNotFoundError, ValueError):
        # If write-back fails (e.g. mismatch), fall back to serving original
        os.unlink(tmp.name)
        tmp_name = file_path
    else:
        tmp_name = tmp.name

    download_name = display_name(deck['name'])
    return FileResponse(
        tmp_name,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={'Content-Disposition': f'attachment; filename="{download_name}.pptx"'},
    )


@router.get('/api/logs')
async def get_logs(request: Request):
    """API: return recent in-memory application log records.

    Captures land in the ring buffer attached during ``create_app``; this
    endpoint is read-only triage that surfaces server-internal state
    (including admin-action audit lines), so it is admin-gated as well as
    Bearer-gated. Allowed in view-only deployments because it has no
    mutating side effects. Bearer tokens themselves are scrubbed by
    ``install_authorization_scrub`` before they reach the buffer.

    Query params: ``limit`` (1-2000, default 200), ``level`` (DEBUG/INFO/
    WARNING/ERROR/CRITICAL), ``since`` (record id from a prior poll),
    ``logger_prefix`` (filter by logger name prefix, e.g. ``aippt``).
    """
    if not _extract_bearer_token(request):
        return JSONResponse(
            {"error": "Microsoft sign-in required to view logs."},
            status_code=401,
        )
    denied = _require_admin(request, action="get_logs")
    if denied is not None:
        return denied

    buffer = getattr(request.app.state, "log_buffer", None)
    if buffer is None:
        return JSONResponse(
            {"error": "Log buffer not configured."},
            status_code=503,
        )

    try:
        limit = int(request.query_params.get("limit", "200"))
    except ValueError:
        limit = 200
    limit = max(1, min(limit, buffer.capacity))

    level = request.query_params.get("level")
    logger_prefix = request.query_params.get("logger_prefix")
    since_raw = request.query_params.get("since")
    since: int | None = None
    if since_raw:
        try:
            since = int(since_raw)
        except ValueError:
            since = None

    records = buffer.snapshot(
        limit=limit, level=level, since=since,
        logger_prefix=logger_prefix,
    )
    next_cursor = records[-1]["id"] if records else since
    return {
        "capacity": buffer.capacity,
        "count": len(records),
        "next_cursor": next_cursor,
        "records": records,
    }


@router.delete('/api/decks/{deck_id}')
async def remove_deck(deck_id: int, request: Request):
    """API: Delete a deck (cascade) and purge its rendered images.

    Admin-gated (NTID allowlist in ``gateway.yaml`` ``admin_ntids:``) on top
    of the Bearer requirement. ``view_only`` deployments reject with 403.
    ``purge_images=false`` keeps the PNG dir; default purges.
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse(
            {"error": "Deck delete is disabled in view-only mode"},
            status_code=403,
        )
    if not _extract_bearer_token(request):
        return JSONResponse(
            {"error": "Microsoft sign-in required to delete decks."},
            status_code=401,
        )
    denied = _require_admin(request, action=f"delete_deck:{deck_id}")
    if denied is not None:
        return denied

    db_path = request.app.state.db_path
    images_base = request.app.state.images_dir

    deck = get_deck_by_id(deck_id, db_path)
    if deck is None:
        return JSONResponse({"error": "Deck not found"}, status_code=404)

    purge = request.query_params.get("purge_images", "true").lower() not in (
        "false", "0", "no",
    )

    info = delete_deck(deck_id, db_path=db_path)
    if info is None:
        return JSONResponse({"error": "Deck not found"}, status_code=404)

    purged_dir = None
    if purge:
        # Per-deck images live at {images_dir}/{deck_name}/. delete_deck
        # already removed the DB rows; clean the matching PNG dir so a
        # subsequent re-upload doesn't pile up under an ephemeral volume.
        candidate = os.path.join(images_base, deck["name"])
        if os.path.isdir(candidate):
            import shutil as _shutil
            _shutil.rmtree(candidate, ignore_errors=True)
            purged_dir = candidate

    return {
        "deck_id": deck_id,
        "name": info["name"],
        "display_name": display_name(info["name"]),
        "slide_count": info["slide_count"],
        "tag_count": info["tag_count"],
        "section_count": info["section_count"],
        "images_purged": purged_dir is not None,
        "images_dir": purged_dir,
    }


# ---------------------------------------------------------------------------
# Live preview endpoints  (Task 6 + Task 7)
# ---------------------------------------------------------------------------


def _preview_registry(request: Request):
    registry = getattr(request.app.state, "preview_registry", None)
    if registry is None:
        raise RuntimeError("Preview registry not initialised")
    return registry


@router.post("/api/preview/sessions")
async def create_preview_session(request: Request):
    """Create a preview session for a slides-as-code script.

    Body: ``{"script": "<path>"}``

    Returns ``{"token": "...", "ws_url": "/ws/preview/<token>", "script": "..."}``
    or 403 in view-only mode / when the script is outside the allow-list.
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Live preview is disabled in view-only mode"}, status_code=403)

    body = await request.json()
    script = body.get("script", "").strip()
    if not script:
        return JSONResponse({"error": "script path is required"}, status_code=400)

    registry = _preview_registry(request)
    try:
        session = await registry.create(script)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=403)
    except FileNotFoundError as exc:
        return JSONResponse({"error": str(exc)}, status_code=404)

    return {
        "token": session.token,
        "ws_url": f"{_base_prefix()}/ws/preview/{session.token}",
        "script": session.script_path,
    }


@router.get("/api/preview/sessions")
async def find_preview_session_by_script(request: Request, script: str = None):
    """Find a preview session by its script path. Used by the inline chat preview."""
    if not script:
        return JSONResponse({"error": "script query param required"}, status_code=400)
    registry = _preview_registry(request)
    resolved = str(Path(script).resolve())
    script_name = Path(script).name
    sessions = list(registry._sessions.values())
    # Exact resolved-path match first.
    for session in sessions:
        if str(Path(session.script_path).resolve()) == resolved:
            return {
                "token": session.token,
                "ws_url": f"{_base_prefix()}/ws/preview/{session.token}",
                "script": session.script_path,
            }
    # Fallback: same filename (handles uploads copy vs output copy mismatch).
    for session in sessions:
        if Path(session.script_path).name == script_name:
            return {
                "token": session.token,
                "ws_url": f"{_base_prefix()}/ws/preview/{session.token}",
                "script": session.script_path,
            }
    return JSONResponse({"error": "No active preview session for that script"}, status_code=404)


@router.get("/api/preview/sessions/{token}")
async def get_preview_session(token: str, request: Request):
    """Return the current state of a preview session (polling fallback)."""
    registry = _preview_registry(request)
    session = registry.get(token)
    if session is None:
        return JSONResponse({"error": "Session not found"}, status_code=404)
    return {"token": token, "script": session.script_path, "state": session.last_state}


@router.delete("/api/preview/sessions/{token}")
async def delete_preview_session(token: str, request: Request):
    """Tear down a preview session and stop its watcher."""
    registry = _preview_registry(request)
    session = registry.get(token)
    if session is None:
        return JSONResponse({"error": "Session not found"}, status_code=404)
    await registry.delete(token)
    return {"status": "stopped", "token": token}


@router.get("/api/preview/sessions/{token}/pptx")
async def get_preview_pptx(token: str, request: Request):
    """Serve the last rendered .pptx for a preview session (consumed by PptxViewJS)."""
    registry = _preview_registry(request)
    session = registry.get(token)
    if session is None:
        return JSONResponse({"error": "Session not found"}, status_code=404)
    pptx_path = getattr(session, "last_pptx_path", None)
    if not pptx_path or not Path(pptx_path).exists():
        return JSONResponse({"error": "No rendered .pptx available yet"}, status_code=404)
    script_stem = Path(session.script_path).stem if session.script_path else "deck"
    filename = f"{script_stem}.pptx"
    return FileResponse(
        pptx_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Cache-Control": "no-cache",
            "Content-Disposition": f'attachment; filename="{filename}"',
        },
    )


@router.post("/api/preview/sessions/{token}/catalog")
async def catalog_preview_session(token: str, request: Request):
    """Catalog the last rendered PPTX from a preview session into the deck library.

    Sets source_script_path so the deck shows up as a 'script' origin and
    a chat conversation can be created with the correct script path wired in.
    Returns the deck id and whether it was newly cataloged or already existed.
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Catalog is disabled in view-only mode"}, status_code=403)
    registry = _preview_registry(request)
    session = registry.get(token)
    if session is None:
        return JSONResponse({"error": "Session not found"}, status_code=404)
    pptx_path = getattr(session, "last_pptx_path", None)
    if not pptx_path or not Path(pptx_path).exists():
        return JSONResponse({"error": "No rendered .pptx available yet — force a render first"}, status_code=404)

    # Optional client-captured slide images (Approach B): the browser renders
    # each slide via PptxViewJS and posts them so cataloged preview decks get
    # per-slide thumbnails without a Microsoft Graph render at catalog time.
    # Absent/empty body is fine — the deck just keeps placeholder cards.
    images_payload = []
    try:
        body = await request.json()
        if isinstance(body, dict):
            payload = body.get("images")
            if isinstance(payload, list):
                images_payload = payload
    except Exception:
        images_payload = []

    db_path = request.app.state.db_path
    try:
        deck_id = catalog_deck(
            deck_path=pptx_path,
            db_path=db_path,
            source_script_path=session.script_path,
            source_engine="pptxgenjs",
        )
        # Stage a writable copy of the script so chat "edit" patches (which write
        # the source back in place) don't hit the read-only /app filesystem when
        # the script came from the bundled examples/ tree. Falls back to the
        # original path if staging fails (local dev / already-writable trees).
        uploads_dir = getattr(request.app.state, "uploads_dir", "uploads")
        project_root = getattr(request.app.state, "project_root", os.getcwd())
        staged_script = _stage_writable_script(
            session.script_path, deck_id, uploads_dir, project_root
        )
        # catalog_deck early-returns on hash match without updating source_script_path.
        # Always ensure it points at the staged (writable) path so the chat
        # conversation and its patch/revert writes target the copy.
        conn = get_db(db_path)
        try:
            conn.execute(
                "UPDATE decks SET source_script_path = ? WHERE id = ? AND (source_script_path IS NULL OR source_script_path != ?)",
                (staged_script, deck_id, staged_script),
            )
            conn.commit()
        finally:
            conn.close()
    except Exception as exc:
        logger.error("catalog_preview_session error: %s", exc)
        return JSONResponse({"error": str(exc)}, status_code=500)

    # Best-effort thumbnail generation: never turns a successful catalog into a
    # 500. On any failure the deck keeps placeholder cards, exactly as before.
    thumb_count = 0
    if images_payload:
        try:
            images_dir = _per_deck_images_dir(request, pptx_path)
            thumb_count = store_client_images(
                deck_id,
                images_payload,
                out_dir=images_dir,
                db_path=db_path,
                base_dir=project_root,
            )
            # Write-through to durable object storage (no-op in fs mode).
            persist_tree(request.app.state, images_dir)
        except Exception:
            logger.exception("Thumbnail generation failed for deck %s", deck_id)

    return {"deck_id": deck_id, "script_path": staged_script, "thumbnails": thumb_count}


@router.get("/api/preview/scripts")
async def list_preview_scripts(request: Request):
    """Return discoverable slides-as-code scripts under output/ and examples/.

    Used by the Live Preview UI to populate the script picker dropdown.
    """
    project_root = getattr(request.app.state, "project_root", os.getcwd())
    search_dirs = [
        os.path.join(project_root, "output"),
        os.path.join(project_root, "examples"),
    ]
    results = []
    seen = set()
    for base in search_dirs:
        base_path = Path(base)
        if not base_path.is_dir():
            continue
        for ext in ("*.js", "*.mjs", "*.py"):
            for p in sorted(base_path.rglob(ext)):
                # Skip preview artifacts and node_modules
                parts = set(p.parts)
                if parts & {".preview", "node_modules", "__pycache__"}:
                    continue
                key = str(p.resolve())
                if key in seen:
                    continue
                seen.add(key)
                suffix = p.suffix.lower()
                engine = "node" if suffix in (".js", ".mjs") else "python"
                try:
                    mtime = p.stat().st_mtime
                except OSError:
                    mtime = 0
                results.append({
                    "path": str(p),
                    "name": p.name,
                    "engine": engine,
                    "mtime": mtime,
                })
    return results


@router.websocket("/ws/preview/{token}")
async def preview_websocket(token: str, websocket: WebSocket):
    """WebSocket endpoint for live preview events.

    On connect: replays the last known state so a fresh tab gets immediate
    context. Subsequently forwards ``render_started``, ``render_complete``,
    ``render_failed``, and ``watch_changed`` events.

    Client may send ``{"action": "force_render"}`` or ``{"action": "cancel"}``.
    """
    registry = getattr(websocket.app.state, "preview_registry", None)
    if registry is None:
        await websocket.close(code=1011)
        return

    session = registry.get(token)
    if session is None:
        await websocket.close(code=4004)
        return

    await websocket.accept()
    session.add_client(websocket)

    # Replay last known state immediately
    last = session.last_state
    if last.get("event") and last["event"] != "idle":
        try:
            await websocket.send_json(last)
        except Exception:
            pass

    try:
        while True:
            try:
                data = await asyncio.wait_for(websocket.receive_json(), timeout=30)
                action = (data or {}).get("action", "")
                if action == "force_render":
                    await session.force_render()
                # "cancel" is a no-op for now; the watcher handles coalescing
            except asyncio.TimeoutError:
                # Send a keepalive ping
                try:
                    await websocket.send_json({"event": "ping"})
                except Exception:
                    break
            except WebSocketDisconnect:
                break
            except Exception:
                break
    finally:
        session.remove_client(websocket)


# ---------------------------------------------------------------------------
# Chat-with-a-Deck endpoints
# ---------------------------------------------------------------------------

# Active cancel tokens keyed by conversation_id.
# Safe for single-worker uvicorn only — would leak/misfire across multiple workers.
_chat_cancel_tokens: dict = {}


@router.get("/api/chat/conversations")
async def list_chat_conversations(request: Request, deck_id: int):
    """List conversations for a deck."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Chat is disabled in view-only mode"}, status_code=403)
    from aippt.chat import ChatService
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    try:
        svc = ChatService(conn, llm=None)
        return svc.list_conversations(deck_id)
    finally:
        conn.close()


@router.post("/api/chat/conversations")
async def create_chat_conversation(request: Request):
    """Create a new conversation for a deck."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Chat is disabled in view-only mode"}, status_code=403)
    body = await request.json()
    deck_id = body.get("deck_id")
    title = body.get("title", "New conversation")
    if not deck_id:
        return JSONResponse({"error": "deck_id is required"}, status_code=400)
    from aippt.chat import ChatService
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    try:
        svc = ChatService(conn, llm=None)
        # Look up the deck's source script path so the LLM gets the exact path
        deck_row = conn.execute(
            "SELECT source_script_path FROM decks WHERE id=?", (deck_id,)
        ).fetchone()
        script_path = deck_row["source_script_path"] if deck_row else None
        conv_id = svc.create_conversation(deck_id, title, script_path=script_path)
        return {"id": conv_id, "deck_id": deck_id, "title": title}
    finally:
        conn.close()


@router.get("/api/chat/conversations/{conversation_id}")
async def get_chat_conversation(conversation_id: int, request: Request):
    """Get conversation metadata and messages."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Chat is disabled in view-only mode"}, status_code=403)
    from aippt.chat import ChatService
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    try:
        svc = ChatService(conn, llm=None)
        conv = svc.get_conversation(conversation_id)
        if not conv:
            return JSONResponse({"error": "Conversation not found"}, status_code=404)
        messages = svc.get_messages(conversation_id)
        return {"conversation": conv, "messages": messages}
    finally:
        conn.close()


@router.patch("/api/chat/conversations/{conversation_id}")
async def rename_chat_conversation(conversation_id: int, request: Request):
    """Rename a conversation."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Chat is disabled in view-only mode"}, status_code=403)
    body = await request.json()
    title = body.get("title", "")
    if not title:
        return JSONResponse({"error": "title is required"}, status_code=400)
    from aippt.chat import ChatService
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    try:
        svc = ChatService(conn, llm=None)
        ok = svc.rename_conversation(conversation_id, title)
        if not ok:
            return JSONResponse({"error": "Conversation not found"}, status_code=404)
        return {"id": conversation_id, "title": title}
    finally:
        conn.close()


@router.delete("/api/chat/conversations/{conversation_id}")
async def delete_chat_conversation(conversation_id: int, request: Request):
    """Delete a conversation and all its messages."""
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Chat is disabled in view-only mode"}, status_code=403)
    from aippt.chat import ChatService
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    try:
        svc = ChatService(conn, llm=None)
        ok = svc.delete_conversation(conversation_id)
        if not ok:
            return JSONResponse({"error": "Conversation not found"}, status_code=404)
        return {"deleted": conversation_id}
    finally:
        conn.close()


@router.post("/api/chat/conversations/{conversation_id}/messages")
async def chat_send_message(conversation_id: int, request: Request):
    """Send a user message and stream the assistant reply via SSE.

    Request body (JSON):
        message   (str, required)  — the user's text
        slide_id  (int, optional)  — scope to a specific slide
        model     (str, optional)  — override the LLM model
        ntid      (str, optional)  — user NTID for gateway

    Response: ``text/event-stream`` with ``data:`` lines.

    Event types:
        data: {"type":"chunk","text":"..."}
        data: {"type":"patch_proposed","message_id":N,"patches":[...]}
        data: {"type":"cancelled"}
        data: {"type":"done"}
        data: {"type":"error","message":"..."}
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Chat is disabled in view-only mode"}, status_code=403)

    body = await request.json()
    user_message = body.get("message", "").strip()
    slide_id = body.get("slide_id")
    mode = body.get("mode", "ask")
    if mode not in ("ask", "edit", "review"):
        mode = "ask"
    # In view-only mode, force ask (read-only) regardless of what the client sent
    if getattr(request.app.state, "view_only", False):
        mode = "ask"
    model_override = body.get("model")
    user_ntid = body.get("ntid") or request.headers.get("X-User-NTID")

    if not user_message:
        return JSONResponse({"error": "message is required"}, status_code=400)

    db_path = request.app.state.db_path

    # Verify conversation exists
    from aippt.chat import ChatService, CancelToken
    conn_check = get_db(db_path)
    svc_check = ChatService(conn_check, llm=None)
    if not svc_check.get_conversation(conversation_id):
        return JSONResponse({"error": "Conversation not found"}, status_code=404)
    conn_check.close()

    try:
        llm = _make_llm_client(request, "improve", model_override, user_ntid)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=400)

    cancel_token = CancelToken()
    _chat_cancel_tokens[conversation_id] = cancel_token

    def _sse_event(payload: dict) -> str:
        return f"data: {json.dumps(payload)}\n\n"

    def generate():
        conn = get_db(db_path)
        svc = ChatService(conn, llm, db_cwd=str(Path(db_path).parent))
        try:
            for chunk in svc.stream_reply(
                conversation_id,
                user_message,
                slide_id=slide_id,
                mode=mode,
                cancel_token=cancel_token,
            ):
                if chunk == "[CANCELLED]":
                    yield _sse_event({"type": "cancelled"})
                    return
                elif chunk.startswith("[PATCH_PROPOSED:"):
                    # Format: [PATCH_PROPOSED:msg_id:json_array]
                    rest = chunk[len("[PATCH_PROPOSED:"):-1]
                    colon = rest.index(":")
                    msg_id = int(rest[:colon])
                    patch_json = rest[colon + 1:]
                    yield _sse_event({"type": "patch_proposed", "message_id": msg_id, "patches": json.loads(patch_json)})
                elif chunk.startswith("[REVIEW_COMPLETE:"):
                    # Format: [REVIEW_COMPLETE:msg_id:json_object]
                    rest = chunk[len("[REVIEW_COMPLETE:"):-1]
                    colon = rest.index(":")
                    msg_id = int(rest[:colon])
                    findings_json = rest[colon + 1:]
                    yield _sse_event({"type": "review_complete", "message_id": msg_id, "findings": json.loads(findings_json)})
                else:
                    yield _sse_event({"type": "chunk", "text": chunk})
            yield _sse_event({"type": "done"})
        except Exception as exc:
            logger.error("Chat stream error: %s", exc)
            yield _sse_event({"type": "error", "message": str(exc)})
        finally:
            conn.close()
            _chat_cancel_tokens.pop(conversation_id, None)

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/chat/conversations/{conversation_id}/cancel")
async def cancel_chat_stream(conversation_id: int, request: Request):
    """Cancel an in-progress streaming reply."""
    token = _chat_cancel_tokens.get(conversation_id)
    if token:
        token.cancel()
        return {"cancelled": True}
    return {"cancelled": False}


@router.post("/api/chat/conversations/{conversation_id}/messages/{message_id}/apply")
async def apply_chat_patch(conversation_id: int, message_id: int, request: Request):
    """Apply the patch proposed in an assistant message.

    The user must explicitly call this after reviewing the proposed diff.
    Returns 403 in view-only mode.
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Patch apply is disabled in view-only mode"}, status_code=403)

    from aippt.chat import ChatService
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    svc = ChatService(conn, llm=None, db_cwd=str(Path(db_path).parent))
    try:
        ok, reason = svc.apply_message_patch(message_id)
    finally:
        conn.close()

    if ok:
        return {"applied": True}
    return JSONResponse({"error": reason}, status_code=400)


@router.post("/api/chat/conversations/{conversation_id}/messages/{message_id}/revert")
async def revert_chat_patch(conversation_id: int, message_id: int, request: Request):
    """Revert a previously applied patch on an assistant message.

    Returns 403 in view-only mode.
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Patch revert is disabled in view-only mode"}, status_code=403)

    from aippt.chat import ChatService
    db_path = request.app.state.db_path
    conn = get_db(db_path)
    svc = ChatService(conn, llm=None, db_cwd=str(Path(db_path).parent))
    try:
        ok, reason = svc.revert_message_patch(message_id)
    finally:
        conn.close()

    if ok:
        return {"reverted": True}
    return JSONResponse({"error": reason}, status_code=400)


@router.post("/api/chat/slides/{slide_id}/revert")
async def revert_slide_field(slide_id: int, request: Request):
    """Revert the most recent LLM-applied edit on a slide field.

    Request body (JSON):
        field  (str, required)  — "title", "content_text", or "notes"
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Revert is disabled in view-only mode"}, status_code=403)
    body = await request.json()
    field = body.get("field", "").strip()
    if not field:
        return JSONResponse({"error": "field is required"}, status_code=400)

    from aippt.patch import revert_last, PATCHABLE_FIELDS
    if field not in PATCHABLE_FIELDS:
        return JSONResponse(
            {"error": f"field must be one of {sorted(PATCHABLE_FIELDS)}"}, status_code=400
        )

    db_path = request.app.state.db_path
    conn = get_db(db_path)
    try:
        ok, msg = revert_last(slide_id, field, conn, source="chat-revert",
                              cwd=str(Path(db_path).parent))
        if not ok:
            return JSONResponse({"error": msg}, status_code=404)
        return {"reverted": True, "message": msg}
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# Deck Review & Auto-Refine endpoints
# ---------------------------------------------------------------------------

# Per-deck threading.Event used to unblock the refine loop when the browser
# posts fresh thumbnails.  Keyed by deck_id.
_refine_events: dict[int, "threading.Event"] = {}


@router.post("/api/decks/{deck_id}/review")
async def review_deck(deck_id: int, request: Request):
    """Stream a full-deck visual + notes-flow review pass (no auto-apply).

    Creates a new chat conversation for this review run, sends a review-mode
    message for each slide, runs the notes-flow pass, then emits the merged
    R4 findings report.

    Response: ``text/event-stream``

    Event types::

        data: {"type":"progress","step":"review","slide":N,"total":M}
        data: {"type":"findings","data":{...ReviewResult...}}
        data: {"type":"done"}
        data: {"type":"error","message":"..."}
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Review is disabled in view-only mode"}, status_code=403)

    body = await request.json()
    model_override = body.get("model")
    user_ntid = body.get("ntid") or request.headers.get("X-User-NTID")

    db_path = request.app.state.db_path

    # Verify deck exists and has slides
    conn_check = get_db(db_path)
    deck_row = conn_check.execute("SELECT id, name FROM decks WHERE id = ?", (deck_id,)).fetchone()
    conn_check.close()
    if not deck_row:
        return JSONResponse({"error": "Deck not found"}, status_code=404)

    try:
        llm = _make_llm_client(request, "improve", model_override, user_ntid)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=400)

    def _sse(payload: dict) -> str:
        return f"data: {json.dumps(payload)}\n\n"

    def generate():
        import threading as _threading
        from aippt.chat import ChatService
        from aippt.deck_review import _review_all_slides, ReviewResult

        conn = get_db(db_path)
        try:
            # Create a dedicated review conversation
            svc = ChatService(conn, llm, db_cwd=str(Path(db_path).parent))
            conv_id = svc.create_conversation(deck_id, title="Review run")

            slides = conn.execute(
                "SELECT id, position FROM slides WHERE deck_id = ? ORDER BY position",
                (deck_id,),
            ).fetchall()
            total = len(slides)

            all_findings = []
            for i, slide in enumerate(slides, 1):
                yield _sse({"type": "progress", "step": "review", "slide": i, "total": total,
                            "detail": f"Reviewing slide {i} of {total}…"})
                try:
                    chunks = []
                    for chunk in svc.stream_reply(
                        conv_id,
                        f"Review slide {slide['position']}",
                        slide_id=slide["id"],
                        mode="review",
                    ):
                        if not chunk.startswith("[REVIEW_COMPLETE:") and not chunk.startswith("[CANCELLED]"):
                            chunks.append(chunk)
                    from aippt.deck_review import parse_review_output
                    sr = parse_review_output("".join(chunks))
                    for f in sr.findings:
                        if f.slide_num is None:
                            f.slide_num = slide["position"]
                    all_findings.extend(sr.findings)
                except Exception as slide_exc:
                    logger.warning("Review slide %s error: %s — skipping", slide["position"], slide_exc)
                    from aippt.deck_review import Finding, Severity
                    all_findings.append(Finding(
                        slide_num=slide["position"],
                        severity=Severity.LOW,
                        category="content",
                        description="Review could not complete for this slide (network error)",
                        actionable=False,
                    ))

            # Notes-flow pass
            yield _sse({"type": "progress", "step": "notes", "slide": total, "total": total})
            notes_result = svc.run_notes_flow_pass(conv_id)
            all_findings.extend(notes_result.findings)

            from aippt.deck_review import ReviewResult
            merged = ReviewResult(findings=all_findings)
            yield _sse({"type": "findings", "data": merged.to_dict()})
            yield _sse({"type": "done"})
        except Exception as exc:
            logger.error("Review stream error for deck %s: %s", deck_id, exc)
            yield _sse({"type": "error", "message": str(exc)})
        finally:
            conn.close()

    return StreamingResponse(generate(), media_type="text/event-stream")


@router.post("/api/decks/{deck_id}/refine")
async def refine_deck(deck_id: int, request: Request):
    """Stream the bounded auto-refine loop for a deck.

    After each re-render round the server emits ``recapture_needed``; the
    browser should re-render the updated deck via PptxViewJS, POST fresh
    thumbnails to ``/api/decks/{id}/thumbnails``, then POST to
    ``/api/decks/{id}/thumbnails-ready`` to unblock the next round.

    Response: ``text/event-stream``

    Event types::

        data: {"type":"progress","step":"review"|"patch"|"render"|"recapture","detail":"..."}
        data: {"type":"round_complete","round":N,"patches_applied":M,"findings":[...]}
        data: {"type":"recapture_needed","deck_id":N}
        data: {"type":"thumbnails_ready"}
        data: {"type":"complete","summary":{...}}
        data: {"type":"error","message":"..."}
    """
    if getattr(request.app.state, "view_only", False):
        return JSONResponse({"error": "Refine is disabled in view-only mode"}, status_code=403)

    body = await request.json()
    max_rounds = int(body.get("max_rounds", 2))
    model_override = body.get("model")
    user_ntid = body.get("ntid") or request.headers.get("X-User-NTID")

    db_path = request.app.state.db_path
    project_root = str(Path(request.app.state.db_path).parent)

    conn_check = get_db(db_path)
    deck_row = conn_check.execute("SELECT id FROM decks WHERE id = ?", (deck_id,)).fetchone()
    conn_check.close()
    if not deck_row:
        return JSONResponse({"error": "Deck not found"}, status_code=404)

    try:
        llm = _make_llm_client(request, "improve", model_override, user_ntid)
    except ValueError as exc:
        return JSONResponse({"error": str(exc)}, status_code=400)

    import threading as _threading
    evt = _threading.Event()
    _refine_events[deck_id] = evt

    import queue as _queue
    event_q: _queue.Queue = _queue.Queue()

    def _progress(step: str, detail: str) -> None:
        event_q.put(("progress", {"step": step, "detail": detail}))
        if step == "recapture":
            event_q.put(("recapture_needed", {"deck_id": deck_id}))

    def _worker():
        from aippt.deck_review import run_auto_refine
        try:
            result = run_auto_refine(
                deck_id=deck_id,
                conversation_id=_get_or_create_refine_conversation(deck_id, llm, db_path),
                db_path=db_path,
                project_root=project_root,
                llm_client=llm,
                max_rounds=max_rounds,
                progress_callback=_progress,
                thumbnails_ready_event=evt,
            )
            event_q.put(("complete", result.to_dict()))
        except Exception as exc:
            logger.error("Refine worker error for deck %s: %s", deck_id, exc)
            event_q.put(("error", {"message": str(exc)}))
        finally:
            _refine_events.pop(deck_id, None)

    def _sse(payload: dict) -> str:
        return f"data: {json.dumps(payload)}\n\n"

    async def generate():
        import asyncio
        loop = asyncio.get_running_loop()
        fut = loop.run_in_executor(None, _worker)
        try:
            while True:
                try:
                    kind, payload = event_q.get(timeout=0.1)
                except Exception:
                    if fut.done():
                        # Drain any remaining events
                        while not event_q.empty():
                            kind, payload = event_q.get_nowait()
                            yield _sse({"type": kind, **payload})
                        break
                    await asyncio.sleep(0.05)
                    continue

                yield _sse({"type": kind, **payload})
                if kind in ("complete", "error"):
                    break

                if kind == "recapture_needed":
                    # Signal browser then wait for thumbnails_ready event
                    # The event is set by POST /api/decks/{id}/thumbnails-ready
                    yield _sse({"type": "thumbnails_ready"})  # placeholder — real signal from browser
        finally:
            await fut
    return StreamingResponse(generate(), media_type="text/event-stream")


def _get_or_create_refine_conversation(deck_id: int, llm, db_path: str) -> int:
    """Return an existing 'Auto-refine' conversation or create a new one."""
    from aippt.chat import ChatService
    conn = get_db(db_path)
    try:
        svc = ChatService(conn, llm, db_cwd=str(Path(db_path).parent))
        convs = svc.list_conversations(deck_id)
        for c in convs:
            if c.get("title", "").startswith("Auto-refine"):
                return c["id"]
        # Look up script path
        row = conn.execute("SELECT source_script_path FROM decks WHERE id = ?", (deck_id,)).fetchone()
        script_path = row["source_script_path"] if row else None
        return svc.create_conversation(deck_id, title="Auto-refine", script_path=script_path)
    finally:
        conn.close()


@router.post("/api/decks/{deck_id}/thumbnails")
async def post_deck_thumbnails(deck_id: int, request: Request):
    """Accept fresh slide images from the browser (no session token needed).

    Request body (JSON)::

        {"images": [{"position": 1, "data": "<base64-png>"}, ...]}

    Used by the auto-refine loop's recapture step when there is no live
    preview session open.
    """
    db_path = request.app.state.db_path
    uploads_dir = request.app.state.uploads_dir

    conn_check = get_db(db_path)
    deck_row = conn_check.execute("SELECT id FROM decks WHERE id = ?", (deck_id,)).fetchone()
    conn_check.close()
    if not deck_row:
        return JSONResponse({"error": "Deck not found"}, status_code=404)

    body = await request.json()
    images = body.get("images", [])
    if not images:
        return JSONResponse({"error": "images is required"}, status_code=400)

    from aippt.thumbnails import store_client_images
    import os

    images_dir = os.path.join(uploads_dir, "images", str(deck_id))
    os.makedirs(images_dir, exist_ok=True)

    try:
        count = store_client_images(
            deck_id=deck_id,
            images=images,
            out_dir=images_dir,
            db_path=db_path,
            base_dir=str(Path(db_path).parent),
        )
        return {"stored": count}
    except Exception as exc:
        logger.error("Failed to store deck thumbnails for deck %s: %s", deck_id, exc)
        return JSONResponse({"error": str(exc)}, status_code=500)


@router.post("/api/decks/{deck_id}/thumbnails-ready")
async def thumbnails_ready(deck_id: int, request: Request):
    """Signal the refine loop that fresh thumbnails have been posted.

    Called by the browser after it POSTs new slide images following a
    ``recapture_needed`` event.  Sets the per-deck threading.Event so
    the refine loop can proceed to the next round.
    """
    evt = _refine_events.get(deck_id)
    if evt is None:
        return JSONResponse({"error": "No active refine loop for this deck"}, status_code=404)
    evt.set()
    return {"acknowledged": True}
