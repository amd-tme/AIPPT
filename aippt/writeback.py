"""PPTX notes write-back: flush DB-edited notes into a PowerPoint file."""
import logging
import os
import shutil
from dataclasses import dataclass, field
from datetime import datetime
from typing import Optional

from pptx import Presentation

from aippt.catalog import get_db, get_deck_slides

logger = logging.getLogger(__name__)


@dataclass
class WritebackResult:
    deck_id: int
    slides_written: int = 0
    slides_skipped: int = 0
    slides_total: int = 0
    backup_path: Optional[str] = None
    warnings: list = field(default_factory=list)


def create_backup(deck_path: str) -> str:
    """Copy a deck to a timestamped .pptx.bak file alongside the original.

    Args:
        deck_path: Absolute or relative path to the PowerPoint file.

    Returns:
        The path to the backup file.

    Raises:
        FileNotFoundError: If deck_path does not exist.
    """
    if not os.path.exists(deck_path):
        raise FileNotFoundError(f"Deck not found: {deck_path}")

    timestamp = datetime.now().strftime("%Y-%m-%dT%H-%M-%S")
    backup_path = f"{deck_path}.{timestamp}.pptx.bak"
    shutil.copy2(deck_path, backup_path)
    logger.info("Created backup: %s", backup_path)
    return backup_path


def write_notes_to_pptx(
    deck_path: str,
    db_path: str = "slides.db",
    deck_id: Optional[int] = None,
    output_path: Optional[str] = None,
) -> WritebackResult:
    """Write notes stored in the DB back into the PPTX file.

    For each slide in the deck, if the DB record has non-empty notes the
    function overwrites the notes text frame in the presentation.  Slides
    whose DB notes field is empty or NULL are skipped.

    Args:
        deck_path: Path to the PowerPoint file.
        db_path: Path to the SQLite database.
        deck_id: Database deck ID.  When omitted the deck is looked up by its
            absolute file path.
        output_path: Where to write the modified presentation.  When omitted
            the original file is overwritten in-place.

    Returns:
        WritebackResult with counts and metadata.

    Raises:
        FileNotFoundError: If deck_path does not exist.
        ValueError: If the deck is not found in the database, or if the slide
            count in the PPTX differs from the slide count recorded in the DB.
    """
    if not os.path.exists(deck_path):
        raise FileNotFoundError(f"Deck not found: {deck_path}")

    abs_deck_path = os.path.abspath(deck_path)
    rel_deck_path = os.path.relpath(abs_deck_path)

    # Resolve deck_id from DB when not supplied
    if deck_id is None:
        conn = get_db(db_path)
        try:
            # Try relative path first (new format), then absolute (legacy)
            row = conn.execute(
                "SELECT id FROM decks WHERE file_path = ? OR file_path = ?",
                (rel_deck_path, abs_deck_path),
            ).fetchone()
        finally:
            conn.close()
        if row is None:
            raise ValueError(
                f"Deck not found in database: {abs_deck_path}"
            )
        deck_id = row["id"]
    else:
        # Verify the supplied deck_id actually exists
        conn = get_db(db_path)
        try:
            row = conn.execute(
                "SELECT id FROM decks WHERE id = ?", (deck_id,)
            ).fetchone()
        finally:
            conn.close()
        if row is None:
            raise ValueError(
                f"Deck ID {deck_id} not found in database"
            )

    db_slides = get_deck_slides(deck_id, db_path=db_path)
    prs = Presentation(deck_path)

    if len(prs.slides) != len(db_slides):
        raise ValueError(
            f"Slide count mismatch: PPTX has {len(prs.slides)} slides "
            f"but DB has {len(db_slides)} slides for deck_id={deck_id}"
        )

    result = WritebackResult(deck_id=deck_id, slides_total=len(db_slides))

    for db_slide in db_slides:
        position = db_slide["position"]  # 1-based
        notes = db_slide.get("notes") or ""

        if not notes.strip():
            result.slides_skipped += 1
            logger.debug("Slide %d: no notes, skipping", position)
            continue

        pptx_slide = prs.slides[position - 1]
        # Accessing .notes_slide auto-creates the notes slide in python-pptx
        pptx_slide.notes_slide.notes_text_frame.text = notes
        result.slides_written += 1
        logger.debug("Slide %d: wrote notes (%d chars)", position, len(notes))

    save_path = output_path if output_path else deck_path
    prs.save(save_path)
    logger.info(
        "Saved presentation to %s (%d written, %d skipped)",
        save_path,
        result.slides_written,
        result.slides_skipped,
    )

    return result
