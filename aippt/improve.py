"""Slide improvement pipeline.

Analyzes existing slides with multimodal feedback, rewrites content
via LLM, and applies improvements back to the PPTX.
"""

import difflib
import logging
import os
import re
import subprocess
from typing import Dict, Optional, Tuple

from pptx import Presentation

from aippt.analyze import analyze_slide
from aippt.layouts import _apply_bullets_to_text_frame

logger = logging.getLogger(__name__)


REWRITE_SYSTEM_PROMPT = (
    "You are a presentation content expert rewriting slide content based on "
    "expert feedback. You improve specificity, organization, and detail while "
    "keeping the same topic.\n\n"
    "Constraints:\n"
    "- Return ONLY the improved bullet content, nothing else\n"
    "- One bullet per line starting with '- '\n"
    "- Use '  - ' (2-space indent) for sub-bullets\n"
    "- Use numbered items (1., 2., 3.) if the content is sequential\n"
    "- Use 'Keyword: description' format for bold lead-ins where appropriate\n"
    "- Keep total content to 4-8 bullets (with sub-bullets as needed)\n"
    "- Focus on: specificity, concrete examples, technical accuracy, "
    "logical organization\n"
    "- Do NOT suggest colors, fonts, icons, shapes, or visual design changes\n"
    "- Do NOT include any preamble, explanation, or commentary"
)

FOCUS_GUIDANCE = {
    "accuracy": "Focus on technical accuracy. Replace vague claims with specific, verifiable statements.",
    "detail": "Add concrete examples, metrics, and specifics. Expand abbreviated points.",
    "brevity": "Prioritize conciseness. Remove redundant qualifiers and combine overlapping points.",
    "structure": "Improve logical organization. Group related points and establish clear hierarchy.",
    "general": "",  # No additional guidance
}

AUDIENCE_REWRITE_PROMPTS = {
    "engineers": (
        "\n\nTarget audience: ENGINEERS. Preserve technical depth and specifics. "
        "Include technology names, version numbers, and implementation details. "
        "Use 5-8 detailed bullets with sub-bullets for depth."
    ),
    "executives": (
        "\n\nTarget audience: EXECUTIVES. Frame every point as business impact. "
        "Replace technical jargon with outcomes and metrics. "
        "Use 3-5 concise, punchy bullets — one insight per bullet."
    ),
    "product": (
        "\n\nTarget audience: PRODUCT MANAGERS. Focus on features, user impact, "
        "and adoption metrics. Use product terminology. "
        "Use 4-6 feature-oriented bullets."
    ),
    "mixed": (
        "\n\nTarget audience: MIXED. Write clear explanations accessible to both "
        "technical and non-technical readers. Define technical terms on first use. "
        "Use 4-7 bullets at accessible depth."
    ),
}


VALIDATION_SYSTEM_PROMPT = (
    "You are a presentation quality reviewer. Compare improved slide content "
    "against the original expert feedback and determine whether the feedback "
    "was adequately addressed.\n\n"
    "Return your evaluation in EXACTLY this format:\n"
    "VERDICT: PASS | PARTIAL | FAIL\n"
    "ADDRESSED: [comma-separated list of feedback points that were addressed]\n"
    "UNADDRESSED: [comma-separated list of feedback points NOT addressed]\n"
    "SUGGESTION: [brief guidance for retry, if any]\n\n"
    "Use PASS when all key feedback points are addressed.\n"
    "Use PARTIAL when some but not all key points are addressed.\n"
    "Use FAIL when the rewrite missed the main feedback points.\n\n"
    "Also check for information preservation: if original slide content is "
    "provided, verify that the rewrite did not drop important facts, data, "
    "or claims from the original. Mark as PARTIAL or FAIL if significant "
    "information was lost."
)

# Signal words for adaptive focus selection
_FOCUS_SIGNALS = {
    "detail": ["vague", "unclear", "unspecific", "generic", "lacks specificity",
               "not specific", "need more detail", "missing examples"],
    "brevity": ["verbose", "redundant", "wordy", "repetitive", "too long",
                "too much text", "overly detailed"],
    "accuracy": ["inaccurate", "incorrect", "misleading", "unsupported",
                 "wrong", "factual error", "unverified"],
    "structure": ["disorganized", "jumbled", "no hierarchy", "poor flow",
                  "illogical order", "unstructured", "bad organization"],
}


def select_focus(feedback: str) -> str:
    """Auto-select improvement focus area from analysis feedback text.

    Counts signal-word hits per focus area and returns the area with the
    most evidence.  Ties are broken by ``_FOCUS_SIGNALS`` insertion order.

    Args:
        feedback: Analysis feedback text

    Returns:
        Focus area string: 'detail', 'brevity', 'accuracy', 'structure', or 'general'
    """
    if not feedback:
        return "general"
    feedback_lower = feedback.lower()
    counts: Dict[str, int] = {}
    for focus, signals in _FOCUS_SIGNALS.items():
        count = sum(1 for signal in signals if signal in feedback_lower)
        if count > 0:
            counts[focus] = count
    if not counts:
        return "general"
    return max(counts, key=counts.get)


def parse_validation_response(response: str) -> Dict[str, str]:
    """Parse structured validation response from LLM.

    Extracts VERDICT, ADDRESSED, UNADDRESSED, and SUGGESTION fields.
    Uses regex with multiline matching so multi-line field values
    (e.g. a bulleted UNADDRESSED list) are fully captured.

    Args:
        response: Raw LLM validation response

    Returns:
        Dict with 'verdict', 'addressed', 'unaddressed', 'suggestion' keys
    """
    result = {
        'verdict': 'FAIL',
        'addressed': '',
        'unaddressed': '',
        'suggestion': '',
    }
    fields = ['VERDICT', 'ADDRESSED', 'UNADDRESSED', 'SUGGESTION']
    for i, field in enumerate(fields):
        if i < len(fields) - 1:
            pattern = rf'{field}:\s*(.*?)(?={fields[i+1]}:|\Z)'
        else:
            pattern = rf'{field}:\s*(.*?)(?:\Z)'
        m = re.search(pattern, response, re.DOTALL)
        if m:
            value = m.group(1).strip()
            if field == 'VERDICT':
                value = value.upper()
                if value not in ('PASS', 'PARTIAL', 'FAIL'):
                    value = 'FAIL'
            result[field.lower()] = value
    return result


def has_converged(original: str, improved: str) -> bool:
    """Check if content has converged (rewrite produced negligible change).

    Returns True when:
    - Content is identical
    - Content similarity ratio > 0.95 (less than 5% change)

    Args:
        original: Previous content version
        improved: New content version

    Returns:
        True if content has converged
    """
    if original == improved:
        return True
    if not original and not improved:
        return True
    if not original or not improved:
        return False
    ratio = difflib.SequenceMatcher(None, original, improved).ratio()
    return ratio > 0.95


def build_rewrite_prompt(title: str, current_content: str, feedback: str,
                         is_two_column: bool = False,
                         keep_titles: bool = False) -> str:
    """Build the prompt for LLM content rewrite.

    Args:
        title: Slide title
        current_content: Current bullet content as newline-separated text
        feedback: Structured improvement feedback from analysis
        is_two_column: Whether the slide uses two-column layout
        keep_titles: If True, skip title rewrite instructions

    Returns:
        Prompt string for LLM
    """
    if is_two_column:
        format_instruction = (
            "Return content in EXACTLY this format:\n"
            "## Left Column\n- bullet\n## Right Column\n- bullet\n"
            "Keep both columns balanced (3-5 bullets each). "
            "No preamble, no explanation."
        )
    else:
        format_instruction = "Return ONLY improved bullet content — no preamble, no explanation."

    if keep_titles:
        title_instruction = ""
    else:
        title_instruction = (
            "\n\nAlso provide an improved title if the expert feedback suggests "
            "the current title is generic or could be more insight-driven.\n\n"
            "Format your response as:\n"
            "TITLE: [improved title, or KEEP if the current title is good]\n"
            "CONTENT:\n"
            "- bullet 1\n"
            "- bullet 2\n"
        )

    return f"""Rewrite this slide's content to address the expert feedback below.
{format_instruction}{title_instruction}

Slide title: {title}

Current content:
{current_content}

Expert feedback:
{feedback}
"""


GENERIC_TITLE_PATTERNS = re.compile(
    r'^(The )?(Problem|Solution|Architecture|Overview|Summary|Introduction|'
    r'Background|Conclusion|Results|Agenda|Outline|Next Steps|Questions|'
    r'Key (Findings|Takeaways|Points))$',
    re.IGNORECASE,
)


def is_generic_title(title: str) -> bool:
    """Check if a title matches common generic patterns.

    Generic titles are category labels like "The Problem" or "Architecture"
    rather than insight-driven titles like "Engineers Waste 3 Hours/Day".

    Args:
        title: Slide title text

    Returns:
        True if the title is generic/label-style
    """
    if not title or not title.strip():
        return False
    return bool(GENERIC_TITLE_PATTERNS.match(title.strip()))


def parse_rewrite_response(response: str) -> Tuple[Optional[str], str]:
    """Parse LLM rewrite response into (new_title, improved_content).

    Extracts the TITLE: line (if present) and the remaining content.
    Returns (None, content) if title should be kept (TITLE: KEEP or absent).

    Args:
        response: Raw LLM response text

    Returns:
        Tuple of (new_title_or_None, content_string)
    """
    lines = response.strip().split('\n')
    new_title = None
    content_start = 0

    # Check first line for TITLE:
    if lines and lines[0].startswith('TITLE:'):
        title_value = lines[0][len('TITLE:'):].strip()
        if title_value.upper() != 'KEEP':
            new_title = title_value
        content_start = 1
        # Skip CONTENT: marker if present
        if content_start < len(lines) and lines[content_start].strip().upper() == 'CONTENT:':
            content_start += 1

    content_lines = lines[content_start:]
    content = '\n'.join(content_lines).strip()

    return new_title, content


def parse_rewritten_content(response: str) -> str:
    """Parse LLM response to extract only bullet/numbered content lines.

    Strips any preamble text the LLM may add before the actual bullets.

    Args:
        response: Raw LLM response text

    Returns:
        Cleaned content string with only bullet/numbered lines
    """
    lines = response.strip().split('\n')
    content_lines = []
    for line in lines:
        stripped = line.strip()
        # Keep lines that start with bullet markers, numbered items, or indented sub-bullets
        if (stripped.startswith(('-', '•', '*')) or
                re.match(r'^\d+\.', stripped) or
                (line.startswith('  ') and stripped.startswith(('-', '•', '*')))):
            content_lines.append(line)
    return '\n'.join(content_lines) if content_lines else response.strip()


def parse_two_column_rewritten_content(response: str) -> Tuple[str, str]:
    """Parse LLM two-column response into (left, right) tuple.

    Expects '## Left Column' and '## Right Column' section markers.
    Falls back to all-left if markers absent.
    """
    if '## Right Column' in response:
        parts = response.split('## Right Column', 1)
        left_raw = parts[0].replace('## Left Column', '').strip()
        right_raw = parts[1].strip()
    else:
        left_raw = response.replace('## Left Column', '').strip()
        right_raw = ''
    left = parse_rewritten_content(left_raw)
    right = parse_rewritten_content(right_raw) if right_raw else ''
    return left, right


def extract_slide_content(slide) -> Tuple[str, str]:
    """Extract title and body text from a PPTX slide.

    Detects two-column layouts (2+ content placeholders) and labels
    sections with ``## Left Column`` / ``## Right Column`` headers.

    Args:
        slide: A pptx.slide.Slide object

    Returns:
        Tuple of (title, body_text)
    """
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text

    # Collect content placeholders (idx > 0), sorted by idx
    content_phs = sorted(
        [(ph.placeholder_format.idx, ph)
         for ph in slide.placeholders
         if ph.placeholder_format.idx > 0],
        key=lambda x: x[0]
    )

    def _extract_lines(text_frame) -> list:
        lines = []
        for para in text_frame.paragraphs:
            text = para.text.strip()
            if text:
                prefix = "  - " if para.level >= 1 else "- "
                lines.append(f"{prefix}{text}")
        return lines

    if len(content_phs) >= 2:
        # Two-column: label with section headers
        left_lines = _extract_lines(content_phs[0][1].text_frame)
        right_lines = _extract_lines(content_phs[1][1].text_frame)
        if left_lines or right_lines:
            body_parts = []
            if left_lines:
                body_parts.append("## Left Column")
                body_parts.extend(left_lines)
            if right_lines:
                body_parts.append("## Right Column")
                body_parts.extend(right_lines)
            return title, '\n'.join(body_parts)
    elif len(content_phs) == 1:
        lines = _extract_lines(content_phs[0][1].text_frame)
        if lines:
            return title, '\n'.join(lines)

    # Fallback: scan all shapes with text frames (textboxes, etc.)
    body_lines = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.has_text_frame and shape.text_frame.text.strip() == title.strip():
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                prefix = "  - " if para.level >= 1 else "- "
                body_lines.append(f"{prefix}{text}")
    return title, '\n'.join(body_lines)


def _do_rewrite(client, title, current_content, feedback, is_two_column,
                keep_titles, audience, focus, unaddressed_hint=None):
    """Execute a single rewrite pass and return (new_title, improved_content, raw)."""
    combined_feedback = feedback
    if unaddressed_hint:
        combined_feedback = (
            f"{feedback}\n\n"
            f"IMPORTANT — Previous rewrite missed these issues:\n{unaddressed_hint}"
        )
    prompt = build_rewrite_prompt(title, current_content, combined_feedback,
                                  is_two_column=is_two_column,
                                  keep_titles=keep_titles)
    system_prompt = REWRITE_SYSTEM_PROMPT
    audience_suffix = AUDIENCE_REWRITE_PROMPTS.get(audience, AUDIENCE_REWRITE_PROMPTS['mixed'])
    system_prompt = system_prompt + audience_suffix
    focus_text = FOCUS_GUIDANCE.get(focus, "")
    if focus_text:
        system_prompt = system_prompt + "\n\n" + focus_text
    raw_response = client.generate_text(
        prompt=prompt,
        system_prompt=system_prompt,
        max_tokens=1000,
        temperature=0.3,
    )
    new_title = None
    if not keep_titles:
        new_title, improved = parse_rewrite_response(raw_response)
    else:
        improved = raw_response
    if not is_two_column:
        improved = parse_rewritten_content(improved)
    return new_title, improved, raw_response


def _verdict_rank(verdict: str) -> int:
    """Return a numeric rank for validation verdicts (higher = better)."""
    return {'PASS': 2, 'PARTIAL': 1, 'FAIL': 0}.get(verdict, 0)


def _do_validate(client, feedback, improved_content, original_content=None):
    """Execute a validation pass and return parsed result dict."""
    validate_prompt = f"Original expert feedback:\n{feedback}\n\n"
    if original_content:
        validate_prompt += f"Original slide content:\n{original_content}\n\n"
    validate_prompt += (
        f"Rewritten content:\n{improved_content}\n\n"
        "Did the rewrite address the feedback without losing important "
        "information? Evaluate each point."
    )
    raw = client.generate_text(
        prompt=validate_prompt,
        system_prompt=VALIDATION_SYSTEM_PROMPT,
        max_tokens=500,
        temperature=0.2,
    )
    return parse_validation_response(raw)


def improve_slide(slide, image_path: Optional[str], client, dry_run: bool = False,
                  focus: str = "general", audience: str = "mixed",
                  keep_titles: bool = False, max_retries: int = 2,
                  no_validate: bool = False):
    """Run the improve pipeline on a single slide.

    Args:
        slide: A pptx.slide.Slide object
        image_path: Path to slide PNG image (or None for text-only)
        client: LLMClient instance
        dry_run: If True, return changes without applying
        focus: Focus area for improvements (accuracy, detail, brevity, structure, general)
        audience: Target audience type
        keep_titles: If True, skip title rewriting
        max_retries: Max validation retries per slide (default: 2)
        no_validate: If True, skip validation pass (old behavior)

    Returns:
        Dict with 'title', 'original', 'improved', 'feedback', 'applied',
        and optionally 'title_rewritten', 'original_title', 'new_title',
        'validation', 'focus_source', 'focus_used'
    """
    title, current_content = extract_slide_content(slide)

    if not current_content.strip():
        logger.info(f"Skipping slide '{title}' — no body content")
        return {'title': title, 'original': '', 'improved': '', 'feedback': '',
                'applied': False, 'status': 'no_content'}

    # Detect placeholder count for two-column awareness
    content_phs = sorted(
        [(ph.placeholder_format.idx, ph)
         for ph in slide.placeholders if ph.placeholder_format.idx > 0],
        key=lambda x: x[0]
    )
    is_two_column = len(content_phs) >= 2

    # Step 1: Analyze
    feedback = analyze_slide(
        client=client,
        image_path=image_path,
        mode='improvements',
        title=title,
        content_text=current_content,
    )

    # Adaptive focus selection
    if focus == "general":
        auto_focus = select_focus(feedback)
        effective_focus = auto_focus
        focus_source = "auto"
    else:
        effective_focus = focus
        focus_source = "user"

    # Step 2: Rewrite (with optional validation loop)
    new_title, improved, _ = _do_rewrite(
        client, title, current_content, feedback, is_two_column,
        keep_titles, audience, effective_focus,
    )

    validation_result = None
    if not no_validate:
        retries = 0
        converged = False
        passed = False
        prev_improved = None
        best_improved = improved
        best_val = None

        while retries <= max_retries:
            # Validate
            val = _do_validate(client, feedback, improved, original_content=current_content)
            # Track best attempt
            if best_val is None or _verdict_rank(val['verdict']) > _verdict_rank(best_val['verdict']):
                best_improved = improved
                best_val = val
            if val['verdict'] == 'PASS':
                passed = True
                break
            if val['verdict'] == 'PARTIAL' and retries > 0:
                passed = True  # Good enough after at least one attempt
                break
            if retries == max_retries:
                break

            # Retry: check for convergence first
            retries += 1
            prev_improved = improved
            unaddressed_hint = val.get('unaddressed', '') or val.get('suggestion', '')
            # Rewrite from original content (not previous attempt) to avoid error
            # accumulation. The unaddressed_hint steers toward missed points.
            new_title_retry, improved, _ = _do_rewrite(
                client, title, current_content, feedback, is_two_column,
                keep_titles, audience, effective_focus,
                unaddressed_hint=unaddressed_hint,
            )
            # Use the latest title if retry produced one
            if new_title_retry:
                new_title = new_title_retry

            if has_converged(prev_improved, improved):
                converged = True
                logger.info(f"  Content converged after {retries} retry(ies)")
                break

        # Use best attempt when max retries exhausted without passing
        if not passed and not converged:
            improved = best_improved

        validation_result = {
            'passed': passed,
            'retries': retries,
            'converged': converged,
            'unaddressed': val.get('unaddressed', ''),
        }

    result = {
        'title': title,
        'original': current_content,
        'improved': improved,
        'feedback': feedback,
        'applied': False,
        'status': 'pending',
        'focus_source': focus_source,
        'focus_used': effective_focus,
    }

    if validation_result is not None:
        result['validation'] = validation_result

    # Track title changes
    if new_title:
        result['title_rewritten'] = True
        result['original_title'] = title
        result['new_title'] = new_title
    else:
        result['title_rewritten'] = False
        if not keep_titles:
            result['original_title'] = title

    if dry_run:
        result['status'] = 'dry_run'
        return result

    # Step 3: Apply body content
    if is_two_column:
        left_content, right_content = parse_two_column_rewritten_content(improved)
        content_phs[0][1].text_frame.clear()
        content_phs[1][1].text_frame.clear()
        _apply_bullets_to_text_frame(content_phs[0][1].text_frame, left_content)
        _apply_bullets_to_text_frame(content_phs[1][1].text_frame, right_content)
        result['applied'] = True
        result['status'] = 'applied'
    elif len(content_phs) == 1:
        tf = content_phs[0][1].text_frame
        tf.clear()
        _apply_bullets_to_text_frame(tf, improved)
        result['applied'] = True
        result['status'] = 'applied'
    else:
        logger.warning(f"No body placeholder found for slide '{title}' — skipping apply")
        result['status'] = 'no_placeholder'

    # Step 3b: Apply title
    if new_title and slide.shapes.title:
        # Preserve formatting from original title
        title_shape = slide.shapes.title
        if title_shape.text_frame.paragraphs:
            original_para = title_shape.text_frame.paragraphs[0]
            # Save run-level formatting before overwrite
            original_runs = original_para.runs
            saved_font = None
            if original_runs:
                r = original_runs[0]
                saved_font = {
                    'bold': r.font.bold,
                    'size': r.font.size,
                    'color_rgb': r.font.color.rgb if r.font.color and r.font.color.type is not None else None,
                    'name': r.font.name,
                }
        # Set new title text
        title_shape.text = new_title
        # Re-apply formatting
        if saved_font and title_shape.text_frame.paragraphs:
            for run in title_shape.text_frame.paragraphs[0].runs:
                if saved_font['bold'] is not None:
                    run.font.bold = saved_font['bold']
                if saved_font['size'] is not None:
                    run.font.size = saved_font['size']
                if saved_font['color_rgb'] is not None:
                    run.font.color.rgb = saved_font['color_rgb']
                if saved_font['name'] is not None:
                    run.font.name = saved_font['name']
        logger.info(f"  Title rewritten: '{title}' → '{new_title}'")

    # Step 4: Update speaker notes with revision history + structured metadata
    try:
        notes_tf = slide.notes_slide.notes_text_frame
        existing_notes = notes_tf.text
        separator = "\n\n" if existing_notes.strip() else ""
        notes_tf.text = (existing_notes + separator +
                         f"--- Revision ---\nOriginal:\n{current_content}\n\nImproved:\n{improved}")

        # Append structured metadata
        from aippt.metadata import append_metadata, content_hash
        meta_kwargs = dict(
            model=getattr(client, 'model', 'unknown'),
            focus=effective_focus,
            focus_source=focus_source,
            audience=audience,
            changes_summary=f"Revised from {len(current_content.splitlines())} to {len(improved.splitlines())} lines",
            original_content_hash=content_hash(current_content),
        )
        if new_title:
            meta_kwargs['title_rewritten'] = True
            meta_kwargs['original_title'] = title
            meta_kwargs['new_title'] = new_title
        if validation_result is not None:
            meta_kwargs['validation'] = validation_result
        append_metadata(slide, "improve", **meta_kwargs)
    except Exception as e:
        logger.warning(f"Could not update notes for slide '{title}': {e}")

    return result


def _try_reexport_images(pptx_path: str, images_dir: str) -> bool:
    """Attempt to re-export slide images after improvement.

    Uses the PowerShell export script if available. Degrades gracefully
    on non-Windows systems or when PowerPoint is not installed.

    Args:
        pptx_path: Path to the saved PPTX file
        images_dir: Directory to export images into

    Returns:
        True if export succeeded, False otherwise
    """
    script_dir = os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "scripts"
    )
    ps_script = os.path.join(script_dir, "Export-SlidesToImages.ps1")

    if not os.path.exists(ps_script):
        logger.debug("Image export script not found, skipping re-export")
        return False

    # Find PowerShell executable
    for ps_name in ("pwsh", "powershell", "powershell.exe"):
        try:
            subprocess.run([ps_name, "-Version"], capture_output=True, timeout=5)
            ps_exe = ps_name
            break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    else:
        logger.debug("PowerShell not found, skipping image re-export")
        return False

    try:
        result = subprocess.run(
            [ps_exe, "-ExecutionPolicy", "Bypass", "-File", ps_script,
             "-PptxPath", os.path.abspath(pptx_path),
             "-OutDir", os.path.abspath(images_dir)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0:
            logger.info("Re-exported slide images after improvement pass")
            return True
        else:
            logger.debug(f"Image re-export failed: {result.stderr.strip()}")
            return False
    except (subprocess.TimeoutExpired, OSError) as e:
        logger.debug(f"Image re-export unavailable: {e}")
        return False


def improve_deck(pptx_path: str, output_path: Optional[str] = None,
                 images_dir: Optional[str] = None, slides_filter: Optional[list] = None,
                 passes: int = 1, dry_run: bool = False, client=None,
                 focus: str = "general", audience: str = "mixed",
                 keep_titles: bool = False, max_retries: int = 2,
                 no_validate: bool = False):
    """Run the improve pipeline on an entire deck.

    Args:
        pptx_path: Path to the PPTX file
        output_path: Output path (default: overwrite in-place)
        images_dir: Directory with slide images
        slides_filter: List of 1-based slide numbers to improve (None = all)
        passes: Number of improvement passes
        dry_run: If True, show changes without modifying
        client: LLMClient instance
        focus: Focus area for improvements (accuracy, detail, brevity, structure, general)
        keep_titles: If True, skip title rewriting
        max_retries: Max validation retries per slide (default: 2)
        no_validate: If True, skip validation pass

    Returns:
        List of result dicts, one per slide processed
    """
    prs = Presentation(pptx_path)
    save_path = output_path or pptx_path
    all_results = []

    for pass_num in range(1, passes + 1):
        if passes > 1:
            logger.info(f"=== Pass {pass_num}/{passes} ===")

        for i, slide in enumerate(prs.slides, 1):
            if slides_filter and i not in slides_filter:
                continue

            # Find image path (matches export-images convention: Slide1.PNG)
            image_path = None
            if images_dir:
                for ext in ('.png', '.PNG', '.jpg', '.jpeg'):
                    candidate = os.path.join(images_dir, f"Slide{i}{ext}")
                    if os.path.exists(candidate):
                        image_path = candidate
                        break

            logger.info(f"Improving slide {i}/{len(prs.slides)}")

            try:
                result = improve_slide(slide, image_path, client, dry_run=dry_run,
                                       focus=focus, audience=audience,
                                       keep_titles=keep_titles,
                                       max_retries=max_retries,
                                       no_validate=no_validate)
                result['slide_num'] = i
                result['pass'] = pass_num
                all_results.append(result)

                if result['applied']:
                    logger.info(f"  Applied improvements to: {result['title']}")
                elif dry_run and result['improved']:
                    logger.info(f"  [DRY RUN] Would improve: {result['title']}")
            except Exception as e:
                logger.error(f"Error improving slide {i}: {e}")
                all_results.append({
                    'title': f'Slide {i}', 'original': '', 'improved': '',
                    'feedback': '', 'applied': False, 'status': 'error',
                    'slide_num': i, 'pass': pass_num,
                })
                continue

        # Save after each pass
        if not dry_run:
            prs.save(save_path)
            logger.info(f"Saved improved deck to: {save_path}")

            # Re-export images so subsequent passes analyze fresh screenshots
            if images_dir and pass_num < passes:
                exported = _try_reexport_images(save_path, images_dir)
                if not exported:
                    logger.warning(
                        "Could not re-export images — subsequent passes may "
                        "analyze stale screenshots"
                    )

    # Deck-level summary
    applied = sum(1 for r in all_results if r.get('applied'))
    retried = sum(1 for r in all_results if r.get('validation', {}).get('retries', 0) > 0)
    maxed = sum(1 for r in all_results
                if r.get('validation', {}).get('retries', 0) >= max_retries
                and not r.get('validation', {}).get('passed'))
    logger.info(f"Improvement summary: {applied}/{len(all_results)} slides improved, "
                f"{retried} needed retries, {maxed} hit max retries")

    return all_results
