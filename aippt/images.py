"""Image generation and handling for slide diagrams.

This module provides functions for:
- Setting up image directories
- Generating SVG diagrams via LLM
- Generating images via DALL-E
- Converting SVG to PNG
- Applying diagram layouts to slides
"""

import base64
import datetime
import hashlib
import logging
import os
import re
import shutil
from typing import Dict, Optional

import requests
from pptx.util import Inches

from aippt.llm import LLMClient

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Directory management
# ---------------------------------------------------------------------------

def setup_image_directory(output_pptx: str) -> str:
    """Create and return path to image directory for this presentation.

    Creates a directory structure: images/<presentation-name>-<timestamp>/

    Args:
        output_pptx: Path to the output PowerPoint file

    Returns:
        Path to the created image directory
    """
    base_name = os.path.splitext(os.path.basename(output_pptx))[0]
    timestamp = datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
    image_dir = os.path.join('images', f"{base_name}-{timestamp}")
    os.makedirs(image_dir, exist_ok=True)
    return image_dir


# ---------------------------------------------------------------------------
# SVG generation
# ---------------------------------------------------------------------------

def generate_diagram_svg(description: str, client: LLMClient) -> str:
    """Generate SVG diagram using LLM.

    Args:
        description: Text description of the diagram to create
        client: Configured LLMClient instance

    Returns:
        SVG content as a string
    """
    prompt = f"""Please create a simple SVG diagram based on this description:
{description}

The SVG should:
1. Use basic shapes (rect, circle, path, line)
2. Include minimal styling
3. Be clean and professional
4. Use a viewBox of 0 0 800 600
5. Use neutral colors

Respond ONLY with the SVG code, starting with <svg> and ending with </svg>."""

    system_prompt = "You are an expert at creating simple, clean SVG diagrams."

    svg_content = client.generate_text(
        prompt=prompt,
        system_prompt=system_prompt,
        max_tokens=2000,
        temperature=0.2,
    )

    # Extract SVG code if there are surrounding explanations
    svg_match = re.search(r'<svg.*?</svg>', svg_content, re.DOTALL)
    if svg_match:
        return svg_match.group(0)

    return svg_content


def save_diagram(content: str, image_dir: str, slide_num: int, file_type: str) -> str:
    """Save diagram to local storage with organized naming.

    Args:
        content: The diagram content (SVG string or binary data)
        image_dir: Directory to save the file
        slide_num: Slide number for naming
        file_type: File extension ('svg' or 'png')

    Returns:
        Path to the saved file
    """
    filename = f"slide{slide_num:02d}-diagram.{file_type}"
    filepath = os.path.join(image_dir, filename)

    mode = 'w' if file_type == 'svg' else 'wb'
    with open(filepath, mode) as f:
        f.write(content)

    return filepath


# ---------------------------------------------------------------------------
# SVG to PNG conversion
# ---------------------------------------------------------------------------

def convert_svg_to_png(svg_path: str) -> str:
    """Convert SVG file to PNG using cairosvg.

    Args:
        svg_path: Path to the SVG file

    Returns:
        Path to the created PNG file

    Raises:
        ImportError: If cairosvg is not installed
        Exception: If conversion fails
    """
    try:
        import cairosvg

        png_path = svg_path.replace('.svg', '.png')

        cairosvg.svg2png(
            url=svg_path,
            write_to=png_path,
            output_width=1024,
            output_height=768
        )

        logger.debug(f"Converted {svg_path} to {png_path}")
        return png_path

    except ImportError:
        logger.error(
            "cairosvg not installed. Please install with: pip install cairosvg"
        )
        raise
    except Exception as e:
        logger.error(f"Error converting SVG to PNG: {e}")
        raise


# ---------------------------------------------------------------------------
# DALL-E image generation
# ---------------------------------------------------------------------------

def generate_dalle_image(
    description: str,
    image_dir: str,
    slide_num: int,
    client: LLMClient,
) -> str:
    """Generate diagram using OpenAI compatible image API.

    Args:
        description: Text description of the image to generate
        image_dir: Directory to save the generated image
        slide_num: Slide number for naming
        client: Configured LLMClient instance with image support

    Returns:
        Path to the saved image file

    Raises:
        ValueError: If client is not provided
        Exception: If image generation fails
    """
    if not client:
        raise ValueError("LLM client is required for image generation")

    prompt = f"""Create a professional technical diagram or flowchart showing:
{description}

Requirements:
- Clean, minimalist design
- Professional appearance suitable for technical presentations
- White or light background
- Clear labels and arrows
- High contrast for visibility
- Use of professional color scheme (blues, grays)
- Appropriate for business/technical audience"""

    image_url = client.generate_image(prompt=prompt, size="1024x1024")
    image_data = requests.get(image_url).content

    image_path = os.path.join(image_dir, f"slide{slide_num:02d}-dalle.png")
    with open(image_path, 'wb') as f:
        f.write(image_data)

    logger.debug(f"Generated image: {image_path}")
    return image_path


# ---------------------------------------------------------------------------
# Diagram layout application
# ---------------------------------------------------------------------------

def apply_diagram_layout(
    slide,
    content: str,
    suggestions: Dict,
    image_dir: str,
    slide_num: int,
    client: Optional[LLMClient] = None,
    image_gen: str = 'none',
):
    """Apply diagram layout with generated image to slide.

    Args:
        slide: A python-pptx Slide object
        content: Text content for key points below the diagram
        suggestions: Dict with 'VISUALS' key containing description
        image_dir: Directory for storing generated images
        slide_num: Current slide number
        client: LLMClient instance for image generation
        image_gen: Image generation mode ('none', 'claude', 'dalle', 'openai')
    """
    try:
        if not suggestions.get('VISUALS'):
            logger.warning("No visual suggestions provided")
            return

        if image_gen in ('dalle', 'openai'):
            image_path = generate_dalle_image(
                suggestions['VISUALS'],
                image_dir,
                slide_num,
                client
            )
            logger.debug(f"Generated DALL-E image: {image_path}")
        else:
            # Use LLM for SVG generation
            svg_content = generate_diagram_svg(suggestions['VISUALS'], client)
            svg_path = save_diagram(svg_content, image_dir, slide_num, 'svg')
            image_path = convert_svg_to_png(svg_path)
            logger.debug(f"Generated SVG diagram: {svg_path}")

        # Add image to slide
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(4)
        slide.shapes.add_picture(image_path, left, top, width, height)

        logger.debug(f"Added image to slide {slide_num}: {image_path}")

        # Add key points below diagram
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5),
            Inches(8), Inches(1.5)
        )
        tf = content_box.text_frame

        for line in content.split('\n'):
            stripped = line.strip()
            if stripped and stripped.startswith(('-', '•', '*')):
                p = tf.add_paragraph()
                p.level = 1
                p.text = stripped.lstrip('-•* ')

    except Exception as e:
        logger.error(f"Error in apply_diagram_layout: {e}")
        raise


# ---------------------------------------------------------------------------
# MCP image generation
# ---------------------------------------------------------------------------

def _mcp_cache_key(prompt: str, model: str, aspect_ratio: str) -> str:
    """Generate a 12-char hex cache key from prompt parameters."""
    data = f"{model}:{prompt}:{aspect_ratio}"
    return hashlib.sha256(data.encode()).hexdigest()[:12]


async def generate_mcp_image(
    prompt: str,
    output_dir: str,
    slide_num: int,
    mcp_manager,
    server_name: str = "txt2img",
    model: str = "gemini-2.0-flash-preview-image-generation",
    classification: str = "internal",
    aspect_ratio: str = "16:9",
    cache_dir: str | None = None,
) -> str | None:
    """Generate an image via MCP txt2img server.

    Args:
        prompt: Text description of the image to generate.
        output_dir: Directory to save the generated image.
        slide_num: Slide number for file naming.
        mcp_manager: MCPManager instance for server communication.
        server_name: MCP server name from config (default: txt2img).
        model: Image generation model name.
        classification: Content classification (default: internal).
        aspect_ratio: Image aspect ratio (default: 16:9).
        cache_dir: Directory for prompt-hash cache. None disables caching.

    Returns:
        Path to the saved image file, or None on failure.
    """
    cache_key = _mcp_cache_key(prompt, model, aspect_ratio)

    # Check cache
    if cache_dir:
        os.makedirs(cache_dir, exist_ok=True)
        cached_path = os.path.join(cache_dir, f"{cache_key}.png")
        if os.path.exists(cached_path):
            output_path = os.path.join(output_dir, f"slide_{slide_num:02d}_gen.png")
            shutil.copy2(cached_path, output_path)
            logger.info(f"Cache hit for slide {slide_num}: {cache_key}")
            return output_path

    try:
        result = await mcp_manager.call_tool(
            server_name,
            "generate_slide_image",
            {
                "prompt": prompt,
                "model": model,
                "classification": classification,
                "aspect_ratio": aspect_ratio,
            },
        )

        # Extract base64 image data from MCP response
        if not result.content or not hasattr(result.content[0], 'text') or not result.content[0].text:
            logger.warning(f"MCP returned empty content for slide {slide_num}")
            return None
        image_b64 = result.content[0].text
        image_data = base64.b64decode(image_b64)

        # Save to output directory
        output_path = os.path.join(output_dir, f"slide_{slide_num:02d}_gen.png")
        with open(output_path, 'wb') as f:
            f.write(image_data)

        # Save to cache
        if cache_dir:
            cached_path = os.path.join(cache_dir, f"{cache_key}.png")
            with open(cached_path, 'wb') as f:
                f.write(image_data)
            logger.debug(f"Cached image for slide {slide_num}: {cache_key}")

        logger.info(f"Generated MCP image for slide {slide_num}: {output_path}")
        return output_path

    except Exception as e:
        logger.warning(f"MCP image generation failed for slide {slide_num}: {e}")
        return None
