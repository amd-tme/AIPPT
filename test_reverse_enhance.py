import sys
import tempfile
from pathlib import Path
from pptx import Presentation
from aippt.cli import build_parser, cmd_reverse

# Create a minimal test PPTX
with tempfile.TemporaryDirectory() as tmpdir:
    tmp_path = Path(tmpdir)
    
    # Create test PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test"
    
    pptx_path = tmp_path / "test.pptx"
    prs.save(str(pptx_path))
    
    # Try reverse with --enhance
    parser = build_parser()
    args = parser.parse_args(['reverse', str(pptx_path), str(tmp_path / "output.md"), '--enhance'])
    
    try:
        result = cmd_reverse(args)
        print(f"Result: {result}")
    except TypeError as e:
        print(f"ERROR: {e}")
        sys.exit(1)
