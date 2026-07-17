"""
clean_deck.py — Post-process a generated PPTX to match AMD corporate standards.

Fixes:
  1. Font normalization: replaces stray non-Klavika/Arial fonts with Arial
  2. Comment cleanup: removes all author comments embedded in slides
  3. Reports a summary table of changes made

Usage:
    python scripts/clean_deck.py output.pptx [--in-place] [--dry-run]
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    sys.exit("python-pptx not installed — run: pip install python-pptx")

# Fonts that should stay as-is; everything else becomes FALLBACK_FONT
ALLOWED_FONTS = {"Klavika", "Arial", "Consolas", "Calibri", None}
FALLBACK_FONT = "Arial"


def _normalize_fonts(prs: Presentation) -> dict:
    stats = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    name = font.name
                    if name not in ALLOWED_FONTS:
                        stats[name] = stats.get(name, 0) + 1
                        font.name = FALLBACK_FONT
    return stats


def _remove_comments(prs: Presentation) -> int:
    count = 0
    for slide in prs.slides:
        # Access the slide's comments part if it exists
        if hasattr(slide, "comments") and slide.comments is not None:
            comment_list = slide.comments.comment_lst
            n = len(comment_list)
            for _ in range(n):
                comment_list.remove(comment_list[0])
            count += n
    return count


def main():
    parser = argparse.ArgumentParser(description="Clean an AMD-generated PPTX deck.")
    parser.add_argument("input", help="Path to .pptx file")
    parser.add_argument("--in-place", action="store_true", help="Overwrite input file")
    parser.add_argument("--dry-run", action="store_true", help="Report changes without saving")
    args = parser.parse_args()

    src = Path(args.input)
    if not src.exists():
        sys.exit(f"File not found: {src}")

    prs = Presentation(src)

    font_stats = _normalize_fonts(prs)
    comment_count = _remove_comments(prs)

    print(f"\nclean_deck.py — {src.name}")
    print("-" * 50)
    if font_stats:
        print("Font replacements → Arial:")
        for name, n in sorted(font_stats.items(), key=lambda x: -x[1]):
            print(f"  {name!r}: {n} run(s)")
    else:
        print("  No non-standard fonts found.")

    print(f"Comments removed: {comment_count}")

    if args.dry_run:
        print("\n[dry-run] No file written.")
        return

    out = src if args.in_place else src.with_stem(src.stem + "-clean")
    prs.save(out)
    print(f"\nSaved: {out}")


if __name__ == "__main__":
    main()
